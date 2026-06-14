import json
import os
import shutil
import subprocess
import sys
import zipfile
from io import BytesIO
from pathlib import Path, PurePosixPath

from lxml.etree import SubElement
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

_OUTPUTS_VIRTUAL_PREFIX = "/mnt/user-data/outputs/"

_IMAGE_DEPENDENT_LAYOUTS = {"full_bleed_image", "content_image"}


def _pptxgenjs_enabled() -> bool:
    return os.getenv("SOPHIA_PPTXGENJS", "").strip().lower() in {"1", "true", "yes"}


def _candidate_js_runtime_dirs() -> list[Path]:
    configured = os.getenv("SOPHIA_ARTIFACT_JS_RUNTIME")
    candidates: list[Path] = []
    if configured:
        candidates.append(Path(configured))
    candidates.append(Path("/app/backend/packages/harness/deerflow/sophia/js"))
    script_path = Path(__file__).resolve()
    for parent in script_path.parents:
        candidates.append(parent / "backend/packages/harness/deerflow/sophia/js")
    unique: list[Path] = []
    seen: set[Path] = set()
    for candidate in candidates:
        resolved = candidate.resolve()
        if resolved not in seen:
            unique.append(resolved)
            seen.add(resolved)
    return unique


def _pptxgenjs_script() -> Path | None:
    for runtime_dir in _candidate_js_runtime_dirs():
        script = runtime_dir / "compile_pptx.mjs"
        if script.exists():
            return script
    return None


def _generate_with_pptxgenjs(plan_file: str, slide_images: list[str], output_file: str) -> str:
    script = _pptxgenjs_script()
    node = shutil.which("node")
    if script is None or node is None:
        raise RuntimeError("PptxGenJS runtime is unavailable")
    command = [
        node,
        str(script),
        "--plan-file",
        plan_file,
        "--output-file",
        output_file,
    ]
    if slide_images:
        command.append("--slide-images")
        command.extend(slide_images)
    result = subprocess.run(command, check=False, text=True, capture_output=True)
    if result.stdout:
        print(result.stdout.rstrip())
    if result.stderr:
        print(result.stderr.rstrip(), file=sys.stderr)
    if result.returncode != 0:
        raise RuntimeError(f"PptxGenJS compiler failed with exit code {result.returncode}")
    return "Successfully generated presentation with PptxGenJS"


def generate_ppt(
    plan_file: str,
    slide_images: list[str],
    output_file: str,
) -> str:
    """
    Generate a PowerPoint presentation from a plan, optionally with slide images.

    Args:
        plan_file: Path to JSON file containing presentation plan
        slide_images: Optional list of paths to slide images in order
        output_file: Path to output PPTX file

    Returns:
        Status message
    """
    if _pptxgenjs_enabled():
        return _generate_with_pptxgenjs(plan_file, slide_images, output_file)

    plan = load_plan(plan_file)
    _emit_plan_lint(plan)
    slide_width, slide_height = slide_dimensions(plan)
    prs = Presentation()
    prs.slide_width = slide_width
    prs.slide_height = slide_height

    blank_layout = prs.slide_layouts[6]  # Blank layout
    slides_info = plan.get("slides", [])
    if not slides_info:
        slides_info = [{"slide_number": 1, "type": "title", "title": plan.get("title", "Presentation")}]

    if slide_images and len(slide_images) != len(slides_info):
        print(
            "Slide image count does not match plan slide count; extra slides "
            "will use plan-referenced visuals or text layouts",
            file=sys.stderr,
        )

    for index, slide_info in enumerate(slides_info):
        slide = prs.slides.add_slide(blank_layout)
        image_from_cli = index < len(slide_images)
        image_path = slide_images[index] if image_from_cli else slide_visual_image_path(slide_info, plan_file, output_file)
        if image_path and not os.path.exists(image_path):
            if image_from_cli:
                raise FileNotFoundError(f"Slide image not found: {image_path}")
            # Degrade to the text layout instead of aborting the whole deck
            # (or leaving an empty visual card) on a missing file.
            print(
                f"Slide image missing, using text layout: {image_path}",
                file=sys.stderr,
            )
            image_path = None
        layout_name = resolve_layout(slide_info, image_path)
        if image_path is None and layout_name in _IMAGE_DEPENDENT_LAYOUTS:
            # Explicit image layout without a usable image: re-resolve as if
            # no layout was requested. Title-typed slides keep the title
            # treatment; everything else falls back to content_text.
            print(f"Layout '{layout_name}' requires an image, using text layout", file=sys.stderr)
            layout_name = "title" if str(slide_info.get("type") or "").lower() == "title" else "content_text"
        LAYOUT_DISPATCH[layout_name](slide, slide_info, plan, image_path, slide_width, slide_height)
        add_speaker_notes(slide, slide_info)

    picture_count = save_and_validate_pptx(prs, output_file)
    return f"Successfully generated presentation with {len(prs.slides)} slides (picture_count={picture_count})"


THEMES = {
    "boardroom": {
        "background": RGBColor(0x0E, 0x1A, 0x2B),
        "title": RGBColor(0xF8, 0xFA, 0xF5),
        "body": RGBColor(0xB8, 0xC4, 0xD6),
        "accent": RGBColor(0x3B, 0x82, 0xC4),
        "card": RGBColor(0x14, 0x24, 0x3A),
        "border": RGBColor(0x2C, 0x40, 0x5C),
        "accent_deep": RGBColor(0x12, 0x2B, 0x47),
        "overlay": RGBColor(0x0A, 0x14, 0x20),
        "overlay_text": RGBColor(0xF8, 0xFA, 0xF5),
        "title_font": "Georgia",
        "body_font": "Calibri",
    },
    "daylight": {
        "background": RGBColor(0xF8, 0xFA, 0xFC),
        "title": RGBColor(0x12, 0x18, 0x26),
        "body": RGBColor(0x2D, 0x37, 0x48),
        "accent": RGBColor(0x1C, 0x7E, 0xD6),
        "card": RGBColor(0xFF, 0xFF, 0xFF),
        "border": RGBColor(0xCB, 0xD5, 0xE1),
        "accent_deep": RGBColor(0x18, 0x64, 0xAB),
        "overlay": RGBColor(0x10, 0x2A, 0x43),
        "overlay_text": RGBColor(0xFF, 0xFF, 0xFF),
        "title_font": "Calibri",
        "body_font": "Calibri",
    },
    "ember": {
        "background": RGBColor(0x1C, 0x14, 0x10),
        "title": RGBColor(0xF8, 0xF1, 0xE7),
        "body": RGBColor(0xD9, 0xC8, 0xB4),
        "accent": RGBColor(0xE0, 0x8A, 0x00),
        "card": RGBColor(0x28, 0x1D, 0x16),
        "border": RGBColor(0x4A, 0x38, 0x26),
        "accent_deep": RGBColor(0x33, 0x21, 0x14),
        "overlay": RGBColor(0x14, 0x0E, 0x0A),
        "overlay_text": RGBColor(0xF8, 0xF1, 0xE7),
        "title_font": "Georgia",
        "body_font": "Calibri",
    },
    "mist": {
        "background": RGBColor(0xF1, 0xF5, 0xF9),
        "title": RGBColor(0x1E, 0x29, 0x3B),
        "body": RGBColor(0x33, 0x41, 0x55),
        "accent": RGBColor(0x47, 0x55, 0x69),
        "card": RGBColor(0xFF, 0xFF, 0xFF),
        "border": RGBColor(0xCB, 0xD5, 0xE1),
        "accent_deep": RGBColor(0x33, 0x41, 0x55),
        "overlay": RGBColor(0x1E, 0x29, 0x3B),
        "overlay_text": RGBColor(0xF8, 0xFA, 0xFC),
        "title_font": "Arial",
        "body_font": "Arial",
    },
    # Sage & Terracotta (60/30/10: warm paper dominant, paper-white cards
    # supporting, sage accent ~10% with terracotta reserved for highlights).
    "terra": {
        "background": RGBColor(0xF4, 0xF1, 0xDE),
        "title": RGBColor(0x2C, 0x2C, 0x2C),
        "body": RGBColor(0x44, 0x40, 0x3A),
        "accent": RGBColor(0x87, 0xA9, 0x6B),
        "card": RGBColor(0xFA, 0xF7, 0xEC),
        "border": RGBColor(0xDC, 0xD5, 0xBF),
        "accent_deep": RGBColor(0xE0, 0x7A, 0x5F),
        "overlay": RGBColor(0x2C, 0x2C, 0x2C),
        "overlay_text": RGBColor(0xF4, 0xF1, 0xDE),
        "title_font": "Georgia",
        "body_font": "Calibri",
    },
    # Black & Gold (60/30/10: near-black dominant, charcoal cards supporting,
    # gold accent ~10%; pure-black divider/closing slides).
    "noir": {
        "background": RGBColor(0x11, 0x11, 0x11),
        "title": RGBColor(0xF4, 0xF6, 0xF6),
        "body": RGBColor(0xC8, 0xC8, 0xC3),
        "accent": RGBColor(0xBF, 0x9A, 0x4A),
        "card": RGBColor(0x1C, 0x1C, 0x1C),
        "border": RGBColor(0x3A, 0x33, 0x21),
        "accent_deep": RGBColor(0x00, 0x00, 0x00),
        "overlay": RGBColor(0x00, 0x00, 0x00),
        "overlay_text": RGBColor(0xF4, 0xF6, 0xF6),
        "title_font": "Georgia",
        "body_font": "Calibri",
    },
}

_STYLE_ALIASES = {
    "dark-premium": "boardroom",
    "keynote": "boardroom",
    "glassmorphism": "boardroom",
    "business": "daylight",
    "minimal": "mist",
    "academic": "mist",
    "creative": "ember",
    "earthy": "terra",
    "organic": "terra",
    "luxury": "noir",
    "premium-dark": "noir",
}


def slide_theme(plan: dict) -> dict:
    key = str(plan.get("theme") or plan.get("style") or "").strip().lower()
    key = _STYLE_ALIASES.get(key, key)
    return THEMES.get(key, THEMES["daylight"])


def resolve_layout(slide_info: dict, image_path: str | None) -> str:
    explicit = str(slide_info.get("layout") or "").strip().lower()
    if explicit in LAYOUT_DISPATCH:
        return explicit
    if explicit:
        print(f"Unknown layout '{explicit}', inferring", file=sys.stderr)
    if str(slide_info.get("type") or "").lower() == "title":
        return "title"
    return "content_image" if image_path else "content_text"


def tint_color(color, factor: float) -> RGBColor:
    return RGBColor(*(int(channel + (255 - channel) * factor) for channel in color))


def set_fill_transparency(shape, alpha_pct: int) -> None:
    # alpha_pct is the remaining opacity percent: 40 -> 40% opaque / 60% transparent.
    srgb = shape.fill._xPr.find(f"{qn('a:solidFill')}/{qn('a:srgbClr')}")
    if srgb is None:
        return
    for existing in srgb.findall(qn("a:alpha")):
        srgb.remove(existing)
    SubElement(srgb, qn("a:alpha")).set("val", str(int(alpha_pct) * 1000))


def apply_slide_background(slide, theme: dict) -> None:
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = theme["background"]


def _add_accent_shape(slide, theme: dict, shape_type, left, top, width, height):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["accent"]
    shape.line.color.rgb = theme["accent"]
    return shape


def _motif_rule(slide, theme: dict, slide_width) -> None:
    # Thin accent underline below the title block.
    _add_accent_shape(slide, theme, MSO_SHAPE.RECTANGLE, Inches(0.78), Inches(1.42), Inches(1.9), Inches(0.045))


def _motif_corner(slide, theme: dict, slide_width) -> None:
    # L-shaped bracket in the top-left corner (two thin accent rects).
    _add_accent_shape(slide, theme, MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(0.3), Inches(0.85), Inches(0.05))
    _add_accent_shape(slide, theme, MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(0.3), Inches(0.05), Inches(0.85))


def _motif_dot_grid(slide, theme: dict, slide_width) -> None:
    # 3x3 grid of 4px accent dots, top-right corner.
    dot = int(Pt(4))
    spacing = int(Inches(0.18))
    grid_extent = 2 * spacing + dot
    left_base = int(slide_width) - int(Inches(0.7)) - grid_extent
    top_base = int(Inches(0.4))
    for row in range(3):
        for col in range(3):
            _add_accent_shape(slide, theme, MSO_SHAPE.OVAL, left_base + col * spacing, top_base + row * spacing, dot, dot)


_MOTIF_RENDERERS = {
    "rule": _motif_rule,
    "corner": _motif_corner,
    "dot_grid": _motif_dot_grid,
}


def apply_motif(slide, motif, theme: dict, slide_width) -> None:
    """Render the plan-level repeated accent motif. Absent/unknown motif is a no-op."""
    renderer = _MOTIF_RENDERERS.get(str(motif or "").strip().lower())
    if renderer is None:
        return
    renderer(slide, theme, int(slide_width) if slide_width else int(Inches(13.333)))


def add_title_and_subtitle(slide, slide_info: dict, plan: dict, theme: dict) -> None:
    title = str(slide_info.get("title") or plan.get("title") or "Untitled")
    subtitle = str(slide_info.get("subtitle") or "")

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(12), Inches(1.05))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(34 if len(title) < 50 else 28)
    p.font.color.rgb = theme["title"]
    p.font.name = theme["title_font"]

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.55), Inches(11.6), Inches(0.55))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(17)
        p.font.color.rgb = theme["body"]
        p.font.name = theme["body_font"]


def add_text_layout_slide(slide, slide_info: dict, plan: dict, image_path: str | None = None, slide_width=None, slide_height=None) -> None:
    theme = slide_theme(plan)
    apply_slide_background(slide, theme)
    add_title_and_subtitle(slide, slide_info, plan, theme)
    apply_motif(slide, plan.get("motif"), theme, slide_width)
    points = slide_points(slide_info)

    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.25), Inches(0.08), Inches(4.35))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = theme["accent"]
    accent_bar.line.color.rgb = theme["accent"]
    body_box = slide.shapes.add_textbox(Inches(1.05), Inches(2.2), Inches(11.2), Inches(4.6))
    body_frame = body_box.text_frame
    body_frame.word_wrap = True
    body_frame.clear()
    if not points:
        points = [str(slide_info.get("summary") or slide_info.get("body") or "")]
    for idx, point in enumerate(points[:6]):
        paragraph = body_frame.paragraphs[0] if idx == 0 else body_frame.add_paragraph()
        paragraph.text = str(point)
        paragraph.level = 0
        paragraph.font.size = Pt(20 if len(points) <= 4 else 17)
        paragraph.font.color.rgb = theme["body"]
        paragraph.font.name = theme["body_font"]
        paragraph.space_after = Pt(8)


def add_content_with_image_slide(slide, slide_info: dict, plan: dict, image_path: str, slide_width, slide_height) -> None:
    theme = slide_theme(plan)
    apply_slide_background(slide, theme)
    add_title_and_subtitle(slide, slide_info, plan, theme)
    apply_motif(slide, plan.get("motif"), theme, slide_width)
    points = slide_points(slide_info)
    if not points:
        points = [str(slide_info.get("summary") or slide_info.get("body") or "")]

    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.25), Inches(0.08), Inches(4.35))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = theme["accent"]
    accent_bar.line.color.rgb = theme["accent"]

    body_box = slide.shapes.add_textbox(Inches(1.05), Inches(2.15), Inches(5.3), Inches(4.8))
    body_frame = body_box.text_frame
    body_frame.word_wrap = True
    body_frame.clear()
    for idx, point in enumerate(points[:5]):
        paragraph = body_frame.paragraphs[0] if idx == 0 else body_frame.add_paragraph()
        paragraph.text = str(point)
        paragraph.level = 0
        paragraph.font.size = Pt(18 if len(points) <= 4 else 16)
        paragraph.font.color.rgb = theme["body"]
        paragraph.font.name = theme["body_font"]
        paragraph.space_after = Pt(8)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.05), Inches(5.75), Inches(4.85))
    card.fill.solid()
    card.fill.fore_color.rgb = theme["card"]
    card.line.color.rgb = theme["border"]
    card.line.width = Pt(1.25)
    add_picture_in_box(slide, image_path, Inches(7.05), Inches(2.3), Inches(5.25), Inches(4.35))


def add_title_slide(slide, slide_info: dict, plan: dict, image_path: str | None, slide_width, slide_height) -> None:
    theme = slide_theme(plan)
    apply_slide_background(slide, theme)
    title = str(slide_info.get("title") or plan.get("title") or "Untitled")
    subtitle = str(slide_info.get("subtitle") or "")
    if image_path:
        add_full_bleed_picture(slide, image_path, slide_width, slide_height)
        add_overlay_band(slide, theme, slide_width, slide_height, title, subtitle)
        return

    text_width = int(slide_width) - int(Inches(1.8))
    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.35), text_width, Inches(1.4))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = theme["title"]
    p.font.name = theme["title_font"]

    rule_width = Inches(2.2)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int((int(slide_width) - int(rule_width)) / 2), Inches(4.0), rule_width, Inches(0.05))
    rule.fill.solid()
    rule.fill.fore_color.rgb = theme["accent"]
    rule.line.color.rgb = theme["accent"]

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.9), Inches(4.3), text_width, Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_frame.clear()
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.color.rgb = theme["body"]
        p.font.name = theme["body_font"]


def add_full_bleed_image_slide(slide, slide_info: dict, plan: dict, image_path: str, slide_width, slide_height) -> None:
    theme = slide_theme(plan)
    apply_slide_background(slide, theme)
    # Picture first so the overlay band sits above it in z-order.
    add_full_bleed_picture(slide, image_path, slide_width, slide_height)
    add_overlay_band(slide, theme, slide_width, slide_height, str(slide_info.get("title") or ""), str(slide_info.get("subtitle") or ""))


def add_section_divider_slide(slide, slide_info: dict, plan: dict, image_path: str | None, slide_width, slide_height) -> None:
    theme = slide_theme(plan)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = theme["accent_deep"]

    number = str(slide_info.get("section_number") or slide_info.get("slide_number") or "")
    if number:
        number_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(6), Inches(2.3))
        number_frame = number_box.text_frame
        number_frame.clear()
        p = number_frame.paragraphs[0]
        p.text = number
        p.font.bold = True
        p.font.size = Pt(110)
        p.font.color.rgb = tint_color(theme["accent_deep"], 0.18)
        p.font.name = theme["title_font"]

    title_box = slide.shapes.add_textbox(Inches(0.85), Inches(3.45), int(slide_width) - int(Inches(1.7)), Inches(1.3))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = str(slide_info.get("title") or plan.get("title") or "")
    p.font.bold = True
    p.font.size = Pt(40)
    p.font.color.rgb = theme["overlay_text"]
    p.font.name = theme["title_font"]

    subtitle = str(slide_info.get("subtitle") or "")
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.9), Inches(4.8), int(slide_width) - int(Inches(1.8)), Inches(0.7))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_frame.clear()
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.color.rgb = tint_color(theme["accent_deep"], 0.65)
        p.font.name = theme["body_font"]


def add_quote_slide(slide, slide_info: dict, plan: dict, image_path: str | None, slide_width, slide_height) -> None:
    theme = slide_theme(plan)
    apply_slide_background(slide, theme)
    points = slide_points(slide_info)
    quote = str(slide_info.get("quote") or (points[0] if points else "") or slide_info.get("title") or "")
    attribution = str(slide_info.get("attribution") or slide_info.get("subtitle") or "")

    mark_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.45), Inches(2.6), Inches(2.1))
    mark_frame = mark_box.text_frame
    mark_frame.clear()
    p = mark_frame.paragraphs[0]
    p.text = "“"
    p.font.bold = True
    p.font.size = Pt(120)
    p.font.color.rgb = theme["accent"]
    p.font.name = theme["title_font"]

    quote_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.4), int(slide_width) - int(Inches(3.0)), Inches(2.9))
    quote_frame = quote_box.text_frame
    quote_frame.word_wrap = True
    quote_frame.clear()
    p = quote_frame.paragraphs[0]
    p.text = quote
    p.font.italic = True
    p.font.size = Pt(32)
    p.font.color.rgb = theme["title"]
    p.font.name = theme["title_font"]

    if attribution:
        attribution_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), int(slide_width) - int(Inches(3.0)), Inches(0.6))
        attribution_frame = attribution_box.text_frame
        attribution_frame.word_wrap = True
        attribution_frame.clear()
        p = attribution_frame.paragraphs[0]
        p.text = f"— {attribution}"
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(16)
        p.font.color.rgb = theme["body"]
        p.font.name = theme["body_font"]


def add_two_column_slide(slide, slide_info: dict, plan: dict, image_path: str | None, slide_width, slide_height) -> None:
    theme = slide_theme(plan)
    apply_slide_background(slide, theme)
    add_title_and_subtitle(slide, slide_info, plan, theme)
    apply_motif(slide, plan.get("motif"), theme, slide_width)

    columns = [column for column in (slide_info.get("columns") or []) if isinstance(column, dict)]
    if not columns:
        points = slide_points(slide_info)
        midpoint = (len(points) + 1) // 2
        columns = [{"points": points[:midpoint]}, {"points": points[midpoint:]}]
    columns = columns[:2]

    margin = int(Inches(0.75))
    gutter = int(Inches(0.6))
    column_width = int((int(slide_width) - 2 * margin - gutter) / 2)
    top = int(Inches(2.3))
    for index, column in enumerate(columns):
        left = margin + index * (column_width + gutter)
        rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, column_width, Inches(0.05))
        rule.fill.solid()
        rule.fill.fore_color.rgb = theme["accent"]
        rule.line.color.rgb = theme["accent"]

        column_box = slide.shapes.add_textbox(left, top + int(Inches(0.25)), column_width, Inches(4.3))
        column_frame = column_box.text_frame
        column_frame.word_wrap = True
        column_frame.clear()
        wrote_first = False
        heading = str(column.get("heading") or "")
        if heading:
            paragraph = column_frame.paragraphs[0]
            wrote_first = True
            paragraph.text = heading
            paragraph.font.bold = True
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = theme["accent"]
            paragraph.font.name = theme["title_font"]
            paragraph.space_after = Pt(8)
        column_points = [str(point) for point in (column.get("points") or []) if str(point).strip()]
        for point in column_points[:6]:
            paragraph = column_frame.add_paragraph() if wrote_first else column_frame.paragraphs[0]
            wrote_first = True
            paragraph.text = point
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = theme["body"]
            paragraph.font.name = theme["body_font"]
            paragraph.space_after = Pt(6)


def add_closing_slide(slide, slide_info: dict, plan: dict, image_path: str | None, slide_width, slide_height) -> None:
    theme = slide_theme(plan)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = theme["accent_deep"]

    text_width = int(slide_width) - int(Inches(1.8))
    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.7), text_width, Inches(1.3))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = str(slide_info.get("title") or "Thank You")
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = theme["overlay_text"]
    p.font.name = theme["title_font"]

    subtitle = str(slide_info.get("subtitle") or "")
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.9), Inches(4.15), text_width, Inches(0.7))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_frame.clear()
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.color.rgb = tint_color(theme["accent_deep"], 0.65)
        p.font.name = theme["body_font"]


def _add_stat_column(slide, theme: dict, stat: dict, left, column_width) -> None:
    value_box = slide.shapes.add_textbox(left, Inches(2.85), column_width, Inches(1.3))
    value_frame = value_box.text_frame
    value_frame.word_wrap = True
    value_frame.clear()
    p = value_frame.paragraphs[0]
    p.text = str(stat.get("value") or "")
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(54)
    p.font.color.rgb = theme["accent"]
    p.font.name = theme["title_font"]

    label_box = slide.shapes.add_textbox(left, Inches(4.2), column_width, Inches(0.7))
    label_frame = label_box.text_frame
    label_frame.word_wrap = True
    label_frame.clear()
    p = label_frame.paragraphs[0]
    p.text = str(stat.get("label") or "")
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.color.rgb = theme["body"]
    p.font.name = theme["body_font"]


def add_stat_band_slide(slide, slide_info: dict, plan: dict, image_path: str | None, slide_width, slide_height) -> None:
    stats = [stat for stat in (slide_info.get("stats") or []) if isinstance(stat, dict)][:4]
    if not stats:
        # No stats to feature: degrade to the two_column treatment.
        add_two_column_slide(slide, slide_info, plan, image_path, slide_width, slide_height)
        return
    theme = slide_theme(plan)
    apply_slide_background(slide, theme)
    if slide_info.get("title"):
        add_title_and_subtitle(slide, slide_info, plan, theme)
    apply_motif(slide, plan.get("motif"), theme, slide_width)

    margin = int(Inches(0.75))
    column_width = int((int(slide_width) - 2 * margin) / len(stats))
    for index, stat in enumerate(stats):
        _add_stat_column(slide, theme, stat, margin + index * column_width, column_width)


LAYOUT_DISPATCH = {
    "title": add_title_slide,
    "content_text": add_text_layout_slide,
    "content_image": add_content_with_image_slide,
    "full_bleed_image": add_full_bleed_image_slide,
    "section_divider": add_section_divider_slide,
    "quote": add_quote_slide,
    "two_column": add_two_column_slide,
    "closing": add_closing_slide,
    "stat_band": add_stat_band_slide,
}


_VISUAL_ANCHOR_LAYOUTS = {"full_bleed_image", "content_image", "section_divider", "quote", "stat_band"}


def _has_image_ref(slide_info: dict) -> bool:
    raw = slide_info.get("image") or slide_info.get("chart_path") or slide_info.get("visual_path")
    return isinstance(raw, str) and bool(raw.strip())


def _lint_resolved_layout(slide_info: dict) -> str:
    # Silent mirror of resolve_layout (no stderr) using the raw image ref as the
    # image signal — lint runs before any on-disk existence checks.
    explicit = str(slide_info.get("layout") or "").strip().lower()
    if explicit in LAYOUT_DISPATCH:
        return explicit
    if str(slide_info.get("type") or "").lower() == "title":
        return "title"
    return "content_image" if _has_image_ref(slide_info) else "content_text"


def _lint_slide_fact(slide_info: dict, position: int) -> dict:
    layout = _lint_resolved_layout(slide_info)
    return {
        "number": slide_info.get("slide_number") or position,
        "layout": layout,
        "anchor": layout in _VISUAL_ANCHOR_LAYOUTS or _has_image_ref(slide_info),
        "title": str(slide_info.get("title") or ""),
    }


def _lint_consecutive_layouts(facts: list[dict]) -> list[str]:
    return [
        f"slides {prev['number']} and {curr['number']} both use layout '{curr['layout']}' — vary the rhythm: reorder so a different layout (section_divider, quote, stat_band, or an image slide) sits between them, or merge the content"
        for prev, curr in zip(facts, facts[1:])
        if prev["layout"] == curr["layout"]
    ]


def _dry_stretch_warning(run: list[dict]) -> str:
    return f"slides {run[0]['number']}-{run[-1]['number']} are a text-only stretch with no visual anchor (full_bleed_image, content_image, section_divider, quote, or stat_band) — aim for at least one anchor every 3 slides"


def _lint_visual_cadence(facts: list[dict]) -> list[str]:
    warnings: list[str] = []
    run: list[dict] = []
    for fact in facts:
        if fact["anchor"]:
            if len(run) >= 3:
                warnings.append(_dry_stretch_warning(run))
            run = []
        else:
            run.append(fact)
    if len(run) >= 3:
        warnings.append(_dry_stretch_warning(run))
    return warnings


def _lint_title_lengths(facts: list[dict]) -> list[str]:
    return [f"slide {fact['number']} title is {len(fact['title'])} chars (>60) — long titles wrap in most layouts; tighten it" for fact in facts if len(fact["title"]) > 60]


def lint_plan(plan: dict) -> list[str]:
    """Cadence/design warnings for a presentation plan. Advisory only — never fails the build."""
    slides = [slide for slide in (plan.get("slides") or []) if isinstance(slide, dict)]
    if not slides:
        return []
    facts = [_lint_slide_fact(slide, position) for position, slide in enumerate(slides, start=1)]
    return _lint_consecutive_layouts(facts) + _lint_visual_cadence(facts) + _lint_title_lengths(facts)


def _emit_plan_lint(plan: dict) -> None:
    warnings = lint_plan(plan)
    for warning in warnings:
        print(f"PLAN_LINT: {warning}", file=sys.stderr)
    if warnings:
        print(f"PLAN_LINT: {len(warnings)} warning(s)", file=sys.stderr)


def add_overlay_band(slide, theme: dict, slide_width, slide_height, title: str, subtitle: str) -> None:
    band_height = int(int(slide_height) * 0.28)
    band_top = int(slide_height) - band_height
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, band_top, int(slide_width), band_height)
    band.fill.solid()
    band.fill.fore_color.rgb = theme["overlay"]
    band.line.fill.background()
    set_fill_transparency(band, 40)

    text_left = int(Inches(0.7))
    text_width = int(slide_width) - 2 * text_left
    if title:
        title_box = slide.shapes.add_textbox(text_left, band_top + int(Inches(0.2)), text_width, Inches(0.75))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.clear()
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(30)
        p.font.color.rgb = theme["overlay_text"]
        p.font.name = theme["title_font"]
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(text_left, band_top + int(Inches(1.0)), text_width, Inches(0.55))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_frame.clear()
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.color.rgb = theme["overlay_text"]
        p.font.name = theme["body_font"]


def add_full_bleed_picture(slide, image_path: str, slide_width, slide_height) -> None:
    if not os.path.exists(image_path):
        raise FileNotFoundError("Slide image not found")
    with Image.open(image_path) as img:
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")
        img_bytes = fitted_image_cover_payload(img, slide_width, slide_height)
        slide.shapes.add_picture(img_bytes, 0, 0, int(slide_width), int(slide_height))


def slide_points(slide_info: dict) -> list[str]:
    for key in ("key_points", "bullets", "points"):
        value = slide_info.get(key)
        if isinstance(value, list):
            return [str(item) for item in value if str(item).strip()]
    return []


def slide_visual_image_path(slide_info: dict, plan_file: str, output_file: str) -> str | None:
    raw = slide_info.get("image") or slide_info.get("chart_path") or slide_info.get("visual_path")
    if not isinstance(raw, str) or not raw.strip():
        return None
    return resolve_image_path(raw.strip(), plan_file, output_file)


def infer_outputs_root(output_file: str) -> Path | None:
    output_path = Path(output_file).resolve()
    for candidate in (output_path.parent, *output_path.parents):
        if candidate.name == "outputs":
            return candidate
    return output_path.parent


def resolve_image_path(raw_path: str, plan_file: str, output_file: str) -> str:
    normalized = raw_path.replace("\\", "/").strip()
    if normalized.startswith(_OUTPUTS_VIRTUAL_PREFIX):
        relative = normalized[len(_OUTPUTS_VIRTUAL_PREFIX) :].strip("/")
        if ".." in PurePosixPath(relative).parts:
            raise ValueError("Slide image path may not contain '..'")
        outputs_root = infer_outputs_root(output_file)
        if outputs_root is not None:
            return str((outputs_root / relative).resolve())
    path = Path(normalized)
    if path.is_absolute():
        return str(path)
    return str((Path(plan_file).resolve().parent / path).resolve())


def load_plan(plan_file: str) -> dict:
    with open(plan_file, encoding="utf-8") as f:
        plan = json.load(f)
    if not isinstance(plan, dict):
        raise ValueError("Invalid presentation plan: top-level JSON must be an object")
    slides = plan.get("slides")
    if slides is not None and not isinstance(slides, list):
        raise ValueError("Invalid presentation plan: slides must be a list")
    if isinstance(slides, list):
        for index, slide in enumerate(slides, start=1):
            if not isinstance(slide, dict):
                raise ValueError(f"Invalid presentation plan: slide {index} must be an object")
    return plan


def slide_dimensions(plan: dict):
    aspect_ratio = plan.get("aspect_ratio", "16:9")
    if aspect_ratio == "4:3":
        return Inches(10), Inches(7.5)
    return Inches(13.333), Inches(7.5)


def add_slide_image(slide, image_path: str, slide_width, slide_height) -> None:
    if not os.path.exists(image_path):
        raise FileNotFoundError("Slide image not found")
    with Image.open(image_path) as img:
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")
        img_bytes, left, top, width, height = fitted_image_payload(
            img,
            slide_width,
            slide_height,
        )
        slide.shapes.add_picture(img_bytes, left, top, width, height)


def add_picture_in_box(slide, image_path: str, left, top, width, height) -> None:
    if not os.path.exists(image_path):
        raise FileNotFoundError("Slide content image not found")
    with Image.open(image_path) as img:
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")
        img_bytes, image_width, image_height = fitted_image_payload_for_box(img, width, height)
        offset_left = left + int((int(width) - int(image_width)) / 2)
        offset_top = top + int((int(height) - int(image_height)) / 2)
        slide.shapes.add_picture(img_bytes, offset_left, offset_top, image_width, image_height)


def fitted_image_payload(img, slide_width, slide_height):
    img_width, img_height = img.size
    img_aspect = img_width / img_height
    slide_width_emu = int(slide_width)
    slide_height_emu = int(slide_height)
    slide_aspect = slide_width / slide_height
    if img_aspect > slide_aspect:
        new_width_emu = slide_width_emu
        new_height_emu = int(slide_width_emu / img_aspect)
        left = Inches(0)
        top = Inches((slide_height_emu - new_height_emu) / 914400)
    else:
        new_height_emu = slide_height_emu
        new_width_emu = int(slide_height_emu * img_aspect)
        left = Inches((slide_width_emu - new_width_emu) / 914400)
        top = Inches(0)
    img_bytes = BytesIO()
    img.save(img_bytes, format="JPEG", quality=95)
    img_bytes.seek(0)
    return img_bytes, left, top, Inches(new_width_emu / 914400), Inches(new_height_emu / 914400)


def fitted_image_payload_for_box(img, box_width, box_height):
    img_width, img_height = img.size
    img_aspect = img_width / img_height
    box_width_emu = int(box_width)
    box_height_emu = int(box_height)
    box_aspect = box_width_emu / box_height_emu
    if img_aspect > box_aspect:
        new_width_emu = box_width_emu
        new_height_emu = int(box_width_emu / img_aspect)
    else:
        new_height_emu = box_height_emu
        new_width_emu = int(box_height_emu * img_aspect)
    img_bytes = BytesIO()
    img.save(img_bytes, format="PNG")
    img_bytes.seek(0)
    return img_bytes, new_width_emu, new_height_emu


def fitted_image_cover_payload(img, slide_width, slide_height):
    # Center-crop to the slide aspect, then resize to slide pixel size
    # (1280 wide -> 720 tall for 16:9) so full-bleed embeds stay small.
    slide_aspect = int(slide_width) / int(slide_height)
    img_width, img_height = img.size
    img_aspect = img_width / img_height
    if img_aspect > slide_aspect:
        crop_width = max(1, int(round(img_height * slide_aspect)))
        offset = (img_width - crop_width) // 2
        box = (offset, 0, offset + crop_width, img_height)
    else:
        crop_height = max(1, int(round(img_width / slide_aspect)))
        offset = (img_height - crop_height) // 2
        box = (0, offset, img_width, offset + crop_height)
    target_width = 1280
    target_height = max(1, round(target_width / slide_aspect))
    resized = img.crop(box).resize((target_width, target_height))
    img_bytes = BytesIO()
    resized.save(img_bytes, format="JPEG", quality=90)
    img_bytes.seek(0)
    return img_bytes


def add_speaker_notes(slide, slide_info) -> None:
    notes = speaker_notes(slide_info)
    if not notes:
        return
    text_frame = slide.notes_slide.notes_text_frame
    if text_frame is not None:
        text_frame.text = "\n".join(notes)


def speaker_notes(slide_info) -> list[str]:
    notes = []
    if slide_info.get("title"):
        notes.append(f"Title: {slide_info['title']}")
    if slide_info.get("subtitle"):
        notes.append(f"Subtitle: {slide_info['subtitle']}")
    if slide_info.get("key_points"):
        notes.append("Key Points:")
        notes.extend(f"  • {point}" for point in slide_info["key_points"])
    return notes


def pptx_picture_count(prs) -> int:
    return sum(
        1
        for slide in prs.slides
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    )


def save_and_validate_pptx(prs, output_file: str) -> int:
    os.makedirs(os.path.dirname(output_file) or ".", exist_ok=True)
    picture_count = pptx_picture_count(prs)
    prs.save(output_file)
    if not os.path.exists(output_file) or os.path.getsize(output_file) <= 0:
        raise RuntimeError("PPTX save did not produce output bytes")
    with zipfile.ZipFile(output_file, "r") as archive:
        names = set(archive.namelist())
    required = {"[Content_Types].xml", "_rels/.rels", "ppt/presentation.xml"}
    missing = sorted(required.difference(names))
    if missing:
        raise RuntimeError(f"PPTX save missing required Office entries: {','.join(missing)}")
    return picture_count


def pptx_package_picture_count(output_file: str) -> int:
    with zipfile.ZipFile(output_file, "r") as archive:
        return sum(1 for name in archive.namelist() if name.startswith("ppt/media/"))


def main() -> int:
    import argparse
    parser = argparse.ArgumentParser(
        description="Generate PowerPoint presentation from slide images"
    )
    parser.add_argument(
        "--plan-file",
        required=True,
        help="Absolute path to JSON presentation plan file",
    )
    parser.add_argument(
        "--slide-images",
        nargs="*",
        default=[],
        help="Optional absolute paths to slide images in order (space-separated)",
    )
    parser.add_argument(
        "--output-file",
        required=True,
        help="Output path for generated PPTX file",
    )

    args = parser.parse_args()

    try:
        message = generate_ppt(
            args.plan_file,
            args.slide_images,
            args.output_file,
        )
        size = os.path.getsize(args.output_file)
        picture_count = pptx_package_picture_count(args.output_file)
        print(message)
        print(
            "PPT generation diagnostics: "
            f"slide_image_count={len(args.slide_images)} picture_count={picture_count} "
            f"output_ext=.pptx bytes={size}",
            file=sys.stderr,
        )
        return 0
    except Exception as e:
        print(
            f"Error while generating presentation: {type(e).__name__}: {e}",
            file=sys.stderr,
        )
        return 1


if __name__ == "__main__":
    sys.exit(main())
