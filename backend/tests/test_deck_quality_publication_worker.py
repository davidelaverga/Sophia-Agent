from __future__ import annotations

import asyncio
import hashlib
import io
import json
import threading
from datetime import UTC, datetime, timedelta
from decimal import Decimal
from pathlib import Path
from types import SimpleNamespace

import pytest
from pptx import Presentation
from pptx.util import Inches

from app.gateway.workers.deck_quality_publication_worker import (
    DeckQualityPublicationWorker,
    build_configured_deck_quality_publication_worker,
    get_deck_quality_publication_worker,
    get_deck_quality_publication_worker_or_none,
    install_deck_quality_publication_worker,
    start_deck_quality_publication_worker,
    stop_deck_quality_publication_worker,
)
from deerflow.config.deck_quality_config import DeckQualityConfig
from deerflow.sophia.deck_quality.canonical import canonical_json_bytes, canonical_sha256
from deerflow.sophia.deck_quality.idempotency import derive_quality_run_id
from deerflow.sophia.deck_quality.publication_persistence import (
    PublicationErrorCode,
    PublicationLease,
    PublicationOperationKind,
    PublicationRecord,
    PublicationRequest,
    PublicationState,
    expected_publication_source_pack_path,
)
from deerflow.sophia.deck_quality.publisher import (
    DeckQualitySourceHashes,
    DeckQualitySourcePack,
    deck_quality_immutable_artifact_snapshot_path,
    deck_quality_producer_archive_path,
    encode_deck_quality_producer_bundle,
)
from deerflow.sophia.deck_quality.schemas import BlindBrief, QualityInstrumentLock
from deerflow.sophia.storage.supabase_artifact_store import ArtifactObjectSizeError

CANARY_USER = "canary-user"
THREAD_ID = "thread-01"
BUILD_ID = "build-01"
ARTIFACT_VERSION_ID = "artifact-version-01"


def _instrument(**overrides: object) -> QualityInstrumentLock:
    values: dict[str, object] = {
        "rubric_version": "deck-rubric-v2",
        "rubric_hash": "a" * 64,
        "prompt_hashes": {
            "blind_visual": "b" * 64,
            "plan_realization": "c" * 64,
        },
        "judge_plan_hash": "d" * 64,
        "judge_profile_version": "v2",
        "evidence_preprocessor_version": "deck-evidence-v4",
        "judge_invoker_version": "deck-judge-invoker-v4",
        "assessment_schema_versions": {
            "blind_visual": "v4",
            "plan_realization": "v4",
        },
        "adjudication_policy_hash": "e" * 64,
    }
    values.update(overrides)
    return QualityInstrumentLock.model_validate(values)


def _config(*, users: frozenset[str] = frozenset({CANARY_USER})) -> DeckQualityConfig:
    return DeckQualityConfig(
        enabled=True,
        mode="shadow",
        scope="canary",
        canary_user_ids=users,
        max_quality_cost_usd=Decimal("0.60"),
    )


def _pptx_bytes() -> bytes:
    output = io.BytesIO()
    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    first = presentation.slides.add_slide(blank)
    first.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1)).text = (
        "Observe, appraise, act"
    )
    second = presentation.slides.add_slide(blank)
    second.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1)).text = (
        "Close the loop"
    )
    presentation.save(output)
    return output.getvalue()


def _pack(
    artifact_bytes: bytes,
    *,
    instrument: QualityInstrumentLock | None = None,
    task_id: str = "task-01",
    user_id: str = CANARY_USER,
    thread_id: str = THREAD_ID,
    artifact_version_id: str = ARTIFACT_VERSION_ID,
) -> DeckQualitySourcePack:
    instrument = instrument or _instrument()
    quality_run_id = derive_quality_run_id(
        artifact_version_id=artifact_version_id,
        campaign_id="DQ-1",
        instrument=instrument,
    )
    creative = {"subject_materials": ["control loops"], "signature": "cycle"}
    design = {"palette": ["ink", "paper"], "rhythm": "setup-mechanism-close"}
    build = {
        "build_id": BUILD_ID,
        "slides": [{"selector": "slide:1"}, {"selector": "slide:2"}],
    }
    brief = BlindBrief(
        request="Explain the control loop.",
        subject="Control loops",
        audience="AI engineers",
        goal="Explain the mechanism",
    )
    mechanical = {
        "checks": {
            "authoritative_gate": True,
            "native_editability": True,
            "render_success": True,
        }
    }
    return DeckQualitySourcePack(
        quality_run_id=quality_run_id,
        instrument=instrument,
        instrument_identity_hash=canonical_sha256(instrument),
        user_id=user_id,
        thread_id=thread_id,
        task_id=task_id,
        build_id=BUILD_ID,
        builder_run_id="builder-run-01",
        parent_builder_trace_id="019f675a-dcc1-7053-80dc-c6f572fb4d87",
        logical_artifact_id="artifact-01",
        artifact_version_id=artifact_version_id,
        manifest_revision=1,
        artifact_virtual_path="/mnt/user-data/outputs/deck.pptx",
        accepted_delivery_object_path=deck_quality_immutable_artifact_snapshot_path(
            user_id=user_id,
            thread_id=thread_id,
            build_id=BUILD_ID,
            logical_artifact_id="artifact-01",
            artifact_version_id=artifact_version_id,
            artifact_sha256=hashlib.sha256(artifact_bytes).hexdigest(),
            artifact_virtual_path="/mnt/user-data/outputs/deck.pptx",
        ),
        immutable_snapshot_object_path=deck_quality_immutable_artifact_snapshot_path(
            user_id=user_id,
            thread_id=thread_id,
            build_id=BUILD_ID,
            logical_artifact_id="artifact-01",
            artifact_version_id=artifact_version_id,
            artifact_sha256=hashlib.sha256(artifact_bytes).hexdigest(),
            artifact_virtual_path="/mnt/user-data/outputs/deck.pptx",
        ),
        artifact_sha256=hashlib.sha256(artifact_bytes).hexdigest(),
        creative_plan=creative,
        design_plan=design,
        build_record=build,
        blind_brief=brief,
        mechanical_record=mechanical,
        source_hashes=DeckQualitySourceHashes(
            creative_plan=canonical_sha256(creative),
            design_plan=canonical_sha256(design),
            build_record=canonical_sha256(build),
            blind_brief=canonical_sha256(brief),
            mechanical_record=canonical_sha256(mechanical),
        ),
    )


def _validated_update(record: PublicationRecord, **updates: object) -> PublicationRecord:
    payload = record.model_dump(mode="python")
    payload.update(updates)
    return PublicationRecord.model_validate(payload)


def _pending_record(
    pack: DeckQualitySourcePack,
    encoded: bytes,
    *,
    now: datetime,
    row_task_id: str = "task-01",
) -> PublicationRecord:
    source_hash = hashlib.sha256(encoded).hexdigest()
    instrument = pack.instrument
    return PublicationRecord(
        quality_run_id=pack.quality_run_id,
        campaign_id="DQ-1",
        scope_kind="canary",
        instrument_schema_version=instrument.schema_version,
        instrument_identity_hash=canonical_sha256(instrument),
        rubric_version=instrument.rubric_version,
        rubric_hash=instrument.rubric_hash,
        prompt_hashes=instrument.prompt_hashes,
        judge_plan_hash=instrument.judge_plan_hash,
        judge_profile_version=instrument.judge_profile_version,
        evidence_preprocessor_version=instrument.evidence_preprocessor_version,
        judge_invoker_version=instrument.judge_invoker_version,
        assessment_schema_versions=instrument.assessment_schema_versions,
        adjudication_policy_hash=instrument.adjudication_policy_hash,
        user_id=CANARY_USER,
        thread_id=THREAD_ID,
        task_id=row_task_id,
        build_id=BUILD_ID,
        builder_run_id="builder-run-01",
        parent_builder_trace_id="019f675a-dcc1-7053-80dc-c6f572fb4d87",
        logical_artifact_id="artifact-01",
        artifact_version_id=pack.artifact_version_id,
        manifest_revision=1,
        artifact_object_path=pack.artifact_storage_object_path,
        artifact_hash=pack.artifact_sha256,
        source_pack_object_path=expected_publication_source_pack_path(
            user_id=CANARY_USER,
            thread_id=THREAD_ID,
            build_id=BUILD_ID,
            quality_run_id=pack.quality_run_id,
        ),
        source_pack_hash=source_hash,
        state=PublicationState.PENDING,
        attempt_count=0,
        max_attempts=3,
        error_count=0,
        next_attempt_at=now,
        deadline_at=now + timedelta(minutes=3),
        quality_max_attempts=5,
        quality_run_deadline_at=now + timedelta(minutes=15),
        lease_epoch=0,
        requested_at=now,
        updated_at=now,
    )


class FakePublicationStore:
    def __init__(self, record: PublicationRecord, *, now: datetime) -> None:
        self.record = record
        self.now = now
        self.claim_calls: list[dict[str, object]] = []
        self.retry_calls: list[dict[str, object]] = []
        self.fail_calls: list[dict[str, object]] = []
        self.promote_calls: list[dict[str, object]] = []
        self.probe_calls = 0
        self.close_calls = 0
        self.get_calls = 0
        self.lose_claim_response_once = False
        self.lose_promote_response_once = False

    async def probe(self) -> None:
        self.probe_calls += 1

    async def aclose(self) -> None:
        self.close_calls += 1

    async def claim(
        self,
        *,
        lease_owner: str,
        claim_token: str,
        lease_seconds: int = 60,
        limit: int = 1,
    ) -> tuple[PublicationRecord, ...]:
        arguments = {
            "lease_owner": lease_owner,
            "claim_token": claim_token,
            "lease_seconds": lease_seconds,
            "limit": limit,
        }
        self.claim_calls.append(arguments)
        if self.record.state is PublicationState.RUNNING:
            if (
                self.record.lease_owner == lease_owner
                and self.record.claim_token == claim_token
            ):
                return (self.record,)
            return ()
        if self.record.state not in {
            PublicationState.PENDING,
            PublicationState.RETRY_WAIT,
        }:
            return ()
        claim_hash = canonical_sha256(arguments)
        self.record = _validated_update(
            self.record,
            state=PublicationState.RUNNING,
            attempt_count=self.record.attempt_count + 1,
            lease_owner=lease_owner,
            lease_epoch=self.record.lease_epoch + 1,
            lease_expires_at=self.now + timedelta(seconds=lease_seconds),
            claim_token=claim_token,
            claim_hash=claim_hash,
            started_at=self.record.started_at or self.now,
            updated_at=self.now,
            last_operation_kind=None,
            last_operation_token=None,
            last_operation_hash=None,
        )
        if self.lose_claim_response_once:
            self.lose_claim_response_once = False
            raise RuntimeError("simulated lost claim response")
        return (self.record,)

    async def retry(
        self,
        lease: PublicationLease,
        *,
        operation_token: str,
        error_code: PublicationErrorCode,
        error_stage: str,
        delay_seconds: int = 15,
    ) -> PublicationRecord:
        assert lease.quality_run_id == self.record.quality_run_id
        self.retry_calls.append(
            {
                "operation_token": operation_token,
                "error_code": error_code,
                "error_stage": error_stage,
                "delay_seconds": delay_seconds,
            }
        )
        terminal = self.record.attempt_count >= self.record.max_attempts
        self.record = _validated_update(
            self.record,
            state=(PublicationState.FAILED if terminal else PublicationState.RETRY_WAIT),
            error_count=self.record.error_count + 1,
            next_attempt_at=self.now + timedelta(seconds=delay_seconds),
            lease_owner=None,
            lease_expires_at=None,
            claim_token=None,
            claim_hash=None,
            last_operation_kind=PublicationOperationKind.RETRY,
            last_operation_token=operation_token,
            last_operation_hash="f" * 64,
            last_error_code=(
                PublicationErrorCode.ATTEMPT_LIMIT_EXHAUSTED
                if terminal
                else error_code
            ),
            last_error_stage=error_stage,
            last_error_at=self.now,
            finished_at=(self.now if terminal else None),
            updated_at=self.now,
        )
        return self.record

    async def fail(
        self,
        lease: PublicationLease,
        *,
        operation_token: str,
        error_code: PublicationErrorCode,
        error_stage: str,
    ) -> PublicationRecord:
        assert lease.quality_run_id == self.record.quality_run_id
        self.fail_calls.append(
            {
                "operation_token": operation_token,
                "error_code": error_code,
                "error_stage": error_stage,
            }
        )
        self.record = _validated_update(
            self.record,
            state=PublicationState.FAILED,
            error_count=self.record.error_count + 1,
            lease_owner=None,
            lease_expires_at=None,
            claim_token=None,
            claim_hash=None,
            last_operation_kind=PublicationOperationKind.FAIL,
            last_operation_token=operation_token,
            last_operation_hash="f" * 64,
            last_error_code=error_code,
            last_error_stage=error_stage,
            last_error_at=self.now,
            finished_at=self.now,
            updated_at=self.now,
        )
        return self.record

    async def promote(
        self,
        lease: PublicationLease,
        *,
        operation_token: str,
        input_manifest_object_path: str,
        input_manifest_hash: str,
    ) -> PublicationRecord:
        assert lease.quality_run_id == self.record.quality_run_id
        self.promote_calls.append(
            {
                "operation_token": operation_token,
                "input_manifest_object_path": input_manifest_object_path,
                "input_manifest_hash": input_manifest_hash,
            }
        )
        self.record = _validated_update(
            self.record,
            state=PublicationState.PUBLISHED,
            input_manifest_object_path=input_manifest_object_path,
            input_manifest_hash=input_manifest_hash,
            lease_owner=None,
            lease_expires_at=None,
            claim_token=None,
            claim_hash=None,
            last_operation_kind=PublicationOperationKind.PROMOTE,
            last_operation_token=operation_token,
            last_operation_hash="f" * 64,
            finished_at=self.now,
            updated_at=self.now,
        )
        if self.lose_promote_response_once:
            self.lose_promote_response_once = False
            raise RuntimeError("simulated lost promotion response")
        return self.record

    async def get(self, quality_run_id: str) -> PublicationRecord | None:
        self.get_calls += 1
        return self.record if quality_run_id == self.record.quality_run_id else None


class FakeObjects:
    def __init__(self, objects: dict[str, bytes]) -> None:
        self.objects = dict(objects)
        self.reads: list[tuple[str, int]] = []
        self.creates: list[tuple[str, str]] = []
        self.list_calls: list[tuple[str, int]] = []
        self.deletes: list[str] = []
        self.fail_reads: set[str] = set()
        self.oversized_paths: set[str] = set()
        self.fail_listing = False
        self.lose_create_response_once: set[str] = set()
        self.lose_delete_response_once: set[str] = set()
        self.fail_delete_once: set[str] = set()

    def read_bounded(self, object_path: str, *, max_bytes: int) -> bytes | None:
        self.reads.append((object_path, max_bytes))
        if object_path in self.fail_reads:
            raise RuntimeError("simulated read outage")
        value = self.objects.get(object_path)
        if value is None:
            return None
        if object_path in self.oversized_paths:
            raise ArtifactObjectSizeError("simulated oversize object")
        if value is not None and len(value) > max_bytes:
            raise ArtifactObjectSizeError("simulated oversize object")
        return bytes(value)

    def read(self, object_path: str) -> bytes | None:
        value = self.objects.get(object_path)
        return bytes(value) if value is not None else None

    def create_if_absent(
        self,
        object_path: str,
        content: bytes,
        *,
        content_type: str,
    ) -> str:
        self.creates.append((object_path, content_type))
        if object_path in self.objects:
            return "exists"
        self.objects[object_path] = bytes(content)
        if object_path in self.lose_create_response_once:
            self.lose_create_response_once.remove(object_path)
            raise RuntimeError("simulated lost create response")
        return "created"

    def list_flat_page(
        self,
        prefix: str,
        *,
        limit: int,
    ) -> list[str]:
        self.list_calls.append((prefix, limit))
        if self.fail_listing:
            raise RuntimeError("simulated listing outage")
        prefix_with_separator = f"{prefix.rstrip('/')}/"
        paths = sorted(
            path
            for path in self.objects
            if path.startswith(prefix_with_separator)
            and "/" not in path.removeprefix(prefix_with_separator)
        )
        return paths[:limit]

    def delete_if_present(self, object_path: str) -> str:
        self.deletes.append(object_path)
        if object_path in self.fail_delete_once:
            self.fail_delete_once.remove(object_path)
            raise RuntimeError("simulated delete failure before commit")
        outcome = "deleted" if self.objects.pop(object_path, None) is not None else "missing"
        self.oversized_paths.discard(object_path)
        if object_path in self.lose_delete_response_once:
            self.lose_delete_response_once.remove(object_path)
            raise RuntimeError("simulated lost delete response")
        return outcome


def _producer_bundle(
    pack: DeckQualitySourcePack,
    artifact_bytes: bytes,
) -> tuple[str, bytes, bytes]:
    source_pack_bytes = canonical_json_bytes(pack)
    bundle_bytes, descriptor = encode_deck_quality_producer_bundle(
        pack=pack,
        source_pack_bytes=source_pack_bytes,
    )
    return descriptor.object_path, bundle_bytes, source_pack_bytes


class FakeProducerReconciliationStore:
    """RPC-shaped fake that preserves commit-before-response-loss semantics."""

    def __init__(self, pack: DeckQualitySourcePack, *, now: datetime) -> None:
        self.pack = pack
        self.packs = {pack.quality_run_id: pack}
        self.records: dict[str, PublicationRecord] = {}
        self.now = now
        self.record: PublicationRecord | None = None
        self.request_calls: list[
            tuple[PublicationRequest, str, str]
        ] = []
        self.claim_calls: list[dict[str, object]] = []
        self.get_calls: list[str] = []
        self.lose_request_response_once = False
        self.fail_request_before_commit_remaining = 0
        self.probe_calls = 0
        self.close_calls = 0

    def add_pack(self, pack: DeckQualitySourcePack) -> None:
        self.packs[pack.quality_run_id] = pack

    async def probe(self) -> None:
        self.probe_calls += 1

    async def aclose(self) -> None:
        self.close_calls += 1

    async def get(self, quality_run_id: str) -> PublicationRecord | None:
        self.get_calls.append(quality_run_id)
        return self.records.get(quality_run_id)

    async def request_ready(
        self,
        request: PublicationRequest,
        *,
        source_pack_object_path: str,
        source_pack_hash: str,
    ) -> PublicationRecord:
        self.request_calls.append(
            (request, source_pack_object_path, source_pack_hash)
        )
        if self.fail_request_before_commit_remaining > 0:
            self.fail_request_before_commit_remaining -= 1
            raise RuntimeError("simulated request-ready failure before commit")
        expected_pack = self.packs[request.quality_run_id]
        assert request.instrument == expected_pack.instrument
        assert request.user_id == expected_pack.user_id
        assert request.thread_id == expected_pack.thread_id
        assert request.artifact_object_path == expected_pack.immutable_snapshot_object_path
        assert request.artifact_hash == expected_pack.artifact_sha256
        record = self.records.get(request.quality_run_id)
        if record is None:
            template = _pending_record(
                expected_pack,
                b"source hash replaced by atomic request",
                now=self.now,
            )
            record = _validated_update(
                template,
                source_pack_object_path=source_pack_object_path,
                source_pack_hash=source_pack_hash,
                artifact_object_path=request.artifact_object_path,
                artifact_hash=request.artifact_hash,
                deadline_at=request.deadline_at,
                quality_run_deadline_at=request.quality_run_deadline_at,
                next_attempt_at=self.now,
                requested_at=self.now,
                updated_at=self.now,
            )
            self.records[request.quality_run_id] = record
            self.record = record
        if self.lose_request_response_once:
            self.lose_request_response_once = False
            raise RuntimeError("simulated lost request-ready response")
        return record

    async def claim(
        self,
        *,
        lease_owner: str,
        claim_token: str,
        lease_seconds: int = 60,
        limit: int = 1,
    ) -> tuple[PublicationRecord, ...]:
        self.claim_calls.append(
            {
                "lease_owner": lease_owner,
                "claim_token": claim_token,
                "lease_seconds": lease_seconds,
                "limit": limit,
            }
        )
        return ()


def _producer_setup(
    tmp_path: Path,
    *,
    pack: DeckQualitySourcePack | None = None,
    artifact_bytes: bytes | None = None,
    claim_token: str = "producer-claim-token",
) -> tuple[
    DeckQualityPublicationWorker,
    FakeProducerReconciliationStore,
    FakeObjects,
    DeckQualitySourcePack,
    bytes,
    str,
    bytes,
]:
    now = datetime.now(UTC).replace(microsecond=0)
    if artifact_bytes is None:
        artifact_bytes = _pptx_bytes()
    pack = pack or _pack(artifact_bytes)
    object_path, bundle_bytes, source_pack_bytes = _producer_bundle(
        pack,
        artifact_bytes,
    )
    objects = FakeObjects(
        {
            object_path: bundle_bytes,
            expected_publication_source_pack_path(
                user_id=pack.user_id,
                thread_id=pack.thread_id,
                build_id=pack.build_id,
                quality_run_id=pack.quality_run_id,
            ): source_pack_bytes,
            pack.artifact_storage_object_path: artifact_bytes,
        }
    )
    store = FakeProducerReconciliationStore(pack, now=now)
    worker = DeckQualityPublicationWorker(
        config=_config(),
        instrument=_instrument(),
        store=store,  # type: ignore[arg-type]
        object_store=objects,
        lease_owner="producer-worker-1",
        materialization_root=tmp_path / "producer-materialized",
        claim_token_factory=lambda: claim_token,
        clock=lambda: store.now,
    )
    return (
        worker,
        store,
        objects,
        pack,
        artifact_bytes,
        object_path,
        source_pack_bytes,
    )


def _setup(
    tmp_path: Path,
    *,
    pack_task_id: str = "task-01",
    row_task_id: str = "task-01",
    config: DeckQualityConfig | None = None,
) -> tuple[
    DeckQualityPublicationWorker,
    FakePublicationStore,
    FakeObjects,
    DeckQualitySourcePack,
]:
    now = datetime.now(UTC).replace(microsecond=0)
    artifact = _pptx_bytes()
    pack = _pack(artifact, task_id=pack_task_id)
    encoded = canonical_json_bytes(pack)
    record = _pending_record(pack, encoded, now=now, row_task_id=row_task_id)
    assert record.source_pack_object_path is not None
    objects = FakeObjects(
        {
            record.source_pack_object_path: encoded,
            record.artifact_object_path: artifact,
        }
    )
    store = FakePublicationStore(record, now=now)
    tokens = iter(("claim-token-1", "claim-token-2", "claim-token-3"))
    worker = DeckQualityPublicationWorker(
        config=config or _config(),
        instrument=_instrument(),
        store=store,
        object_store=objects,
        lease_owner="publication-worker-1",
        materialization_root=tmp_path / "materialized",
        claim_token_factory=lambda: next(tokens),
        clock=lambda: store.now,
    )
    return worker, store, objects, pack


@pytest.mark.anyio
async def test_worker_materializes_v2_inputs_and_atomically_promotes(
    tmp_path: Path,
) -> None:
    worker, store, objects, pack = _setup(tmp_path)
    source_path = store.record.source_pack_object_path
    assert source_path is not None

    result = await worker.run_once()

    assert result.claimed == 1
    assert result.published == 1
    assert store.record.state is PublicationState.PUBLISHED
    assert store.claim_calls[0] == {
        "lease_owner": "publication-worker-1",
        "claim_token": "claim-token-1",
        "lease_seconds": 120,
        "limit": 2,
    }
    assert objects.reads[:2] == [
        (source_path, 8 * 1024 * 1024),
        (pack.artifact_storage_object_path, 32 * 1024 * 1024),
    ]
    assert len(store.promote_calls) == 1
    promoted = store.promote_calls[0]
    assert str(promoted["operation_token"]).startswith("dq1-pub-promote:")
    manifest_path = str(promoted["input_manifest_object_path"])
    manifest_hash = str(promoted["input_manifest_hash"])
    assert manifest_path == store.record.expected_input_manifest_object_path
    assert hashlib.sha256(objects.objects[manifest_path]).hexdigest() == manifest_hash
    manifest = json.loads(objects.objects[manifest_path])
    assert manifest["schema_version"] == "deck-quality-pre-render-input-manifest/v2"
    assert [item["role"] for item in manifest["objects"]] == [
        "accepted_artifact",
        "creative_plan",
        "design_plan",
        "build_record",
        "blind_brief",
        "mechanical_record",
    ]
    assert all(item["media_type"] != "application/pdf" for item in manifest["objects"])
    materialized = tmp_path / "materialized"
    assert materialized.is_dir() and list(materialized.iterdir()) == []


@pytest.mark.anyio
async def test_worker_rejects_source_identity_before_artifact_read(
    tmp_path: Path,
) -> None:
    worker, store, objects, _pack_value = _setup(
        tmp_path,
        pack_task_id="different-task",
        row_task_id="task-01",
    )
    source_path = store.record.source_pack_object_path

    result = await worker.run_once()

    assert result.failed == 1
    assert store.fail_calls[0]["error_code"] is PublicationErrorCode.ARTIFACT_SNAPSHOT_STALE
    assert store.fail_calls[0]["error_stage"] == "source_identity"
    assert objects.reads == [(source_path, 8 * 1024 * 1024)]
    assert store.promote_calls == []


@pytest.mark.anyio
async def test_worker_rejects_noncanary_row_before_any_object_read(
    tmp_path: Path,
) -> None:
    worker, store, objects, _pack_value = _setup(
        tmp_path,
        config=_config(users=frozenset({"another-canary"})),
    )

    result = await worker.run_once()

    assert result.failed == 1
    assert store.fail_calls[0]["error_stage"] == "row_identity"
    assert objects.reads == []


@pytest.mark.anyio
async def test_worker_rejects_noncurrent_instrument_before_any_object_read(
    tmp_path: Path,
) -> None:
    _worker, store, objects, _pack_value = _setup(tmp_path)
    current_instrument = _instrument(rubric_hash="9" * 64)
    worker = DeckQualityPublicationWorker(
        config=_config(),
        instrument=current_instrument,
        store=store,
        object_store=objects,
        lease_owner="publication-worker-1",
        materialization_root=tmp_path / "instrument-mismatch",
        claim_token_factory=lambda: "instrument-claim-token",
        clock=lambda: store.now,
    )

    result = await worker.run_once()

    assert result.failed == 1
    assert store.fail_calls[0]["error_stage"] == "row_identity"
    assert objects.reads == []


@pytest.mark.anyio
async def test_worker_rejects_source_hash_before_artifact_reference_read(
    tmp_path: Path,
) -> None:
    worker, store, objects, _pack_value = _setup(tmp_path)
    source_path = store.record.source_pack_object_path
    assert source_path is not None
    objects.objects[source_path] += b"\n"

    result = await worker.run_once()

    assert result.failed == 1
    assert store.fail_calls[0]["error_code"] is PublicationErrorCode.ARTIFACT_SNAPSHOT_STALE
    assert store.fail_calls[0]["error_stage"] == "source_hash"
    assert objects.reads == [(source_path, 8 * 1024 * 1024)]


@pytest.mark.anyio
async def test_worker_retries_transient_source_read_without_following_artifact(
    tmp_path: Path,
) -> None:
    worker, store, objects, _pack_value = _setup(tmp_path)
    source_path = store.record.source_pack_object_path
    assert source_path is not None
    objects.fail_reads.add(source_path)

    result = await worker.run_once()

    assert result.retry_scheduled == 1
    assert store.retry_calls[0]["error_code"] is PublicationErrorCode.PERSISTENCE_ERROR
    assert store.retry_calls[0]["error_stage"] == "source_read"
    assert objects.reads == [(source_path, 8 * 1024 * 1024)]
    assert store.promote_calls == []


@pytest.mark.anyio
async def test_worker_fails_deterministic_artifact_hash_mismatch(
    tmp_path: Path,
) -> None:
    worker, store, objects, pack = _setup(tmp_path)
    objects.objects[pack.artifact_storage_object_path] = b"different artifact"

    result = await worker.run_once()

    assert result.failed == 1
    assert store.fail_calls[0]["error_code"] is PublicationErrorCode.ARTIFACT_VERIFICATION_FAILED
    assert store.fail_calls[0]["error_stage"] == "artifact_verify"
    assert len(objects.reads) == 2
    assert store.promote_calls == []


@pytest.mark.anyio
async def test_claim_and_promotion_response_loss_reconcile_without_duplicate_attempt(
    tmp_path: Path,
) -> None:
    worker, store, _objects, _pack_value = _setup(tmp_path)
    store.lose_claim_response_once = True
    store.lose_promote_response_once = True

    result = await worker.run_once()

    assert result.published == 1
    assert len(store.claim_calls) == 2
    assert store.claim_calls[0] == store.claim_calls[1]
    assert store.record.attempt_count == 1
    assert len(store.promote_calls) == 1
    assert store.get_calls == 1


@pytest.mark.anyio
async def test_worker_lifecycle_and_configuration_helpers(tmp_path: Path) -> None:
    worker, store, objects, _pack_value = _setup(tmp_path)
    app = SimpleNamespace(state=SimpleNamespace())

    install_deck_quality_publication_worker(app, worker)
    assert get_deck_quality_publication_worker(app) is worker
    assert get_deck_quality_publication_worker_or_none(app) is worker
    await worker.probe()
    assert store.probe_calls == 1

    # Make the row terminal so the background lifecycle only polls an empty
    # bounded claim and cannot perform another materialization.
    store.record = _validated_update(
        store.record,
        state=PublicationState.FAILED,
        last_error_code=PublicationErrorCode.ARTIFACT_SNAPSHOT_STALE,
        last_error_stage="test_terminal",
        last_error_at=store.now,
        finished_at=store.now,
    )
    await start_deck_quality_publication_worker(worker)
    assert worker.running
    await stop_deck_quality_publication_worker(worker)
    assert not worker.running
    assert store.probe_calls == 2
    assert store.close_calls == 1

    disabled = DeckQualityConfig()
    assert build_configured_deck_quality_publication_worker(
        config=disabled,
        instrument=_instrument(),
        store=store,  # type: ignore[arg-type]
        object_store=objects,  # type: ignore[arg-type]
    ) is None
    with pytest.raises(ValueError, match="instrument"):
        DeckQualityPublicationWorker(
            config=_config(),
            instrument=_instrument(evidence_preprocessor_version="different-preprocessor"),
            store=store,
            object_store=objects,
        )
    with pytest.raises(ValueError, match="claim limit"):
        DeckQualityPublicationWorker(
            config=_config(),
            instrument=_instrument(),
            store=store,
            object_store=objects,
            claim_limit=3,
        )


@pytest.mark.anyio
async def test_worker_reconciles_durable_producer_bundle_before_claim(
    tmp_path: Path,
) -> None:
    (
        worker,
        store,
        objects,
        pack,
        artifact_bytes,
        object_path,
        source_pack_bytes,
    ) = _producer_setup(tmp_path)

    result = await worker.run_once()

    assert result.producer_seen == 1
    assert result.producer_reconciled == 1
    assert result.producer_quarantined == 0
    assert result.producer_failed == 0
    assert result.claimed == 0
    assert objects.list_calls == [
        ("dq1/producer-inbox/v1", 32),
        ("dq1/producer-failures/v1", 32),
        ("dq1/producer-rejections/v1", 32),
    ]
    assert len(store.request_calls) == 1
    request, source_path, source_hash = store.request_calls[0]
    assert request.quality_run_id == pack.quality_run_id
    assert request.artifact_object_path == pack.immutable_snapshot_object_path
    assert source_path == expected_publication_source_pack_path(
        user_id=CANARY_USER,
        thread_id=THREAD_ID,
        build_id=BUILD_ID,
        quality_run_id=pack.quality_run_id,
    )
    assert source_hash == hashlib.sha256(source_pack_bytes).hexdigest()
    assert objects.objects[pack.immutable_snapshot_object_path] == artifact_bytes
    assert objects.objects[source_path] == source_pack_bytes
    archive_path = deck_quality_producer_archive_path(pack.quality_run_id)
    assert objects.objects[archive_path]
    assert object_path not in objects.objects
    assert objects.deletes == [object_path]
    assert objects.reads == [
        (object_path, 64 * 1024),
        (archive_path, 64 * 1024),
        (source_path, 8 * 1024 * 1024),
        (pack.immutable_snapshot_object_path, 32 * 1024 * 1024),
        (archive_path, 64 * 1024),
        (object_path, 64 * 1024),
    ]


@pytest.mark.anyio
async def test_producer_reconciliation_fences_create_and_rpc_response_loss_and_restart(
    tmp_path: Path,
) -> None:
    (
        worker,
        store,
        objects,
        pack,
        _artifact_bytes,
        _object_path,
        _source_pack_bytes,
    ) = _producer_setup(tmp_path)
    objects.lose_create_response_once.add(
        deck_quality_producer_archive_path(pack.quality_run_id)
    )
    objects.lose_create_response_once.add(pack.immutable_snapshot_object_path)
    objects.lose_delete_response_once.add(_object_path)
    store.lose_request_response_once = True

    first = await worker.run_once()

    assert first.producer_reconciled == 1
    assert first.producer_failed == 0
    assert len(store.request_calls) == 2
    assert store.request_calls[0] == store.request_calls[1]
    assert store.record is not None
    assert store.record.state is PublicationState.PENDING
    creates_after_first = list(objects.creates)

    restarted = DeckQualityPublicationWorker(
        config=_config(),
        instrument=_instrument(),
        store=store,  # type: ignore[arg-type]
        object_store=objects,
        lease_owner="producer-worker-after-restart",
        materialization_root=tmp_path / "producer-after-restart",
        claim_token_factory=lambda: "producer-restart-claim-token",
        clock=lambda: store.now,
    )
    replay = await restarted.run_once()

    assert replay.producer_seen == 0
    assert replay.producer_reconciled == 0
    assert replay.producer_failed == 0
    assert len(store.request_calls) == 2
    assert objects.creates == creates_after_first


@pytest.mark.anyio
async def test_request_ready_double_failure_recovers_exact_committed_row(
    tmp_path: Path,
) -> None:
    worker, store, objects, pack, _artifact, inbox_path, source_pack_bytes = (
        _producer_setup(tmp_path)
    )
    source_path = expected_publication_source_pack_path(
        user_id=pack.user_id,
        thread_id=pack.thread_id,
        build_id=pack.build_id,
        quality_run_id=pack.quality_run_id,
    )
    exact = _validated_update(
        _pending_record(pack, source_pack_bytes, now=store.now),
        source_pack_object_path=source_path,
        source_pack_hash=hashlib.sha256(source_pack_bytes).hexdigest(),
    )
    store.records[pack.quality_run_id] = exact
    store.fail_request_before_commit_remaining = 2

    result = await worker.run_once()

    assert result.producer_reconciled == 1
    assert result.producer_failed == 0
    assert store.get_calls == [pack.quality_run_id]
    assert inbox_path not in objects.objects


@pytest.mark.anyio
async def test_request_ready_identity_conflict_is_quarantined_without_starvation(
    tmp_path: Path,
) -> None:
    worker, store, objects, pack, _artifact, inbox_path, source_pack_bytes = (
        _producer_setup(tmp_path)
    )
    source_path = expected_publication_source_pack_path(
        user_id=pack.user_id,
        thread_id=pack.thread_id,
        build_id=pack.build_id,
        quality_run_id=pack.quality_run_id,
    )
    mismatched = _validated_update(
        _pending_record(pack, source_pack_bytes, now=store.now),
        task_id="conflicting-task",
        source_pack_object_path=source_path,
        source_pack_hash=hashlib.sha256(source_pack_bytes).hexdigest(),
    )
    store.records[pack.quality_run_id] = mismatched
    store.fail_request_before_commit_remaining = 2

    result = await worker.run_once()

    assert result.producer_reconciled == 0
    assert result.producer_quarantined == 1
    assert result.producer_failed == 0
    assert result.producer_failure_evidence == 1
    assert store.get_calls == [pack.quality_run_id]
    assert inbox_path not in objects.objects


@pytest.mark.anyio
async def test_archive_before_db_failure_leaves_restart_safe_inbox(
    tmp_path: Path,
) -> None:
    worker, store, objects, pack, _artifact, inbox_path, _source = (
        _producer_setup(tmp_path)
    )
    store.fail_request_before_commit_remaining = 2

    first = await worker.run_once()

    archive_path = deck_quality_producer_archive_path(pack.quality_run_id)
    assert first.producer_reconciled == 0
    assert first.producer_failed == 1
    assert inbox_path in objects.objects
    assert archive_path not in objects.objects
    assert store.records == {}

    restarted = DeckQualityPublicationWorker(
        config=_config(),
        instrument=_instrument(),
        store=store,  # type: ignore[arg-type]
        object_store=objects,
        lease_owner="archive-db-restart",
        materialization_root=tmp_path / "archive-db-restart",
        claim_token_factory=lambda: "archive-db-restart-token",
        clock=lambda: store.now,
    )
    replay = await restarted.run_once()

    assert replay.producer_reconciled == 1
    assert replay.producer_failed == 0
    assert inbox_path not in objects.objects
    assert len(store.records) == 1
    assert len(store.request_calls) == 3


@pytest.mark.anyio
async def test_db_commit_before_delete_failure_replays_without_duplicate_row(
    tmp_path: Path,
) -> None:
    worker, store, objects, pack, _artifact, inbox_path, _source = (
        _producer_setup(tmp_path)
    )
    objects.fail_delete_once.add(inbox_path)

    first = await worker.run_once()

    assert first.producer_reconciled == 0
    assert first.producer_failed == 1
    assert inbox_path in objects.objects
    assert deck_quality_producer_archive_path(pack.quality_run_id) in objects.objects
    assert len(store.records) == 1
    assert len(store.request_calls) == 1

    restarted = DeckQualityPublicationWorker(
        config=_config(),
        instrument=_instrument(),
        store=store,  # type: ignore[arg-type]
        object_store=objects,
        lease_owner="db-delete-restart",
        materialization_root=tmp_path / "db-delete-restart",
        claim_token_factory=lambda: "db-delete-restart-token",
        clock=lambda: store.now,
    )
    replay = await restarted.run_once()

    assert replay.producer_reconciled == 1
    assert replay.producer_failed == 0
    assert inbox_path not in objects.objects
    assert len(store.records) == 1
    assert len(store.request_calls) == 2


@pytest.mark.anyio
async def test_worker_rejects_tampered_producer_bundle_without_materialization(
    tmp_path: Path,
) -> None:
    worker, store, objects, _pack_value, *_rest = _producer_setup(tmp_path)
    producer_path = next(iter(objects.objects))
    bundle = objects.objects[producer_path]
    objects.objects[producer_path] = bundle[:-1] + bytes((bundle[-1] ^ 1,))

    result = await worker.run_once()

    assert result.producer_seen == 1
    assert result.producer_reconciled == 0
    assert result.producer_quarantined == 1
    assert result.producer_failed == 0
    assert result.producer_failure_evidence == 1
    assert store.request_calls == []
    assert len(objects.creates) == 2
    assert objects.creates[0][0].startswith("dq1/producer-quarantine/v1/")
    assert objects.creates[0][1] == "application/octet-stream"
    assert objects.creates[1][0].startswith("dq1/producer-rejections/v1/")
    assert objects.creates[1][1] == "application/json"
    assert producer_path not in objects.objects


@pytest.mark.anyio
async def test_distinct_poison_replay_at_same_inbox_path_gets_distinct_quarantine(
    tmp_path: Path,
) -> None:
    worker, store, objects, _pack_value, _artifact, inbox_path, _source = (
        _producer_setup(tmp_path)
    )
    objects.objects[inbox_path] = b"first malformed bundle"

    first = await worker.run_once()

    assert first.producer_quarantined == 1
    first_quarantine = {
        path
        for path in objects.objects
        if path.startswith("dq1/producer-quarantine/v1/")
    }
    assert len(first_quarantine) == 1

    objects.objects[inbox_path] = b"second distinct malformed bundle"
    restarted = DeckQualityPublicationWorker(
        config=_config(),
        instrument=_instrument(),
        store=store,  # type: ignore[arg-type]
        object_store=objects,
        lease_owner="distinct-poison-restart",
        materialization_root=tmp_path / "distinct-poison-restart",
        claim_token_factory=lambda: "distinct-poison-restart-token",
        clock=lambda: store.now,
    )
    second = await restarted.run_once()

    assert second.producer_quarantined == 1
    all_quarantine = {
        path
        for path in objects.objects
        if path.startswith("dq1/producer-quarantine/v1/")
    }
    assert len(all_quarantine) == 2
    assert first_quarantine < all_quarantine
    assert inbox_path not in objects.objects


@pytest.mark.anyio
async def test_worker_rejects_noncanary_producer_bundle_before_materialization(
    tmp_path: Path,
) -> None:
    artifact_bytes = _pptx_bytes()
    pack = _pack(artifact_bytes, user_id="ordinary-user")
    worker, store, objects, *_rest = _producer_setup(
        tmp_path,
        pack=pack,
        artifact_bytes=artifact_bytes,
    )

    result = await worker.run_once()

    assert result.producer_seen == 1
    assert result.producer_reconciled == 0
    assert result.producer_quarantined == 1
    assert result.producer_failed == 0
    assert result.producer_failure_evidence == 1
    assert store.request_calls == []
    assert len(objects.creates) == 2
    assert objects.creates[0][0].startswith("dq1/producer-quarantine/v1/")
    assert objects.creates[0][1] == "application/octet-stream"
    assert objects.creates[1][0].startswith("dq1/producer-rejections/v1/")
    assert objects.creates[1][1] == "application/json"


@pytest.mark.anyio
async def test_large_artifact_conflict_records_only_content_free_evidence(
    tmp_path: Path,
) -> None:
    worker, store, objects, pack, _artifact, inbox_path, _source = (
        _producer_setup(tmp_path)
    )
    inbox_content = objects.objects[inbox_path]
    conflicting_artifact = b"x" * (2 * 1024 * 1024)
    objects.objects[pack.artifact_storage_object_path] = conflicting_artifact

    result = await worker.run_once()

    assert result.producer_seen == 1
    assert result.producer_reconciled == 0
    assert result.producer_quarantined == 1
    assert result.producer_failed == 0
    assert result.producer_failure_evidence == 1
    assert store.request_calls == []
    assert inbox_path not in objects.objects
    assert sum(content == conflicting_artifact for content in objects.objects.values()) == 1
    rejection_paths = [
        path
        for path in objects.objects
        if path.startswith("dq1/producer-rejections/v1/")
    ]
    assert len(rejection_paths) == 1
    evidence = json.loads(objects.objects[rejection_paths[0]])
    assert evidence == {
        "schema_version": "deck-quality-producer-rejection-evidence/v1",
        "reason": "storage_conflict",
        "inbox_path_sha256": hashlib.sha256(inbox_path.encode("utf-8")).hexdigest(),
        "observed_content_sha256": hashlib.sha256(inbox_content).hexdigest(),
        "observed_size_bytes": len(inbox_content),
        "conflicting_path_sha256": hashlib.sha256(
            pack.artifact_storage_object_path.encode("utf-8")
        ).hexdigest(),
        "conflicting_content_sha256": hashlib.sha256(conflicting_artifact).hexdigest(),
        "conflicting_size_bytes": len(conflicting_artifact),
    }


@pytest.mark.anyio
async def test_quarantined_poison_bundles_do_not_starve_later_valid_bundle_or_restart(
    tmp_path: Path,
) -> None:
    worker, store, objects, pack, *_rest = _producer_setup(tmp_path)
    poison_paths = (
        "dq1/producer-inbox/v1/000-invalid.bin",
        "dq1/producer-inbox/v1/001-invalid.bin",
    )
    for path in poison_paths:
        objects.objects[path] = b"permanently malformed producer object"

    first = await worker.run_once()

    assert first.producer_seen == 3
    assert first.producer_reconciled == 1
    assert first.producer_quarantined == 2
    assert first.producer_failed == 0
    assert len(store.request_calls) == 1
    assert store.request_calls[0][0].quality_run_id == pack.quality_run_id
    error_markers = sorted(
        path
        for path in objects.objects
        if path.startswith("dq1/producer-quarantine/v1/")
    )
    assert len(error_markers) == 2
    creates_after_first = list(objects.creates)

    restarted = DeckQualityPublicationWorker(
        config=_config(),
        instrument=_instrument(),
        store=store,  # type: ignore[arg-type]
        object_store=objects,
        lease_owner="producer-poison-restart",
        materialization_root=tmp_path / "poison-restart",
        claim_token_factory=lambda: "poison-restart-claim-token",
        clock=lambda: store.now,
    )
    replay = await restarted.run_once()

    assert replay.producer_seen == 0
    assert replay.producer_reconciled == 0
    assert replay.producer_quarantined == 0
    assert replay.producer_failed == 0
    assert replay.producer_failure_evidence == 2
    assert len(store.request_calls) == 1
    assert objects.creates == creates_after_first


@pytest.mark.anyio
async def test_full_page_of_oversized_poison_retires_before_later_valid_bundle(
    tmp_path: Path,
) -> None:
    worker, store, objects, pack, _artifact, valid_path, _source = (
        _producer_setup(tmp_path)
    )
    oversized_paths = {
        f"dq1/producer-inbox/v1/{index:03d}-oversized.bin"
        for index in range(32)
    }
    for path in oversized_paths:
        objects.objects[path] = b"bounded test sentinel"
    objects.oversized_paths.update(oversized_paths)

    first = await worker.run_once()

    assert first.producer_seen == 32
    assert first.producer_quarantined == 32
    assert first.producer_failed == 0
    assert valid_path in objects.objects
    assert oversized_paths.isdisjoint(objects.objects)
    assert len(
        [
            path
            for path in objects.objects
            if path.startswith("dq1/producer-quarantine/v1/oversized/")
        ]
    ) == 32

    restarted = DeckQualityPublicationWorker(
        config=_config(),
        instrument=_instrument(),
        store=store,  # type: ignore[arg-type]
        object_store=objects,
        lease_owner="oversize-page-restart",
        materialization_root=tmp_path / "oversize-page-restart",
        claim_token_factory=lambda: "oversize-page-restart-token",
        clock=lambda: store.now,
    )
    second = await restarted.run_once()

    assert second.producer_seen == 1
    assert second.producer_reconciled == 1
    assert second.producer_failed == 0
    assert valid_path not in objects.objects
    assert deck_quality_producer_archive_path(pack.quality_run_id) in objects.objects


@pytest.mark.anyio
async def test_more_than_one_page_of_sequential_producers_never_scans_archive_history(
    tmp_path: Path,
) -> None:
    _initial_worker, store, objects, _initial_pack, *_rest = _producer_setup(
        tmp_path
    )
    objects.objects.clear()
    artifact_bytes = _pptx_bytes()
    tokens = iter(f"sequential-claim-{index}" for index in range(40))
    worker = DeckQualityPublicationWorker(
        config=_config(),
        instrument=_instrument(),
        store=store,  # type: ignore[arg-type]
        object_store=objects,
        lease_owner="sequential-producer-worker",
        materialization_root=tmp_path / "sequential-producers",
        claim_token_factory=lambda: next(tokens),
        clock=lambda: store.now,
    )

    quality_run_ids: list[str] = []
    for index in range(40):
        pack = _pack(
            artifact_bytes,
            artifact_version_id=f"artifact-version-sequential-{index}",
        )
        store.add_pack(pack)
        inbox_path, bundle_bytes, _source_pack = _producer_bundle(
            pack,
            artifact_bytes,
        )
        objects.objects[inbox_path] = bundle_bytes
        objects.objects[
            expected_publication_source_pack_path(
                user_id=pack.user_id,
                thread_id=pack.thread_id,
                build_id=pack.build_id,
                quality_run_id=pack.quality_run_id,
            )
        ] = _source_pack
        objects.objects[pack.artifact_storage_object_path] = artifact_bytes

        result = await worker.run_once()

        assert result.producer_seen == 1
        assert result.producer_reconciled == 1
        assert result.producer_failed == 0
        assert inbox_path not in objects.objects
        assert objects.objects[
            deck_quality_producer_archive_path(pack.quality_run_id)
        ] == bundle_bytes
        quality_run_ids.append(pack.quality_run_id)

    assert len(store.request_calls) == 40
    assert objects.list_calls == [
        call
        for _ in range(40)
        for call in (
            ("dq1/producer-inbox/v1", 32),
            ("dq1/producer-failures/v1", 32),
            ("dq1/producer-rejections/v1", 32),
        )
    ]
    assert len(
        [
            path
            for path in objects.objects
            if path.startswith("dq1/producer-archive/v1/")
        ]
    ) == 40
    assert len(set(quality_run_ids)) == 40


@pytest.mark.anyio
async def test_worker_listing_outage_fails_cycle_before_claim(tmp_path: Path) -> None:
    worker, store, objects, *_rest = _producer_setup(tmp_path)
    objects.fail_listing = True

    with pytest.raises(RuntimeError, match="listing outage"):
        await worker.run_once()

    assert store.claim_calls == []
    assert store.request_calls == []
    assert objects.reads == []


@pytest.mark.anyio
async def test_worker_probe_requires_flat_page_listing_capability(tmp_path: Path) -> None:
    worker, store, objects, *_rest = _producer_setup(tmp_path)
    setattr(objects, "list_flat_page", None)

    with pytest.raises(RuntimeError, match="not flat-page-listable"):
        await worker.probe()

    assert store.probe_calls == 1


@pytest.mark.anyio
async def test_worker_live_readiness_degrades_while_retryable_inbox_failure_remains(
    tmp_path: Path,
) -> None:
    worker, _store, objects, *_rest = _producer_setup(tmp_path)
    producer_path = next(iter(objects.objects))
    objects.fail_reads.add(producer_path)

    worker.start()
    try:
        for _attempt in range(100):
            readiness = worker.readiness()
            if readiness.get("reason") == "cycle_failed":
                break
            await asyncio.sleep(0.01)
        else:
            pytest.fail("worker did not expose the failed cycle")

        assert readiness == {
            "status": "degraded",
            "reason": "cycle_failed",
            "error_type": "RuntimeError",
        }
    finally:
        await worker.stop()


@pytest.mark.anyio
async def test_worker_live_readiness_stays_degraded_for_persisted_failure_evidence(
    tmp_path: Path,
) -> None:
    worker, _store, objects, _pack, _artifact, inbox_path, _source = (
        _producer_setup(tmp_path)
    )
    objects.objects.pop(inbox_path)
    objects.objects["dq1/producer-failures/v1/candidate.json"] = b"{}"

    worker.start()
    try:
        for _attempt in range(100):
            readiness = worker.readiness()
            if readiness.get("reason") == "cycle_failed":
                break
            await asyncio.sleep(0.01)
        else:
            pytest.fail("worker did not expose persistent producer failure evidence")

        assert readiness == {
            "status": "degraded",
            "reason": "cycle_failed",
            "error_type": "RuntimeError",
        }
    finally:
        await worker.stop()


@pytest.mark.anyio
async def test_worker_live_readiness_degrades_when_heartbeat_is_stale(
    tmp_path: Path,
) -> None:
    worker, store, _objects, _pack_value = _setup(tmp_path)

    worker.start()
    try:
        for _attempt in range(200):
            readiness = worker.readiness()
            if readiness.get("status") == "ready":
                break
            await asyncio.sleep(0.01)
        else:
            pytest.fail("worker did not publish its first successful heartbeat")

        store.now += timedelta(seconds=31)
        assert worker.readiness() == {
            "status": "degraded",
            "reason": "heartbeat_stale",
        }
    finally:
        await worker.stop()


@pytest.mark.anyio
async def test_slow_sync_inbox_call_is_single_flight_and_readiness_goes_stale(
    tmp_path: Path,
) -> None:
    worker, store, objects, *_rest = _producer_setup(tmp_path)
    entered = threading.Event()
    release = threading.Event()
    slow_calls = 0
    original_list = objects.list_flat_page

    def slow_list(prefix: str, *, limit: int) -> list[str]:
        nonlocal slow_calls
        slow_calls += 1
        entered.set()
        release.wait(timeout=5)
        return original_list(prefix, limit=limit)

    setattr(objects, "list_flat_page", slow_list)
    worker.start()
    try:
        assert await asyncio.to_thread(entered.wait, 1)
        store.now += timedelta(seconds=31)
        await asyncio.sleep(0.05)
        assert slow_calls == 1
        assert worker.readiness() == {
            "status": "degraded",
            "reason": "heartbeat_stale",
        }
    finally:
        release.set()
        await worker.stop()
