from __future__ import annotations

import asyncio
import hashlib
import io
import json
import threading
import zipfile
from decimal import Decimal
from pathlib import Path
from types import SimpleNamespace
from typing import Any

import anyio
import pytest
from langsmith.run_helpers import get_tracing_context
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from deerflow.sandbox.tools import replace_virtual_path
from deerflow.sophia.build_manifest import (
    DECK_STYLE_ROOT_SELECTOR,
    BuildComponent,
    BuildManifest,
)
from deerflow.sophia.deck_build.html_sanitizer import assemble_compact_slide_html
from deerflow.sophia.deck_build.models import DeckBuildResult
from deerflow.sophia.deck_build.tracing import NATIVE_DECK_COMPILE_MODE
from deerflow.sophia.deck_design_lift import candidate_compiler as candidate_module
from deerflow.sophia.deck_design_lift import materializer as materializer_module
from deerflow.sophia.deck_design_lift.candidate_compiler import (
    BaselineDeckRender,
    BaselineVisualAsset,
    DeckCandidateBaseline,
    DeckCandidateCompilationError,
    DurableCandidateDq1Publisher,
    ProductionDeckCandidateCompiler,
    baseline_from_authenticated_snapshot,
)
from deerflow.sophia.deck_design_lift.materializer import (
    DeckCandidateCompileRequest,
    DeckCandidateSource,
)
from deerflow.sophia.deck_design_lift.schemas import (
    DeckRepairProgram,
    RepairRenderEvidence,
    SelectorRepair,
    SkillRef,
)
from deerflow.sophia.deck_quality.canonical import canonical_json_bytes, canonical_sha256
from deerflow.sophia.deck_quality.idempotency import derive_quality_run_id
from deerflow.sophia.deck_quality.instrument import DeckQualityRuntimeInstrument
from deerflow.sophia.deck_quality.publisher import (
    DeckQualityProducerBundleReceipt,
    DeckQualitySourcePack,
    PreparedDeckQualityPublication,
    capture_deck_quality_source_pack,
    deck_quality_immutable_artifact_snapshot_path,
    deck_quality_producer_archive_path,
    deck_quality_producer_bundle_path,
)
from deerflow.sophia.deck_quality.schemas import (
    AdjudicationPolicy,
    QualityInstrumentLock,
    VisibleTextSlide,
)

BUILD_ID = "build_psi_candidate_01"
USER_ID = "user_canary_01"
THREAD_ID = "thread_canary_01"
INITIAL_ARTIFACT_ID = "artifact_initial_01"
CANDIDATE_ARTIFACT_ID = "artifact_candidate_01"
INITIAL_QUALITY_RUN_ID = "quality_initial_01"
TRANSACTION_ID = "transaction_candidate_01"
OPERATION_ID = "operation_candidate_01"
FIXED_TIME = "2026-07-20T00:00:00+00:00"


def _sha(content: bytes) -> str:
    return hashlib.sha256(content).hexdigest()


def _png(color: str) -> bytes:
    output = io.BytesIO()
    Image.new("RGB", (160, 90), color).save(output, format="PNG")
    return output.getvalue()


def _jpg(color: str) -> bytes:
    output = io.BytesIO()
    Image.new("RGB", (160, 90), color).save(output, format="JPEG", quality=95)
    return output.getvalue()


def _pptx() -> bytes:
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def _retimestamp_pptx(content: bytes, timestamp: tuple[int, int, int, int, int, int]) -> bytes:
    output = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(content), mode="r") as source:
        with zipfile.ZipFile(output, mode="w", compression=zipfile.ZIP_DEFLATED) as destination:
            for original in source.infolist():
                info = zipfile.ZipInfo(original.filename, date_time=timestamp)
                info.compress_type = original.compress_type
                info.create_system = original.create_system
                info.external_attr = original.external_attr
                destination.writestr(info, source.read(original))
    return output.getvalue()


def _instrument() -> DeckQualityRuntimeInstrument:
    policy = AdjudicationPolicy(
        critical_score_floor=4,
        min_weighted_score=Decimal("3.5"),
    )
    lock = QualityInstrumentLock(
        rubric_version="deck-rubric-v2",
        rubric_hash="a" * 64,
        prompt_hashes={
            "blind_visual": "b" * 64,
            "plan_realization": "c" * 64,
        },
        judge_plan_hash="d" * 64,
        judge_profile_version="deck-visual-judge-v2",
        evidence_preprocessor_version="deck-evidence-v4",
        judge_invoker_version="deck-judge-invoker-v4",
        assessment_schema_versions={
            "blind_visual": "deck-quality-blind-assessment/v4",
            "mechanical": "deck-quality-mechanical-projection/v1",
            "plan_realization": "deck-quality-plan-assessment/v4",
        },
        adjudication_policy_hash=canonical_sha256(policy),
    )
    return DeckQualityRuntimeInstrument.model_construct(lock=lock)


def _program(instrument: DeckQualityRuntimeInstrument) -> DeckRepairProgram:
    skill = SkillRef(
        path="skills/public/hands-on-deck/designing-slides.md",
        source_hash="e" * 64,
        excerpt_hash="f" * 64,
    )
    repair = SelectorRepair(
        selector="slide:1",
        failure_codes=("weak_visual_hierarchy",),
        render_evidence=(
            RepairRenderEvidence(
                selector="slide:1",
                path="renders/slide-1.png",
                sha256="1" * 64,
            ),
        ),
        instruction="Improve hierarchy without changing any text.",
        retained_content=("all factual slide text",),
    )
    payload = {
        "schema_version": "sophia-deck-repair-program/v1",
        "build_id": BUILD_ID,
        "initial_quality_run_id": INITIAL_QUALITY_RUN_ID,
        "initial_manifest_revision": 1,
        "repair_attempt": 1,
        "plan_revision_allowed": False,
        "authorized_selectors": ("slide:1",),
        "authorized_source_roles": {"slide:1": ("slide_css",)},
        "deck_instruction": "Apply exactly one bounded hierarchy repair.",
        "selector_repairs": (repair,),
        "must_preserve": ("all facts and slide count",),
        "must_not": ("revise plans or generate assets",),
        "skill_refs": (skill,),
        "expected_improvements": ("weak_visual_hierarchy",),
        "forbidden_regressions": ("content_fidelity",),
        "rubric_version": instrument.lock.rubric_version,
        "instrument_hash": canonical_sha256(instrument.lock),
    }
    return DeckRepairProgram(**payload, program_hash=canonical_sha256(payload))


def _root_program(instrument: DeckQualityRuntimeInstrument) -> DeckRepairProgram:
    skill = SkillRef(
        path="skills/public/hands-on-deck/designing-slides.md",
        source_hash="e" * 64,
        excerpt_hash="f" * 64,
    )
    repair = SelectorRepair(
        selector=DECK_STYLE_ROOT_SELECTOR,
        failure_codes=("repetitive_structure",),
        render_evidence=(
            RepairRenderEvidence(
                selector=f"slide:{index}",
                path=f"renders/slide-{index}.png",
                sha256=str(index) * 64,
            )
            for index in (1, 2)
        ),
        instruction="Improve the shared hierarchy without changing text.",
        retained_content=("all factual slide text",),
    )
    payload = {
        "schema_version": "sophia-deck-repair-program/v1",
        "build_id": BUILD_ID,
        "initial_quality_run_id": INITIAL_QUALITY_RUN_ID,
        "initial_manifest_revision": 1,
        "repair_attempt": 1,
        "plan_revision_allowed": False,
        "authorized_selectors": (DECK_STYLE_ROOT_SELECTOR,),
        "authorized_source_roles": {DECK_STYLE_ROOT_SELECTOR: ("deck_css",)},
        "deck_instruction": "Apply one shared stylesheet repair.",
        "selector_repairs": (repair,),
        "must_preserve": ("all facts and slide count",),
        "must_not": ("revise plans or generate assets",),
        "skill_refs": (skill,),
        "expected_improvements": ("repetitive_structure",),
        "forbidden_regressions": ("content_fidelity",),
        "rubric_version": instrument.lock.rubric_version,
        "instrument_hash": canonical_sha256(instrument.lock),
    }
    return DeckRepairProgram(**payload, program_hash=canonical_sha256(payload))


def _source_values() -> dict[tuple[str, str], str]:
    deck_css = ".slide-root { background: #001122; color: #ffffff; }"
    values = {(DECK_STYLE_ROOT_SELECTOR, "deck_css"): deck_css}
    for index in (1, 2):
        selector = f"slide:{index}"
        body = f'<h1 data-deck-id="title-{index}" data-deck-role="title" data-deck-required="true">Title {index}</h1><p data-deck-id="narrative-{index}" data-deck-role="narrative" data-deck-required="true">Narrative {index}</p>'
        css = f"#title-{index} {{ font-size: 52px; }}"
        values[(selector, "body")] = body
        values[(selector, "slide_css")] = css
        values[(selector, "notes")] = f"Notes {index}"
        values[(selector, "assembled")] = assemble_compact_slide_html(
            deck_stylesheet=deck_css,
            html_body=body,
            slide_css=css,
        )
    return values


def _component(
    selector: str,
    *,
    index: int,
    values: dict[tuple[str, str], str],
) -> BuildComponent:
    roles = (
        ("deck_css",)
        if selector == DECK_STYLE_ROOT_SELECTOR
        else (
            "body",
            "slide_css",
            "notes",
            "assembled",
        )
    )
    component_id = "component_style" if selector == DECK_STYLE_ROOT_SELECTOR else f"component_slide_{index}"
    version_id = f"component_version_{index}"
    source_roles = {role: f"artifacts/{USER_ID}/{THREAD_ID}/foundation/.builder/builds/{BUILD_ID}/components/{component_id}/versions/{version_id}/{role}.txt" for role in roles}
    source_hashes = {role: _sha(values[(selector, role)].encode()) for role in roles}
    if selector != DECK_STYLE_ROOT_SELECTOR:
        source_hashes["deck_css"] = _sha(values[(DECK_STYLE_ROOT_SELECTOR, "deck_css")].encode())
    return BuildComponent(
        id=component_id,
        selector=selector,
        type="deck_style" if selector == DECK_STYLE_ROOT_SELECTOR else "slide",
        index=index,
        source_path=source_roles[roles[0]],
        source_roles=source_roles,
        source_hashes=source_hashes,
        shared_dependencies=[] if selector == DECK_STYLE_ROOT_SELECTOR else [DECK_STYLE_ROOT_SELECTOR],
        asset_paths=[],
        status="gated",
        gate_results={"mechanical_passed": True},
        current_version_id=version_id,
        provenance={"authored_by": "fresh"},
    )


def _manifest(values: dict[tuple[str, str], str]) -> BuildManifest:
    components = [
        _component(DECK_STYLE_ROOT_SELECTOR, index=0, values=values),
        _component("slide:1", index=1, values=values),
        _component("slide:2", index=2, values=values),
    ]
    return BuildManifest(
        manifest_revision=1,
        build_id=BUILD_ID,
        user_id=USER_ID,
        thread_id=THREAD_ID,
        format="pptx",
        status="complete",
        logical_artifact_id="logical_deck_01",
        current_artifact_version_id=INITIAL_ARTIFACT_ID,
        deliverable_path="/mnt/user-data/outputs/deck.pptx",
        components=components,
        format_extensions={
            "deck": {
                "current_pptx_hash": "2" * 64,
                "artifact_storage_object_path": "artifacts/original.pptx",
            }
        },
        created_at=FIXED_TIME,
        updated_at=FIXED_TIME,
    )


def _inventory() -> dict[str, Any]:
    return {
        f"slide:{index}": {
            "native_slide_index": index - 1,
            "shape_count": 2,
            "full_slide_picture_count": 0,
            "shapes": [
                {
                    "name": f"title-{index}",
                    "type": "TEXT_BOX",
                    "text_preview": f"Title {index}",
                    "full_slide": False,
                },
                {
                    "name": f"narrative-{index}",
                    "type": "TEXT_BOX",
                    "text_preview": f"Narrative {index}",
                    "full_slide": False,
                },
            ],
        }
        for index in (1, 2)
    }


def _plans() -> tuple[dict[str, Any], dict[str, Any]]:
    design = {
        "source": "frozen",
        "subject": "PSI control loop",
        "audience": "technical leaders",
        "goal": "explain the mechanism",
    }
    creative = {
        "subject": "PSI control loop",
        "audience": "technical leaders",
        "goal": "explain the mechanism",
        "image_strategy": "diagram_native",
        "image_assets": [],
        "design_plan": design,
    }
    return creative, design


def _baseline(instrument: DeckQualityRuntimeInstrument) -> DeckCandidateBaseline:
    creative, design = _plans()
    build_record = {
        "build_id": BUILD_ID,
        "deck_title": "PSI Control Loop",
        "register": "professional_technical",
        "visual_policy": "auto",
        "deck_authoring_contract": "compact_model_html_v1",
        "style_profile": {},
        "native_editability_score": 0.9,
        "full_slide_picture_count": 0,
        "native_shape_inventory": _inventory(),
        "slides": [
            {
                "selector": f"slide:{index}",
                "title": f"Title {index}",
                "narrative": f"Narrative {index}",
                "role": "cover" if index == 1 else "closing",
                "layout_kind": "cover_hero" if index == 1 else "closing_summary",
                "claim": None,
                "visual_prompt": None,
                "speaker_notes": f"Notes {index}",
            }
            for index in (1, 2)
        ],
    }
    renders = (_png("navy"), _png("navy"))
    return DeckCandidateBaseline(
        build_id=BUILD_ID,
        user_id=USER_ID,
        thread_id=THREAD_ID,
        task_id="task-01",
        builder_run_id="builder-run-01",
        parent_builder_trace_id="trace-root-01",
        initial_quality_run_id=INITIAL_QUALITY_RUN_ID,
        logical_artifact_id="logical_deck_01",
        initial_artifact_version_id=INITIAL_ARTIFACT_ID,
        initial_manifest_revision=1,
        task_brief="Build a two-slide deck explaining the PSI control loop.",
        build_record=build_record,
        creative_plan_record=creative,
        design_plan_record=design,
        instrument=instrument,
        visible_text=tuple(
            VisibleTextSlide(
                selector=f"slide:{index}",
                text=f"Title {index}\nNarrative {index}",
                source_hash=canonical_sha256(
                    {
                        "selector": f"slide:{index}",
                        "text": f"Title {index}\nNarrative {index}",
                    }
                ),
            )
            for index in (1, 2)
        ),
        renders=tuple(
            BaselineDeckRender(
                selector=f"slide:{index}",
                content=content,
                sha256=_sha(content),
            )
            for index, content in enumerate(renders, start=1)
        ),
    )


def _visible_text_slide(selector: str, text: str) -> VisibleTextSlide:
    return VisibleTextSlide(
        selector=selector,
        text=text,
        source_hash=canonical_sha256({"selector": selector, "text": text}),
    )


def _baseline_with_first_slide_text(
    instrument: DeckQualityRuntimeInstrument,
    text: str,
) -> DeckCandidateBaseline:
    baseline = _baseline(instrument)
    return baseline.model_copy(
        update={
            "visible_text": (
                _visible_text_slide("slide:1", text),
                baseline.visible_text[1],
            )
        }
    )


def _stub_candidate_visible_text(
    monkeypatch: pytest.MonkeyPatch,
    *,
    baseline: DeckCandidateBaseline,
    first_slide_text: str,
) -> None:
    candidate = (
        _visible_text_slide("slide:1", first_slide_text),
        baseline.visible_text[1],
    )
    monkeypatch.setattr(
        candidate_module,
        "_candidate_visible_text",
        lambda _content: candidate,
    )


def _request(
    instrument: DeckQualityRuntimeInstrument,
) -> DeckCandidateCompileRequest:
    values = _source_values()
    manifest = _manifest(values)
    repaired_css = "#title-1 { font-size: 64px; font-weight: 700; }"
    sources: list[DeckCandidateSource] = []
    for component in manifest.components:
        changed = component.selector == "slide:1"
        for role in component.source_roles:
            content = repaired_css if (component.selector, role) == ("slide:1", "slide_css") else values[(component.selector, role)]
            sources.append(
                DeckCandidateSource(
                    selector=component.selector,
                    source_role=role,
                    object_path=component.source_roles[role],
                    source_hash=_sha(content.encode()),
                    content=content.encode(),
                    model_authored=(component.selector, role) == ("slide:1", "slide_css"),
                    component_version_changed=changed,
                )
            )
    return DeckCandidateCompileRequest(
        transaction_id=TRANSACTION_ID,
        operation_id=OPERATION_ID,
        build_id=BUILD_ID,
        user_id=USER_ID,
        thread_id=THREAD_ID,
        candidate_manifest_revision=2,
        artifact_version_id=CANDIDATE_ARTIFACT_ID,
        candidate_hash="3" * 64,
        baseline_manifest=manifest,
        program=_program(instrument),
        sources=tuple(sources),
        derived_source_targets=(("slide:1", "assembled"),),
    )


def _root_request(
    instrument: DeckQualityRuntimeInstrument,
) -> DeckCandidateCompileRequest:
    values = _source_values()
    manifest = _manifest(values)
    repaired_css = ".slide-root { background: #102030; color: #ffffff; }"
    sources: list[DeckCandidateSource] = []
    for component in manifest.components:
        for role in component.source_roles:
            content = repaired_css if (component.selector, role) == (DECK_STYLE_ROOT_SELECTOR, "deck_css") else values[(component.selector, role)]
            sources.append(
                DeckCandidateSource(
                    selector=component.selector,
                    source_role=role,
                    object_path=component.source_roles[role],
                    source_hash=_sha(content.encode()),
                    content=content.encode(),
                    model_authored=(component.selector, role) == (DECK_STYLE_ROOT_SELECTOR, "deck_css"),
                    component_version_changed=True,
                )
            )
    return DeckCandidateCompileRequest(
        transaction_id=TRANSACTION_ID,
        operation_id=OPERATION_ID,
        build_id=BUILD_ID,
        user_id=USER_ID,
        thread_id=THREAD_ID,
        candidate_manifest_revision=2,
        artifact_version_id=CANDIDATE_ARTIFACT_ID,
        candidate_hash="5" * 64,
        baseline_manifest=manifest,
        program=_root_program(instrument),
        sources=tuple(sources),
        derived_source_targets=(("slide:1", "assembled"), ("slide:2", "assembled")),
    )


class _BaselineLoader:
    def __init__(self, baseline: DeckCandidateBaseline) -> None:
        self.baseline = baseline
        self.calls = 0

    async def load(self, request: DeckCandidateCompileRequest) -> DeckCandidateBaseline:
        assert request.build_id == BUILD_ID
        self.calls += 1
        return self.baseline


class _Publisher:
    def __init__(self) -> None:
        self.calls: list[Any] = []
        self.source_pack_bytes: list[bytes] = []

    async def publish(
        self,
        *,
        prepared,
        instrument,
        pptx_bytes,
        source_pack,
        source_pack_bytes,
    ):
        self.calls.append((prepared, instrument, pptx_bytes))
        captured_pack, captured_bytes = capture_deck_quality_source_pack(
            prepared=prepared,
            instrument=instrument,
        )
        assert captured_pack == source_pack
        assert captured_bytes == source_pack_bytes
        self.source_pack_bytes.append(source_pack_bytes)
        quality_run_id = derive_quality_run_id(
            artifact_version_id=prepared.artifact_version_id,
            campaign_id="DQ-1",
            instrument=instrument.lock,
        )
        return DeckQualityProducerBundleReceipt(
            quality_run_id=quality_run_id,
            bundle_object_path=deck_quality_producer_bundle_path(quality_run_id),
            bundle_hash="4" * 64,
            bundle_size_bytes=100,
        )


class _MemoryPublicationStore:
    def __init__(self) -> None:
        self.objects: dict[str, bytes] = {}
        self.creates: list[str] = []

    def create_if_absent(self, object_path, content, *, content_type):
        assert content_type
        self.creates.append(object_path)
        if object_path in self.objects:
            return "exists"
        self.objects[object_path] = content
        return "created"

    def read(self, object_path):
        return self.objects.get(object_path)

    def read_bounded(self, object_path, *, max_bytes):
        value = self.read(object_path)
        if value is not None and len(value) > max_bytes:
            raise RuntimeError("oversized")
        return value


class _ScriptedSnapshotStore(_MemoryPublicationStore):
    def __init__(self, *snapshot_events: str) -> None:
        super().__init__()
        self.snapshot_events = list(snapshot_events)
        self.snapshot_create_count = 0

    def create_if_absent(self, object_path, content, *, content_type):
        if content_type.endswith("presentationml.presentation"):
            self.snapshot_create_count += 1
            if self.snapshot_events:
                event = self.snapshot_events.pop(0)
                self.creates.append(object_path)
                if event == "raise_after_commit":
                    self.objects[object_path] = content
                    raise RuntimeError("synthetic create response loss")
                if event == "raise_before_commit":
                    raise RuntimeError("synthetic pre-commit failure")
                if event == "invalid_outcome":
                    return "invalid"
                raise AssertionError(f"unknown snapshot event: {event}")
        return super().create_if_absent(
            object_path,
            content,
            content_type=content_type,
        )


class _FakeDeckService:
    def __init__(
        self,
        *,
        mutate_plan: bool = False,
        archive_datetime: tuple[int, int, int, int, int, int] | None = None,
        build_timestamp: str | None = None,
        retained_build_value: str = "stable",
    ) -> None:
        self.mutate_plan = mutate_plan
        self.archive_datetime = archive_datetime
        self.build_timestamp = build_timestamp
        self.retained_build_value = retained_build_value
        self.foundation_config = None
        self.tracing_enabled = None
        self.raw_pptx_bytes: bytes | None = None
        self.thread_ids: list[int] = []

    def prepare_and_build(self, **kwargs) -> DeckBuildResult:
        self.thread_ids.append(threading.get_ident())
        runtime = kwargs["runtime"]
        slides = kwargs["slides"]
        output_path = kwargs["output_path"]
        self.foundation_config = runtime.context["build_foundation_config"]
        self.tracing_enabled = get_tracing_context()["enabled"]
        output_host = Path(replace_virtual_path(output_path, runtime.state["thread_data"]))
        output_host.parent.mkdir(parents=True, exist_ok=True)
        presentation = Presentation()
        blank = presentation.slide_layouts[6]
        for index, slide_input in enumerate(slides, start=1):
            slide = presentation.slides.add_slide(blank)
            title = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
            title.name = f"title-{index}"
            title.text_frame.text = slide_input["title"]
            narrative = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
            narrative.name = f"narrative-{index}"
            narrative.text_frame.text = slide_input["narrative"]
        presentation.save(output_host)
        raw_pptx = output_host.read_bytes()
        if self.archive_datetime is not None:
            raw_pptx = _retimestamp_pptx(raw_pptx, self.archive_datetime)
            output_host.write_bytes(raw_pptx)
        self.raw_pptx_bytes = raw_pptx

        outputs = Path(runtime.state["thread_data"]["outputs_path"])
        deck_root = outputs / "deck_build"
        deck_root.mkdir(parents=True, exist_ok=True)
        creative = dict(kwargs["creative_plan"])
        if self.mutate_plan:
            creative["goal"] = "silently revised"
        (deck_root / "creative_plan.json").write_text(json.dumps(creative), encoding="utf-8")
        (deck_root / "design_plan.json").write_text(json.dumps(kwargs["design_plan"]), encoding="utf-8")
        assembled_slides = []
        for index, slide in enumerate(slides, start=1):
            assembled_slides.append(
                {
                    **slide,
                    "selector": f"slide:{index}",
                    "html_source": assemble_compact_slide_html(
                        deck_stylesheet=kwargs["deck_stylesheet"],
                        html_body=slide["html_body"],
                        slide_css=slide["slide_css"],
                    ),
                }
            )
        build_record = {
            "build_id": BUILD_ID,
            "deck_title": kwargs["deck_title"],
            "register": kwargs["register"],
            "visual_policy": kwargs["visual_policy"],
            "deck_authoring_contract": kwargs["authoring_contract"],
            "style_profile": kwargs["style_profile"],
            "native_editability_score": 0.9,
            "full_slide_picture_count": 0,
            "candidate_retained_provenance": self.retained_build_value,
            "native_shape_inventory": _inventory(),
            "slides": assembled_slides,
        }
        if self.build_timestamp is not None:
            build_record.update(
                {
                    "created_at": self.build_timestamp,
                    "updated_at": self.build_timestamp,
                    "service_elapsed_ms": int(self.archive_datetime[-1]) if self.archive_datetime else 0,
                    "langsmith_trace_ids": [f"trace-{self.build_timestamp}"],
                }
            )
        build_path = deck_root / "build.json"
        build_path.write_text(json.dumps(build_record), encoding="utf-8")
        render_root = outputs / ".builder" / "deck_native" / "rendered"
        render_root.mkdir(parents=True, exist_ok=True)
        for index in range(2):
            (render_root / f"slide-{index}.jpg").write_bytes(_jpg("navy"))
        return DeckBuildResult(
            success=True,
            build_id=BUILD_ID,
            deck_build_path="/mnt/user-data/outputs/deck_build/build.json",
            pptx_path=output_path,
            deck_compile_mode=NATIVE_DECK_COMPILE_MODE,
            deck_authoring_contract=kwargs["authoring_contract"],
            native_required=True,
            legacy_screenshot_debug=False,
            native_editability_score=0.9,
            full_slide_picture_count=0,
            slide_count=2,
            expected_visual_count=0,
            successful_visual_count=0,
            missing_visual_count=0,
            mechanical_gate_results={"passed": True},
            source_retention_report={"passed": True},
            native_contrast_report={"passed": True},
            native_mechanical_report={
                "render_success": True,
                "lint_fix_success": True,
                "lint_residue_count": 0,
            },
        )


def test_content_proof_accepts_visible_token_boundary_regrouping(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    instrument = _instrument()
    baseline = _baseline_with_first_slide_text(
        instrument,
        "PSI \ufb01delity:\ncontrol loop.",
    )
    _stub_candidate_visible_text(
        monkeypatch,
        baseline=baseline,
        first_slide_text="PSI fidelity: control\nloop.",
    )

    proof = candidate_module._content_proof(
        baseline=baseline,
        candidate_pptx=b"candidate",
        selectors=("slide:1", "slide:2"),
        candidate_editability_score=0.9,
    )

    assert proof.required_content_preserved is True


def test_content_proof_rejects_visible_token_reordering(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    instrument = _instrument()
    baseline = _baseline_with_first_slide_text(
        instrument,
        "PSI control loop.",
    )
    _stub_candidate_visible_text(
        monkeypatch,
        baseline=baseline,
        first_slide_text="PSI loop control.",
    )

    with pytest.raises(DeckCandidateCompilationError, match="content_changed"):
        candidate_module._content_proof(
            baseline=baseline,
            candidate_pptx=b"candidate",
            selectors=("slide:1", "slide:2"),
            candidate_editability_score=0.9,
        )


@pytest.mark.parametrize(
    "candidate_text",
    ("PSI control loop now.", "PSI control."),
    ids=("addition", "removal"),
)
def test_content_proof_rejects_visible_token_addition_or_removal(
    monkeypatch: pytest.MonkeyPatch,
    candidate_text: str,
) -> None:
    instrument = _instrument()
    baseline = _baseline_with_first_slide_text(
        instrument,
        "PSI control loop.",
    )
    _stub_candidate_visible_text(
        monkeypatch,
        baseline=baseline,
        first_slide_text=candidate_text,
    )

    with pytest.raises(DeckCandidateCompilationError, match="content_changed"):
        candidate_module._content_proof(
            baseline=baseline,
            candidate_pptx=b"candidate",
            selectors=("slide:1", "slide:2"),
            candidate_editability_score=0.9,
        )


@pytest.mark.parametrize(
    "candidate_text",
    ("psi control loop.", "PSI control loop!"),
    ids=("case", "punctuation"),
)
def test_content_proof_rejects_visible_token_case_or_punctuation_changes(
    monkeypatch: pytest.MonkeyPatch,
    candidate_text: str,
) -> None:
    instrument = _instrument()
    baseline = _baseline_with_first_slide_text(
        instrument,
        "PSI control loop.",
    )
    _stub_candidate_visible_text(
        monkeypatch,
        baseline=baseline,
        first_slide_text=candidate_text,
    )

    with pytest.raises(DeckCandidateCompilationError, match="content_changed"):
        candidate_module._content_proof(
            baseline=baseline,
            candidate_pptx=b"candidate",
            selectors=("slide:1", "slide:2"),
            candidate_editability_score=0.9,
        )


def test_failed_native_mechanical_service_maps_to_specific_compiler_code() -> None:
    instrument = _instrument()
    request = _request(instrument)
    result = DeckBuildResult(
        success=False,
        build_id=BUILD_ID,
        deck_build_path="/mnt/user-data/outputs/deck_build/build.json",
        failure_code="deck_mechanical_gate_failed",
    )

    with pytest.raises(DeckCandidateCompilationError) as error:
        candidate_module._validate_service_result(
            result,
            request=request,
            slide_count=2,
            output_path="/mnt/user-data/outputs/candidate.pptx",
            authoring_contract="compact_model_html_v2",
        )

    assert error.value.code == "mechanical_gate_failed"


def test_authenticated_parent_thread_snapshot_projects_builder_manifest_owner() -> None:
    instrument = _instrument()
    expected = _baseline(instrument)
    manifest = _manifest(_source_values())
    parent_thread_id = "parent_companion_thread_01"
    evidence_manifest = SimpleNamespace(
        quality_run_id=INITIAL_QUALITY_RUN_ID,
        thread_id=parent_thread_id,
        task_id=THREAD_ID,
        selectors=tuple(item.selector for item in expected.renders),
        render_hashes={item.selector: item.sha256 for item in expected.renders},
    )
    authenticated = SimpleNamespace(
        row=SimpleNamespace(
            quality_run_id=INITIAL_QUALITY_RUN_ID,
            build_id=BUILD_ID,
            user_id=USER_ID,
            thread_id=parent_thread_id,
            task_id=THREAD_ID,
            builder_run_id=expected.builder_run_id,
            parent_builder_trace_id=expected.parent_builder_trace_id,
            logical_artifact_id=expected.logical_artifact_id,
            artifact_version_id=expected.initial_artifact_version_id,
            manifest_revision=expected.initial_manifest_revision,
        ),
        manifest=manifest,
        evidence_manifest=evidence_manifest,
        evidence_bundle=SimpleNamespace(
            build_record=expected.build_record,
            snapshot=SimpleNamespace(
                build_id=BUILD_ID,
                brief=SimpleNamespace(request=expected.task_brief),
                creative_plan=expected.creative_plan_record,
                design_plan=expected.design_plan_record,
                visible_text=expected.visible_text,
            ),
        ),
    )

    projected = baseline_from_authenticated_snapshot(
        authenticated,
        instrument=instrument,
        render_contents={item.selector: item.content for item in expected.renders},
    )

    assert projected.thread_id == manifest.thread_id == THREAD_ID
    assert projected.task_id == THREAD_ID


def test_authenticated_snapshot_rejects_unrelated_parent_and_task_owner() -> None:
    instrument = _instrument()
    expected = _baseline(instrument)
    manifest = _manifest(_source_values())
    authenticated = SimpleNamespace(
        row=SimpleNamespace(
            quality_run_id=INITIAL_QUALITY_RUN_ID,
            build_id=BUILD_ID,
            user_id=USER_ID,
            thread_id="unrelated_parent_thread_01",
            task_id="unrelated_builder_thread_01",
            builder_run_id=expected.builder_run_id,
            parent_builder_trace_id=expected.parent_builder_trace_id,
            logical_artifact_id=expected.logical_artifact_id,
            artifact_version_id=expected.initial_artifact_version_id,
            manifest_revision=expected.initial_manifest_revision,
        ),
        manifest=manifest,
        evidence_manifest=SimpleNamespace(
            quality_run_id=INITIAL_QUALITY_RUN_ID,
            thread_id="unrelated_parent_thread_01",
            task_id="unrelated_builder_thread_01",
            selectors=tuple(item.selector for item in expected.renders),
            render_hashes={item.selector: item.sha256 for item in expected.renders},
        ),
        evidence_bundle=SimpleNamespace(
            build_record=expected.build_record,
            snapshot=SimpleNamespace(
                build_id=BUILD_ID,
                brief=SimpleNamespace(request=expected.task_brief),
                creative_plan=expected.creative_plan_record,
                design_plan=expected.design_plan_record,
                visible_text=expected.visible_text,
            ),
        ),
    )

    with pytest.raises(DeckCandidateCompilationError, match="baseline_invalid"):
        baseline_from_authenticated_snapshot(
            authenticated,
            instrument=instrument,
            render_contents={item.selector: item.content for item in expected.renders},
        )


def test_candidate_compiler_rebuilds_sources_proves_locality_and_publishes_dq1(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    instrument = _instrument()
    loader = _BaselineLoader(_baseline(instrument))
    publisher = _Publisher()
    service = _FakeDeckService()
    compiler = ProductionDeckCandidateCompiler(
        baseline_loader=loader,
        publisher=publisher,
        service_factory=lambda _batch, _single: service,
    )
    blocking_threads: list[int] = []
    original_canonicalize = candidate_module._canonicalize_candidate_artifact
    original_render_proof = candidate_module._render_proof

    def canonicalize(*args: Any, **kwargs: Any) -> bytes:
        blocking_threads.append(threading.get_ident())
        return original_canonicalize(*args, **kwargs)

    def render_proof(*args: Any, **kwargs: Any):
        blocking_threads.append(threading.get_ident())
        return original_render_proof(*args, **kwargs)

    monkeypatch.setattr(candidate_module, "_canonicalize_candidate_artifact", canonicalize)
    monkeypatch.setattr(candidate_module, "_render_proof", render_proof)

    compilation = asyncio.run(compiler.compile(_request(instrument)))

    assert loader.calls == 1
    assert len(publisher.calls) == 1
    assert service.foundation_config.enabled is False
    assert service.foundation_config.manifest_mode == "off"
    assert service.tracing_enabled is False
    assert service.thread_ids and service.thread_ids[0] != threading.get_ident()
    assert len(blocking_threads) == 2
    assert all(thread_id != threading.get_ident() for thread_id in blocking_threads)
    assert compilation.mechanical.status == "passed"
    assert compilation.locality.authorized_selectors == ("slide:1",)
    assert compilation.locality.changed_component_versions == ("slide:1",)
    assert compilation.locality.unchanged_component_versions == (
        DECK_STYLE_ROOT_SELECTOR,
        "slide:2",
    )
    assert compilation.content.required_content_preserved is True
    assert compilation.creative_plan_record["plan_revision_changed"] is False
    assert compilation.design_plan_record["plan_revision_changed"] is False
    assert compilation.native_record["verified"] is True
    assert compilation.render_collateral_record["compared_selectors"] == ["slide:2"]
    assert tuple((item.selector, item.source_role) for item in compilation.derived_sources) == (("slide:1", "assembled"),)
    assert compilation.dq1_publication_metadata["quality_run_id"].startswith("quality_")
    for record in (
        compilation.build_record,
        compilation.creative_plan_record,
        compilation.design_plan_record,
        compilation.mechanical_record,
        compilation.native_record,
        compilation.render_collateral_record,
        compilation.dq1_publication_metadata,
    ):
        assert materializer_module._safe_record(record) == record


def test_candidate_compiler_is_byte_stable_across_distinct_save_times() -> None:
    instrument = _instrument()
    publisher = _Publisher()
    service_inputs = iter(
        (
            ((2026, 7, 17, 12, 0, 0), "2026-07-17T12:00:00+00:00"),
            ((2026, 7, 17, 12, 0, 4), "2026-07-17T12:00:04+00:00"),
        )
    )
    services: list[_FakeDeckService] = []

    def service_factory(_batch, _single):
        archive_datetime, build_timestamp = next(service_inputs)
        service = _FakeDeckService(
            archive_datetime=archive_datetime,
            build_timestamp=build_timestamp,
        )
        services.append(service)
        return service

    compiler = ProductionDeckCandidateCompiler(
        baseline_loader=_BaselineLoader(_baseline(instrument)),
        publisher=publisher,
        service_factory=service_factory,
    )

    first = asyncio.run(compiler.compile(_request(instrument)))
    second = asyncio.run(compiler.compile(_request(instrument)))

    raw_first = services[0].raw_pptx_bytes
    raw_second = services[1].raw_pptx_bytes
    assert raw_first is not None and raw_second is not None
    assert _sha(raw_first) != _sha(raw_second)
    with zipfile.ZipFile(io.BytesIO(raw_first)) as first_package:
        first_members = {name: first_package.read(name) for name in first_package.namelist()}
    with zipfile.ZipFile(io.BytesIO(raw_second)) as second_package:
        second_members = {name: second_package.read(name) for name in second_package.namelist()}
    assert first_members == second_members

    assert first == second
    assert publisher.calls[0][2] == publisher.calls[1][2] == first.pptx_bytes
    assert publisher.source_pack_bytes[0] == publisher.source_pack_bytes[1]
    assert candidate_module._canonicalize_pptx_package(first.pptx_bytes) == first.pptx_bytes
    with zipfile.ZipFile(io.BytesIO(first.pptx_bytes)) as package:
        assert package.namelist() == sorted(package.namelist())
        assert all(info.date_time == candidate_module._CANONICAL_ZIP_DATETIME for info in package.infolist())
        assert {name: package.read(name) for name in package.namelist()} == first_members
    presentation = Presentation(io.BytesIO(first.pptx_bytes))
    assert len(presentation.slides) == 2
    assert [shape.name for shape in presentation.slides[0].shapes] == ["title-1", "narrative-1"]
    assert first.build_record["authoritative_build_record_hash"] == second.build_record["authoritative_build_record_hash"]


def test_candidate_compiler_retries_after_crash_with_one_exact_dq1_publication() -> None:
    instrument = _instrument()
    store = _MemoryPublicationStore()
    publisher = DurableCandidateDq1Publisher(store_factory=lambda: store)
    service_inputs = iter(
        (
            ((2026, 7, 17, 12, 0, 0), "2026-07-17T12:00:00+00:00"),
            ((2026, 7, 17, 12, 0, 6), "2026-07-17T12:00:06+00:00"),
        )
    )

    def service_factory(_batch, _single):
        archive_datetime, build_timestamp = next(service_inputs)
        return _FakeDeckService(
            archive_datetime=archive_datetime,
            build_timestamp=build_timestamp,
        )

    compiler = ProductionDeckCandidateCompiler(
        baseline_loader=_BaselineLoader(_baseline(instrument)),
        publisher=publisher,
        service_factory=service_factory,
    )
    request = _request(instrument)

    first = asyncio.run(compiler.compile(request))

    # Simulate a process crash after DQ-1 publication but before the DQ-2 stage
    # record is stored.  The gateway may also acknowledge and archive the live
    # outbox before the materializer retries the identical compilation.
    quality_run_id = first.dq1_publication_metadata["quality_run_id"]
    live_bundle_path = deck_quality_producer_bundle_path(quality_run_id)
    archive_bundle_path = deck_quality_producer_archive_path(quality_run_id)
    store.objects[archive_bundle_path] = store.objects.pop(live_bundle_path)

    recovered = asyncio.run(compiler.compile(request))

    assert first == recovered
    assert _sha(first.pptx_bytes) == _sha(recovered.pptx_bytes)
    snapshot_paths = [path for path in store.objects if path.endswith("/candidate.pptx")]
    source_pack_paths = [path for path in store.objects if path.endswith("/publication/source_pack/manifest.json")]
    stored_bundle_paths = [path for path in store.objects if path in {live_bundle_path, archive_bundle_path}]
    assert len(snapshot_paths) == len(source_pack_paths) == len(stored_bundle_paths) == 1
    assert store.objects[snapshot_paths[0]] == recovered.pptx_bytes
    assert recovered.dq1_publication_metadata["bundle_object_path"] == live_bundle_path
    assert recovered.dq1_publication_metadata["bundle_archive_object_path"] == archive_bundle_path
    source_pack = json.loads(store.objects[source_pack_paths[0]])
    assert {
        "created_at",
        "updated_at",
        "service_elapsed_ms",
        "langsmith_trace_ids",
    }.isdisjoint(source_pack["build_record"])


def test_candidate_compiler_rejects_stale_bundle_when_current_source_pack_changes() -> None:
    instrument = _instrument()
    store = _MemoryPublicationStore()
    publisher = DurableCandidateDq1Publisher(store_factory=lambda: store)
    retained_values = iter(("first-retained-value", "changed-retained-value"))

    def service_factory(_batch, _single):
        return _FakeDeckService(
            archive_datetime=(2026, 7, 17, 12, 0, 0),
            retained_build_value=next(retained_values),
        )

    compiler = ProductionDeckCandidateCompiler(
        baseline_loader=_BaselineLoader(_baseline(instrument)),
        publisher=publisher,
        service_factory=service_factory,
    )
    request = _request(instrument)
    first = asyncio.run(compiler.compile(request))
    quality_run_id = first.dq1_publication_metadata["quality_run_id"]
    live_bundle_path = deck_quality_producer_bundle_path(quality_run_id)
    archive_bundle_path = deck_quality_producer_archive_path(quality_run_id)
    store.objects[archive_bundle_path] = store.objects.pop(live_bundle_path)

    with pytest.raises(DeckCandidateCompilationError, match="publication_failed") as exc_info:
        asyncio.run(compiler.compile(request))

    assert exc_info.value.code == "publication_failed"
    snapshot_paths = [path for path in store.objects if path.endswith("/candidate.pptx")]
    source_pack_paths = [path for path in store.objects if path.endswith("/publication/source_pack/manifest.json")]
    assert len(snapshot_paths) == len(source_pack_paths) == 1
    assert store.objects[snapshot_paths[0]] == first.pptx_bytes
    stored_source_pack = json.loads(store.objects[source_pack_paths[0]])
    assert stored_source_pack["build_record"]["candidate_retained_provenance"] == "first-retained-value"
    assert archive_bundle_path in store.objects


def test_candidate_compiler_rejects_plan_revision_before_publication() -> None:
    instrument = _instrument()
    publisher = _Publisher()
    compiler = ProductionDeckCandidateCompiler(
        baseline_loader=_BaselineLoader(_baseline(instrument)),
        publisher=publisher,
        service_factory=lambda _batch, _single: _FakeDeckService(mutate_plan=True),
    )

    with pytest.raises(DeckCandidateCompilationError, match="plan_changed") as exc_info:
        asyncio.run(compiler.compile(_request(instrument)))

    assert exc_info.value.code == "plan_changed"
    assert publisher.calls == []


def test_candidate_compiler_proves_shared_style_dependency_closure() -> None:
    instrument = _instrument()
    compiler = ProductionDeckCandidateCompiler(
        baseline_loader=_BaselineLoader(_baseline(instrument)),
        publisher=_Publisher(),
        service_factory=lambda _batch, _single: _FakeDeckService(),
    )

    compilation = asyncio.run(compiler.compile(_root_request(instrument)))

    assert compilation.locality.authorized_selectors == (DECK_STYLE_ROOT_SELECTOR,)
    assert compilation.locality.changed_component_versions == (
        DECK_STYLE_ROOT_SELECTOR,
        "slide:1",
        "slide:2",
    )
    assert compilation.locality.unchanged_component_versions == ()
    assert compilation.locality.shared_dependency_changed is True
    assert compilation.render_collateral_record["compared_selectors"] == []
    assert tuple(item.selector for item in compilation.derived_sources) == (
        "slide:1",
        "slide:2",
    )


def test_candidate_compiler_rejects_source_hash_mismatch_before_build() -> None:
    instrument = _instrument()
    request = _request(instrument)
    bad_source = request.sources[0].model_copy(update={"source_hash": "9" * 64})
    request = request.model_copy(update={"sources": (bad_source, *request.sources[1:])})
    service_calls = 0

    def service_factory(_batch, _single):
        nonlocal service_calls
        service_calls += 1
        return _FakeDeckService()

    compiler = ProductionDeckCandidateCompiler(
        baseline_loader=_BaselineLoader(_baseline(instrument)),
        publisher=_Publisher(),
        service_factory=service_factory,
    )

    with pytest.raises(DeckCandidateCompilationError, match="source_hash_mismatch"):
        asyncio.run(compiler.compile(request))

    assert service_calls == 0


def test_reuse_only_image_runner_copies_exact_frozen_asset_bytes(tmp_path: Path) -> None:
    content = _png("teal")
    asset = BaselineVisualAsset(
        asset_id="hero-texture",
        selector="slide:1",
        content=content,
        sha256=_sha(content),
    )
    runner = candidate_module._ReuseOnlyImageRunner((asset,))
    outputs = tmp_path / "outputs"
    manifest = outputs / "assets" / "slide-visuals.manifest.json"
    manifest.parent.mkdir(parents=True)
    manifest.write_text(
        json.dumps(
            {
                "items": [
                    {
                        "slide_index": 1,
                        "output_file": "/mnt/user-data/outputs/assets/slide-01.png",
                    }
                ]
            }
        ),
        encoding="utf-8",
    )
    runtime = SimpleNamespace(
        state={
            "thread_data": {
                "outputs_path": str(outputs),
                "workspace_path": str(tmp_path / "workspace"),
                "uploads_path": str(tmp_path / "uploads"),
            }
        }
    )

    summary = runner.run_batch(
        "/mnt/user-data/outputs/assets/slide-visuals.manifest.json",
        runtime,
    )

    assert summary["complete"] is True
    assert summary["source"] == "immutable_baseline_reuse"
    assert (outputs / "assets" / "slide-01.png").read_bytes() == content
    assert summary["items"][0]["reused_asset_hash"] == _sha(content)


def test_candidate_snapshot_reconciles_lost_create_response() -> None:
    store = _ScriptedSnapshotStore("raise_after_commit")
    content = b"exact candidate"
    path = "artifacts/test/candidate.pptx"

    DurableCandidateDq1Publisher._persist_snapshot(
        store=store,
        object_path=path,
        pptx_bytes=content,
    )

    assert store.snapshot_create_count == 1
    assert store.objects[path] == content


def test_candidate_snapshot_retries_once_after_precommit_failure() -> None:
    store = _ScriptedSnapshotStore("raise_before_commit")
    content = b"exact candidate"
    path = "artifacts/test/candidate.pptx"

    DurableCandidateDq1Publisher._persist_snapshot(
        store=store,
        object_path=path,
        pptx_bytes=content,
    )

    assert store.snapshot_create_count == 2
    assert store.objects[path] == content


def test_candidate_snapshot_reconciles_lost_retry_response() -> None:
    store = _ScriptedSnapshotStore(
        "raise_before_commit",
        "raise_after_commit",
    )
    content = b"exact candidate"
    path = "artifacts/test/candidate.pptx"

    DurableCandidateDq1Publisher._persist_snapshot(
        store=store,
        object_path=path,
        pptx_bytes=content,
    )

    assert store.snapshot_create_count == 2
    assert store.objects[path] == content


def test_candidate_snapshot_rejects_conflicting_existing_bytes() -> None:
    store = _ScriptedSnapshotStore()
    path = "artifacts/test/candidate.pptx"
    store.objects[path] = b"other candidate"

    with pytest.raises(DeckCandidateCompilationError, match="publication_failed"):
        DurableCandidateDq1Publisher._persist_snapshot(
            store=store,
            object_path=path,
            pptx_bytes=b"exact candidate",
        )

    assert store.snapshot_create_count == 1
    assert store.objects[path] == b"other candidate"


def test_candidate_snapshot_does_not_retry_ambiguous_conflict() -> None:
    path = "artifacts/test/candidate.pptx"

    class _AmbiguousConflictStore(_ScriptedSnapshotStore):
        def create_if_absent(self, object_path, content, *, content_type):
            self.snapshot_create_count += 1
            self.objects[object_path] = b"other candidate"
            raise RuntimeError("synthetic response loss")

    store = _AmbiguousConflictStore()
    with pytest.raises(DeckCandidateCompilationError, match="publication_failed"):
        DurableCandidateDq1Publisher._persist_snapshot(
            store=store,
            object_path=path,
            pptx_bytes=b"exact candidate",
        )

    assert store.snapshot_create_count == 1
    assert store.objects[path] == b"other candidate"


def test_candidate_snapshot_fails_closed_after_one_exhausted_retry() -> None:
    store = _ScriptedSnapshotStore(
        "raise_before_commit",
        "raise_before_commit",
    )
    path = "artifacts/test/candidate.pptx"

    with pytest.raises(DeckCandidateCompilationError, match="publication_failed"):
        DurableCandidateDq1Publisher._persist_snapshot(
            store=store,
            object_path=path,
            pptx_bytes=b"exact candidate",
        )

    assert store.snapshot_create_count == 2
    assert path not in store.objects


def test_candidate_snapshot_logs_no_storage_exception_details(
    caplog: pytest.LogCaptureFixture,
) -> None:
    private_sentinel = "private-storage-error-sentinel"
    private_path = "artifacts/private-user/private-thread/candidate.pptx"

    class _PrivateFailureStore(_MemoryPublicationStore):
        def create_if_absent(self, object_path, content, *, content_type):
            raise RuntimeError(private_sentinel)

    with caplog.at_level("INFO", logger=candidate_module.__name__):
        with pytest.raises(DeckCandidateCompilationError, match="publication_failed"):
            DurableCandidateDq1Publisher._persist_snapshot(
                store=_PrivateFailureStore(),
                object_path=private_path,
                pptx_bytes=b"exact candidate",
            )

    messages = "\n".join(record.getMessage() for record in caplog.records)
    assert messages == (
        "DQ2 candidate publication failed "
        "stage=immutable_snapshot code=retry_create_failed"
    )
    assert private_sentinel not in messages
    assert private_path not in messages
    assert all(record.exc_info is None for record in caplog.records)


def test_candidate_snapshot_does_not_retry_invalid_store_outcome() -> None:
    store = _ScriptedSnapshotStore("invalid_outcome")
    path = "artifacts/test/candidate.pptx"

    with pytest.raises(DeckCandidateCompilationError, match="publication_failed"):
        DurableCandidateDq1Publisher._persist_snapshot(
            store=store,
            object_path=path,
            pptx_bytes=b"exact candidate",
        )

    assert store.snapshot_create_count == 1
    assert path not in store.objects


@pytest.mark.parametrize("inject_store", [False, True])
def test_candidate_publisher_selects_producer_protocol(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
    inject_store: bool,
) -> None:
    store = _MemoryPublicationStore()
    monkeypatch.setattr(
        candidate_module,
        "SupabaseImmutableObjectStore",
        lambda: store,
    )
    captured: dict[str, Any] = {}
    quality_run_id = "quality_" + "9" * 64
    receipt = DeckQualityProducerBundleReceipt(
        quality_run_id=quality_run_id,
        bundle_object_path=deck_quality_producer_bundle_path(quality_run_id),
        bundle_hash="8" * 64,
        bundle_size_bytes=100,
    )

    def persist(**kwargs: Any) -> DeckQualityProducerBundleReceipt:
        captured.update(kwargs)
        return receipt

    monkeypatch.setattr(
        candidate_module,
        "persist_deck_quality_producer_bundle",
        persist,
    )
    pptx_bytes = b"candidate bytes"
    artifact_hash = _sha(pptx_bytes)
    virtual_path = "/mnt/user-data/outputs/candidate.pptx"
    snapshot_path = deck_quality_immutable_artifact_snapshot_path(
        user_id=USER_ID,
        thread_id=THREAD_ID,
        build_id=BUILD_ID,
        logical_artifact_id="logical_deck_01",
        artifact_version_id=CANDIDATE_ARTIFACT_ID,
        artifact_sha256=artifact_hash,
        artifact_virtual_path=virtual_path,
    )
    prepared = PreparedDeckQualityPublication.model_construct(
        outputs_root=tmp_path,
        artifact_virtual_path=virtual_path,
        artifact_storage_object_path=snapshot_path,
        artifact_sha256=artifact_hash,
        logical_artifact_id="logical_deck_01",
        artifact_version_id=CANDIDATE_ARTIFACT_ID,
        manifest_revision=2,
        build_id=BUILD_ID,
        user_id=USER_ID,
        thread_id=THREAD_ID,
    )
    source_pack = DeckQualitySourcePack.model_construct(
        artifact_sha256=artifact_hash,
        artifact_version_id=CANDIDATE_ARTIFACT_ID,
        manifest_revision=2,
    )
    source_pack_bytes = canonical_json_bytes(source_pack)
    publisher = (
        DurableCandidateDq1Publisher(store_factory=lambda: store)
        if inject_store
        else DurableCandidateDq1Publisher()
    )

    result = publisher._publish_sync(
        prepared=prepared,
        instrument=_instrument(),
        pptx_bytes=pptx_bytes,
        source_pack=source_pack,
        source_pack_bytes=source_pack_bytes,
    )

    assert result == receipt
    assert store.objects[snapshot_path] == pptx_bytes
    if inject_store:
        assert captured["object_store"] is store
    else:
        assert "object_store" not in captured


def test_durable_candidate_publisher_uploads_snapshot_and_replays_bundle(
    tmp_path: Path,
) -> None:
    outputs = tmp_path / "outputs"
    deck_root = outputs / "deck_build"
    deck_root.mkdir(parents=True)
    for filename, payload in (
        ("creative_plan.json", {"image_strategy": "diagram_native"}),
        ("design_plan.json", {"style_lane": "technical"}),
        ("build.json", {"build_id": BUILD_ID, "slide_count": 1}),
    ):
        (deck_root / filename).write_text(json.dumps(payload), encoding="utf-8")
    pptx_bytes = _pptx()
    artifact_hash = _sha(pptx_bytes)
    virtual_path = f"/mnt/user-data/outputs/.builder/builds/{BUILD_ID}/artifacts/{CANDIDATE_ARTIFACT_ID}/candidate.pptx"
    snapshot_path = deck_quality_immutable_artifact_snapshot_path(
        user_id=USER_ID,
        thread_id=THREAD_ID,
        build_id=BUILD_ID,
        logical_artifact_id="logical_deck_01",
        artifact_version_id=CANDIDATE_ARTIFACT_ID,
        artifact_sha256=artifact_hash,
        artifact_virtual_path=virtual_path,
    )
    prepared = PreparedDeckQualityPublication(
        outputs_root=outputs,
        artifact_virtual_path=virtual_path,
        artifact_storage_object_path=snapshot_path,
        artifact_sha256=artifact_hash,
        artifact_id=CANDIDATE_ARTIFACT_ID,
        logical_artifact_id="logical_deck_01",
        artifact_version_id=CANDIDATE_ARTIFACT_ID,
        manifest_revision=2,
        build_id=BUILD_ID,
        user_id=USER_ID,
        thread_id=THREAD_ID,
        task_id="task-01",
        builder_run_id=OPERATION_ID,
        parent_builder_trace_id="trace-root-01",
        task_brief="Build a concise PSI control-loop deck.",
        mechanical_gate_results={"passed": True},
        source_retention_report={"passed": True},
        native_contrast_report={"passed": True},
        native_mechanical_report={
            "render_success": True,
            "lint_fix_success": True,
            "lint_residue_count": 0,
        },
        native_editability_score=0.9,
        missing_expected_visual_count=0,
    )
    store = _MemoryPublicationStore()
    publisher = DurableCandidateDq1Publisher(store_factory=lambda: store)
    instrument = _instrument()
    source_pack, source_pack_bytes = capture_deck_quality_source_pack(
        prepared=prepared,
        instrument=instrument,
    )

    first = asyncio.run(
        publisher.publish(
            prepared=prepared,
            instrument=instrument,
            pptx_bytes=pptx_bytes,
            source_pack=source_pack,
            source_pack_bytes=source_pack_bytes,
        )
    )
    replay = asyncio.run(
        publisher.publish(
            prepared=prepared,
            instrument=instrument,
            pptx_bytes=pptx_bytes,
            source_pack=source_pack,
            source_pack_bytes=source_pack_bytes,
        )
    )

    assert first == replay
    assert store.objects[snapshot_path] == pptx_bytes
    assert first.quality_run_id.startswith("quality_")


@pytest.mark.anyio
async def test_durable_candidate_publisher_worker_is_not_abandoned_on_cancellation(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    publisher = DurableCandidateDq1Publisher(store_factory=lambda: object())
    entered = threading.Event()
    release = threading.Event()
    child_done = anyio.Event()
    cancel_scope: list[anyio.CancelScope] = []
    worker_threads: list[int] = []
    durable_writes: list[str] = []
    receipt = DeckQualityProducerBundleReceipt(
        quality_run_id="quality_" + "a" * 64,
        bundle_object_path=deck_quality_producer_bundle_path("quality_" + "a" * 64),
        bundle_hash="b" * 64,
        bundle_size_bytes=100,
    )

    def blocking_publish(**_kwargs: Any) -> DeckQualityProducerBundleReceipt:
        worker_threads.append(threading.get_ident())
        entered.set()
        assert release.wait(timeout=2)
        durable_writes.append("committed")
        return receipt

    monkeypatch.setattr(publisher, "_publish_sync", blocking_publish)

    async def publish() -> None:
        with anyio.CancelScope() as scope:
            cancel_scope.append(scope)
            try:
                await publisher.publish(
                    prepared=object(),  # type: ignore[arg-type]
                    instrument=object(),  # type: ignore[arg-type]
                    pptx_bytes=b"candidate",
                    source_pack=object(),  # type: ignore[arg-type]
                    source_pack_bytes=b"source-pack",
                )
            finally:
                child_done.set()

    async with anyio.create_task_group() as task_group:
        task_group.start_soon(publish)
        while not entered.is_set():
            await anyio.sleep(0)
        cancel_scope[0].cancel()
        await anyio.sleep(0.05)
        assert not child_done.is_set()
        assert durable_writes == []
        release.set()
        await child_done.wait()

    assert worker_threads == [worker_threads[0]]
    assert worker_threads[0] != threading.get_ident()
    assert durable_writes == ["committed"]
    await anyio.sleep(0.05)
    assert durable_writes == ["committed"]
