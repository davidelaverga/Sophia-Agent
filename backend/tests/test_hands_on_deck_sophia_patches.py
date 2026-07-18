from __future__ import annotations

import importlib.util
import json
import sys
from pathlib import Path
from types import SimpleNamespace

import pytest
from PIL import ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt

from deerflow.sophia.deck_build.mechanical_gates import _compiled_typography_issues
from deerflow.sophia.deck_native import DeckNativeService

PROJECT_ROOT = Path(__file__).resolve().parents[2]
HTML2PATCH_PATH = (
    PROJECT_ROOT
    / "third_party/hands_on_deck/skills/hands-on-deck/scripts/html2patch.py"
)
INVENTORY_PATH = (
    PROJECT_ROOT
    / "third_party/hands_on_deck/skills/hands-on-deck/scripts/inventory.py"
)


def _html2patch_module():
    spec = importlib.util.spec_from_file_location("sophia_html2patch", HTML2PATCH_PATH)
    assert spec is not None and spec.loader is not None
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def _inventory_module():
    spec = importlib.util.spec_from_file_location("sophia_hands_on_deck_inventory", INVENTORY_PATH)
    assert spec is not None and spec.loader is not None
    module = importlib.util.module_from_spec(spec)
    sys.modules[spec.name] = module
    spec.loader.exec_module(module)
    return module


def test_inventory_measurement_converts_points_to_pixels_and_preserves_weight(monkeypatch) -> None:
    module = _inventory_module()
    captured: list[tuple[str, int, bool, bool]] = []

    def fake_font(font_name: str, size: int, *, bold: bool = False, italic: bool = False):
        captured.append((font_name, size, bold, italic))
        return ImageFont.load_default(size=size)

    monkeypatch.setattr(module, "load_measure_font", fake_font)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(0.6))
    shape.text_frame.margin_left = 0
    shape.text_frame.margin_right = 0
    shape.text_frame.margin_top = 0
    shape.text_frame.margin_bottom = 0
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.text = "Habitat is fragmented, not absent"
    run = paragraph.runs[0]
    run.font.name = "Georgia"
    run.font.size = Pt(37.5)
    run.font.bold = True

    module.ShapeData(shape, slide=slide)

    assert captured == [("Georgia", 50, True, False)]


def test_inventory_fit_scale_keeps_largest_fitting_size() -> None:
    module = _inventory_module()
    shape_data = object.__new__(module.ShapeData)
    shape_data._frame_overflow_inches = lambda *, font_scale=1.0: (  # type: ignore[method-assign]
        None if font_scale <= 0.84 else 0.6
    )

    assert shape_data.font_scale_to_fit(0.6) == 0.84


def test_inventory_treats_required_cambria_heading_font_as_serif() -> None:
    module = _inventory_module()

    assert "cambria" in module.SERIF_HINTS
    fallbacks = module._fallback_font_paths(
        serif=True,
        bold=True,
        italic=False,
        font_name="Cambria",
    )
    assert fallbacks[0].endswith("/Caladea-Bold.ttf")


def test_inventory_measures_later_mixed_runs_with_their_own_font_metrics(monkeypatch) -> None:
    module = _inventory_module()
    captured: list[tuple[str, int, bool, bool]] = []

    def fake_font(font_name: str, size: int, *, bold: bool = False, italic: bool = False):
        captured.append((font_name, size, bold, italic))
        return ImageFont.load_default(size=size)

    monkeypatch.setattr(module, "load_measure_font", fake_font)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.6))
    shape.text_frame.margin_left = 0
    shape.text_frame.margin_right = 0
    shape.text_frame.margin_top = 0
    shape.text_frame.margin_bottom = 0
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.clear()
    lead = paragraph.add_run()
    lead.text = "Signal "
    lead.font.name = "Arial"
    lead.font.size = Pt(12)
    emphasis = paragraph.add_run()
    emphasis.text = "critical corridor failure"
    emphasis.font.name = "Cambria"
    emphasis.font.size = Pt(40)
    emphasis.font.bold = True

    shape_data = module.ShapeData(shape, slide=slide)

    assert ("Arial", 16, False, False) in captured
    assert ("Cambria", 53, True, False) in captured
    assert shape_data.frame_overflow_bottom is not None


def test_one_source_element_maps_deterministically_to_box_and_text_shapes(tmp_path: Path) -> None:
    module = _html2patch_module()
    source = {
        "sourceId": "system-title",
        "sourceRole": "title",
        "sourceRequired": True,
    }
    extract = {
        "items": [
            {
                **source,
                "type": "box",
                "box": {"x": 100, "y": 80, "w": 900, "h": 120},
                "rotation": 0,
                "fill": "FFFFFF",
                "fillAlpha": 1,
                "gradient": None,
                "border": {"color": "111827", "w": 2, "dashed": False},
                "partialBorders": [],
                "radiusPx": 0,
            },
            {
                **source,
                "type": "text",
                "box": {"x": 120, "y": 90, "w": 860, "h": 90},
                "rotation": 0,
                "vanchor": "middle",
                "paragraphs": [
                    {
                        "lineHeightPx": 42,
                        "runs": [
                            {
                                "text": "System title",
                                "style": {
                                    "font": "Arial",
                                    "sizePx": 36,
                                    "color": "111827",
                                    "alpha": 1,
                                    "bold": True,
                                    "italic": False,
                                    "underline": False,
                                    "link": None,
                                },
                            }
                        ],
                    }
                ],
                "meta": {
                    "align": "left",
                    "lineHeightPx": 42,
                    "fontSizePx": 36,
                    "padding": [0, 0, 0, 0],
                },
            },
        ]
    }
    source_map: dict = {"schema_version": "sophia-deck-source-map/v1", "slides": {}}

    operations = module.compile_page(
        extract,
        0,
        tmp_path / "slide.html",
        tmp_path,
        "s1",
        [],
        source_map,
    )

    names = source_map["slides"]["slide:1"]["elements"]["system-title"]["shape_names"]
    assert names == ["s1-system-title-box-1", "s1-system-title-text-1"]
    assert [operation["name"] for operation in operations] == names


def test_descendant_shapes_map_to_direct_and_all_semantic_ancestors(tmp_path: Path) -> None:
    module = _html2patch_module()

    def box(source_id: str, x: int) -> dict:
        direct = {
            "sourceId": source_id,
            "sourceRole": "detail",
            "sourceRequired": True,
        }
        return {
            **direct,
            "sourceRefs": [
                direct,
                {
                    "sourceId": "cluster",
                    "sourceRole": "diagram",
                    "sourceRequired": True,
                },
                {
                    "sourceId": "system",
                    "sourceRole": "architecture",
                    "sourceRequired": True,
                },
            ],
            "type": "box",
            "box": {"x": x, "y": 160, "w": 300, "h": 180},
            "rotation": 0,
            "fill": "E2E8F0",
            "fillAlpha": 1,
            "gradient": None,
            "border": None,
            "partialBorders": [],
            "radiusPx": 0,
        }

    extract = {"items": [box("node-a", 100), box("node-b", 440)]}
    source_map: dict = {"schema_version": "sophia-deck-source-map/v1", "slides": {}}

    operations = module.compile_page(
        extract,
        0,
        tmp_path / "slide.html",
        tmp_path,
        "s1",
        [],
        source_map,
    )

    operation_names = [operation["name"] for operation in operations]
    assert operation_names == ["s1-node-a-box-1", "s1-node-b-box-1"]
    elements = source_map["slides"]["slide:1"]["elements"]
    assert elements["node-a"]["shape_names"] == ["s1-node-a-box-1"]
    assert elements["node-b"]["shape_names"] == ["s1-node-b-box-1"]
    assert elements["cluster"]["shape_names"] == operation_names
    assert elements["system"]["shape_names"] == operation_names
    assert elements["cluster"]["source_required"] is True
    assert elements["system"]["source_role"] == "architecture"


def test_table_item_deduplicates_descendant_refs_onto_one_native_table(tmp_path: Path) -> None:
    module = _html2patch_module()
    panel = {
        "sourceId": "compare-table-panel",
        "sourceRole": "comparison",
        "sourceRequired": False,
    }
    action = {
        "sourceId": "row-action-psi",
        "sourceRole": "content",
        "sourceRequired": True,
    }
    style = {
        "bold": False,
        "italic": False,
        "sizePx": 24,
        "font": "Arial",
        "color": "111827",
        "fill": None,
        "align": "left",
    }
    extract = {
        "items": [
            {
                **panel,
                "sourceRefs": [panel, action, panel, action],
                "type": "table",
                "box": {"x": 96, "y": 160, "w": 1200, "h": 180},
                "rows": [["Action selection"]],
                "cellStyles": [[style]],
                "colWidths": [1200],
                "fontSizePx": 24,
            }
        ]
    }
    source_map: dict = {"schema_version": "sophia-deck-source-map/v1", "slides": {}}

    operations = module.compile_page(
        extract,
        0,
        tmp_path / "slide.html",
        tmp_path,
        "h2p-1",
        [],
        source_map,
    )

    table = next(operation for operation in operations if operation["op"] == "add-table")
    elements = source_map["slides"]["slide:1"]["elements"]
    assert elements["compare-table-panel"]["shape_names"] == [table["name"]]
    assert elements["row-action-psi"]["shape_names"] == [table["name"]]
    assert elements["row-action-psi"]["source_required"] is True


def _wide_base_deck(path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(20)
    presentation.slide_height = Inches(11.25)
    presentation.save(path)


def _html_to_patch_or_skip(tmp_path: Path, source: str, *, stem: str):
    if importlib.util.find_spec("playwright") is None:
        pytest.skip("Python Playwright is not installed in this backend env")
    base = tmp_path / f"{stem}-base.pptx"
    html = tmp_path / f"{stem}.html"
    patch = tmp_path / f"{stem}.patch.json"
    _wide_base_deck(base)
    html.write_text(source, encoding="utf-8")
    result = DeckNativeService().html_to_patch(
        html_paths=[str(html)],
        base_deck_path=str(base),
        output_patch_path=str(patch),
    )
    if not result.success and any(
        token in error.lower()
        for error in result.errors
        for token in ("browser", "chromium", "playwright")
    ):
        pytest.skip("Playwright is installed but its browser runtime is unavailable")
    assert result.success is True, result.errors
    return (
        json.loads(patch.read_text(encoding="utf-8")),
        json.loads(Path(result.source_map_path or "").read_text(encoding="utf-8")),
    )


def test_table_extraction_maps_only_visible_represented_descendant_ids(tmp_path: Path) -> None:
    patch, source_map = _html_to_patch_or_skip(
        tmp_path,
        """<!doctype html><html><body style="margin:0;width:1920px;height:1080px;background:#FFFFFF;color:#111827">
        <table data-deck-id="compare-table-panel" data-deck-role="comparison"
               style="position:absolute;left:96px;top:120px;width:1500px;border-collapse:collapse;font:20px Arial">
          <tr><th style="height:80px">Dimension</th><th>PSI agent</th></tr>
          <tr><td>Action</td><td data-deck-id="row-action-psi" data-deck-role="content" data-deck-required="true">Weighted action</td></tr>
          <tr><td>Arbitration</td><td data-deck-id="row-arb-psi" data-deck-role="content" data-deck-required="true">
            <span data-deck-id="visible-inline" data-deck-role="label">Motive wins</span>
            <span data-deck-id="hidden-inline" data-deck-role="label" style="display:none">HIDDEN COPY</span>
          </td></tr>
          <tr><td>Debug</td><td data-deck-id="row-debug-psi" data-deck-role="content" data-deck-required="true">
            Inspectable
            <span data-deck-id="nonrepresented-chip"
                  style="display:inline-block;width:12px;height:12px;background:#EF4444"></span>
          </td></tr>
        </table>
        </body></html>""",
        stem="semantic-table",
    )

    table = next(operation for operation in patch["ops"] if operation["op"] == "add-table")
    assert "HIDDEN COPY" not in json.dumps(table["rows"])
    cell_updates = [
        operation
        for operation in patch["ops"]
        if operation.get("op") == "set-text" and operation.get("shape") == table["name"]
    ]
    assert {
        float(operation["text"][0]["font_size"])
        for operation in cell_updates
    } == {18.0}
    assert len(cell_updates) == sum(len(row) for row in table["rows"])
    elements = source_map["slides"]["slide:1"]["elements"]
    expected = {
        "compare-table-panel",
        "row-action-psi",
        "row-arb-psi",
        "row-debug-psi",
        "visible-inline",
    }
    assert expected <= set(elements)
    assert "hidden-inline" not in elements
    assert "nonrepresented-chip" not in elements
    for source_id in expected:
        assert elements[source_id]["shape_names"] == [table["name"]]

    # The source map makes the single native table shape required, so every
    # represented cell must satisfy the required 24px/18pt whole-shape floor.
    base = tmp_path / "semantic-table-apply-base.pptx"
    patch_path = tmp_path / "semantic-table-apply.patch.json"
    output = tmp_path / "semantic-table-applied.pptx"
    _wide_base_deck(base)
    patch_path.write_text(json.dumps(patch), encoding="utf-8")
    applied = DeckNativeService().apply_patch(
        base_deck_path=str(base),
        patch_path=str(patch_path),
        output_path=str(output),
        fix=False,
    )
    assert applied.success is True, applied.errors
    issues = _compiled_typography_issues(
        SimpleNamespace(source_element_map=source_map, slides=[]),
        native_pptx_path=output,
    )
    assert issues == []


def test_extraction_omits_transparent_and_fully_clipped_text_but_keeps_visibility_override(
    tmp_path: Path,
) -> None:
    patch, source_map = _html_to_patch_or_skip(
        tmp_path,
        """<!doctype html><html><body style="margin:0;width:1920px;height:1080px;background:#FFFFFF;color:#111827;font:24px Arial">
        <p data-deck-id="ordinary-copy" data-deck-role="narrative" data-deck-required="true"
           style="position:absolute;left:96px;top:100px;width:1200px;margin:0">
          Visible ordinary copy
          <span data-deck-id="transparent-inline" style="color:transparent">TRANSPARENT SECRET</span>
        </p>
        <div style="position:absolute;left:96px;top:220px;width:800px;height:60px;visibility:hidden">
          <span data-deck-id="restored-visible" data-deck-role="label"
                style="visibility:visible">Restored visible child</span>
        </div>
        <div style="position:absolute;left:96px;top:320px;width:60px;height:30px;overflow:hidden">
          <span data-deck-id="fully-clipped" data-deck-role="label"
                style="position:absolute;left:200px;top:0;white-space:nowrap">CLIPPED SECRET</span>
        </div>
        <table data-deck-id="visible-table"
               style="position:absolute;left:96px;top:430px;width:1200px;border-collapse:collapse;font:24px Arial">
          <tr><td data-deck-id="visible-cell">Visible table copy
            <span data-deck-id="transparent-table-inline" style="color:rgba(1,2,3,0)">TABLE SECRET</span>
          </td></tr>
        </table>
        </body></html>""",
        stem="hidden-text",
    )

    serialized = json.dumps(patch)
    assert "Visible ordinary copy" in serialized
    assert "Restored visible child" in serialized
    assert "Visible table copy" in serialized
    assert "TRANSPARENT SECRET" not in serialized
    assert "CLIPPED SECRET" not in serialized
    assert "TABLE SECRET" not in serialized

    elements = source_map["slides"]["slide:1"]["elements"]
    assert {"ordinary-copy", "restored-visible", "visible-table", "visible-cell"} <= set(elements)
    assert "transparent-inline" not in elements
    assert "fully-clipped" not in elements
    assert "transparent-table-inline" not in elements


def _text_operation(patch: dict, source_id: str) -> dict:
    return next(
        operation
        for operation in patch["ops"]
        if operation.get("op") == "add-shape"
        and f"-{source_id}-text-" in str(operation.get("name") or "")
    )


def _single_run_font_size(operation: dict) -> float:
    paragraph = operation["text"][0]
    if "font_size" in paragraph:
        return float(paragraph["font_size"])
    return float(paragraph["runs"][0]["font_size"])


def test_browser_typography_floor_reflows_before_geometry_and_records_provenance(tmp_path: Path) -> None:
    patch, source_map = _html_to_patch_or_skip(
        tmp_path,
        """<!doctype html><html><body style="margin:0;width:1920px;height:1080px;background:#FFFFFF;color:#111827;font-family:Arial">
        <p data-deck-id="required-copy" data-deck-role="narrative" data-deck-required="true"
           style="position:absolute;left:96px;top:100px;width:1200px;margin:0;font-size:16px;line-height:12px">Required narrative</p>
        <p data-deck-id="optional-label" data-deck-role="label"
           style="position:absolute;left:96px;top:250px;width:800px;margin:0;font-size:12px;line-height:10px">Optional label</p>
        <p data-deck-id="role-required" data-deck-role="title"
           style="position:absolute;left:96px;top:400px;width:1200px;margin:0;font-size:18px;line-height:18px">Required by role</p>
        <section data-deck-id="required-wrapper" data-deck-role="content" data-deck-required="true"
                 style="position:absolute;left:96px;top:550px;width:1200px;height:100px">
          <span data-deck-id="nested-copy" data-deck-role="label"
                style="display:block;font-size:10px;line-height:9px">Nested required copy</span>
        </section>
        <p data-deck-id="hidden-copy" data-deck-role="narrative" data-deck-required="true"
           style="display:none;font-size:8px;line-height:8px">Hidden copy</p>
        </body></html>""",
        stem="typography-floor",
    )

    required = _text_operation(patch, "required-copy")
    optional = _text_operation(patch, "optional-label")
    role_required = _text_operation(patch, "role-required")
    nested = _text_operation(patch, "nested-copy")
    assert _single_run_font_size(required) == 18.0
    assert _single_run_font_size(optional) == 15.0
    assert _single_run_font_size(role_required) == 18.0
    assert _single_run_font_size(nested) == 18.0
    assert required["text"][0]["line_spacing"] == 21.6
    assert optional["text"][0]["line_spacing"] == 18.0
    assert required["size"][1] == pytest.approx(28.8 / 96, abs=0.001)
    assert optional["size"][1] == pytest.approx(24 / 96, abs=0.001)

    slide_map = source_map["slides"]["slide:1"]
    normalizations = slide_map["typography_normalizations"]
    by_source_id = {item["sourceId"]: item for item in normalizations}
    assert slide_map["typography_normalization_count"] == 4
    assert slide_map["typography_normalizations_truncated"] is False
    assert "hidden-copy" not in by_source_id
    assert by_source_id["required-copy"] == {
        "selector": "slide:1",
        "sourceId": "required-copy",
        "sourceRole": "narrative",
        "elementTag": "p",
        "required": True,
        "minimumPx": 24,
        "oldFontPx": 16,
        "newFontPx": 24,
        "oldLineHeightPx": 12,
        "newLineHeightPx": 28.8,
        "fontChanged": True,
        "lineHeightChanged": True,
    }
    assert by_source_id["optional-label"]["minimumPx"] == 20
    assert by_source_id["optional-label"]["newFontPx"] == 20
    assert by_source_id["role-required"]["required"] is True
    assert by_source_id["nested-copy"]["required"] is True
