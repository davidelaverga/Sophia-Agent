from __future__ import annotations

import json
from copy import deepcopy
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pydantic import ValidationError
from test_deck_build_service import _creative_plan, _slides

from deerflow.agents.sophia_agent.middlewares.builder_artifact import _pptx_contains_visual_evidence
from deerflow.agents.sophia_agent.middlewares.builder_task import BuilderTaskMiddleware
from deerflow.sophia.deck_build.compiler_capabilities import (
    SUPPORTED_CSS_FEATURES,
    lossy_css_in_html,
    rejected_css_in_html,
    unsupported_tags_in_html,
)
from deerflow.sophia.deck_build.models import DeckSlideCompositionPlan, DeckSlideSpec
from deerflow.sophia.deck_build.native_contrast import evaluate_native_contrast
from deerflow.sophia.deck_build.source_retention import evaluate_source_retention, retention_summary
from deerflow.sophia.deck_build.tool_contract import (
    DeckCreativePlanInput,
    PrepareDeckBuildInput,
    prepare_deck_build_validation_summary,
)


def test_prepare_contract_normalizes_exactly_one_json_layer() -> None:
    model = PrepareDeckBuildInput.model_validate(
        {
            "deck_title": "Technical Deck",
            "slides": json.dumps(_slides()),
            "output_path": "/mnt/user-data/outputs/deck.pptx",
            "creative_plan": json.dumps(_creative_plan()),
        }
    )

    assert len(model.slides) == 3
    assert model.creative_plan.skill_refs[0] == "hands-on-deck/designing-slides"


def test_prepare_contract_rejects_double_encoded_json() -> None:
    with pytest.raises(ValidationError) as exc:
        PrepareDeckBuildInput.model_validate(
            {
                "deck_title": "Technical Deck",
                "slides": json.dumps(json.dumps(_slides())),
                "output_path": "/mnt/user-data/outputs/deck.pptx",
                "creative_plan": json.dumps(_creative_plan()),
            }
        )

    assert exc.value.errors()[0]["loc"] == ("slides",)


def test_prepare_contract_rejects_mixed_compact_and_legacy_sources() -> None:
    slides = _slides()
    slides[0]["html_body"] = "<section>compact</section>"
    with pytest.raises(ValidationError) as exc:
        PrepareDeckBuildInput.model_validate(
            {
                "deck_title": "Technical Deck",
                "slides": slides,
                "output_path": "/mnt/user-data/outputs/deck.pptx",
                "creative_plan": _creative_plan(),
                "deck_stylesheet": "main { background: #101828; }",
            }
        )

    assert "slides.0" in str(exc.value)


def _compact_slide() -> dict:
    slide = deepcopy(_slides()[0])
    slide.pop("html_source")
    slide["html_body"] = '<h1 data-deck-id="title" data-deck-role="title">Compact</h1>'
    return slide


def test_prepare_contract_rejects_document_tags_and_oversized_fragments() -> None:
    slide = _compact_slide()
    slide["html_body"] = "<body>not a fragment</body>"
    with pytest.raises(ValidationError) as tag_error:
        PrepareDeckBuildInput.model_validate(
            {
                "deck_title": "Technical Deck",
                "slides": [slide],
                "output_path": "/mnt/user-data/outputs/deck.pptx",
                "creative_plan": _creative_plan(),
                "deck_stylesheet": "main { background: #101828; }",
            }
        )
    assert tag_error.value.errors()[0]["loc"] == ("slides", 0)

    slide["html_body"] = "x" * (16 * 1024 + 1)
    with pytest.raises(ValidationError) as size_error:
        PrepareDeckBuildInput.model_validate(
            {
                "deck_title": "Technical Deck",
                "slides": [slide],
                "output_path": "/mnt/user-data/outputs/deck.pptx",
                "creative_plan": _creative_plan(),
                "deck_stylesheet": "main { background: #101828; }",
            }
        )
    assert size_error.value.errors()[0]["loc"] == ("slides", 0)


def test_prepare_contract_rejects_total_authoring_payload_over_128_kib() -> None:
    slides = []
    for index in range(8):
        slide = _compact_slide()
        slide["title"] = f"Compact {index}"
        slide["html_body"] = "x" * (16 * 1024)
        slides.append(slide)

    with pytest.raises(ValidationError, match="131072-byte limit"):
        PrepareDeckBuildInput.model_validate(
            {
                "deck_title": "Technical Deck",
                "slides": slides,
                "output_path": "/mnt/user-data/outputs/deck.pptx",
                "creative_plan": _creative_plan(),
                "deck_stylesheet": "main { background: #101828; }",
            }
        )


def test_compiler_capabilities_reject_svg_and_lossy_semantic_css() -> None:
    source = """<html><style>.x { opacity: .5; filter: blur(2px); transform: translateX(4px); }</style>
    <body><svg><path d="M0 0" /></svg></body></html>"""

    assert unsupported_tags_in_html(source) == ["path", "svg"]
    assert rejected_css_in_html(source) == ["filter", "transform"]
    assert lossy_css_in_html(source) == ["opacity"]
    assert {"grid", "rotate-transform", "solid-fill"}.issubset(SUPPORTED_CSS_FEATURES)


def test_compiler_capabilities_do_not_confuse_text_transform_with_transform() -> None:
    source = """<html><style>
    main { background: #101828; text-transform: uppercase; }
    </style><body><main style="width:1920px;height:1080px"></main></body></html>"""

    assert rejected_css_in_html(source) == []
    assert lossy_css_in_html(source) == []


def test_compiler_capabilities_recurse_into_nested_css_at_rules() -> None:
    source = """<html><style>
    @media all { .hidden { opacity: 0; } }
    @media screen { .shifted { transform: translateX(4px); } }
    @supports (display: grid) { @layer deck { .blurred { filter: blur(2px); } } }
    </style><body><main></main></body></html>"""

    assert lossy_css_in_html(source) == ["opacity"]
    assert rejected_css_in_html(source) == ["filter", "transform"]


def test_compact_v2_profile_is_required_in_model_schema_and_bounded() -> None:
    schema = PrepareDeckBuildInput.model_json_schema()
    assert "authoring_contract" in schema["required"]
    assert schema["properties"]["authoring_contract"]["const"] == "compact_model_html_v2"
    body = schema["$defs"]["DeckSlideInput"]["properties"]["html_body"]
    string_variant = next(item for item in body["anyOf"] if item.get("type") == "string")
    assert string_variant["maxLength"] == 3 * 1024

    slide = _compact_slide()
    slide["html_body"] = "x" * (3 * 1024 + 1)
    with pytest.raises(ValidationError, match="compact-v2 3072-byte limit"):
        PrepareDeckBuildInput.model_validate(
            {
                "deck_title": "Technical Deck",
                "slides": [slide],
                "output_path": "/mnt/user-data/outputs/deck.pptx",
                "creative_plan": _creative_plan(),
                "authoring_contract": "compact_model_html_v2",
                "deck_stylesheet": "main { background: #101828; }",
            }
        )

    # Queued/internal compact-v1 payloads retain the previous 16 KiB body limit.
    model = PrepareDeckBuildInput.model_validate(
        {
            "deck_title": "Technical Deck",
            "slides": [slide],
            "output_path": "/mnt/user-data/outputs/deck.pptx",
            "creative_plan": _creative_plan(),
            "authoring_contract": "compact_model_html_v1",
            "deck_stylesheet": "main { background: #101828; }",
        }
    )
    assert model.authoring_contract == "compact_model_html_v1"


def _compact_v2_args(*, body_sizes: list[int] | None = None) -> dict:
    slides = []
    for index, body_size in enumerate(body_sizes or [256], start=1):
        slide = _compact_slide()
        slide["title"] = f"Compact {index}"
        slide["html_body"] = "x" * body_size
        slides.append(slide)
    return {
        "deck_title": "Compact Limit Diagnostics",
        "slides": slides,
        "output_path": "/mnt/user-data/outputs/deck.pptx",
        "creative_plan": _creative_plan(),
        "authoring_contract": "compact_model_html_v2",
        "deck_stylesheet": "main { background: #101828; }",
    }


def test_prepare_validation_summary_enumerates_all_compact_v2_body_limits() -> None:
    args = _compact_v2_args(body_sizes=[1173, 3220, 2896, 3442, 2625])

    summary = prepare_deck_build_validation_summary(args)

    slide_two = "slides[1].html_body is 3220 bytes; compact-v2 limit is 3072 bytes"
    slide_four = "slides[3].html_body is 3442 bytes; compact-v2 limit is 3072 bytes"
    assert summary.count(slide_two) == 1
    assert summary.count(slide_four) == 1
    assert summary.index(slide_two) < summary.index(slide_four)
    assert "xxxxxxxx" not in summary
    assert len(summary) <= 1200


def test_prepare_validation_summary_matches_size_semantics_and_deduplicates() -> None:
    whitespace_args = _compact_v2_args()
    whitespace_args["deck_stylesheet"] = (" " * 200) + ("x" * 8000)
    assert prepare_deck_build_validation_summary(whitespace_args) == ""

    unicode_args = _compact_v2_args()
    unicode_args["slides"][0]["slide_css"] = "é" * 513
    unicode_summary = prepare_deck_build_validation_summary(unicode_args)
    assert unicode_summary == "slides[0].slide_css is 1026 bytes; compact-v2 limit is 1024 bytes"

    plan_args = _compact_v2_args()
    plan_args["creative_plan"]["story_arc"] = "x" * 13_000
    normalized_plan = DeckCreativePlanInput.model_validate(plan_args["creative_plan"]).model_dump(mode="json")
    plan_bytes = len(json.dumps(normalized_plan, separators=(",", ":"), ensure_ascii=False).encode("utf-8"))
    plan_summary = prepare_deck_build_validation_summary(plan_args)
    assert f"creative_plan is {plan_bytes} bytes; compact-v2 limit is 12288 bytes" in plan_summary

    total_args = _compact_v2_args()
    total_args["style_profile"] = {"padding": "x" * 50_000}
    raw_bytes = len(json.dumps(total_args, separators=(",", ":"), ensure_ascii=False).encode("utf-8"))
    total_summary = prepare_deck_build_validation_summary(total_args)
    total_target = f"prepare_deck_build arguments are {raw_bytes} bytes; compact-v2 limit is 49152 bytes"
    assert total_summary.count(total_target) == 1
    assert total_summary.count("prepare_deck_build arguments") == 1


def test_source_retention_reports_missing_required_ids() -> None:
    composition = DeckSlideCompositionPlan(
        selector="slide:1",
        slide_role="architecture",
        headline_intent="Show the system",
        layout_name="system-map",
        composition_rationale="Connect the system parts.",
        native_elements=["title", "diagram"],
        image_asset_ids=[],
        required_element_ids=["title", "diagram"],
        structural_fingerprint="title-over-wide-system-map",
    )
    slide = DeckSlideSpec(
        selector="slide:1",
        index=1,
        role="architecture",
        layout_kind="single_visual_focus",
        title="System",
        narrative="A native system diagram.",
        composition_plan=composition,
    )
    reports = evaluate_source_retention(
        slides=[slide],
        native_shape_inventory={"slide:1": {"shapes": [{"name": "s1-title-text"}]}},
        source_element_map={
            "slides": {
                "slide:1": {
                    "elements": {
                        "title": {"shape_names": ["s1-title-text"]},
                        "diagram": {"shape_names": ["s1-diagram-box"]},
                    }
                }
            }
        },
    )

    summary = retention_summary(reports)
    assert summary["passed"] is False
    assert summary["missing_required"] == [{"selector": "slide:1", "source_id": "diagram"}]


def _contrast_deck(path: Path, *, explicit: bool, foreground: RGBColor | None = None) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    shape.name = "s1-title-text"
    shape.text_frame.paragraphs[0].text = "Deterministic contrast"
    if explicit:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x0A, 0x0E, 0x14)
        run = shape.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(28)
        run.font.color.rgb = foreground or RGBColor(0xEE, 0xF4, 0xFB)
    presentation.save(path)


def test_native_contrast_passes_explicit_required_text_and_fails_indeterminate(tmp_path: Path) -> None:
    source_map = {
        "slides": {
            "slide:1": {
                "elements": {
                    "title": {
                        "source_required": True,
                        "shape_names": ["s1-title-text"],
                    }
                }
            }
        }
    }
    passing = tmp_path / "passing.pptx"
    indeterminate = tmp_path / "indeterminate.pptx"
    _contrast_deck(passing, explicit=True)
    _contrast_deck(indeterminate, explicit=False)

    assert evaluate_native_contrast(pptx_path=passing, source_element_map=source_map)["passed"] is True
    failed = evaluate_native_contrast(pptx_path=indeterminate, source_element_map=source_map)
    assert failed["passed"] is False
    assert failed["indeterminate_required_count"] == 1

    low_contrast = tmp_path / "low-contrast.pptx"
    _contrast_deck(low_contrast, explicit=True, foreground=RGBColor(0x11, 0x18, 0x27))
    low = evaluate_native_contrast(pptx_path=low_contrast, source_element_map=source_map)
    assert low["passed"] is False
    assert low["issues"][0]["reason"] == "contrast_below_threshold"


def test_native_shape_only_pptx_counts_as_visual_evidence(tmp_path: Path) -> None:
    path = tmp_path / "native-shape.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(4), Inches(2))
    presentation.save(path)

    assert _pptx_contains_visual_evidence(path) is True


def test_presentation_inventory_exposes_only_deck_scoped_design_skills(monkeypatch) -> None:
    class _Skill:
        def __init__(self, name: str) -> None:
            self.name = name
            self.description = f"{name} description"

        def get_container_file_path(self, base: str) -> str:
            return f"{base}/{self.name}/SKILL.md"

    names = [
        "hands-on-deck",
        "deck-impeccable",
        "deck-hallmark",
        "ppt-generation",
        "image-generation",
        "visual-design",
        "hallmark",
        "pdf-report",
    ]
    monkeypatch.setattr("deerflow.skills.load_skills", lambda **_kwargs: [_Skill(name) for name in names])

    block = BuilderTaskMiddleware._build_skills_inventory_block(presentation_design_mode=True)

    assert block is not None
    for expected in names[:5]:
        assert f"<name>{expected}</name>" in block
    for excluded in names[5:]:
        assert f"<name>{excluded}</name>" not in block
