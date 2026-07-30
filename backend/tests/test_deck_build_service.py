from __future__ import annotations

import hashlib
import json
import subprocess
from pathlib import Path
from types import SimpleNamespace

import pytest
from langchain_core.messages import ToolMessage
from langgraph.types import Command
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from deerflow.agents.sophia_agent.builder_tools import (
    assert_deck_tool_contract,
    build_builder_tools_for_task_type,
    deck_build_service_enabled,
)
from deerflow.agents.sophia_agent.middlewares import builder_artifact as builder_artifact_module
from deerflow.agents.sophia_agent.middlewares.builder_artifact import BuilderArtifactMiddleware
from deerflow.sandbox.tools import replace_virtual_path
from deerflow.sophia.deck_build import service as deck_service
from deerflow.sophia.deck_build.mechanical_gates import MechanicalGateIssue, MechanicalGateResult
from deerflow.sophia.deck_build.service import DeckBuildService
from deerflow.sophia.deck_build.storage import load_deck_build
from deerflow.sophia.deck_build.tool_contract import DeckCreativePlanInput, PrepareDeckBuildInput
from deerflow.sophia.deck_build.tracing import NATIVE_DECK_COMPILE_MODE
from deerflow.sophia.deck_native.models import (
    NativeDeckInspectResult,
    NativeDeckLintFixResult,
    NativeDeckPatchResult,
    NativeDeckPreflight,
    NativeDeckRenderResult,
)
from deerflow.sophia.tools.prepare_deck_build import prepare_deck_build

_OUTPUTS = "/mnt/user-data/outputs/"


def _runtime(outputs: Path, *, user_request: str = "Build a visual 3 slide deck") -> SimpleNamespace:
    outputs.mkdir(parents=True, exist_ok=True)
    workspace = outputs.parent / "workspace"
    uploads = outputs.parent / "uploads"
    workspace.mkdir(parents=True, exist_ok=True)
    uploads.mkdir(parents=True, exist_ok=True)
    return SimpleNamespace(
        state={
            "thread_id": "builder-thread",
            "parent_thread_id": "companion-thread",
            "user_id": "user-1",
            "task_id": "task-1",
            "builder_pptx_requested_slide_count": 3,
            "builder_artifact_target_path": f"{_OUTPUTS}deck.pptx",
            "delegation_context": {"request": user_request},
            "thread_data": {
                "outputs_path": str(outputs),
                "workspace_path": str(workspace),
                "uploads_path": str(uploads),
            },
        },
        context={"thread_id": "builder-thread"},
        config={},
    )


def test_deck_service_reserves_terminal_cleanup_from_shared_deadline(monkeypatch) -> None:
    runtime = SimpleNamespace(
        state={
            "builder_deadline_epoch_ms": 1_300_000,
            "builder_budget": {"terminal_reserve_seconds": 30},
        },
        context={},
        config={},
    )
    monkeypatch.setattr(deck_service.time, "time", lambda: 1_000.0)

    assert deck_service._service_deadline_epoch_ms(runtime) == 1_270_000
    assert deck_service._remaining_deadline_seconds(runtime) == 270


def test_deck_service_v2_body_pool_matches_typed_contract() -> None:
    deck = SimpleNamespace(
        deck_stylesheet="main{font-family:Calibri,Arial,sans-serif}",
        deck_authoring_contract="compact_model_html_v2",
    )

    deck_service._validate_v2_authoring_sizes(
        deck,
        [{"html_body": "x" * 6144}, {"html_body": "x" * 2048}],
    )

    with pytest.raises(deck_service.DeckBuildFailure, match="aggregate budget is 8192 bytes"):
        deck_service._validate_v2_authoring_sizes(
            deck,
            [{"html_body": "x" * 6144}, {"html_body": "x" * 2049}],
        )
    with pytest.raises(deck_service.DeckBuildFailure, match="hard 6144-byte limit"):
        deck_service._validate_v2_authoring_sizes(
            deck,
            [{"html_body": "x" * 6145}, *[{"html_body": "x"} for _ in range(4)]],
        )


def _slide_html(index: int, title: str, narrative: str, *, include_asset: bool = False) -> str:
    asset = f'<figure class="asset"><img src="../assets/slide-{index:02d}.png" alt="" /></figure>' if include_asset else ""
    return f"""<!doctype html>
<html><head><meta charset="utf-8"><style>
html, body {{ margin: 0; padding: 0; width: 1920px; height: 1080px; background: #0A0E14; }}
body {{ overflow: hidden; color: #EEF4FB; font-family: Aptos, Arial, sans-serif; }}
.canvas {{ position: relative; width: 1920px; height: 1080px; background: #0A0E14; }}
h1 {{ position: absolute; left: 120px; top: 90px; width: 1250px; margin: 0; font-size: 64px; line-height: 1.05; }}
.narrative {{ position: absolute; left: 120px; top: 830px; width: 1320px; font-size: 30px; line-height: 1.3; color: #A7B4C2; }}
.diagram {{ position: absolute; left: 120px; top: 280px; width: 1500px; height: 420px; border: 3px solid #38BDF8; background: #111827; }}
.node {{ position: absolute; top: 130px; width: 300px; height: 120px; border: 2px solid #38BDF8; }}
.node.a {{ left: 120px; }} .node.b {{ left: 600px; }} .node.c {{ left: 1080px; }}
.asset {{ position: absolute; inset: 0; margin: 0; }}
.asset img {{ width: 100%; height: 100%; object-fit: cover; }}
</style></head><body><main class="canvas">
{asset}
	<h1 data-deck-id="title-{index}" data-deck-role="title" data-deck-required="true">{title}</h1>
	<section class="diagram" data-deck-id="diagram-{index}" data-deck-role="diagram"><div class="node a"></div><div class="node b"></div><div class="node c"></div></section>
	<p class="narrative" data-deck-id="narrative-{index}" data-deck-role="narrative" data-deck-required="true">{narrative}</p>
</main></body></html>"""


def _slides(count: int = 3, *, include_asset: bool = True) -> list[dict]:
    roles = ["cover", "architecture", "closing"]
    layouts = ["cover_hero", "single_visual_focus", "closing_summary"]
    slides = []
    for index in range(1, count + 1):
        title = f"Slide {index} System Story"
        narrative = "A concise technical narrative explains the point with calm professional framing."
        slides.append(
            {
                "title": f"Slide {index} System Story",
                "narrative": narrative,
                "role": roles[index - 1],
                "layout_kind": layouts[index - 1],
                "visual_prompt": f"Professional technical visual metaphor for slide {index}",
                "speaker_notes": "Optional notes.",
                "html_source": _slide_html(index, title, narrative, include_asset=include_asset and index == 1),
            }
        )
    return slides


def _compact_slides(*, visible_eyebrow_on_slide: int | None = None) -> list[dict]:
    slides = _slides(include_asset=False)
    for index, slide in enumerate(slides, start=1):
        eyebrow = '<div class="eyebrow">SECTION</div>' if visible_eyebrow_on_slide == index else ""
        slide.pop("html_source")
        slide["visual_prompt"] = ""
        slide["html_body"] = (
            f'{eyebrow}<section id="t{index}" data-deck-id="title-{index}" data-deck-role="title" '
            f'data-deck-required="true"><h1>{slide["title"]}</h1></section>'
            f'<section class="diagram" data-deck-id="diagram-{index}" data-deck-role="diagram"></section>'
            f'<div id="n{index}" class="narrative" data-deck-id="narrative-{index}" data-deck-role="narrative" '
            f'data-deck-required="true"><p>{slide["narrative"]}</p></div>'
        )
        slide["repair_anchor_ids"] = [f"t{index}", f"n{index}"]
    return slides


def _compact_stylesheet(*extra: str, font_family: str = "Calibri,Arial,sans-serif") -> str:
    geometry = "".join(
        (
            f"#t{index}{{position:absolute;left:120px;top:80px;width:1200px;height:120px;"
            "box-sizing:border-box;margin:0}"
            f"#n{index}{{position:absolute;left:120px;top:820px;width:1320px;height:120px;"
            "box-sizing:border-box;margin:0}"
        )
        for index in range(1, 4)
    )
    return (
        f"main{{width:1920px;height:1080px;background:#F7F1E1;color:#2B2926;font-family:{font_family}}}"
        + geometry
        + "".join(extra)
    )


def _creative_plan(*, include_asset: bool = True) -> dict:
    image_assets = []
    if include_asset:
        image_assets.append(
            {
                "asset_id": "cover-texture",
                "slide_selector": "slide:1",
                "role": "hero_background",
                "reason": "Establishes subject atmosphere without carrying semantic text.",
                "prompt": "Dark technical abstract system texture, no readable text, no labels.",
                "aspect_ratio": "16:9",
                "integration": "full_bleed_background",
                "no_baked_text": True,
            }
        )
    return {
        "subject": "Technical Deck",
        "audience": "technical stakeholders",
        "goal": "explain the system clearly",
        "viewing_context": "Presented live to technical stakeholders on a 16:9 display.",
        "subject_materials": ["system topology", "signal flow", "operational evidence"],
        "story_arc": "Frame the system, explain the architecture, close with synthesis.",
        "design_plan": {
            "source": "test",
            "subject": "Technical Deck",
            "audience": "technical stakeholders",
            "goal": "explain the system clearly",
            "style_lane": "technical_blueprint",
            "palette": [
                {"name": "background", "hex": "#0A0E14", "role": "slide substrate"},
                {"name": "surface", "hex": "#111827", "role": "panel"},
                {"name": "ink", "hex": "#EEF4FB", "role": "text"},
                {"name": "muted", "hex": "#A7B4C2", "role": "secondary text"},
                {"name": "accent", "hex": "#38BDF8", "role": "linework"},
            ],
            "typography": {"display": "Aptos Display", "body": "Aptos", "utility": "Aptos"},
            "grid": {"slide_width_px": 1920, "slide_height_px": 1080},
            "signature": "dark native technical diagram language",
            "rhythm": "cover, architecture, synthesis",
            "anti_slop_profile": ["native text", "structural variety"],
            "requested_style_terms": ["dark_technical"],
        },
        "image_strategy": "hybrid" if include_asset else "diagram_native",
        "image_strategy_rationale": "Use one atmospheric asset while keeping semantic system structure native.",
        "image_assets": image_assets,
        "slide_compositions": [
            {
                "selector": f"slide:{index}",
                "slide_role": role,
                "headline_intent": f"Explain slide {index}",
                "layout_name": layout,
                "composition_rationale": "Use native structure with clear hierarchy.",
                "native_elements": ["title", "narrative", "diagram"],
                "image_asset_ids": ["cover-texture"] if include_asset and index == 1 else [],
                "required_element_ids": [f"title-{index}", f"diagram-{index}", f"narrative-{index}"],
                "structural_fingerprint": f"{role}-{layout}-slide-{index}",
                "risk_notes": [],
            }
            for index, (role, layout) in enumerate(
                [
                    ("cover", "cover_with_texture"),
                    ("architecture", "architecture_native_diagram"),
                    ("closing", "closing_synthesis"),
                ],
                start=1,
            )
        ],
        "skill_refs": ["hands-on-deck/designing-slides", "deck-impeccable/critique"],
        "plan_critique": {
            "initial_scores": {
                "philosophy": 3,
                "hierarchy": 4,
                "execution_feasibility": 4,
                "specificity": 3,
                "restraint": 4,
                "variety": 3,
            },
            "weakest_point": "The initial structures were too similar.",
            "revision_made": "Assigned a distinct spatial fingerprint to each narrative role.",
            "final_scores": {
                "philosophy": 4,
                "hierarchy": 4,
                "execution_feasibility": 4,
                "specificity": 4,
                "restraint": 4,
                "variety": 4,
            },
        },
        "anti_slop_commitments": ["structural variety across slides", "no screenshot substrate"],
    }


def _fake_compiler(calls: list[dict]):
    def compile_deck(runtime, output_path: str, title: str, slides_dir: str) -> dict:
        calls.append({"output_path": output_path, "title": title, "slides_dir": slides_dir})
        host = Path(replace_virtual_path(output_path, runtime.state["thread_data"]))
        host.parent.mkdir(parents=True, exist_ok=True)
        host.write_bytes(b"fake pptx")
        slide_count = len(list((host.parent / "slides").glob("*.html")))
        return {
            "success": True,
            "pptx_path": output_path,
            "size_bytes": host.stat().st_size,
            "engine": "fake",
            "slide_count": slide_count,
            "overflow_slides": [],
        }

    return compile_deck


def _fake_batch(runtime: SimpleNamespace, *, create_outputs: bool = True, complete: bool = True):
    def run_batch(manifest_path: str, tool_runtime) -> dict:
        manifest_host = Path(replace_virtual_path(manifest_path, tool_runtime.state["thread_data"]))
        manifest = json.loads(manifest_host.read_text(encoding="utf-8"))
        items = []
        for item in manifest["items"]:
            output_file = item["output_file"]
            if create_outputs:
                host = Path(replace_virtual_path(output_file, tool_runtime.state["thread_data"]))
                host.parent.mkdir(parents=True, exist_ok=True)
                host.write_bytes(b"png")
            items.append({"output_file": output_file, "success": create_outputs, "error_class": None if create_outputs else "api_error"})
        return {
            "summary_present": True,
            "complete": complete and create_outputs,
            "requested": len(items),
            "images_generated": len(items) if create_outputs else 0,
            "failed": 0 if create_outputs else len(items),
            "items": items,
            "error_class_histogram": {},
        }

    return run_batch


class _FakeNativeService:
    def __init__(
        self,
        calls: list[dict] | None = None,
        *,
        full_slide_picture_count: int = 0,
        narrative_font_pt: float = 20.0,
        compact_label_font_pt: float | None = None,
    ) -> None:
        self.calls = calls if calls is not None else []
        self.full_slide_picture_count = full_slide_picture_count
        self.narrative_font_pt = narrative_font_pt
        self.compact_label_font_pt = compact_label_font_pt
        self.slide_count = 0
        self.source_map_path: str | None = None
        self.inventory_path: str | None = None

    def preflight(self) -> NativeDeckPreflight:
        return NativeDeckPreflight(True, True, True, True, [])

    def html_to_patch(self, *, html_paths: list[str], base_deck_path: str, output_patch_path: str) -> NativeDeckPatchResult:
        self.slide_count = len(html_paths)
        self.calls.append(
            {
                "stage": "html_to_patch",
                "html_basenames": [Path(path).name for path in html_paths],
                "base_file": Path(base_deck_path).name,
                "patch_file": Path(output_patch_path).name,
            }
        )
        patch = Path(output_patch_path)
        patch.parent.mkdir(parents=True, exist_ok=True)
        patch.write_text(json.dumps({"ops": []}), encoding="utf-8")
        source_map = {
            "schema_version": "sophia-deck-source-map/v1",
            "slides": {
                f"slide:{index}": {
                    "elements": {
                        f"title-{index}": {
                            "source_role": "title",
                            "source_required": True,
                            "shape_names": [f"s{index}-title-{index}-text"],
                        },
                        f"diagram-{index}": {
                            "source_role": "diagram",
                            "source_required": False,
                            "shape_names": [f"s{index}-diagram-{index}-box"],
                        },
                        f"narrative-{index}": {
                            "source_role": "narrative",
                            "source_required": True,
                            "shape_names": [f"s{index}-narrative-{index}-text"],
                        },
                    }
                }
                for index in range(1, self.slide_count + 1)
            },
        }
        source_map_path = patch.with_suffix(".source-map.json")
        source_map_path.write_text(json.dumps(source_map), encoding="utf-8")
        self.source_map_path = str(source_map_path)
        return NativeDeckPatchResult(
            success=True,
            output_pptx_path=None,
            patch_path=str(patch),
            patch_op_count=0,
            validation_error_count=0,
            errors=[],
            source_map_path=self.source_map_path,
        )

    def apply_patch(self, *, base_deck_path: str, patch_path: str, output_path: str, fix: bool = True) -> NativeDeckPatchResult:
        self.calls.append({"stage": "apply_patch", "output_path": output_path, "fix": fix})
        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)
        presentation = Presentation()
        presentation.slide_width = Inches(13.333333)
        presentation.slide_height = Inches(7.5)
        blank = presentation.slide_layouts[6]
        inventory: dict[str, dict] = {}
        for index in range(1, self.slide_count + 1):
            slide = presentation.slides.add_slide(blank)
            background = slide.background.fill
            background.solid()
            background.fore_color.rgb = RGBColor(0x0A, 0x0E, 0x14)
            title = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(10), Inches(0.8))
            title.name = f"s{index}-title-{index}-text"
            title.text_frame.paragraphs[0].text = f"Slide {index} System Story"
            title_run = title.text_frame.paragraphs[0].runs[0]
            title_run.font.size = Pt(32)
            title_run.font.color.rgb = RGBColor(0xEE, 0xF4, 0xFB)
            diagram = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2), Inches(10), Inches(2.5))
            diagram.name = f"s{index}-diagram-{index}-box"
            diagram.fill.solid()
            diagram.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x27)
            narrative = slide.shapes.add_textbox(Inches(0.8), Inches(5.7), Inches(10), Inches(0.6))
            narrative.name = f"s{index}-narrative-{index}-text"
            narrative.text_frame.paragraphs[0].text = "A concise technical narrative explains the point."
            narrative_run = narrative.text_frame.paragraphs[0].runs[0]
            narrative_run.font.size = Pt(self.narrative_font_pt)
            narrative_run.font.color.rgb = RGBColor(0xEE, 0xF4, 0xFB)
            if index == 1 and self.compact_label_font_pt is not None:
                label = slide.shapes.add_textbox(Inches(0.8), Inches(5.1), Inches(4), Inches(0.4))
                label.name = "s1-status-label-text"
                label.text_frame.paragraphs[0].text = "Success measure"
                label_run = label.text_frame.paragraphs[0].runs[0]
                label_run.font.size = Pt(self.compact_label_font_pt)
                label_run.font.color.rgb = RGBColor(0xEE, 0xF4, 0xFB)
            inventory[f"slide:{index}"] = {
                "native_slide_index": index - 1,
                "title": title.name,
                "body": narrative.name,
                "visual": diagram.name,
                "shapes": [
                    {"name": title.name},
                    {"name": diagram.name},
                    {"name": narrative.name},
                ],
            }
        presentation.save(output)
        inventory_path = output.with_suffix(".shape-inventory.json")
        inventory_path.write_text(json.dumps({"slides": inventory}), encoding="utf-8")
        self.inventory_path = str(inventory_path)
        return NativeDeckPatchResult(True, str(output), patch_path, 0, 0, [])

    def inspect(self, _pptx_path: str) -> NativeDeckInspectResult:
        native_text_shapes = max(1, self.slide_count) * 2
        return NativeDeckInspectResult(
            True,
            slide_count=self.slide_count,
            shape_count=max(1, self.slide_count) * 3 + self.full_slide_picture_count,
            native_text_shape_count=native_text_shapes,
            picture_shape_count=self.full_slide_picture_count,
            full_slide_picture_count=self.full_slide_picture_count,
            native_editability_score=0.9,
            shape_inventory_path=self.inventory_path,
            raw_json_path=None,
            errors=[],
        )

    def lint_fix(self, *, pptx_path: str, touched_slides: list[int] | None = None) -> NativeDeckLintFixResult:
        return NativeDeckLintFixResult(True, 0, 0, 0, len(touched_slides or []), [], [])

    def render(self, *, pptx_path: str, output_dir: str, slides: list[int] | None = None) -> NativeDeckRenderResult:
        Path(output_dir).mkdir(parents=True, exist_ok=True)
        for slide in slides or []:
            (Path(output_dir) / f"slide-{slide}.jpg").write_bytes(b"jpg")
        return NativeDeckRenderResult(True, output_dir, len(slides or []), [])

    def diff(self, *, before_path: str, after_path: str) -> dict:
        return {"success": True, "changed": True, "errors": []}


class _InvalidPptxNativeService(_FakeNativeService):
    def apply_patch(
        self,
        *,
        base_deck_path: str,
        patch_path: str,
        output_path: str,
        fix: bool = True,
    ) -> NativeDeckPatchResult:
        result = super().apply_patch(
            base_deck_path=base_deck_path,
            patch_path=patch_path,
            output_path=output_path,
            fix=fix,
        )
        Path(output_path).write_bytes(b"not a valid pptx package")
        return result


def test_deck_build_service_required_deck_writes_manifest_html_pptx_and_build_json(tmp_path: Path) -> None:
    runtime = _runtime(tmp_path / "outputs")
    native_calls: list[dict] = []
    service = DeckBuildService(
        image_batch_runner=_fake_batch(runtime),
        native_service=_FakeNativeService(native_calls),
    )

    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="Technical Deck",
        slides=_slides(),
        output_path=f"{_OUTPUTS}deck.pptx",
        creative_plan=_creative_plan(),
    )

    assert result.success is True
    assert result.pptx_path == f"{_OUTPUTS}deck.pptx"
    assert result.deck_route == "deck_creative_html_native"
    assert result.deck_compile_mode == NATIVE_DECK_COMPILE_MODE
    assert result.native_editability_score == 0.9
    assert result.native_text_shape_count >= 6
    assert result.full_slide_picture_count == 0
    assert result.expected_visual_count == 1
    assert result.successful_visual_count == 1
    assert result.referenced_visual_count == 1
    assert result.generated_asset_count == 1
    assert result.native_html_slide_count == 2
    assert result.hybrid_slide_count == 1
    assert native_calls[0]["html_basenames"] == ["01-cover.html", "02-architecture.html", "03-closing.html"]
    assert next(call for call in native_calls if call["stage"] == "apply_patch")["fix"] is False
    outputs = tmp_path / "outputs"
    prompt_files = sorted((outputs / "assets" / "prompts").glob("slide-*.json"))
    assert len(prompt_files) == 1
    prompt_payload = json.loads(prompt_files[0].read_text(encoding="utf-8"))
    assert prompt_payload["style"]["visual_style"] == "asset_only_supporting_visual"
    assert prompt_payload["technical"]["deck_asset"] is True
    assert prompt_payload["technical"]["slide_visual"] is False
    assert prompt_payload["technical"]["slide_index"] == 1
    assert "supporting asset" in prompt_payload["constraints"][0]
    assert "handwritten" not in json.dumps(prompt_payload).lower()
    manifest = json.loads((outputs / "assets" / "slide-visuals.manifest.json").read_text(encoding="utf-8"))
    assert manifest["manifest_author"] == "DeckBuildService"
    assert len(manifest["items"]) == 1
    assert manifest["items"][0]["deck_asset"] is True
    assert manifest["items"][0]["slide_visual"] is False
    assert len(list((outputs / "slides").glob("*.html"))) == 3
    html = (outputs / "slides" / "02-architecture.html").read_text(encoding="utf-8")
    assert 'class="diagram"' in html
    assert 'data-deck-id="title-2"' in html
    assert ">Slide 2 System Story</h1>" in html
    build = json.loads((outputs / "deck_build" / "build.json").read_text(encoding="utf-8"))
    assert build["schema_version"] == "sophia-deck-build/v1"
    assert build["status"] == "evaluated"
    assert build["deck_route"] == "deck_creative_html_native"
    assert build["deck_compile_mode"] == NATIVE_DECK_COMPILE_MODE
    assert build["native_editability_score"] == 0.9
    assert build["image_generation_status"] == "success"
    assert build["primary_image_batch_status"] == "success"
    assert build["design_plan_path"] == f"{_OUTPUTS}deck_build/design_plan.json"
    assert build["creative_plan_path"] == f"{_OUTPUTS}deck_build/creative_plan.json"
    assert build["asset_policy_path"] == f"{_OUTPUTS}deck_build/asset_policy.json"
    assert build["generated_asset_count"] == 1
    assert build["native_html_slide_count"] == 2
    assert "Dark technical abstract system texture" in build["slides"][0]["visual_prompt"]
    loaded = load_deck_build(result.deck_build_path, runtime)
    assert loaded is not None
    assert loaded.build_id == result.build_id
    assert loaded.slides[0].selector == "slide:1"


def test_deck_build_service_scopes_native_lint_but_renders_every_slide(
    tmp_path: Path,
) -> None:
    runtime = _runtime(tmp_path / "outputs")
    lint_calls: list[list[int]] = []
    render_calls: list[list[int]] = []

    class RecordingNativeService(_FakeNativeService):
        def lint_fix(
            self,
            *,
            pptx_path: str,
            touched_slides: list[int] | None = None,
        ) -> NativeDeckLintFixResult:
            lint_calls.append(list(touched_slides or []))
            return super().lint_fix(
                pptx_path=pptx_path,
                touched_slides=touched_slides,
            )

        def render(
            self,
            *,
            pptx_path: str,
            output_dir: str,
            slides: list[int] | None = None,
        ) -> NativeDeckRenderResult:
            render_calls.append(list(slides or []))
            return super().render(
                pptx_path=pptx_path,
                output_dir=output_dir,
                slides=slides,
            )

    result = DeckBuildService(
        image_batch_runner=_fake_batch(runtime),
        native_service=RecordingNativeService([]),
    ).prepare_and_build(
        runtime=runtime,
        deck_title="Technical Deck",
        slides=_slides(),
        output_path=f"{_OUTPUTS}deck.pptx",
        creative_plan=_creative_plan(),
        native_lint_slide_indices=(0, 2),
    )

    assert result.success is True
    assert lint_calls == [[0, 2]]
    assert render_calls == [[0, 1, 2]]


@pytest.mark.parametrize(
    "scope",
    [
        (),
        (0, 0),
        (False,),
        (-1,),
        (3,),
        ("0",),
    ],
)
def test_deck_build_service_rejects_invalid_native_lint_scope(
    tmp_path: Path,
    scope: tuple[object, ...],
) -> None:
    runtime = _runtime(tmp_path / "outputs")
    lint_calls: list[list[int]] = []

    class RecordingNativeService(_FakeNativeService):
        def lint_fix(
            self,
            *,
            pptx_path: str,
            touched_slides: list[int] | None = None,
        ) -> NativeDeckLintFixResult:
            lint_calls.append(list(touched_slides or []))
            return super().lint_fix(
                pptx_path=pptx_path,
                touched_slides=touched_slides,
            )

    result = DeckBuildService(
        image_batch_runner=_fake_batch(runtime),
        native_service=RecordingNativeService([]),
    ).prepare_and_build(
        runtime=runtime,
        deck_title="Technical Deck",
        slides=_slides(),
        output_path=f"{_OUTPUTS}deck.pptx",
        creative_plan=_creative_plan(),
        native_lint_slide_indices=scope,  # type: ignore[arg-type]
    )

    assert result.success is False
    assert result.failure_code == "invalid_deck_ir"
    assert lint_calls == []


def test_deck_build_service_resolves_manifest_owner_from_runtime_config(tmp_path: Path) -> None:
    runtime = _runtime(tmp_path / "outputs")
    runtime.state.pop("thread_id")
    runtime.state.pop("user_id")
    runtime.state.pop("parent_thread_id")
    runtime.context = {}
    runtime.config = {
        "configurable": {
            "thread_id": "configured-builder-thread",
            "user_id": "configured-user",
            "parent_thread_id": "configured-companion-thread",
        }
    }
    service = DeckBuildService(
        image_batch_runner=_fake_batch(runtime),
        native_service=_FakeNativeService([]),
    )

    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="Configured Identity Deck",
        slides=_slides(),
        output_path=f"{_OUTPUTS}deck.pptx",
        creative_plan=_creative_plan(),
    )

    assert result.success is True
    build = json.loads((tmp_path / "outputs" / "deck_build" / "build.json").read_text(encoding="utf-8"))
    assert build["thread_id"] == "configured-builder-thread"
    assert build["user_id"] == "configured-user"
    assert build["parent_thread_id"] == "configured-companion-thread"


def test_deck_build_service_contrast_analyzer_failure_is_clean_terminal_result(tmp_path: Path) -> None:
    runtime = _runtime(tmp_path / "outputs")
    result = DeckBuildService(
        image_batch_runner=_fake_batch(runtime),
        native_service=_InvalidPptxNativeService(),
    ).prepare_and_build(
        runtime=runtime,
        deck_title="Technical Deck",
        slides=_slides(),
        output_path=f"{_OUTPUTS}deck.pptx",
        creative_plan=_creative_plan(),
    )

    assert result.success is False
    assert result.failure_code == "deck_native_contrast_analysis_failed"
    assert result.retryable is False


def test_deck_build_service_clears_stale_slide_html_before_compile(tmp_path: Path) -> None:
    runtime = _runtime(tmp_path / "outputs")
    slides_dir = tmp_path / "outputs" / "slides"
    slides_dir.mkdir(parents=True)
    (slides_dir / "99-stale.html").write_text("<html>stale</html>", encoding="utf-8")
    native_calls: list[dict] = []

    service = DeckBuildService(
        image_batch_runner=_fake_batch(runtime),
        native_service=_FakeNativeService(native_calls),
    )

    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="Technical Deck",
        slides=_slides(),
        output_path=f"{_OUTPUTS}deck.pptx",
        creative_plan=_creative_plan(),
    )

    assert result.success is True
    assert "99-stale.html" not in native_calls[0]["html_basenames"]
    assert native_calls[0]["html_basenames"] == ["01-cover.html", "02-architecture.html", "03-closing.html"]


def test_deck_build_service_nested_output_path_evaluates_against_outputs_root(tmp_path: Path) -> None:
    runtime = _runtime(tmp_path / "outputs")
    service = DeckBuildService(
        image_batch_runner=_fake_batch(runtime),
        native_service=_FakeNativeService(),
    )

    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="Nested Technical Deck",
        slides=_slides(),
        output_path=f"{_OUTPUTS}decks/foo.pptx",
        creative_plan=_creative_plan(),
    )

    assert result.success is True
    assert result.pptx_path == f"{_OUTPUTS}decks/foo.pptx"
    assert (tmp_path / "outputs" / "slides" / "01-cover.html").is_file()
    assert (tmp_path / "outputs" / "decks" / "foo.pptx").is_file()


def test_deck_build_service_full_slide_picture_in_native_deck_warns(tmp_path: Path) -> None:
    runtime = _runtime(tmp_path / "outputs")
    service = DeckBuildService(
        image_batch_runner=_fake_batch(runtime),
        native_service=_FakeNativeService(full_slide_picture_count=1),
    )

    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="Overflow Deck",
        slides=_slides(),
        output_path=f"{_OUTPUTS}deck.pptx",
        creative_plan=_creative_plan(),
    )

    assert result.success is True
    assert result.full_slide_picture_count == 1
    assert result.quality_warning == "native_full_bleed_picture_present"


def test_deck_build_service_missing_visual_prompt_does_not_fail_normal_native_slide(tmp_path: Path) -> None:
    runtime = _runtime(tmp_path / "outputs")
    batch_called = False

    def batch_runner(manifest_path, tool_runtime):
        nonlocal batch_called
        batch_called = True
        return _fake_batch(runtime)(manifest_path, tool_runtime)

    service = DeckBuildService(image_batch_runner=batch_runner, native_service=_FakeNativeService())
    slides = _slides()
    slides[1]["visual_prompt"] = ""

    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="Technical Deck",
        slides=slides,
        output_path=f"{_OUTPUTS}deck.pptx",
        creative_plan=_creative_plan(),
    )

    assert result.success is True
    assert result.expected_visual_count == 1
    assert batch_called is True
    assert (tmp_path / "outputs" / "deck.pptx").exists()


def test_deck_build_service_allows_negated_visual_prompt_guardrails(tmp_path: Path) -> None:
    runtime = _runtime(tmp_path / "outputs")
    service = DeckBuildService(
        image_batch_runner=_fake_batch(runtime),
        native_service=_FakeNativeService(),
    )
    slides = _slides()
    slides[0]["visual_prompt"] = "Professional system visual, no axis labels, without formulas, not neon."

    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="Technical Deck",
        slides=slides,
        output_path=f"{_OUTPUTS}deck.pptx",
        creative_plan=_creative_plan(),
    )

    assert result.success is True


def test_deck_build_service_ignores_hidden_unused_eyebrow_selector(tmp_path: Path) -> None:
    runtime = _runtime(tmp_path / "outputs")
    service = DeckBuildService(
        image_batch_runner=_fake_batch(runtime),
        native_service=_FakeNativeService(),
    )

    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="Technical Deck",
        slides=_compact_slides(),
        output_path=f"{_OUTPUTS}deck.pptx",
        deck_stylesheet=_compact_stylesheet(
            ".eyebrow-none{display:none}"
            "h1{font-size:64px}.diagram{width:1200px;height:500px}.narrative{font-size:30px}"
        ),
        authoring_contract="compact_model_html_v2",
        creative_plan=_creative_plan(include_asset=False),
    )

    assert result.success is True


def test_deck_build_service_allows_emitted_20px_label_and_ignores_unused_tiny_rule(tmp_path: Path) -> None:
    runtime = _runtime(tmp_path / "outputs")
    native_service = _FakeNativeService(
        compact_label_font_pt=15.0,
    )
    service = DeckBuildService(
        image_batch_runner=_fake_batch(runtime),
        native_service=native_service,
    )

    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="Technical Deck",
        slides=_compact_slides(),
        output_path=f"{_OUTPUTS}deck.pptx",
        deck_stylesheet=_compact_stylesheet(
            ".unused-utility{font-size:12px}.status-label{font-size:20px}"
            "h1{font-size:64px}.diagram{width:1200px;height:500px}.narrative{font-size:30px}"
        ),
        authoring_contract="compact_model_html_v2",
        creative_plan=_creative_plan(include_asset=False),
    )

    assert result.success is True
    assert result.quality_status == "passed"
    assert result.quality_warning is None
    assert result.source_quality_report["soft_warnings"] == []


def test_deck_build_service_routes_compiled_required_tiny_text_to_mechanical_repair(tmp_path: Path) -> None:
    runtime = _runtime(tmp_path / "outputs")
    native_service = _FakeNativeService(
        narrative_font_pt=17.25,
    )
    service = DeckBuildService(
        image_batch_runner=_fake_batch(runtime),
        native_service=native_service,
    )

    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="Technical Deck",
        slides=_compact_slides(),
        output_path=f"{_OUTPUTS}deck.pptx",
        deck_stylesheet=_compact_stylesheet(
            "h1{font-size:64px}.diagram{width:1200px;height:500px}.narrative{font-size:23px}"
        ),
        authoring_contract="compact_model_html_v2",
        creative_plan=_creative_plan(include_asset=False),
    )

    assert result.success is False
    assert result.failure_code == "deck_mechanical_gate_failed"
    assert result.retryable is True
    assert result.source_quality_report["passed"] is True
    assert any(
        issue["code"] == "native_required_text_too_small"
        for issue in result.mechanical_gate_results["issues"]
    )
    assert result.repair_instruction is not None
    assert result.repair_instruction["generic_repair_target_count"] == 1
    typography_target = result.repair_instruction["repair_targets"][0]
    assert len(typography_target["typography_occurrences"]) == 3
    assert "24px" in result.repair_instruction["repair_message"]


def test_deck_build_service_still_rejects_visible_eyebrow_chrome(tmp_path: Path) -> None:
    runtime = _runtime(tmp_path / "outputs")
    native_calls: list[dict] = []
    service = DeckBuildService(
        image_batch_runner=_fake_batch(runtime),
        native_service=_FakeNativeService(native_calls),
    )

    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="Technical Deck",
        slides=_compact_slides(visible_eyebrow_on_slide=2),
        output_path=f"{_OUTPUTS}deck.pptx",
        deck_stylesheet=_compact_stylesheet(
            "h1{font-size:64px}.diagram{width:1200px;height:500px}.narrative{font-size:30px}"
        ),
        authoring_contract="compact_model_html_v2",
        creative_plan=_creative_plan(include_asset=False),
    )

    assert result.success is False
    assert result.failure_code == "deck_source_quality_failed"
    assert result.retryable is True
    assert "chrome" in str(result.failure_summary).lower()
    assert native_calls  # Native gates still run so one repair can receive every target.
    assert result.source_quality_report["passed"] is False
    assert result.repair_instruction is not None
    assert result.repair_instruction["source_quality_repair_target_count"] == 1
    assert result.repair_instruction["repair_targets"][0]["selector"] == "slide:2"


def test_deck_build_service_combines_source_and_mechanical_targets_for_one_repair(
    tmp_path: Path,
    monkeypatch,
) -> None:
    runtime = _runtime(tmp_path / "outputs")
    slides = _compact_slides()
    for slide in slides:
        slide["html_body"] = '<div class="eyebrow">SECTION</div>' + slide["html_body"]

    mechanical_attempt = 0

    def evaluate_gates(_deck, *, rendered_dir, native_pptx_path=None):
        nonlocal mechanical_attempt
        del rendered_dir, native_pptx_path
        mechanical_attempt += 1
        if mechanical_attempt > 1:
            return MechanicalGateResult(passed=True)
        return MechanicalGateResult(
            passed=False,
            failure_code="deck_mechanical_gate_failed",
            failure_summary="Slide 2 contains a material overlap.",
            issues=[
                MechanicalGateIssue(
                    code="native_lint_severe_overlap",
                    selector="slide:2",
                    summary="Slide 2 contains a material overlap.",
                    repair_hint="Separate the two source elements.",
                )
            ],
        )

    monkeypatch.setattr(deck_service, "evaluate_mechanical_gates", evaluate_gates)
    service = DeckBuildService(
        image_batch_runner=_fake_batch(runtime),
        native_service=_FakeNativeService(),
    )
    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="Technical Deck",
        slides=slides,
        output_path=f"{_OUTPUTS}deck.pptx",
        deck_stylesheet=_compact_stylesheet(
            "h1{font-size:64px}.diagram{width:1200px;height:500px}.narrative{font-size:30px}"
        ),
        authoring_contract="compact_model_html_v2",
        creative_plan=_creative_plan(include_asset=False),
    )

    assert result.success is False
    assert result.failure_code == "deck_mechanical_gate_failed"
    assert result.retryable is True
    assert "Source quality also failed" in str(result.failure_summary)
    assert result.repair_instruction is not None
    assert result.repair_instruction["source_quality_repair_target_count"] == 1
    assert result.repair_instruction["source_quality_issue_count"] == 3
    assert result.repair_instruction["generic_repair_target_count"] == 1
    quality_selectors = {
        selector
        for target in result.repair_instruction["repair_targets"]
        if target["target_type"] == "quality"
        for selector in target["selectors"]
    }
    assert quality_selectors == {"slide:1", "slide:2", "slide:3"}
    assert "OVERLAP" in result.repair_instruction["repair_message"] or "GATE slide:2" in result.repair_instruction["repair_message"]

    repaired = service.prepare_and_build(
        runtime=runtime,
        deck_title="Technical Deck",
        slides=_compact_slides(),
        output_path=f"{_OUTPUTS}deck.pptx",
        deck_stylesheet=_compact_stylesheet(
            "h1{font-size:64px}.diagram{width:1200px;height:500px}.narrative{font-size:30px}"
        ),
        authoring_contract="compact_model_html_v2",
        creative_plan=_creative_plan(include_asset=False),
    )

    assert repaired.success is True
    assert repaired.source_quality_report["passed"] is True
    assert repaired.mechanical_gate_results["passed"] is True


def test_deck_build_service_rejects_renderer_unsafe_compact_fonts(tmp_path: Path) -> None:
    runtime = _runtime(tmp_path / "outputs")
    service = DeckBuildService(
        image_batch_runner=_fake_batch(runtime),
        native_service=_FakeNativeService(),
    )

    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="Technical Deck",
        slides=_compact_slides(),
        output_path=f"{_OUTPUTS}deck.pptx",
        deck_stylesheet=_compact_stylesheet(
            "h1{font-size:64px}.diagram{width:1200px;height:500px}.narrative{font-size:30px}",
            font_family="Georgia,serif",
        ),
        authoring_contract="compact_model_html_v2",
        creative_plan=_creative_plan(include_asset=False),
    )

    assert result.success is False
    assert result.failure_code == "invalid_deck_ir"
    assert "Cambria" in str(result.failure_summary)


def test_compact_v2_normalizes_nonportable_secondary_font_fallbacks() -> None:
    stylesheet = _compact_stylesheet(
        (
            "h1{font:700 60px/1.1 Cambria,Georgia,serif}"
            ".diagram{width:1200px;height:500px}.narrative{font-size:30px}"
        ),
        font_family="Calibri,Helvetica,Arial,sans-serif",
    )
    deck = SimpleNamespace(
        deck_stylesheet=stylesheet,
        deck_stylesheet_hash=hashlib.sha256(stylesheet.encode("utf-8")).hexdigest(),
        deck_authoring_contract="compact_model_html_v2",
    )

    deck_service._validate_authoring_inputs(deck, _compact_slides())

    assert "Helvetica" not in deck.deck_stylesheet
    assert "Georgia" not in deck.deck_stylesheet
    assert "font-family:Calibri, Arial, sans-serif" in deck.deck_stylesheet
    assert "font:700 60px/1.1 Cambria, serif" in deck.deck_stylesheet
    assert deck.deck_stylesheet_hash == hashlib.sha256(deck.deck_stylesheet.encode("utf-8")).hexdigest()


def test_compact_v2_malformed_stylesheet_is_a_controlled_authoring_failure() -> None:
    stylesheet = _compact_stylesheet(font_family="Calibri,Helvetica,sans-serif") + "}"
    deck = SimpleNamespace(
        deck_stylesheet=stylesheet,
        deck_stylesheet_hash=hashlib.sha256(stylesheet.encode("utf-8")).hexdigest(),
        deck_authoring_contract="compact_model_html_v2",
    )

    with pytest.raises(deck_service.DeckBuildFailure) as exc:
        deck_service._validate_authoring_inputs(deck, _compact_slides())

    assert exc.value.code == "invalid_deck_ir"
    assert "malformed CSS" in exc.value.summary


def test_compact_v2_recomputes_hash_from_canonical_stored_stylesheet() -> None:
    stylesheet = _compact_stylesheet()
    deck = SimpleNamespace(
        deck_stylesheet=stylesheet,
        deck_stylesheet_hash=hashlib.sha256(f"  {stylesheet}  ".encode()).hexdigest(),
        deck_authoring_contract="compact_model_html_v2",
    )

    deck_service._validate_authoring_inputs(deck, _compact_slides())

    assert deck.deck_stylesheet_hash == hashlib.sha256(stylesheet.encode("utf-8")).hexdigest()


def test_compact_v2_rejects_empty_font_family_entries() -> None:
    stylesheet = _compact_stylesheet(font_family="Calibri,,sans-serif")
    deck = SimpleNamespace(
        deck_stylesheet=stylesheet,
        deck_stylesheet_hash=hashlib.sha256(stylesheet.encode("utf-8")).hexdigest(),
        deck_authoring_contract="compact_model_html_v2",
    )

    with pytest.raises(deck_service.DeckBuildFailure) as exc:
        deck_service._validate_authoring_inputs(deck, _compact_slides())

    assert exc.value.code == "invalid_deck_ir"
    assert "empty entries" in exc.value.summary


def test_compact_v2_checks_raw_stylesheet_size_before_font_normalization() -> None:
    oversized_fallbacks = ",".join(["Helvetica"] * 900)
    stylesheet = _compact_stylesheet(font_family=f"Calibri,{oversized_fallbacks},sans-serif")
    deck = SimpleNamespace(
        deck_stylesheet=stylesheet,
        deck_stylesheet_hash=hashlib.sha256(stylesheet.encode("utf-8")).hexdigest(),
        deck_authoring_contract="compact_model_html_v2",
    )

    with pytest.raises(deck_service.DeckBuildFailure) as exc:
        deck_service._validate_authoring_inputs(deck, _compact_slides())

    assert exc.value.code == "invalid_deck_ir"
    assert "compact-v2 8192-byte limit" in exc.value.summary


def test_compact_v2_checks_forbidden_style_close_before_font_normalization() -> None:
    stylesheet = _compact_stylesheet(font_family='Calibri,"</style",sans-serif')
    deck = SimpleNamespace(
        deck_stylesheet=stylesheet,
        deck_stylesheet_hash=hashlib.sha256(stylesheet.encode("utf-8")).hexdigest(),
        deck_authoring_contract="compact_model_html_v2",
    )

    with pytest.raises(deck_service.DeckBuildFailure) as exc:
        deck_service._validate_authoring_inputs(deck, _compact_slides())

    assert exc.value.code == "invalid_deck_ir"
    assert "forbidden closing style tag" in exc.value.summary


def test_compact_v2_rejects_deep_at_rules_without_recursing() -> None:
    stylesheet = "@media all{" * 600 + _compact_stylesheet() + "}" * 600
    deck = SimpleNamespace(
        deck_stylesheet=stylesheet,
        deck_stylesheet_hash=hashlib.sha256(stylesheet.encode("utf-8")).hexdigest(),
        deck_authoring_contract="compact_model_html_v2",
    )

    with pytest.raises(deck_service.DeckBuildFailure) as exc:
        deck_service._validate_authoring_inputs(deck, _compact_slides())

    assert exc.value.code == "invalid_deck_ir"
    assert "at-rule" in exc.value.summary


def test_compact_v2_normalizes_commented_unquoted_primary_font() -> None:
    stylesheet = _compact_stylesheet(font_family="Calibri /* preferred */,Helvetica,sans-serif")
    deck = SimpleNamespace(
        deck_stylesheet=stylesheet,
        deck_stylesheet_hash=hashlib.sha256(stylesheet.encode("utf-8")).hexdigest(),
        deck_authoring_contract="compact_model_html_v2",
    )

    deck_service._validate_authoring_inputs(deck, _compact_slides())

    assert "Helvetica" not in deck.deck_stylesheet
    assert "Calibri /* preferred */, sans-serif" in deck.deck_stylesheet


def test_compact_v2_empty_stylesheet_keeps_absent_hash() -> None:
    deck = SimpleNamespace(
        deck_stylesheet="",
        deck_stylesheet_hash=hashlib.sha256(b"").hexdigest(),
        deck_authoring_contract="compact_model_html_v2",
    )

    with pytest.raises(deck_service.DeckBuildFailure):
        deck_service._validate_authoring_inputs(deck, _compact_slides())

    assert deck.deck_stylesheet_hash is None


@pytest.mark.parametrize(("opening", "closing"), [("(", ")"), ("[", "]")])
def test_compact_v2_deep_component_nesting_is_a_controlled_failure(opening: str, closing: str) -> None:
    nested = opening * 500 + "Helvetica" + closing * 500
    stylesheet = _compact_stylesheet(font_family=f"Calibri,{nested},sans-serif")
    deck = SimpleNamespace(
        deck_stylesheet=stylesheet,
        deck_stylesheet_hash=hashlib.sha256(stylesheet.encode("utf-8")).hexdigest(),
        deck_authoring_contract="compact_model_html_v2",
    )

    with pytest.raises(deck_service.DeckBuildFailure) as exc:
        deck_service._validate_authoring_inputs(deck, _compact_slides())

    assert exc.value.code == "invalid_deck_ir"
    assert "nested or malformed CSS" in exc.value.summary


@pytest.mark.parametrize(
    "stylesheet",
    [
        "main{font-family:Calibri,Arial,sans-serif}h1{font:700 60px/1.1 Georgia,serif}",
        ".unused{font-family:Calibri,Arial,sans-serif}main{color:#1F2A37}",
        "main h1{font-family:Calibri,Arial,sans-serif}",
        "body .unused{font-family:Calibri,Arial,sans-serif}",
        "main{font-family:'Arial Black',sans-serif}",
        "main{font-family:Calibri,Arial,sans-serif}@media all{h1{font-family:Georgia,serif}}",
        "main{font-family:Calibri,Arial,sans-serif}@supports(display:grid){h1{font:700 60px Georgia,serif}}",
        "main{font-family:Calibri,Arial,sans-serif;h1{font-family:Georgia,serif}}",
        "main{font-family:Calibri,Arial,sans-serif}h1{all:initial}",
    ],
)
def test_compact_font_contract_rejects_shorthand_fake_and_unused_safe_families(stylesheet: str) -> None:
    with pytest.raises(deck_service.DeckBuildFailure) as exc:
        deck_service._validate_compact_pptx_font_contract(stylesheet, _compact_slides())

    assert exc.value.code == "invalid_deck_ir"
    assert any(token in exc.value.summary for token in ("Cambria", "nested", "Office-safe"))


def test_compact_font_contract_rejects_unsafe_inline_font_shorthand() -> None:
    slides = _compact_slides()
    slides[0]["html_body"] += '<span style="font:700 40px Georgia,serif">Unsafe</span>'

    with pytest.raises(deck_service.DeckBuildFailure) as exc:
        deck_service._validate_compact_pptx_font_contract(
            "main{font-family:Calibri,Arial,sans-serif}",
            slides,
        )

    assert "inline style" in exc.value.summary


def test_compact_font_contract_rejects_unquoted_inline_font_family() -> None:
    slides = _compact_slides()
    slides[0]["html_body"] += "<span style=font-family:Georgia>Unsafe</span>"

    with pytest.raises(deck_service.DeckBuildFailure) as exc:
        deck_service._validate_compact_pptx_font_contract(
            "main{font-family:Calibri,Arial,sans-serif}",
            slides,
        )

    assert "inline style" in exc.value.summary


def test_compact_font_contract_accepts_safe_font_shorthand() -> None:
    deck_service._validate_compact_pptx_font_contract(
        (
            "main{font:400 24px/1.2 Calibri,Arial,sans-serif}"
            "h1{font:700 60px/1.1 Cambria,serif}"
        ),
        _compact_slides(),
    )


def _source_pair_stylesheet(*, second_position: str = "absolute") -> str:
    return (
        "main{background:#101828;color:#F8FAFC;font-family:Calibri,Arial,sans-serif}"
        "#hero{position:absolute;left:80px;top:80px;width:720px;height:320px;"
        "box-sizing:border-box;margin:0}"
        f"#proof{{position:{second_position};left:920px;top:560px;width:720px;height:320px;"
        "box-sizing:border-box;margin:0}"
        "#hero>div{display:flex;gap:24px}#proof>div{display:grid;grid-template-columns:1fr 1fr}"
    )


def _source_pair_body() -> str:
    return (
        '<section id="hero" data-deck-id="hero" data-deck-role="subject" data-deck-required="true">'
        "<div><strong>Current</strong><span>PSI</span></div></section>"
        '<div id="proof" data-deck-id="proof" data-deck-role="mechanism" data-deck-required="true">'
        "<div><span>Control</span><span>loop</span></div></div>"
    )


def _v51_incomplete_anchor_contract() -> tuple[str, list[dict[str, object]]]:
    pairs = (
        ("hero", "why"),
        ("loopwhy", "loopsteps"),
        ("scenario", "motives"),
        ("tablehead", "tablewrap"),
        ("qwrap", "cta"),
    )
    slides: list[dict[str, object]] = []
    rules = [
        "main{font-family:Calibri,Arial,sans-serif}",
        ".eb{position:absolute;box-sizing:border-box;margin:0}",
    ]
    for index, (first, second) in enumerate(pairs):
        slides.append(
            {
                "html_body": (
                    f'<section id="{first}" data-deck-id="{first}" data-deck-role="title" '
                    f'data-deck-required="true"><h1>Slide {index + 1}</h1></section>'
                    f'<div id="{second}" data-deck-id="{second}" data-deck-role="narrative" '
                    'data-deck-required="true"><p>Inspectable motive arbitration.</p></div>'
                ),
                "slide_css": "",
                "repair_anchor_ids": [first, second],
            }
        )
        rules.extend(
            (
                f"#{first}{{left:120px;top:100px;width:1120px;height:320px}}",
                f"#{second}{{left:120px;top:480px;width:1680px;height:400px}}",
            )
        )
    return "".join(rules), slides


@pytest.mark.parametrize(
    ("body", "stylesheet"),
    [
        (
            '<section data-deck-id="hero" data-deck-role="subject" data-deck-required="true">Current PSI</section>'
            '<div data-deck-id="proof" data-deck-role="mechanism" data-deck-required="true">Control loop</div>',
            _source_pair_stylesheet(),
        ),
        (
            '<section id="hero" data-deck-id="hero" data-deck-role="subject" data-deck-required="true">Current PSI</section>'
            '<div id="other" data-deck-id="other" data-deck-role="detail" data-deck-required="true">Detail</div>',
            _source_pair_stylesheet(),
        ),
        (
            '<section><div id="hero" data-deck-id="hero" data-deck-role="subject" data-deck-required="true">'
            'Current PSI</div><div id="proof" data-deck-id="proof" data-deck-role="mechanism" '
            'data-deck-required="true">Control loop</div></section>',
            _source_pair_stylesheet(),
        ),
        (
            _source_pair_body(),
            _source_pair_stylesheet(second_position="static"),
        ),
    ],
    ids=("missing", "one", "nested", "static"),
)
def test_compact_v2_source_addressability_rejects_unprovable_pairs(
    body: str,
    stylesheet: str,
) -> None:
    with pytest.raises(deck_service.DeckBuildFailure, match="both repair anchors declared") as exc:
        deck_service._validate_compact_source_addressability(
            stylesheet,
            [{"html_body": body, "slide_css": "", "repair_anchor_ids": ["hero", "proof"]}],
        )

    assert exc.value.code == "invalid_deck_ir"


@pytest.mark.parametrize(
    "missing_attribute",
    ("data-deck-id", "data-deck-role", "data-deck-required"),
)
def test_compact_v2_source_addressability_requires_deck_semantics(
    missing_attribute: str,
) -> None:
    authored_attribute = {
        "data-deck-id": ' data-deck-id="proof"',
        "data-deck-role": ' data-deck-role="mechanism"',
        "data-deck-required": ' data-deck-required="true"',
    }[missing_attribute]
    body = _source_pair_body().replace(authored_attribute, "", 1)

    with pytest.raises(deck_service.DeckBuildFailure, match="data-deck-id/data-deck-role"):
        deck_service._validate_compact_source_addressability(
            _source_pair_stylesheet(),
            [{"html_body": body, "slide_css": "", "repair_anchor_ids": ["hero", "proof"]}],
        )


def test_compact_v2_source_addressability_rejects_duplicate_anchor_data_deck_ids() -> None:
    body = _source_pair_body().replace('data-deck-id="proof"', 'data-deck-id="hero"', 1)

    with pytest.raises(deck_service.DeckBuildFailure, match="distinct data-deck-id values") as exc:
        deck_service._validate_compact_source_addressability(
            _source_pair_stylesheet(),
            [{"html_body": body, "slide_css": "", "repair_anchor_ids": ["hero", "proof"]}],
        )

    assert exc.value.code == "invalid_deck_ir"


def test_compact_v2_source_addressability_rejects_duplicate_ids_within_slide() -> None:
    body = _source_pair_body() + '<span id="hero">Duplicate</span>'

    with pytest.raises(deck_service.DeckBuildFailure, match="unique within one slide"):
        deck_service._validate_compact_source_addressability(
            _source_pair_stylesheet(),
            [{"html_body": body, "slide_css": "", "repair_anchor_ids": ["hero", "proof"]}],
        )


def test_compact_v2_source_addressability_accepts_strict_dq2_pair_with_interior_layout() -> None:
    body = _source_pair_body()
    stylesheet = _source_pair_stylesheet()

    deck_service._validate_compact_source_addressability(
        stylesheet,
        [{"html_body": body, "slide_css": "", "repair_anchor_ids": ["hero", "proof"]}],
    )

    from deerflow.sophia.deck_design_lift.repair_author import _strict_geometry_source_witness

    assert _strict_geometry_source_witness(
        body=body,
        baseline_slide_css="",
        deck_css=stylesheet,
        minimum=2,
    ) is not None


def test_compact_v2_normalizes_only_redundant_inline_anchor_geometry() -> None:
    body = _source_pair_body().replace(
        'data-deck-required="true">',
        'data-deck-required="true" '
        'style="position:absolute;left:80px;top:80px;width:720px;height:320px;'
        'box-sizing:border-box;margin:0;color:#F8FAFC">',
        1,
    ).replace(
        'data-deck-required="true">',
        'data-deck-required="true" '
        'style="width:720px;height:320px;background:#182230">',
        1,
    )
    slides = [
        {
            "html_body": body,
            "slide_css": "",
            "repair_anchor_ids": ["hero", "proof"],
        }
    ]

    with pytest.raises(deck_service.DeckBuildFailure, match="both repair anchors declared"):
        deck_service._validate_compact_source_addressability(
            _source_pair_stylesheet(),
            slides,
        )

    normalized, report = deck_service._normalize_compact_v2_anchor_inline_geometry(
        _source_pair_stylesheet(),
        slides,
    )

    assert normalized is not slides
    assert normalized[0] is not slides[0]
    assert report == {
        "normalization_applied": True,
        "normalized_slide_count": 1,
        "normalized_anchor_count": 2,
        "removed_declaration_count": 9,
        "removed_property_names": [
            "box-sizing",
            "height",
            "left",
            "margin",
            "position",
            "top",
            "width",
        ],
        "strict_validator_bypassed": False,
        "candidate_compile_changed": False,
        "raw_content_excluded": True,
    }
    assert "position:absolute" in slides[0]["html_body"]
    assert "color:#F8FAFC" in normalized[0]["html_body"]
    assert "background:#182230" in normalized[0]["html_body"]
    assert "position:absolute" not in normalized[0]["html_body"]
    assert "width:720px" not in normalized[0]["html_body"]
    deck_service._validate_compact_source_addressability(
        _source_pair_stylesheet(),
        normalized,
    )


def test_compact_v2_completes_v51_anchor_invariant_contract_before_strict_validation() -> None:
    stylesheet, slides = _v51_incomplete_anchor_contract()
    original_bodies = [slide["html_body"] for slide in slides]

    with pytest.raises(deck_service.DeckBuildFailure, match="both repair anchors declared"):
        deck_service._validate_compact_source_addressability(stylesheet, slides)

    normalized, report = deck_service._normalize_compact_v2_anchor_invariant_contract(
        stylesheet,
        slides,
    )

    assert normalized != stylesheet
    assert [slide["html_body"] for slide in slides] == original_bodies
    assert report["normalization_applied"] is True
    assert report["normalized_anchor_rule_count"] == 10
    assert report["injected_declaration_count"] == 30
    assert report["html_body_changed"] is False
    assert report["strict_validator_bypassed"] is False
    assert report["candidate_compile_changed"] is False
    assert report["raw_content_excluded"] is True
    assert len(report["carrier_selector_sha256"]) == 64
    assert set(deck_service._compact_shared_id_geometry(normalized)) == {
        identifier
        for slide in slides
        for identifier in slide["repair_anchor_ids"]
    }
    deck_service._validate_compact_source_addressability(normalized, slides)

    second, second_report = deck_service._normalize_compact_v2_anchor_invariant_contract(
        normalized,
        slides,
    )
    assert second == normalized
    assert second_report["normalization_applied"] is False
    assert second_report["normalized_anchor_rule_count"] == 0


@pytest.mark.parametrize(
    "stylesheet_suffix",
    (
        ".other{position:absolute;box-sizing:border-box;margin:0}",
        "section{position:relative}",
        "main #hero{position:relative}",
    ),
)
def test_compact_v2_anchor_invariant_completion_fails_closed_on_ambiguity(
    stylesheet_suffix: str,
) -> None:
    stylesheet, slides = _v51_incomplete_anchor_contract()

    normalized, report = deck_service._normalize_compact_v2_anchor_invariant_contract(
        stylesheet + stylesheet_suffix,
        slides,
    )

    assert normalized == stylesheet + stylesheet_suffix
    assert report["normalization_applied"] is False


def test_compact_v2_anchor_invariant_completion_rejects_used_carrier() -> None:
    stylesheet, slides = _v51_incomplete_anchor_contract()
    slides[0]["html_body"] = str(slides[0]["html_body"]).replace(
        'id="hero"',
        'id="hero" class="eb"',
        1,
    )

    normalized, report = deck_service._normalize_compact_v2_anchor_invariant_contract(
        stylesheet,
        slides,
    )

    assert normalized == stylesheet
    assert report["normalization_applied"] is False


def test_compact_v2_anchor_invariant_completion_rejects_undeclared_cross_slide_id() -> None:
    stylesheet, slides = _v51_incomplete_anchor_contract()
    slides[1]["html_body"] = (
        str(slides[1]["html_body"])
        + '<span id="hero" data-deck-id="unrelated">Unrelated slide-local element</span>'
    )

    normalized, report = deck_service._normalize_compact_v2_anchor_invariant_contract(
        stylesheet,
        slides,
    )

    assert normalized == stylesheet
    assert report["normalization_applied"] is False


@pytest.mark.parametrize(
    "carrier",
    (
        ".eb{position:relative;box-sizing:border-box;margin:0}",
        ".eb{position:absolute;box-sizing:border-box;margin:1px}",
        ".eb{position:absolute;box-sizing:border-box;margin:0;color:#fff}",
        ".eb{position:absolute!important;box-sizing:border-box;margin:0}",
        ".slide-root{position:absolute;box-sizing:border-box;margin:0}",
    ),
)
def test_compact_v2_anchor_invariant_completion_rejects_noncanonical_carrier(
    carrier: str,
) -> None:
    stylesheet, slides = _v51_incomplete_anchor_contract()
    stylesheet = stylesheet.replace(
        ".eb{position:absolute;box-sizing:border-box;margin:0}",
        carrier,
    )

    normalized, report = deck_service._normalize_compact_v2_anchor_invariant_contract(
        stylesheet,
        slides,
    )

    assert normalized == stylesheet
    assert report["normalization_applied"] is False


def test_compact_v2_anchor_invariant_completion_rejects_unsafe_geometry_and_budget() -> None:
    stylesheet, slides = _v51_incomplete_anchor_contract()
    unsafe = stylesheet.replace("width:1120px", "width:calc(100% - 20px)", 1)
    oversized = stylesheet + "/*" + ("x" * (8 * 1024)) + "*/"

    for candidate in (unsafe, oversized):
        normalized, report = deck_service._normalize_compact_v2_anchor_invariant_contract(
            candidate,
            slides,
        )
        assert normalized == candidate
        assert report["normalization_applied"] is False


def test_compact_v2_normalizes_inline_secondary_font_fallbacks_byte_locally() -> None:
    body = (
        '<SVG viewBox="0 0 10 10" aria-label="A &amp; B"><text>Ψ</text></SVG>'
        + _source_pair_body().replace(
            "<strong>",
            '<strong STYLE = "font-family:Cambria,Georgia,serif;color:#F8FAFC">',
            1,
        )
    )
    slides = [{"html_body": body, "slide_css": ""}]

    normalized, report = deck_service._normalize_compact_v2_inline_font_fallbacks(
        slides,
    )

    assert normalized is not slides
    assert normalized[0]["html_body"] == body.replace(
        "Cambria,Georgia,serif",
        "Cambria, serif",
    )
    assert '<SVG viewBox="0 0 10 10" aria-label="A &amp; B">' in normalized[0]["html_body"]
    assert " STYLE = " in normalized[0]["html_body"]
    assert report["normalized_slide_count"] == 1
    assert report["normalized_attribute_count"] == 1
    assert report["normalized_declaration_count"] == 1


def test_compact_v2_normalizes_inline_font_shorthand_secondary_fallback() -> None:
    body = _source_pair_body().replace(
        "<span>Control</span>",
        '<span style="font:700 24px/1.2 Calibri,Helvetica,sans-serif">Control</span>',
        1,
    )

    normalized, _report = deck_service._normalize_compact_v2_inline_font_fallbacks(
        [{"html_body": body, "slide_css": ""}],
    )

    assert "font:700 24px/1.2 Calibri, sans-serif" in normalized[0]["html_body"]
    assert "Helvetica" not in normalized[0]["html_body"]


def test_compact_v2_inline_font_normalization_ignores_style_text_in_other_attributes() -> None:
    body = _source_pair_body().replace(
        "<strong>",
        (
            "<strong "
            "title=\"Narrative style='font-family:Cambria,Georgia,serif' must remain\" "
            "data-note=\"style=&quot;font-family:Calibri,Helvetica,sans-serif&quot;\" "
            'style="font-family:Cambria,Georgia,serif">'
        ),
        1,
    )

    normalized, _report = deck_service._normalize_compact_v2_inline_font_fallbacks(
        [{"html_body": body, "slide_css": ""}],
    )

    assert "Narrative style='font-family:Cambria,Georgia,serif' must remain" in normalized[0]["html_body"]
    assert "style=&quot;font-family:Calibri,Helvetica,sans-serif&quot;" in normalized[0]["html_body"]
    assert 'style="font-family:Cambria, serif"' in normalized[0]["html_body"]


def test_compact_v2_inline_font_normalization_keeps_unsafe_primary_for_validation() -> None:
    body = _source_pair_body().replace(
        "<strong>",
        '<strong style="font-family:Georgia,Cambria,serif">',
        1,
    )
    slides = [{"html_body": body, "slide_css": ""}]

    normalized, report = deck_service._normalize_compact_v2_inline_font_fallbacks(
        slides,
    )

    assert normalized is slides
    assert report["normalization_applied"] is False
    with pytest.raises(deck_service.DeckBuildFailure, match="unsupported PPTX font-family"):
        deck_service._validate_compact_pptx_font_contract(
            _source_pair_stylesheet(),
            normalized,
        )


def test_compact_v2_normalization_does_not_bypass_invalid_shared_geometry() -> None:
    body = _source_pair_body().replace(
        'data-deck-required="true">',
        'data-deck-required="true" style="position:absolute;width:720px">',
        1,
    )
    slides = [
        {
            "html_body": body,
            "slide_css": "",
            "repair_anchor_ids": ["hero", "proof"],
        }
    ]
    stylesheet = _source_pair_stylesheet(second_position="static")

    normalized, report = deck_service._normalize_compact_v2_anchor_inline_geometry(
        stylesheet,
        slides,
    )

    assert normalized is not slides
    assert normalized[0]["html_body"] != slides[0]["html_body"]
    assert "position:absolute" not in normalized[0]["html_body"]
    assert report["normalized_anchor_count"] == 1
    with pytest.raises(deck_service.DeckBuildFailure, match="both repair anchors declared"):
        deck_service._validate_compact_source_addressability(
            stylesheet,
            normalized,
        )


@pytest.mark.parametrize(
    ("candidate_compile", "normalized"),
    [(False, True), (True, False)],
)
def test_deck_service_applies_anchor_normalization_only_to_fresh_authoring(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    candidate_compile: bool,
    normalized: bool,
) -> None:
    runtime = _runtime(tmp_path / "outputs")
    runtime.state["builder_pptx_requested_slide_count"] = 1
    runtime.state["deck_candidate_compile"] = candidate_compile
    body = _source_pair_body().replace(
        'data-deck-required="true">',
        'data-deck-required="true" '
        'style="position:absolute;width:720px;font-family:Cambria,Georgia,serif">',
        1,
    )
    slides = [
        {
            "title": "PSI control",
            "narrative": "Motives arbitrate action.",
            "html_body": body,
            "slide_css": "",
            "repair_anchor_ids": ["hero", "proof"],
        }
    ]
    observed: dict[str, object] = {}
    service = DeckBuildService()

    def stop_after_normalization(
        _deck: object,
        normalized_slides: list[dict[str, object]],
        _output_path: str,
        _runtime: object,
    ) -> None:
        observed["body"] = normalized_slides[0]["html_body"]
        raise deck_service.DeckBuildFailure(
            "test_stop_after_normalization",
            "test stop",
            retryable=False,
        )

    monkeypatch.setattr(service, "_validate_inputs", stop_after_normalization)

    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="PSI Agent Architecture",
        slides=slides,
        output_path=f"{_OUTPUTS}deck.pptx",
        authoring_contract="compact_model_html_v2",
        deck_stylesheet=_source_pair_stylesheet(),
        creative_plan={},
    )

    assert result.failure_code == "test_stop_after_normalization"
    assert ("position:absolute" not in str(observed["body"])) is normalized
    assert ("Georgia" not in str(observed["body"])) is normalized


@pytest.mark.parametrize(
    ("candidate_compile", "normalized"),
    [(False, True), (True, False)],
)
def test_deck_service_completes_anchor_invariants_only_for_fresh_authoring(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    candidate_compile: bool,
    normalized: bool,
) -> None:
    stylesheet, all_slides = _v51_incomplete_anchor_contract()
    slide = {
        **all_slides[0],
        "title": "PSI control",
        "narrative": "Motives arbitrate action.",
    }
    runtime = _runtime(tmp_path / "outputs")
    runtime.state["builder_pptx_requested_slide_count"] = 1
    runtime.state["deck_candidate_compile"] = candidate_compile
    observed: dict[str, object] = {}
    service = DeckBuildService()

    def stop_after_normalization(
        deck: object,
        _normalized_slides: list[dict[str, object]],
        _output_path: str,
        _runtime: object,
    ) -> None:
        observed["stylesheet"] = deck.deck_stylesheet
        observed["stylesheet_hash"] = deck.deck_stylesheet_hash
        raise deck_service.DeckBuildFailure(
            "test_stop_after_normalization",
            "test stop",
            retryable=False,
        )

    monkeypatch.setattr(service, "_validate_inputs", stop_after_normalization)

    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="PSI Agent Architecture",
        slides=[slide],
        output_path=f"{_OUTPUTS}deck.pptx",
        authoring_contract="compact_model_html_v2",
        deck_stylesheet=stylesheet,
        creative_plan={},
    )

    assert result.failure_code == "test_stop_after_normalization"
    assert (str(observed["stylesheet"]) != stylesheet) is normalized
    assert observed["stylesheet_hash"] == hashlib.sha256(
        str(observed["stylesheet"]).encode("utf-8")
    ).hexdigest()


def test_compact_v2_source_addressability_accepts_reversed_declared_pair() -> None:
    deck_service._validate_compact_source_addressability(
        _source_pair_stylesheet(),
        [
            {
                "html_body": _source_pair_body(),
                "slide_css": "",
                "repair_anchor_ids": ["proof", "hero"],
            }
        ],
    )


def test_compact_v2_source_addressability_allows_intervening_eligible_layout_anchor() -> None:
    extra = (
        '<section id="detail" data-deck-id="detail" data-deck-role="evidence" '
        'data-deck-required="true">Additional evidence</section>'
    )
    body = _source_pair_body().replace(
        '</section><div id="proof"',
        f'</section>{extra}<div id="proof"',
        1,
    )
    stylesheet = _source_pair_stylesheet() + (
        "#detail{position:absolute;left:720px;top:440px;width:480px;height:80px;"
        "box-sizing:border-box;margin:0}"
    )

    deck_service._validate_compact_source_addressability(
        stylesheet,
        [{"html_body": body, "slide_css": "", "repair_anchor_ids": ["hero", "proof"]}],
    )


@pytest.mark.parametrize(
    "repair_anchor_ids",
    (None, [], ["hero"], ["hero", "proof", "detail"], ["hero", "hero"], ["Hero", "proof"]),
)
def test_compact_v2_source_addressability_rejects_invalid_declared_pair(
    repair_anchor_ids: object,
) -> None:
    slide = {"html_body": _source_pair_body(), "slide_css": ""}
    if repair_anchor_ids is not None:
        slide["repair_anchor_ids"] = repair_anchor_ids

    with pytest.raises(deck_service.DeckBuildFailure, match="repair_anchor_ids must declare exactly two"):
        deck_service._validate_compact_source_addressability(
            _source_pair_stylesheet(),
            [slide],
        )


def test_compact_v2_source_addressability_rejects_nonempty_slide_css() -> None:
    with pytest.raises(deck_service.DeckBuildFailure, match="slide_css must be empty"):
        deck_service._validate_compact_source_addressability(
            _source_pair_stylesheet(),
            [
                {
                    "html_body": _source_pair_body(),
                    "slide_css": ".unrelated{color:#94A3B8}",
                    "repair_anchor_ids": ["hero", "proof"],
                }
            ],
        )


def test_compact_v2_source_addressability_is_fresh_only_for_trusted_candidate_compile() -> None:
    deck = SimpleNamespace(
        deck_stylesheet=_source_pair_stylesheet(),
        deck_authoring_contract="compact_model_html_v2",
    )
    slides = [{"html_body": _source_pair_body(), "slide_css": "#hero{left:120px}"}]

    with pytest.raises(deck_service.DeckBuildFailure, match="slide_css must be empty"):
        deck_service._validate_authoring_inputs(deck, slides)

    deck_service._validate_authoring_inputs(
        deck,
        slides,
        allow_repair_overlay=True,
    )

    missing_declaration = [{"html_body": _source_pair_body(), "slide_css": ""}]
    with pytest.raises(deck_service.DeckBuildFailure, match="repair_anchor_ids must declare exactly two"):
        deck_service._validate_authoring_inputs(deck, missing_declaration)
    deck_service._validate_authoring_inputs(
        deck,
        missing_declaration,
        allow_repair_overlay=True,
    )


def test_compact_v2_source_addressability_preserves_full_body_cascade() -> None:
    body = '<div class="lead">Lead</div>' + _source_pair_body()
    stylesheet = _source_pair_stylesheet() + "section:nth-child(2){transform:translateX(1px)}"

    from deerflow.sophia.deck_design_lift.repair_author import _strict_geometry_source_witness

    assert _strict_geometry_source_witness(
        body=body,
        baseline_slide_css="",
        deck_css=stylesheet,
        minimum=2,
    ) is None
    with pytest.raises(deck_service.DeckBuildFailure, match="cannot prove"):
        deck_service._validate_compact_source_addressability(
            stylesheet,
            [{"html_body": body, "slide_css": "", "repair_anchor_ids": ["hero", "proof"]}],
        )


def test_compact_v2_source_addressability_rejects_matching_nonzero_margin_rule() -> None:
    stylesheet = _source_pair_stylesheet() + "section{margin:8px}"

    from deerflow.sophia.deck_design_lift.repair_author import _strict_geometry_source_witness

    assert _strict_geometry_source_witness(
        body=_source_pair_body(),
        baseline_slide_css="",
        deck_css=stylesheet,
        minimum=2,
    ) is None
    with pytest.raises(deck_service.DeckBuildFailure, match="not literal zero"):
        deck_service._validate_compact_source_addressability(
            stylesheet,
            [
                {
                    "html_body": _source_pair_body(),
                    "slide_css": "",
                    "repair_anchor_ids": ["hero", "proof"],
                }
            ],
        )


@pytest.mark.parametrize(
    "unsafe_rule",
    (
        "#hero,#proof{margin:auto}",
        "#hero,#proof{margin:0 auto}",
        "#hero{margin-left:1px}",
        "main #hero{margin-inline:0}",
        "#hero{-webkit-margin-start:0}",
        "#hero.foo{margin:auto}",
        "#hero{margin:auto}#hero{margin:0}",
    ),
)
def test_compact_v2_source_addressability_rejects_unsafe_anchor_margin_rules(
    unsafe_rule: str,
) -> None:
    body = _source_pair_body().replace(
        'id="hero"',
        'id="hero" class="foo"',
        1,
    )

    with pytest.raises(deck_service.DeckBuildFailure, match="not literal zero"):
        deck_service._validate_compact_source_addressability(
            _source_pair_stylesheet() + unsafe_rule,
            [
                {
                    "html_body": body,
                    "slide_css": "",
                    "repair_anchor_ids": ["hero", "proof"],
                }
            ],
        )


@pytest.mark.parametrize(
    "safe_rule",
    (
        "#hero,#proof{margin:0}",
        "#hero>*{margin:auto}",
        ".other#hero{margin:auto}",
        "#hero{--anchor-margin:0}",
    ),
)
def test_compact_v2_source_addressability_allows_safe_margin_nonmatches(
    safe_rule: str,
) -> None:
    deck_service._validate_compact_source_addressability(
        _source_pair_stylesheet() + safe_rule,
        [
            {
                "html_body": _source_pair_body(),
                "slide_css": "",
                "repair_anchor_ids": ["hero", "proof"],
            }
        ],
    )


@pytest.mark.parametrize(
    "inline_margin",
    (
        "margin:auto",
        "margin-inline:0",
        "-webkit-margin-start:0",
    ),
)
def test_compact_v2_source_addressability_rejects_unsafe_inline_anchor_margin(
    inline_margin: str,
) -> None:
    body = _source_pair_body().replace(
        'data-deck-required="true">',
        f'data-deck-required="true" style="{inline_margin}">',
        1,
    )

    with pytest.raises(deck_service.DeckBuildFailure, match="not literal zero"):
        deck_service._validate_compact_source_addressability(
            _source_pair_stylesheet(),
            [
                {
                    "html_body": body,
                    "slide_css": "",
                    "repair_anchor_ids": ["hero", "proof"],
                }
            ],
        )


@pytest.mark.parametrize(
    "hero_geometry",
    (
        "left:0;top:0;width:1920px;height:1080px;",
        "left:0;top:0;width:1913px;height:1073px;",
    ),
)
def test_compact_v2_source_addressability_requires_eight_px_clearance(
    hero_geometry: str,
) -> None:
    stylesheet = _source_pair_stylesheet().replace(
        "left:80px;top:80px;width:720px;height:320px;",
        hero_geometry,
        1,
    )

    with pytest.raises(deck_service.DeckBuildFailure, match="at least 8px translation clearance"):
        deck_service._validate_compact_source_addressability(
            stylesheet,
            [
                {
                    "html_body": _source_pair_body(),
                    "slide_css": "",
                    "repair_anchor_ids": ["hero", "proof"],
                }
            ],
        )


@pytest.mark.parametrize(
    "hero_geometry",
    (
        "left:0;top:0;width:1912px;height:1080px;",
        "left:0;top:0;width:1920px;height:1072px;",
    ),
)
def test_compact_v2_source_addressability_accepts_exact_eight_px_clearance(
    hero_geometry: str,
) -> None:
    stylesheet = _source_pair_stylesheet().replace(
        "left:80px;top:80px;width:720px;height:320px;",
        hero_geometry,
        1,
    )

    deck_service._validate_compact_source_addressability(
        stylesheet,
        [
            {
                "html_body": _source_pair_body(),
                "slide_css": "",
                "repair_anchor_ids": ["hero", "proof"],
            }
        ],
    )


def test_compact_v2_source_addressability_rejects_effective_off_canvas_override() -> None:
    body = _source_pair_body().replace(
        'id="hero"',
        'id="hero" class="foo"',
        1,
    )
    stylesheet = _source_pair_stylesheet() + (
        "#hero.foo{left:-8px;top:0;width:1920px;height:1080px}"
    )

    with pytest.raises(deck_service.DeckBuildFailure, match="wholly inside"):
        deck_service._validate_compact_source_addressability(
            stylesheet,
            [
                {
                    "html_body": body,
                    "slide_css": "",
                    "repair_anchor_ids": ["hero", "proof"],
                }
            ],
        )


def test_compact_v2_source_addressability_rejects_non_anchor_witness_pair() -> None:
    body = _source_pair_body() + '<h1 id="x">Fallback title</h1><p id="y">Fallback body</p>'
    stylesheet = _source_pair_stylesheet() + (
        "#hero{transform:translateX(1px)}"
        "#x{position:absolute;left:120px;top:440px;width:600px;height:80px;"
        "box-sizing:border-box;margin:0}"
        "#y{position:absolute;left:920px;top:440px;width:600px;height:80px;"
        "box-sizing:border-box;margin:0}"
    )

    from deerflow.sophia.deck_design_lift.repair_author import _strict_geometry_source_witness

    witness = _strict_geometry_source_witness(
        body=body,
        baseline_slide_css="",
        deck_css=stylesheet,
        minimum=2,
    )
    assert witness is not None
    assert deck_service._compact_witness_anchor_ids(witness) not in ({"hero", "proof"}, None)
    with pytest.raises(deck_service.DeckBuildFailure, match="cannot prove"):
        deck_service._validate_compact_source_addressability(
            stylesheet,
            [{"html_body": body, "slide_css": "", "repair_anchor_ids": ["hero", "proof"]}],
        )


def test_compact_v2_source_addressability_rejects_hidden_descendant_text() -> None:
    body = _source_pair_body().replace("<div><strong>", '<div class="hide"><strong>', 1)
    stylesheet = _source_pair_stylesheet().replace("#hero>div{display:flex;gap:24px}", "") + ".hide{display:none}"

    with pytest.raises(deck_service.DeckBuildFailure, match="both declared compact-v2 repair anchors visible"):
        deck_service._validate_compact_source_addressability(
            stylesheet,
            [{"html_body": body, "slide_css": "", "repair_anchor_ids": ["hero", "proof"]}],
        )


def test_deck_build_service_allows_explicitly_requested_visual_style(tmp_path: Path) -> None:
    runtime = _runtime(
        tmp_path / "outputs",
        user_request="Make a neon cyberpunk pitch deck about resilient infrastructure.",
    )
    service = DeckBuildService(
        image_batch_runner=_fake_batch(runtime),
        native_service=_FakeNativeService(),
    )
    slides = _slides()
    slides[0]["visual_prompt"] = "Neon cyberpunk infrastructure command center with cinematic lighting."

    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="Technical Deck",
        slides=slides,
        output_path=f"{_OUTPUTS}deck.pptx",
        creative_plan=_creative_plan(),
    )

    assert result.success is True


def test_deck_build_service_allows_style_profile_visual_style(tmp_path: Path) -> None:
    runtime = _runtime(tmp_path / "outputs")
    service = DeckBuildService(
        image_batch_runner=_fake_batch(runtime),
        native_service=_FakeNativeService(),
    )
    slides = _slides()
    slides[0]["visual_prompt"] = "Handwritten sketch of the system rollout with marker-like strokes."

    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="Technical Deck",
        slides=slides,
        output_path=f"{_OUTPUTS}deck.pptx",
        style_profile={"visual_style": "handwritten_sketch"},
        creative_plan=_creative_plan(),
    )

    assert result.success is True


def test_deck_build_service_sanitizes_positive_banned_visual_prompt_terms(tmp_path: Path) -> None:
    runtime = _runtime(tmp_path / "outputs")
    batch_called = False

    def batch_runner(_manifest_path, _runtime):
        nonlocal batch_called
        batch_called = True
        return {}

    service = DeckBuildService(image_batch_runner=batch_runner, native_service=_FakeNativeService())
    slides = _slides(include_asset=False)
    slides[0]["visual_prompt"] = "Neon cyberpunk system diagram with dramatic lighting."

    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="Technical Deck",
        slides=slides,
        output_path=f"{_OUTPUTS}deck.pptx",
        creative_plan=_creative_plan(include_asset=False),
    )

    assert result.success is True
    assert result.expected_visual_count == 0
    assert result.image_generation_status == "not_required"
    assert batch_called is False


def test_deck_build_service_missing_batch_summary_fails_without_compile(tmp_path: Path) -> None:
    runtime = _runtime(tmp_path / "outputs")
    native_calls: list[dict] = []
    single_called = False

    def single_runner(_slide, _runtime, _attempt_no):
        nonlocal single_called
        single_called = True
        return {"success": False, "error_class": "should_not_run"}

    service = DeckBuildService(
        image_batch_runner=lambda _manifest_path, _runtime: {"summary_present": False, "complete": False},
        image_single_runner=single_runner,
        native_service=_FakeNativeService(native_calls),
    )

    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="Technical Deck",
        slides=_slides(),
        output_path=f"{_OUTPUTS}deck.pptx",
        creative_plan=_creative_plan(),
    )

    assert result.success is False
    assert result.failure_code == "deck_visual_batch_startup_failed"
    assert single_called is False
    assert native_calls == []
    assert (tmp_path / "outputs" / "slides" / "01-cover.html").exists()


def test_deck_build_service_salvages_partial_timeout_and_repairs_missing_visual(tmp_path: Path) -> None:
    runtime = _runtime(tmp_path / "outputs")
    native_calls: list[dict] = []
    repaired: list[int] = []

    def timeout_batch(manifest_path: str, tool_runtime) -> dict:
        manifest_host = Path(replace_virtual_path(manifest_path, tool_runtime.state["thread_data"]))
        manifest = json.loads(manifest_host.read_text(encoding="utf-8"))
        progress = []
        for index, item in enumerate(manifest["items"], start=1):
            if index == 1:
                continue
            output_file = item["output_file"]
            host = Path(replace_virtual_path(output_file, tool_runtime.state["thread_data"]))
            host.parent.mkdir(parents=True, exist_ok=True)
            host.write_bytes(b"png")
            progress.append({"item_index": index, "output_file": output_file, "success": True, "bytes": 3})
        return {
            "summary_present": False,
            "complete": False,
            "exit_code": 124,
            "error_class": "timeout",
            "items": progress,
        }

    def single_runner(slide, tool_runtime, attempt_no: int) -> dict:
        assert attempt_no == 1
        repaired.append(slide.index)
        host = Path(replace_virtual_path(slide.visual_asset_path, tool_runtime.state["thread_data"]))
        host.parent.mkdir(parents=True, exist_ok=True)
        host.write_bytes(b"png")
        return {"success": True, "bytes": 3}

    service = DeckBuildService(
        image_batch_runner=timeout_batch,
        image_single_runner=single_runner,
        native_service=_FakeNativeService(native_calls),
    )

    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="Technical Deck",
        slides=_slides(),
        output_path=f"{_OUTPUTS}deck.pptx",
        creative_plan=_creative_plan(),
    )

    assert result.success is True
    assert repaired == [1]
    assert result.successful_visual_count == 1
    assert result.image_generation_status == "success_after_repair"
    assert result.primary_image_batch_status == "repaired"
    assert result.serial_repair_count == 1
    assert result.batch_timeout_count == 1
    assert result.partial_batch_salvaged is False
    assert native_calls


def test_deck_build_service_timeout_with_zero_outputs_repairs_selected_assets(tmp_path: Path) -> None:
    runtime = _runtime(tmp_path / "outputs")
    repaired: list[int] = []

    def single_runner(slide, tool_runtime, attempt_no: int) -> dict:
        assert attempt_no == 1
        repaired.append(slide.index)
        host = Path(replace_virtual_path(slide.visual_asset_path, tool_runtime.state["thread_data"]))
        host.parent.mkdir(parents=True, exist_ok=True)
        host.write_bytes(b"png")
        return {"success": True, "bytes": 3}

    service = DeckBuildService(
        image_batch_runner=lambda _manifest_path, _runtime: {
            "summary_present": False,
            "complete": False,
            "exit_code": 124,
            "error_class": "timeout",
            "items": [],
        },
        image_single_runner=single_runner,
        native_service=_FakeNativeService(),
    )

    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="Technical Deck",
        slides=_slides(),
        output_path=f"{_OUTPUTS}deck.pptx",
        creative_plan=_creative_plan(),
    )

    assert result.success is True
    assert repaired == [1]
    assert result.successful_visual_count == 1
    assert result.image_generation_status == "success_after_repair"
    assert result.serial_repair_count == 1
    assert result.partial_batch_salvaged is False


def test_deck_image_batch_timeout_scales_by_manifest_count_and_concurrency(tmp_path: Path, monkeypatch) -> None:
    runtime = _runtime(tmp_path / "outputs")
    manifest_host = tmp_path / "outputs" / "assets" / "slide-visuals.manifest.json"
    manifest_host.parent.mkdir(parents=True)
    manifest_host.write_text(
        json.dumps({"items": [{"output_file": f"{_OUTPUTS}assets/slide-{index:02d}.png"} for index in range(1, 31)]}),
        encoding="utf-8",
    )
    monkeypatch.setenv("SOPHIA_IMAGE_GEN_TIMEOUT", "240")
    monkeypatch.delenv("SOPHIA_IMAGE_GEN_MAX_RETRIES", raising=False)
    monkeypatch.setenv("SOPHIA_IMAGE_GEN_CONCURRENCY", "2")
    monkeypatch.delenv("SOPHIA_DECK_IMAGE_BATCH_TIMEOUT", raising=False)

    timeout = deck_service._deck_image_batch_timeout_seconds(f"{_OUTPUTS}assets/slide-visuals.manifest.json", runtime)

    assert timeout == 7260


def test_deck_image_batch_timeout_override_wins(tmp_path: Path, monkeypatch) -> None:
    runtime = _runtime(tmp_path / "outputs")
    monkeypatch.setenv("SOPHIA_DECK_IMAGE_BATCH_TIMEOUT", "999")

    timeout = deck_service._deck_image_batch_timeout_seconds(f"{_OUTPUTS}missing.manifest.json", runtime)

    assert timeout == 999


def test_deck_image_batch_subprocess_timeout_is_structured(tmp_path: Path, monkeypatch) -> None:
    runtime = _runtime(tmp_path / "outputs")
    script = tmp_path / "generate.py"
    script.write_text("raise SystemExit(0)\n", encoding="utf-8")
    manifest_host = tmp_path / "outputs" / "assets" / "slide-visuals.manifest.json"
    manifest_host.parent.mkdir(parents=True)
    manifest_host.write_text(json.dumps({"items": [{"output_file": f"{_OUTPUTS}assets/slide-01.png"}]}), encoding="utf-8")
    monkeypatch.setattr(deck_service, "_image_script_path", lambda: script)

    def timeout_run(*_args, **_kwargs):
        raise subprocess.TimeoutExpired(cmd=["python"], timeout=10, output="partial stdout", stderr="provider hung")

    monkeypatch.setattr(deck_service, "run_trusted_image_request", timeout_run)

    result = DeckBuildService()._run_image_batch_subprocess(f"{_OUTPUTS}assets/slide-visuals.manifest.json", runtime)

    assert result["summary_present"] is False
    assert result["complete"] is False
    assert result["exit_code"] == 124
    assert result["error_class"] == "timeout"
    assert "timed out" in result["raw_error_excerpt"]


def test_deck_image_batch_timeout_is_capped_by_shared_deadline(
    tmp_path: Path,
    monkeypatch,
) -> None:
    runtime = _runtime(tmp_path / "outputs")
    runtime.state["builder_deadline_epoch_ms"] = 111_000
    runtime.state["builder_budget"] = {"terminal_reserve_seconds": 30}
    monkeypatch.setattr(deck_service.time, "time", lambda: 1.0)
    manifest_host = tmp_path / "outputs" / "assets" / "slide-visuals.manifest.json"
    manifest_host.parent.mkdir(parents=True, exist_ok=True)
    manifest_host.write_text(
        json.dumps({"items": [{"output_file": f"{_OUTPUTS}assets/slide-01.png"}]}),
        encoding="utf-8",
    )

    timeout = deck_service._deck_image_batch_timeout_seconds(
        f"{_OUTPUTS}assets/slide-visuals.manifest.json",
        runtime,
    )

    assert timeout == 80


def test_expired_shared_deadline_returns_non_retryable_failure(
    tmp_path: Path,
    monkeypatch,
) -> None:
    runtime = _runtime(tmp_path / "outputs")
    runtime.state["builder_deadline_epoch_ms"] = 1_000
    monkeypatch.setattr(deck_service.time, "time", lambda: 2.0)

    result = DeckBuildService().prepare_and_build(
        runtime=runtime,
        deck_title="Deadline",
        slides=_slides(),
        output_path=f"{_OUTPUTS}deck.pptx",
        creative_plan=_creative_plan(),
    )

    assert result.success is False
    assert result.failure_code == "deck_deadline_exceeded"
    assert result.retryable is False


def test_deck_build_service_incomplete_visuals_fail_before_compile(tmp_path: Path) -> None:
    runtime = _runtime(tmp_path / "outputs")
    native_calls: list[dict] = []
    repair_attempts: list[tuple[int, int]] = []

    def single_runner(slide, _runtime, attempt_no: int) -> dict:
        repair_attempts.append((slide.index, attempt_no))
        return {"success": False, "error_class": "timeout"}

    service = DeckBuildService(
        image_batch_runner=_fake_batch(runtime, create_outputs=False, complete=False),
        image_single_runner=single_runner,
        native_service=_FakeNativeService(native_calls),
    )

    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="Technical Deck",
        slides=_slides(),
        output_path=f"{_OUTPUTS}deck.pptx",
        creative_plan=_creative_plan(),
    )

    assert result.success is False
    assert result.failure_code == "deck_visuals_incomplete"
    assert result.successful_visual_count == 0
    assert result.missing_visual_count == 1
    assert result.serial_repair_count == 2
    assert repair_attempts == [(1, 1), (1, 2)]
    assert native_calls == []


def test_deck_build_service_terminal_provider_error_does_not_unlock_serial_repair(tmp_path: Path) -> None:
    runtime = _runtime(tmp_path / "outputs")
    single_called = False

    def single_runner(_slide, _runtime, _attempt_no):
        nonlocal single_called
        single_called = True
        return {"success": False, "error_class": "should_not_run"}

    def auth_batch(manifest_path: str, tool_runtime) -> dict:
        manifest_host = Path(replace_virtual_path(manifest_path, tool_runtime.state["thread_data"]))
        manifest = json.loads(manifest_host.read_text(encoding="utf-8"))
        return {
            "summary_present": True,
            "complete": False,
            "requested": len(manifest["items"]),
            "images_generated": 0,
            "failed": len(manifest["items"]),
            "items": [{"output_file": item["output_file"], "success": False, "error_class": "auth_invalid"} for item in manifest["items"]],
            "error_class_histogram": {"auth_invalid": len(manifest["items"])},
        }

    service = DeckBuildService(
        image_batch_runner=auth_batch,
        image_single_runner=single_runner,
        native_service=_FakeNativeService(),
    )

    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="Technical Deck",
        slides=_slides(),
        output_path=f"{_OUTPUTS}deck.pptx",
        creative_plan=_creative_plan(),
    )

    assert result.success is False
    assert result.failure_code == "deck_visuals_incomplete"
    assert result.image_generation_status == "failed"
    assert result.image_generation_reason == "auth_invalid"
    assert result.serial_repair_count == 0
    assert single_called is False


def test_deck_build_service_text_only_requires_explicit_request_and_compiles_without_visuals(tmp_path: Path) -> None:
    runtime = _runtime(tmp_path / "outputs", user_request="Please build a plain text-only 3 slide deck with no visuals.")
    native_calls: list[dict] = []
    service = DeckBuildService(
        image_batch_runner=lambda _manifest_path, _runtime: (_ for _ in ()).throw(AssertionError("no image batch")),
        native_service=_FakeNativeService(native_calls),
    )
    slides = _slides(include_asset=False)
    for slide in slides:
        slide["visual_prompt"] = ""

    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="Text Deck",
        slides=slides,
        output_path=f"{_OUTPUTS}deck.pptx",
        visual_policy="text_only",
        creative_plan=_creative_plan(include_asset=False),
    )

    assert result.success is True
    assert result.expected_visual_count == 0
    assert result.successful_visual_count == 0
    assert result.image_generation_status == "not_required"
    assert native_calls
    assert not (tmp_path / "outputs" / "assets" / "prompts").exists()
    cover_html = (tmp_path / "outputs" / "slides" / "01-cover.html").read_text(encoding="utf-8")
    assert "<img" not in cover_html
    assert 'data-deck-id="title-1"' in cover_html
    assert ">Slide 1 System Story</h1>" in cover_html
    assert 'class="diagram"' in cover_html


def test_deck_build_service_text_only_accepts_delegated_task_brief(tmp_path: Path) -> None:
    runtime = _runtime(tmp_path / "outputs", user_request="")
    runtime.state["delegation_context"] = {"task": "Build a plain text-only deck with no images for the review."}
    native_calls: list[dict] = []
    service = DeckBuildService(
        image_batch_runner=lambda _manifest_path, _runtime: (_ for _ in ()).throw(AssertionError("no image batch")),
        native_service=_FakeNativeService(native_calls),
    )
    slides = _slides(include_asset=False)
    for slide in slides:
        slide["visual_prompt"] = ""

    result = service.prepare_and_build(
        runtime=runtime,
        deck_title="Text Deck",
        slides=slides,
        output_path=f"{_OUTPUTS}deck.pptx",
        visual_policy="text_only",
        creative_plan=_creative_plan(include_asset=False),
    )

    assert result.success is True
    assert native_calls
    assert not (tmp_path / "outputs" / "assets" / "prompts").exists()


def test_prepare_deck_build_tool_schema_excludes_runtime() -> None:
    schema = prepare_deck_build.tool_call_schema.model_json_schema()

    properties = schema.get("properties", {})
    assert "runtime" not in properties
    assert {"deck_title", "slides", "output_path", "register", "visual_policy"}.issubset(properties)
    assert "creative_plan" in schema.get("required", [])
    composition_schema = schema["$defs"]["DeckSlideCompositionInput"]
    grid_schema = schema["$defs"]["DeckGridInput"]
    assert grid_schema["properties"]["footer_policy"]["const"] == "none"
    assert grid_schema["properties"]["eyebrow_policy"]["const"] == "none"
    assert set(composition_schema["required"]) == {
        "selector",
        "slide_role",
        "headline_intent",
        "layout_name",
        "composition_rationale",
        "native_elements",
        "image_asset_ids",
        "required_element_ids",
        "structural_fingerprint",
    }


def test_creative_plan_tool_contract_normalizes_legacy_chrome_policies() -> None:
    payload = _creative_plan(include_asset=False)
    payload["design_plan"]["grid"]["footer_policy"] = "page_numbers"
    payload["design_plan"]["grid"]["eyebrow_policy"] = "only_when_meaningful"

    parsed = DeckCreativePlanInput.model_validate(payload)

    assert parsed.design_plan.grid.footer_policy == "none"
    assert parsed.design_plan.grid.eyebrow_policy == "none"


def test_creative_plan_tool_contract_normalizes_only_direct_aliases() -> None:
    payload = _creative_plan(include_asset=False)
    composition = payload["slide_compositions"][0]
    composition["slide"] = 1
    composition["role"] = composition.pop("slide_role")
    composition["layout"] = composition.pop("layout_name")
    composition.pop("selector")

    parsed = DeckCreativePlanInput.model_validate(payload)
    normalized = parsed.slide_compositions[0]

    assert normalized.selector == "slide:1"
    assert normalized.slide_role == "cover"
    assert normalized.layout_name == "cover_with_texture"
    assert normalized.headline_intent == "Explain slide 1"


def test_presentation_authoring_prompt_forbids_recurring_page_chrome() -> None:
    prompt = builder_artifact_module._PRESENTATION_AUTHORING_SYSTEM_PROMPT

    assert "Do not add eyebrow or kicker labels" in prompt
    assert "footer_policy" in prompt
    assert "eyebrow_policy must both be 'none'" in prompt


def test_presentation_authoring_prompt_matches_compact_v2_body_limit() -> None:
    prompt = builder_artifact_module._PRESENTATION_AUTHORING_SYSTEM_PROMPT

    assert "target each html_body under 4 KiB" in prompt
    assert "combined html_body bytes within 4 KiB times the slide count" in prompt
    assert "slides may borrow unused body budget" in prompt
    assert "each slide capped at the hard 6 KiB ceiling" in prompt
    assert "one slide may borrow" not in prompt
    assert "each html_body under 3 KiB" not in prompt


def test_presentation_authoring_prompt_requires_repair_addressable_anchors() -> None:
    prompt = builder_artifact_module._PRESENTATION_AUTHORING_SYSTEM_PROMPT

    assert "repair_anchor_ids to exactly two short HTML ids" in prompt
    assert "both named repair-addressable layout anchors" in prompt
    assert "non-nested section or div direct children of the service-owned main canvas" in prompt
    assert "HTML id unique within its slide" in prompt
    assert "same two short anchor IDs may be reused in separate slide fragments" in prompt
    assert "shared #id rules scale" in prompt
    assert "[a-z][a-z0-9_-]{0,31}" in prompt
    assert "lowercase ASCII letter followed by at most 31 lowercase ASCII letters" in prompt
    assert "maximum of 32 characters" in prompt
    assert "data-deck-id must be unique within its slide" in prompt
    assert "data-deck-role must be nonempty" in prompt
    assert "data-deck-required must equal true" in prompt
    assert "position:absolute, box-sizing:border-box, margin:0" in prompt
    assert "at least 48x24px and wholly inside the canvas" in prompt
    assert "geometry out of slide_css and inline styles" in prompt
    assert "at least 8px of free canvas" in prompt
    assert "auto, nonzero, or otherwise non-literal-zero physical margin" in prompt
    assert "logical or vendor margin property" in prompt
    assert "rather than overriding it with a later margin:0 reset" in prompt
    assert "Grouped physical margin:0 is safe but unnecessary" in prompt
    assert "reset margins on anchor descendants with separate descendant selectors" in prompt
    assert "Do not use !important, right, bottom, inset, min/max sizing" in prompt
    assert "do not put at-rules or nested rules in authored CSS" in prompt
    assert "Flex and grid remain available inside either anchor" in prompt
    assert "Use real content containers as the two anchors" in prompt
    assert "never duplicate visible content into extra positioned overlay anchors" in prompt
    assert "unrelated visible text-bearing rectangles disjoint with at least a 16px gutter" in prompt
    assert "non-text background with no native text frame" in prompt
    assert "exact-edge connector or background touching remains allowed" in prompt
    assert 'repair_anchor_ids=["hero","proof"]' in prompt
    assert '<section id="hero" data-deck-id="hero"' in prompt
    assert "slide_css omitted or empty for every slide" in prompt
    assert "authenticated repair overlay retains its full budget" in prompt
    assert "only small slide_css overrides" not in prompt


def test_compact_authoring_schema_requires_collision_safe_text_geometry() -> None:
    schema = PrepareDeckBuildInput.model_json_schema()
    slide = schema["$defs"]["DeckSlideInput"]["properties"]["html_body"]["description"]
    stylesheet = schema["properties"]["deck_stylesheet"]["description"]

    assert "real content containers for the anchors" in slide
    assert "never duplicate visible content into extra positioned overlay anchors" in slide
    assert "unrelated visible text-bearing rectangles disjoint with at least a 16px gutter" in stylesheet
    assert "non-text background with no native text frame" in stylesheet
    assert "exact-edge connector or background touching remains allowed" in stylesheet


def test_presentation_authoring_prompt_sets_role_aware_font_floors() -> None:
    prompt = builder_artifact_module._PRESENTATION_AUTHORING_SYSTEM_PROMPT

    assert "required body and narrative text must be at least 24px" in prompt
    assert "Every visible text descendant inside an element marked data-deck-required=true" in prompt
    assert "including nested spans and labels" in prompt
    assert "20-23px is allowed only inside optional elements" in prompt
    assert "No visible text may be below 20px" in prompt
    assert "vertical slack beyond its computed line height" in prompt
    assert "connector bars to exact shared edges or centerlines" in prompt
    assert "Offset connector bars by half their thickness" in prompt
    assert "left=C-W/2" in prompt
    assert "top=C-H/2" in prompt
    assert "Keep every non-bleed shape inside the 1920x1080 canvas" in prompt
    assert "parent-local left/top coordinates" in prompt
    assert "never repeat the parent's slide-global offset on a nested child" in prompt


def test_creative_plan_validation_reports_indexed_nested_path(tmp_path: Path) -> None:
    runtime = _runtime(tmp_path / "outputs")
    plan = _creative_plan(include_asset=False)
    plan["slide_compositions"][0].pop("headline_intent")

    result = DeckBuildService(native_service=_FakeNativeService()).prepare_and_build(
        runtime=runtime,
        deck_title="Technical Deck",
        slides=_slides(include_asset=False),
        output_path=f"{_OUTPUTS}deck.pptx",
        creative_plan=plan,
    )

    assert result.success is False
    assert result.failure_code == "deck_creative_plan_invalid"
    assert result.failure_summary == "creative_plan.slide_compositions[0].headline_intent is required"


def test_presentation_toolset_uses_prepare_deck_build_by_default(monkeypatch) -> None:
    monkeypatch.delenv("SOPHIA_DECK_BUILD_SERVICE_ENABLED", raising=False)

    names = [getattr(tool, "name", "") for tool in build_builder_tools_for_task_type("presentation", vision_enabled=False)]

    assert deck_build_service_enabled() is True
    assert "prepare_deck_build" in names
    assert "prepare_pptx_image_manifest" not in names
    assert "build_deck_from_slides" not in names


def test_presentation_toolset_uses_legacy_only_when_explicitly_disabled(monkeypatch) -> None:
    monkeypatch.setenv("SOPHIA_DECK_BUILD_SERVICE_ENABLED", "false")
    monkeypatch.setenv("SOPHIA_DECK_LEGACY_SCREENSHOT_DEBUG", "true")

    tools = build_builder_tools_for_task_type("presentation", vision_enabled=False)
    names = [getattr(tool, "name", "") for tool in tools]
    contract = assert_deck_tool_contract(tools, task_type="presentation", artifact_target_ext=".pptx")

    assert deck_build_service_enabled() is False
    assert contract is not None
    assert contract["route"] == "legacy_html_slide_to_pptx"
    assert "prepare_deck_build" not in names
    assert "prepare_pptx_image_manifest" in names
    assert "build_deck_from_slides" in names


def test_presentation_toolset_ignores_disabled_flag_without_debug(monkeypatch) -> None:
    monkeypatch.setenv("SOPHIA_DECK_BUILD_SERVICE_ENABLED", "false")
    monkeypatch.delenv("SOPHIA_DECK_LEGACY_SCREENSHOT_DEBUG", raising=False)

    names = [getattr(tool, "name", "") for tool in build_builder_tools_for_task_type("presentation", vision_enabled=False)]

    assert deck_build_service_enabled() is True
    assert "prepare_deck_build" in names
    assert "build_deck_from_slides" not in names


def test_presentation_toolset_forces_deck_service_in_production(monkeypatch) -> None:
    monkeypatch.setenv("SOPHIA_DECK_BUILD_SERVICE_ENABLED", "false")
    monkeypatch.setenv("SOPHIA_DECK_LEGACY_SCREENSHOT_DEBUG", "true")
    monkeypatch.setenv("RENDER", "true")

    names = [getattr(tool, "name", "") for tool in build_builder_tools_for_task_type("presentation", vision_enabled=False)]

    assert deck_build_service_enabled() is True
    assert "prepare_deck_build" in names
    assert "build_deck_from_slides" not in names


def test_pdf_slide_deck_uses_pdf_report_route_even_when_deck_service_default_enabled(monkeypatch) -> None:
    monkeypatch.delenv("SOPHIA_DECK_BUILD_SERVICE_ENABLED", raising=False)

    tools = build_builder_tools_for_task_type(
        "presentation",
        vision_enabled=False,
        artifact_target_ext=".pdf",
    )
    names = [getattr(tool, "name", "") for tool in tools]
    contract = assert_deck_tool_contract(tools, task_type="presentation", artifact_target_ext=".pdf")

    assert contract is None
    assert "prepare_deck_build" not in names
    assert "prepare_pptx_image_manifest" not in names
    assert "build_deck_from_slides" not in names
    assert "render_html_to_pdf" in names


def test_pdf_slide_deck_legacy_tools_are_not_rejected_by_deck_service_guard(monkeypatch) -> None:
    monkeypatch.delenv("SOPHIA_DECK_BUILD_SERVICE_ENABLED", raising=False)
    state = {
        "builder_artifact_target_path": f"{_OUTPUTS}deck.pdf",
        "delegation_context": {"task_type": "presentation"},
    }
    request = SimpleNamespace(
        tool_call={
            "id": "tc-manifest",
            "name": "prepare_pptx_image_manifest",
            "args": {"prompt_files": [f"{_OUTPUTS}assets/prompts/slide-01.json"]},
        },
        state=state,
        runtime=SimpleNamespace(context={}, config={}),
    )

    assert BuilderArtifactMiddleware._deck_build_service_legacy_tool_rejection(request) is None


def test_prepare_deck_build_failure_is_terminal_command(tmp_path: Path, monkeypatch) -> None:
    monkeypatch.setenv("SOPHIA_DECK_BUILD_SERVICE_ENABLED", "true")
    runtime = _runtime(tmp_path / "outputs")
    request = SimpleNamespace(
        tool_call={"id": "tc-deck", "name": "prepare_deck_build", "args": {}},
        state=runtime.state,
        runtime=runtime,
    )
    payload = {
        "success": False,
        "build_id": "deck-1",
        "deck_build_path": f"{_OUTPUTS}deck_build/build.json",
        "failure_code": "deck_visual_batch_startup_failed",
        "failure_summary": "Image batch did not emit IMAGEGEN_BATCH.",
        "retryable": False,
        "slide_count": 3,
        "expected_visual_count": 3,
        "successful_visual_count": 0,
        "referenced_visual_count": 0,
        "missing_visual_count": 3,
        "quality_status": "failed",
    }
    result = ToolMessage(content=json.dumps(payload), tool_call_id="tc-deck", name="prepare_deck_build")
    monkeypatch.setattr(BuilderArtifactMiddleware, "_upload_fallback_and_fire", lambda *args, **kwargs: None)

    command = BuilderArtifactMiddleware()._prepare_deck_build_result_command(request, result)

    assert isinstance(command, Command)
    assert command.goto == "end"
    assert command.update["builder_result"]["artifact_path"] is None
    assert command.update["builder_result"]["failure_code"] == "deck_visual_batch_startup_failed"
    diagnostics = command.update["builder_pptx_diagnostics"]
    assert diagnostics["deck_build_id"] == "deck-1"
    assert diagnostics["missing_expected_visual_count"] == 3


def test_prepare_deck_build_retryable_ir_failure_uses_normal_graph_edge(tmp_path: Path, monkeypatch) -> None:
    monkeypatch.setenv("SOPHIA_DECK_BUILD_SERVICE_ENABLED", "true")
    runtime = _runtime(tmp_path / "outputs")
    request = SimpleNamespace(
        tool_call={"id": "tc-deck", "name": "prepare_deck_build", "args": {}},
        state=runtime.state,
        runtime=runtime,
    )
    payload = {
        "success": False,
        "build_id": "deck-1",
        "deck_build_path": f"{_OUTPUTS}deck_build/build.json",
        "failure_code": "invalid_deck_ir",
        "failure_summary": "Slide 2 narrative is required and must be <= 280 chars.",
        "retryable": True,
        "slide_count": 0,
        "expected_visual_count": 0,
        "successful_visual_count": 0,
        "referenced_visual_count": 0,
        "missing_visual_count": 0,
        "quality_status": "failed",
    }
    result = ToolMessage(content=json.dumps(payload), tool_call_id="tc-deck", name="prepare_deck_build")

    command = BuilderArtifactMiddleware()._prepare_deck_build_result_command(request, result)

    assert isinstance(command, Command)
    assert not command.goto
    assert command.update["builder_deck_ir_repair_attempt_count"] == 1
    assert command.update["builder_last_deck_ir_failure"]["failure_code"] == "invalid_deck_ir"
    assert command.update["messages"] == [result]
    assert command.update["builder_deck_prepare_phase"] == "retry_pending"
    assert "prepare_deck_build exactly once more" in command.update["builder_deck_prepare_repair_message"]

    state = {**runtime.state, **command.update}
    repair_update = BuilderArtifactMiddleware().before_model(state, runtime)
    assert repair_update is not None
    assert "prepare_deck_build exactly once more" in repair_update["messages"][0].content


def test_prepare_deck_build_retryable_source_quality_failure_uses_repair_edge(
    tmp_path: Path,
    monkeypatch,
) -> None:
    monkeypatch.setenv("SOPHIA_DECK_BUILD_SERVICE_ENABLED", "true")
    runtime = _runtime(tmp_path / "outputs")
    request = SimpleNamespace(
        tool_call={"id": "tc-deck", "name": "prepare_deck_build", "args": {}},
        state=runtime.state,
        runtime=runtime,
    )
    repair_message = (
        "Repair every listed source-quality and mechanical issue.\n"
        "1. QUALITY slide:1 [chrome]: remove the eyebrow."
    )
    payload = {
        "success": False,
        "build_id": "deck-1",
        "deck_build_path": f"{_OUTPUTS}deck_build/build.json",
        "failure_code": "deck_source_quality_failed",
        "failure_summary": "chrome on slide:1: remove the eyebrow",
        "retryable": True,
        "slide_count": 3,
        "expected_visual_count": 0,
        "successful_visual_count": 0,
        "referenced_visual_count": 0,
        "missing_visual_count": 0,
        "quality_status": "failed",
        "repair_instruction": {"repair_message": repair_message},
    }
    result = ToolMessage(content=json.dumps(payload), tool_call_id="tc-deck", name="prepare_deck_build")

    command = BuilderArtifactMiddleware()._prepare_deck_build_result_command(request, result)

    assert isinstance(command, Command)
    assert not command.goto
    assert command.update["builder_deck_prepare_phase"] == "retry_pending"
    assert command.update["builder_deck_prepare_repair_message"] == repair_message
    assert command.update["builder_last_deck_creative_failure"]["failure_code"] == "deck_source_quality_failed"


def test_prepare_deck_build_second_source_quality_failure_exhausts_one_repair(
    tmp_path: Path,
    monkeypatch,
) -> None:
    monkeypatch.setenv("SOPHIA_DECK_BUILD_SERVICE_ENABLED", "true")
    runtime = _runtime(tmp_path / "outputs")
    runtime.state["builder_deck_prepare_repair_attempt_count"] = 1
    request = SimpleNamespace(
        tool_call={"id": "tc-deck", "name": "prepare_deck_build", "args": {}},
        state=runtime.state,
        runtime=runtime,
    )
    payload = {
        "success": False,
        "build_id": "deck-1",
        "deck_build_path": f"{_OUTPUTS}deck_build/build.json",
        "failure_code": "deck_source_quality_failed",
        "failure_summary": "chrome on slide:1: remove the eyebrow",
        "retryable": True,
        "slide_count": 3,
        "expected_visual_count": 0,
        "successful_visual_count": 0,
        "referenced_visual_count": 0,
        "missing_visual_count": 0,
        "quality_status": "failed",
    }
    result = ToolMessage(content=json.dumps(payload), tool_call_id="tc-deck", name="prepare_deck_build")
    monkeypatch.setattr(BuilderArtifactMiddleware, "_upload_fallback_and_fire", lambda *args, **kwargs: None)

    command = BuilderArtifactMiddleware()._prepare_deck_build_result_command(request, result)

    assert isinstance(command, Command)
    assert command.goto == "end"
    artifact = command.update["builder_result"]
    assert artifact["failure_code"] == "deck_prepare_retry_exhausted"
    assert artifact["root_failure_code"] == "deck_source_quality_failed"
    assert artifact["last_prepare_failure_code"] == "deck_source_quality_failed"


def test_prepare_deck_build_retryable_ir_second_failure_is_terminal(tmp_path: Path, monkeypatch) -> None:
    monkeypatch.setenv("SOPHIA_DECK_BUILD_SERVICE_ENABLED", "true")
    runtime = _runtime(tmp_path / "outputs")
    runtime.state["builder_deck_ir_repair_attempt_count"] = 1
    request = SimpleNamespace(
        tool_call={"id": "tc-deck", "name": "prepare_deck_build", "args": {}},
        state=runtime.state,
        runtime=runtime,
    )
    payload = {
        "success": False,
        "build_id": "deck-1",
        "deck_build_path": f"{_OUTPUTS}deck_build/build.json",
        "failure_code": "invalid_deck_ir",
        "failure_summary": "Slide 2 narrative is required and must be <= 280 chars.",
        "retryable": True,
        "slide_count": 0,
        "expected_visual_count": 0,
        "successful_visual_count": 0,
        "referenced_visual_count": 0,
        "missing_visual_count": 0,
        "quality_status": "failed",
    }
    result = ToolMessage(content=json.dumps(payload), tool_call_id="tc-deck", name="prepare_deck_build")
    monkeypatch.setattr(BuilderArtifactMiddleware, "_upload_fallback_and_fire", lambda *args, **kwargs: None)

    command = BuilderArtifactMiddleware()._prepare_deck_build_result_command(request, result)

    assert isinstance(command, Command)
    assert command.goto == "end"
    artifact = command.update["builder_result"]
    assert artifact["failure_code"] == "deck_prepare_retry_exhausted"
    assert artifact["root_failure_code"] == "invalid_deck_ir"
    assert artifact["last_prepare_failure_code"] == "invalid_deck_ir"


def test_prepare_deck_build_schema_failure_gets_one_bounded_retry(tmp_path: Path) -> None:
    runtime = _runtime(tmp_path / "outputs")
    runtime.state["builder_pptx_diagnostics"] = {"prepare_emitted_call_count": 1, "prepare_call_count": 1}
    request = SimpleNamespace(
        tool_call={"id": "tc-schema-1", "name": "prepare_deck_build", "args": {}},
        state=runtime.state,
        runtime=runtime,
    )
    result = ToolMessage(
        content="creative_plan.slide_compositions.0.headline_intent: Field required",
        tool_call_id="tc-schema-1",
        name="prepare_deck_build",
        status="error",
    )

    command = BuilderArtifactMiddleware()._prepare_deck_build_result_command(request, result)

    assert isinstance(command, Command)
    assert not command.goto
    assert command.update["builder_deck_prepare_phase"] == "retry_pending"
    diagnostics = command.update["builder_pptx_diagnostics"]
    assert diagnostics["prepare_schema_failure_count"] == 1
    assert diagnostics["prepare_result_count"] == 1
    assert diagnostics["prepare_repair_count"] == 0
    assert diagnostics["prepare_retry_executed"] is True
    assert "prepare_service_call_count" not in diagnostics
    assert diagnostics["deck_root_failure_code"] == "deck_prepare_argument_invalid"
    assert "builder_deck_prepare_repair_attempt_count" not in command.update
    assert "exactly two repair_anchor_ids per slide" in command.update["builder_deck_prepare_repair_message"]


def test_prepare_schema_retry_recovers_all_size_targets_without_tool_metadata(
    tmp_path: Path,
) -> None:
    runtime = _runtime(tmp_path / "outputs")
    runtime.state["builder_pptx_diagnostics"] = {
        "prepare_emitted_call_count": 1,
        "prepare_call_count": 1,
    }
    body_sizes = [1173, 6244, 2896, 6466, 2625]
    template = _compact_slides()[0]
    slides = [
        {
            **template,
            "title": f"Compact {index}",
            "html_body": "x" * body_size,
        }
        for index, body_size in enumerate(body_sizes, start=1)
    ]
    args = {
        "deck_title": "Compact Limit Diagnostics",
        "slides": '<parameter name="_arr">\n' + json.dumps(slides),
        "output_path": "/mnt/user-data/outputs/deck.pptx",
        "creative_plan": _creative_plan(),
        "authoring_contract": "compact_model_html_v2",
        "deck_stylesheet": "main { background: #101828; }",
    }
    request = SimpleNamespace(
        tool_call={"id": "tc-schema-sizes", "name": "prepare_deck_build", "args": args},
        state=runtime.state,
        runtime=runtime,
    )
    result = ToolMessage(
        content="Generic LangChain tool validation failure.",
        tool_call_id="tc-schema-sizes",
        name="prepare_deck_build",
        status="error",
    )

    command = BuilderArtifactMiddleware()._prepare_deck_build_result_command(request, result)

    assert isinstance(command, Command)
    assert not command.goto
    assert command.update["builder_deck_prepare_phase"] == "retry_pending"
    slide_two = "slides[1].html_body is 6244 bytes; compact-v2 hard limit is 6144 bytes"
    slide_four = "slides[3].html_body is 6466 bytes; compact-v2 hard limit is 6144 bytes"
    repair_message = command.update["builder_deck_prepare_repair_message"]
    diagnostics = command.update["builder_pptx_diagnostics"]
    for target in (slide_two, slide_four):
        assert target in repair_message
        assert target in diagnostics["deck_root_failure_summary"]
        assert target in diagnostics["last_prepare_failure_summary"]
    assert "Generic LangChain" not in repair_message


def test_prepare_deck_build_execution_error_is_terminal_without_repair(
    tmp_path: Path,
    monkeypatch,
) -> None:
    runtime = _runtime(tmp_path / "outputs")
    runtime.state["builder_pptx_diagnostics"] = {
        "prepare_emitted_call_count": 1,
        "prepare_call_count": 1,
    }
    request = SimpleNamespace(
        tool_call={"id": "tc-runtime", "name": "prepare_deck_build", "args": {}},
        state=runtime.state,
        runtime=runtime,
    )
    result = ToolMessage(
        content="The authoritative deck build tool failed during execution.",
        tool_call_id="tc-runtime",
        name="prepare_deck_build",
        status="error",
        additional_kwargs={
            "tool_error": {
                "error_class": "TypeError",
                "retryable": False,
                "stage": "tool_execution",
            }
        },
    )
    monkeypatch.setattr(BuilderArtifactMiddleware, "_upload_fallback_and_fire", lambda *args, **kwargs: None)

    command = BuilderArtifactMiddleware()._prepare_deck_build_result_command(request, result)

    assert isinstance(command, Command)
    assert command.goto == "end"
    assert command.update["builder_deck_prepare_phase"] == "terminal"
    assert command.update["builder_result"]["failure_code"] == "deck_prepare_execution_error"
    diagnostics = command.update["builder_pptx_diagnostics"]
    assert diagnostics["prepare_execution_count"] == 1
    assert diagnostics["prepare_result_count"] == 1
    assert "prepare_schema_failure_count" not in diagnostics


def test_prepare_deck_build_second_schema_failure_preserves_root_cause(
    tmp_path: Path,
    monkeypatch,
) -> None:
    runtime = _runtime(tmp_path / "outputs")
    runtime.state["builder_pptx_diagnostics"] = {
        "prepare_emitted_call_count": 2,
        "prepare_call_count": 2,
        "prepare_result_count": 1,
        "prepare_schema_failure_count": 1,
        "deck_root_failure_code": "deck_prepare_argument_invalid",
        "deck_root_failure_summary": "The first call omitted headline_intent.",
    }
    request = SimpleNamespace(
        tool_call={"id": "tc-schema-2", "name": "prepare_deck_build", "args": {}},
        state=runtime.state,
        runtime=runtime,
    )
    result = ToolMessage(
        content="creative_plan.slide_compositions.0.layout_name: Field required",
        tool_call_id="tc-schema-2",
        name="prepare_deck_build",
        status="error",
    )
    monkeypatch.setattr(BuilderArtifactMiddleware, "_upload_fallback_and_fire", lambda *args, **kwargs: None)

    command = BuilderArtifactMiddleware()._prepare_deck_build_result_command(request, result)

    assert isinstance(command, Command)
    assert command.goto == "end"
    artifact = command.update["builder_result"]
    assert artifact["failure_code"] == "deck_prepare_retry_exhausted"
    assert artifact["root_failure_code"] == "deck_prepare_argument_invalid"
    assert artifact["root_failure_summary"] == "The first call omitted headline_intent."
    assert artifact["last_prepare_failure_code"] == "deck_prepare_argument_invalid"
    assert artifact["prepare_repair_count"] == 0
    assert artifact["prepare_schema_failure_count"] == 2
    assert artifact.get("prepare_service_call_count") in {None, 0}


def test_third_prepare_call_is_rejected_before_service_execution(tmp_path: Path, monkeypatch) -> None:
    runtime = _runtime(tmp_path / "outputs")
    runtime.state["builder_deck_prepare_repair_attempt_count"] = 1
    runtime.state["builder_deck_creative_repair_attempt_count"] = 1
    runtime.state["builder_pptx_diagnostics"] = {
        "prepare_emitted_call_count": 3,
        "prepare_call_count": 3,
        "prepare_service_result_count": 1,
        "prepare_repair_count": 1,
        "prepare_retry_executed": True,
        "deck_root_failure_code": "deck_creative_plan_invalid",
        "deck_root_failure_summary": "The first plan failed critique validation.",
    }
    request = SimpleNamespace(
        tool_call={"id": "tc-third", "name": "prepare_deck_build", "args": {}},
        state=runtime.state,
        runtime=runtime,
    )
    monkeypatch.setattr(BuilderArtifactMiddleware, "_upload_fallback_and_fire", lambda *args, **kwargs: None)

    command = BuilderArtifactMiddleware()._prepare_deck_build_exhausted_command(request)

    assert command is not None
    assert command.goto == "end"
    assert command.update["builder_result"]["failure_code"] == "deck_prepare_retry_exhausted"
    assert command.update["builder_result"]["root_failure_code"] == "deck_creative_plan_invalid"
    assert "prepare_call_count" not in command.update["builder_pptx_diagnostics"]


def test_third_prepare_call_after_schema_failure_is_allowed_for_quality_repair(
    tmp_path: Path,
) -> None:
    runtime = _runtime(tmp_path / "outputs")
    runtime.state.update(
        {
            "builder_deck_prepare_repair_attempt_count": 1,
            "builder_deck_creative_repair_attempt_count": 1,
            "builder_pptx_diagnostics": {
                "prepare_emitted_call_count": 3,
                "prepare_call_count": 3,
                "prepare_schema_failure_count": 1,
                "prepare_service_result_count": 1,
                "prepare_repair_count": 1,
                "prepare_retry_executed": True,
            },
        }
    )
    request = SimpleNamespace(
        tool_call={"id": "tc-third", "name": "prepare_deck_build", "args": {}},
        state=runtime.state,
        runtime=runtime,
    )

    command = BuilderArtifactMiddleware()._prepare_deck_build_exhausted_command(request)

    assert command is None


def test_fourth_prepare_call_is_rejected_after_both_corrections(
    tmp_path: Path,
    monkeypatch,
) -> None:
    runtime = _runtime(tmp_path / "outputs")
    runtime.state.update(
        {
            "builder_deck_prepare_repair_attempt_count": 1,
            "builder_deck_creative_repair_attempt_count": 1,
            "builder_pptx_diagnostics": {
                "prepare_emitted_call_count": 4,
                "prepare_call_count": 4,
                "prepare_schema_failure_count": 1,
                "prepare_service_result_count": 1,
                "prepare_repair_count": 1,
                "prepare_retry_executed": True,
                "deck_root_failure_code": "deck_prepare_argument_invalid",
                "deck_root_failure_summary": "The first call failed typed validation.",
            },
        }
    )
    request = SimpleNamespace(
        tool_call={"id": "tc-fourth", "name": "prepare_deck_build", "args": {}},
        state=runtime.state,
        runtime=runtime,
    )
    monkeypatch.setattr(
        BuilderArtifactMiddleware,
        "_upload_fallback_and_fire",
        lambda *args, **kwargs: None,
    )

    command = BuilderArtifactMiddleware()._prepare_deck_build_exhausted_command(request)

    assert command is not None
    assert command.goto == "end"
    artifact = command.update["builder_result"]
    assert artifact["failure_code"] == "deck_prepare_retry_exhausted"
    assert artifact["root_failure_code"] == "deck_prepare_argument_invalid"


def test_second_service_result_is_rejected_after_quality_budget_is_spent(
    tmp_path: Path,
    monkeypatch,
) -> None:
    runtime = _runtime(tmp_path / "outputs")
    runtime.state.update(
        {
            "builder_deck_prepare_repair_attempt_count": 1,
            "builder_deck_creative_repair_attempt_count": 1,
            "builder_pptx_diagnostics": {
                "prepare_emitted_call_count": 3,
                "prepare_call_count": 3,
                "prepare_schema_failure_count": 1,
                "prepare_service_result_count": 2,
                "prepare_repair_count": 1,
                "prepare_retry_executed": True,
            },
        }
    )
    request = SimpleNamespace(
        tool_call={"id": "tc-after-service-cap", "name": "prepare_deck_build", "args": {}},
        state=runtime.state,
        runtime=runtime,
    )
    monkeypatch.setattr(
        BuilderArtifactMiddleware,
        "_upload_fallback_and_fire",
        lambda *args, **kwargs: None,
    )

    command = BuilderArtifactMiddleware()._prepare_deck_build_exhausted_command(request)

    assert command is not None
    assert command.goto == "end"
    assert command.update["builder_result"]["failure_code"] == "deck_prepare_retry_exhausted"


def test_legacy_schema_only_retry_state_does_not_consume_quality_repair() -> None:
    state = {
        "builder_deck_prepare_repair_attempt_count": 1,
        "builder_pptx_diagnostics": {
            "prepare_schema_failure_count": 1,
            "prepare_repair_count": 1,
            "prepare_retry_executed": True,
        },
    }

    repair_count = BuilderArtifactMiddleware._prepare_repair_attempt_count(state)

    assert repair_count == 0


def test_legacy_service_retry_state_still_consumes_quality_repair() -> None:
    state = {
        "builder_pptx_diagnostics": {
            "prepare_repair_count": 1,
            "prepare_retry_executed": True,
        },
    }

    repair_count = BuilderArtifactMiddleware._prepare_repair_attempt_count(state)

    assert repair_count == 1


def test_parallel_prepare_calls_are_rejected_before_service_execution(tmp_path: Path, monkeypatch) -> None:
    runtime = _runtime(tmp_path / "outputs")
    calls = [
        {"id": "tc-parallel-1", "name": "prepare_deck_build", "args": {}},
        {"id": "tc-parallel-2", "name": "prepare_deck_build", "args": {}},
    ]
    update = BuilderArtifactMiddleware._prepare_call_after_model_update(runtime.state, calls)
    runtime.state.update(update)
    request = SimpleNamespace(
        tool_call=calls[0],
        state=runtime.state,
        runtime=runtime,
    )
    monkeypatch.setattr(BuilderArtifactMiddleware, "_upload_fallback_and_fire", lambda *args, **kwargs: None)

    command = BuilderArtifactMiddleware()._prepare_deck_build_exhausted_command(request)

    assert command is not None
    assert command.goto == "end"
    artifact = command.update["builder_result"]
    assert artifact["failure_code"] == "deck_prepare_parallel_calls_forbidden"
    assert artifact["root_failure_code"] == "deck_prepare_parallel_calls_forbidden"
    assert artifact["prepare_parallel_call_count"] == 2
    assert artifact.get("prepare_service_call_count") is None


def test_pptx_terminal_outcome_prefers_deck_failure_payload_counts(monkeypatch) -> None:
    captured: dict[str, dict] = {}

    def capture_span(name: str, **kwargs):
        captured["name"] = name
        captured.update(kwargs)

    monkeypatch.setattr(builder_artifact_module, "_safe_langsmith_span", capture_span)
    state = {
        "builder_artifact_target_path": f"{_OUTPUTS}deck.pptx",
        "builder_pptx_requested_slide_count": 6,
        "builder_pptx_diagnostics": {
            "expected_generated_visual_count": 6,
            "successful_generated_visual_count": 0,
            "referenced_visual_count": 0,
            "missing_expected_visual_count": 6,
            "image_generation_status": "failed",
        },
    }
    artifact = {
        "artifact_path": None,
        "deck_route": "deck_ir_html_raster",
        "deck_compile_mode": "native_html2patch",
        "native_editability_score": 1.0,
        "native_text_shape_count": 12,
        "picture_shape_count": 6,
        "full_slide_picture_count": 1,
        "image_generation_status": "partial",
        "successful_generated_visual_count": 6,
        "referenced_visual_count": 6,
        "missing_expected_visual_count": 0,
        "quality_warning": "native_full_bleed_picture_present",
    }

    builder_artifact_module._trace_pptx_terminal_outcome(
        state=state,
        artifact=artifact,
        status="error",
        failure_code="deck_native_full_slide_picture_forbidden",
    )

    outputs = captured["outputs"]
    assert outputs["image_generation_status"] == "partial"
    assert outputs["successful_generated_visual_count"] == 6
    assert outputs["referenced_visual_count"] == 6
    assert outputs["missing_expected_visual_count"] == 0
    assert outputs["native_editability_score"] == 1.0


def test_prepare_deck_build_missing_success_output_is_terminal_command(tmp_path: Path, monkeypatch) -> None:
    monkeypatch.setenv("SOPHIA_DECK_BUILD_SERVICE_ENABLED", "true")
    runtime = _runtime(tmp_path / "outputs")
    request = SimpleNamespace(
        tool_call={"id": "tc-deck", "name": "prepare_deck_build", "args": {}},
        state=runtime.state,
        runtime=runtime,
    )
    payload = {
        "success": True,
        "build_id": "deck-1",
        "deck_build_path": f"{_OUTPUTS}deck_build/build.json",
        "pptx_path": f"{_OUTPUTS}missing-deck.pptx",
        "slide_count": 3,
        "expected_visual_count": 3,
        "successful_visual_count": 3,
        "referenced_visual_count": 3,
        "missing_visual_count": 0,
        "quality_status": "passed",
    }
    result = ToolMessage(content=json.dumps(payload), tool_call_id="tc-deck", name="prepare_deck_build")
    monkeypatch.setattr(BuilderArtifactMiddleware, "_upload_fallback_and_fire", lambda *args, **kwargs: None)

    command = BuilderArtifactMiddleware()._prepare_deck_build_result_command(request, result)

    assert isinstance(command, Command)
    assert command.goto == "end"
    assert command.update["builder_result"]["artifact_path"] is None
    assert command.update["builder_result"]["failure_code"] == "missing_output"
    diagnostics = command.update["builder_pptx_diagnostics"]
    assert diagnostics["deck_status"] == "failed_terminal"
    assert diagnostics["deck_failure_code"] == "missing_output"
    assert diagnostics["pptx_generator_success_count"] == 0
