from __future__ import annotations

import importlib.util
import json
import shutil
import subprocess
from pathlib import Path

import pytest
from PIL import Image, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt

from deerflow.sophia.deck_native import DeckNativeService, native_mechanical_report
from deerflow.sophia.deck_native import service as native_service_module

PROJECT_ROOT = Path(__file__).resolve().parents[2]
HTML2PATCH_PATH = (
    PROJECT_ROOT
    / "third_party/hands_on_deck/skills/hands-on-deck/scripts/html2patch.py"
)


def _html2patch_module():
    spec = importlib.util.spec_from_file_location("deck_native_html2patch", HTML2PATCH_PATH)
    assert spec is not None and spec.loader is not None
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def _compiler_text_item(
    *,
    x: float,
    y: float,
    width: float,
    height: float = 60,
    align: str = "left",
    font: str = "Arial",
    text: str = "Grid label",
) -> dict:
    return {
        "type": "text",
        "box": {"x": x, "y": y, "w": width, "h": height},
        "rotation": 0,
        "paragraphs": [
            {
                "runs": [
                    {
                        "text": text,
                        "style": {
                            "font": font,
                            "sizePx": 24,
                            "color": "111827",
                            "alpha": 1,
                            "bold": False,
                            "italic": False,
                            "underline": False,
                            "link": None,
                        },
                    }
                ]
            }
        ],
        "meta": {
            "align": align,
            "lineHeightPx": 28,
            "fontSizePx": 24,
            "padding": [0, 0, 0, 0],
        },
    }


def test_deck_native_preflight_reports_missing_scripts(tmp_path: Path) -> None:
    service = DeckNativeService(scripts_dir=tmp_path / "missing-scripts")

    result = service.preflight()

    assert result.success is False
    assert result.scripts_dir_exists is False
    assert result.deck_py_exists is False
    assert result.html2patch_py_exists is False
    assert "deck.py" in "\n".join(result.errors)


def test_deck_native_subprocess_honors_expired_parent_deadline(
    tmp_path: Path,
    monkeypatch,
) -> None:
    script = tmp_path / "deck.py"
    script.write_text("raise SystemExit(0)\n", encoding="utf-8")
    service = DeckNativeService(scripts_dir=tmp_path)
    service.set_deadline_epoch_ms(1_000)
    monkeypatch.setattr(native_service_module.time, "time", lambda: 2.0)

    result = service._run(["python", str(script)])

    assert result.returncode == 124
    assert "deck deadline exceeded" in result.stderr


def test_deck_native_subprocess_timeout_uses_process_group_runner(
    tmp_path: Path,
    monkeypatch,
) -> None:
    script = tmp_path / "deck.py"
    script.write_text("raise SystemExit(0)\n", encoding="utf-8")
    captured: dict = {}

    def _timeout(command, *, timeout, cwd, writable_files, writable_dirs):
        captured.update(
            command=command,
            timeout=timeout,
            cwd=cwd,
            writable_files=writable_files,
            writable_dirs=writable_dirs,
        )
        raise subprocess.TimeoutExpired(command, timeout, output="partial", stderr="hung child")

    monkeypatch.setattr(native_service_module, "run_process_group", _timeout)

    result = DeckNativeService(scripts_dir=tmp_path)._run(["python", str(script)], timeout=30)

    assert captured == {
        "command": ["python", str(script)],
        "timeout": 30,
        "cwd": tmp_path,
        "writable_files": (),
        "writable_dirs": (),
    }
    assert result.returncode == 124
    assert result.stdout == "partial"
    assert "hands-on-deck subprocess timed out after 30s" in result.stderr
    assert "hung child" in result.stderr


def test_deck_native_default_subprocess_timeout_supports_long_decks(
    tmp_path: Path,
    monkeypatch,
) -> None:
    script = tmp_path / "deck.py"
    script.write_text("raise SystemExit(0)\n", encoding="utf-8")
    captured: dict = {}

    def _complete(command, *, timeout, cwd, writable_files, writable_dirs):
        captured.update(
            command=command,
            timeout=timeout,
            cwd=cwd,
            writable_files=writable_files,
            writable_dirs=writable_dirs,
        )
        return subprocess.CompletedProcess(command, 0, stdout="", stderr="")

    monkeypatch.setattr(native_service_module, "run_process_group", _complete)

    result = DeckNativeService(scripts_dir=tmp_path)._run(["python", str(script)])

    assert result.returncode == 0
    assert captured["timeout"] == 600


def test_deck_native_render_clears_stale_images_before_counting_success(
    tmp_path: Path,
    monkeypatch,
) -> None:
    pptx = tmp_path / "deck.pptx"
    pptx.touch()
    render_dir = tmp_path / "rendered"
    render_dir.mkdir()
    stale_render = render_dir / "slide-1.jpg"
    stale_render.write_bytes(b"stale")
    service = DeckNativeService(scripts_dir=tmp_path)
    monkeypatch.setattr(
        service,
        "_run",
        lambda *_args, **_kwargs: subprocess.CompletedProcess([], 0, stdout="", stderr=""),
    )

    result = service.render(pptx_path=str(pptx), output_dir=str(render_dir))

    assert stale_render.exists() is False
    assert result.success is False
    assert result.rendered_slide_count == 0


def test_deck_native_render_requires_every_requested_slide(
    tmp_path: Path,
    monkeypatch,
) -> None:
    pptx = tmp_path / "deck.pptx"
    pptx.touch()
    render_dir = tmp_path / "rendered"
    service = DeckNativeService(scripts_dir=tmp_path)

    def _partial_render(*_args, **_kwargs):
        render_dir.mkdir(exist_ok=True)
        (render_dir / "slide-1.jpg").write_bytes(b"current")
        return subprocess.CompletedProcess([], 0, stdout="", stderr="")

    monkeypatch.setattr(service, "_run", _partial_render)

    result = service.render(
        pptx_path=str(pptx),
        output_dir=str(render_dir),
        slides=[0, 1],
    )

    assert result.success is False
    assert result.rendered_slide_count == 1
    assert result.errors == ["native render incomplete: expected 2 slide image(s), rendered 1"]


def _wide_base_deck(path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(20)
    presentation.slide_height = Inches(11.25)
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def _native_text_patch(path: Path, *, text: str = "Hello Native") -> None:
    path.write_text(
        json.dumps(
            {
                "ops": [
                    {"op": "add-slide", "layout": "Blank"},
                    {
                        "op": "add-shape",
                        "slide": 0,
                        "kind": "textbox",
                        "at": [1, 1],
                        "size": [8, 0.8],
                        "text": [text],
                        "name": "title",
                    },
                ]
            }
        ),
        encoding="utf-8",
    )


def test_deck_native_apply_inspect_render_and_diff(tmp_path: Path) -> None:
    if not shutil.which("soffice") or not shutil.which("pdftoppm"):
        pytest.skip("hands-on-deck render requires LibreOffice and Poppler")
    base = tmp_path / "base.pptx"
    patch = tmp_path / "patch.json"
    output = tmp_path / "out.pptx"
    render_dir = tmp_path / "rendered"
    _wide_base_deck(base)
    _native_text_patch(patch)
    service = DeckNativeService()

    applied = service.apply_patch(base_deck_path=str(base), patch_path=str(patch), output_path=str(output), fix=True)
    inspected = service.inspect(str(output))
    rendered = service.render(pptx_path=str(output), output_dir=str(render_dir), slides=[0])
    diff = service.diff(before_path=str(base), after_path=str(output))

    assert applied.success is True
    assert output.is_file()
    assert inspected.success is True
    assert inspected.native_text_shape_count > 0
    assert inspected.full_slide_picture_count == 0
    assert inspected.native_editability_score >= 0.60
    assert Path(inspected.raw_json_path or "").parent == output.parent / ".builder" / "deck_native" / "inspect"
    assert Path(inspected.shape_inventory_path or "").parent == output.parent / ".builder" / "deck_native" / "inspect"
    assert not output.with_name("out.inspect.json").exists()
    assert not output.with_name("out.shape-inventory.json").exists()
    assert rendered.success is True
    assert rendered.rendered_slide_count == 1
    assert diff["success"] is True
    assert diff["changed"] is True


def test_html2patch_clamps_text_metric_cushion_to_source_layout(tmp_path: Path) -> None:
    module = _html2patch_module()
    columns = [
        _compiler_text_item(x=96, y=200, width=480, text="Baseline"),
        _compiler_text_item(x=576, y=200, width=480, text="Mechanism"),
        _compiler_text_item(x=1056, y=200, width=480, text="Outcome"),
    ]
    centered = _compiler_text_item(x=200, y=40, width=400, align="center", text="Centered")
    right = _compiler_text_item(x=800, y=360, width=400, align="right", text="Right")
    serif = _compiler_text_item(x=200, y=500, width=400, font="Georgia", text="Serif")
    slide_edge = _compiler_text_item(x=1600, y=650, width=320, text="Edge")
    extract = {
        "body": {"w": 1920, "h": 1080},
        "items": [*columns, centered, right, serif, slide_edge],
    }

    operations = module.compile_page(
        extract,
        0,
        tmp_path / "slide.html",
        tmp_path,
        "s1",
        [],
        {"schema_version": "sophia-deck-source-map/v1", "slides": {}},
    )

    first, middle, last, centered_op, right_op, serif_op, edge_op = operations
    assert first["at"][0] + first["size"][0] == pytest.approx(middle["at"][0])
    assert middle["at"][0] + middle["size"][0] == pytest.approx(last["at"][0])
    assert first["size"][0] == pytest.approx(5.0)
    assert middle["size"][0] == pytest.approx(5.0)
    assert last["size"][0] == pytest.approx(5.1)
    assert centered_op["at"][0] + centered_op["size"][0] / 2 == pytest.approx(
        400 / 96,
        abs=0.001,
    )
    assert centered_op["size"][0] == pytest.approx(4.25)
    assert right_op["at"][0] + right_op["size"][0] == pytest.approx(1200 / 96)
    assert right_op["size"][0] == pytest.approx(4.25)
    assert serif_op["size"][0] == pytest.approx(416 / 96, abs=0.001)
    assert edge_op["size"][0] == pytest.approx(320 / 96, abs=0.001)
    assert {
        operation["text"][0]["font_size"]
        for operation in operations
    } == {18.0}

    base = tmp_path / "compiler-base.pptx"
    patch = tmp_path / "compiler.patch.json"
    output = tmp_path / "compiler-output.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(20)
    presentation.slide_height = Inches(11.25)
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(base)
    patch.write_text(json.dumps({"ops": operations}), encoding="utf-8")
    service = DeckNativeService()
    applied = service.apply_patch(
        base_deck_path=str(base),
        patch_path=str(patch),
        output_path=str(output),
        fix=False,
    )
    linted = service.lint_fix(pptx_path=str(output), touched_slides=[0])
    assert applied.success is True
    assert linted.residue_count == 0
    compiled = Presentation(output)
    assert {
        run.font.size.pt
        for shape in compiled.slides[0].shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
    } == {18.0}


def test_html2patch_does_not_increase_source_authored_text_overlap(tmp_path: Path) -> None:
    module = _html2patch_module()
    left = _compiler_text_item(x=100, y=100, width=300, text="Left")
    right = _compiler_text_item(x=390, y=100, width=300, text="Right")
    operations = module.compile_page(
        {"body": {"w": 1920, "h": 1080}, "items": [left, right]},
        0,
        tmp_path / "slide.html",
        tmp_path,
        "s1",
        [],
        {"schema_version": "sophia-deck-source-map/v1", "slides": {}},
    )

    left_op, right_op = operations
    compiled_overlap = (
        left_op["at"][0] + left_op["size"][0] - right_op["at"][0]
    )
    assert compiled_overlap == pytest.approx(10 / 96, abs=0.001)
    assert left_op["size"][0] == pytest.approx(300 / 96, abs=0.001)


def test_html2patch_splits_peer_gap_between_facing_expansions(tmp_path: Path) -> None:
    module = _html2patch_module()
    centered_left = _compiler_text_item(
        x=100,
        y=100,
        width=400,
        align="center",
        text="Centered left",
    )
    centered_right = _compiler_text_item(
        x=506,
        y=100,
        width=400,
        align="center",
        text="Centered right",
    )
    left_facing = _compiler_text_item(
        x=100,
        y=300,
        width=400,
        align="left",
        text="Left aligned",
    )
    right_facing = _compiler_text_item(
        x=506,
        y=300,
        width=400,
        align="right",
        text="Right aligned",
    )
    operations = module.compile_page(
        {
            "body": {"w": 1920, "h": 1080},
            "items": [centered_left, centered_right, left_facing, right_facing],
        },
        0,
        tmp_path / "slide.html",
        tmp_path,
        "s1",
        [],
        {"schema_version": "sophia-deck-source-map/v1", "slides": {}},
    )

    center_left_op, center_right_op, left_op, right_op = operations
    assert center_left_op["at"][0] + center_left_op["size"][0] <= center_right_op["at"][0] + 0.001
    assert left_op["at"][0] + left_op["size"][0] <= right_op["at"][0] + 0.001
    assert {
        operation["text"][0]["font_size"]
        for operation in operations
    } == {18.0}


def test_deck_native_invalid_patch_fails_atomically(tmp_path: Path) -> None:
    base = tmp_path / "base.pptx"
    patch = tmp_path / "invalid.patch.json"
    output = tmp_path / "should-not-exist.pptx"
    _wide_base_deck(base)
    patch.write_text(
        json.dumps({"ops": [{"op": "add-shape", "slide": 99, "kind": "textbox", "at": [1, 1], "size": [1, 1], "text": ["bad"]}]}),
        encoding="utf-8",
    )

    result = DeckNativeService().apply_patch(base_deck_path=str(base), patch_path=str(patch), output_path=str(output))

    assert result.success is False
    assert result.validation_error_count >= 1
    assert not output.exists()


def test_deck_native_lint_fix_reports_residue(tmp_path: Path) -> None:
    base = tmp_path / "base.pptx"
    patch = tmp_path / "covered-text.patch.json"
    output = tmp_path / "covered-text.pptx"
    image = tmp_path / "cover.png"
    _wide_base_deck(base)
    Image.new("RGB", (100, 100), "red").save(image)
    patch.write_text(
        json.dumps(
            {
                "ops": [
                    {"op": "add-slide", "layout": "Blank"},
                    {"op": "add-shape", "slide": 0, "kind": "textbox", "at": [1, 1], "size": [5, 1], "text": ["Hidden text"]},
                    {"op": "add-picture", "slide": 0, "image": str(image), "at": [1, 1], "size": [5, 1]},
                ]
            }
        ),
        encoding="utf-8",
    )
    service = DeckNativeService()
    applied = service.apply_patch(base_deck_path=str(base), patch_path=str(patch), output_path=str(output), fix=False)

    fixed = service.lint_fix(pptx_path=str(output), touched_slides=[0])

    assert applied.success is True
    assert fixed.success is True
    assert fixed.residue_count == 1
    assert fixed.residue[0]["shape"].startswith("s")
    assert fixed.residue_kinds


def test_deck_native_lint_fix_promotes_every_remaining_issue_to_residue(
    tmp_path: Path,
    monkeypatch,
) -> None:
    output = tmp_path / "remaining-issues.pptx"
    Presentation().save(output)
    payload = {
        "fixed": [],
        "residue": [],
        "remaining_issue_shapes": [
            "slide 1 s7: {'frame_overflow_bottom': 0.47}",
            "slide 2 s11: {'misaligned': ['vcenter edge 0.09 off gridline']}",
        ],
        "remaining_issues": [
            {
                "slide": 1,
                "shape": "s7",
                "issues": {"frame_overflow_bottom": 0.47},
            },
            {
                "slide": 2,
                "shape": "s11",
                "issues": {"misaligned": ["vcenter edge 0.09 off gridline"]},
            },
        ],
    }
    service = DeckNativeService()
    monkeypatch.setattr(
        service,
        "_run",
        lambda *_args, **_kwargs: subprocess.CompletedProcess(
            [],
            0,
            stdout=json.dumps(payload),
            stderr="",
        ),
    )

    result = service.lint_fix(pptx_path=str(output), touched_slides=[0])

    assert result.success is True
    assert result.lint_issue_count_before == 2
    assert result.fix_applied_count == 0
    assert result.residue_count == 2
    assert result.residue_kinds == {"frame_overflow": 1, "misaligned": 1}
    assert {(item["slide"], item["shape"]) for item in result.residue} == {
        (1, "s7"),
        (2, "s11"),
    }


def test_deck_native_lint_fix_repairs_compatible_alignment_geometry(
    tmp_path: Path,
) -> None:
    output = tmp_path / "production-alignment.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(20)
    presentation.slide_height = Inches(11.25)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    for index, left in enumerate((1.0, 4.5, 8.0, 11.5)):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left),
            Inches(3.13),
            Inches(2.9),
            Inches(3.88),
        )
        shape.name = f"alignment-peer-{index}"
    target = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(15.0),
        Inches(3.13),
        Inches(2.9),
        Inches(3.8),
    )
    target.name = "production-alignment-target"
    presentation.save(output)

    service = DeckNativeService()
    fixed = service.lint_fix(pptx_path=str(output), touched_slides=[0])

    assert fixed.success is True
    assert fixed.lint_issue_count_before == 1
    assert fixed.fix_applied_count == 1
    assert fixed.issue_kinds == {"align-y": 1}
    assert fixed.residue_count == 0
    repaired = Presentation(output)
    repaired_target = next(shape for shape in repaired.slides[0].shapes if shape.name == "production-alignment-target")
    assert repaired_target.top.inches == pytest.approx(3.13, abs=0.01)
    assert repaired_target.height.inches == pytest.approx(3.88, abs=0.01)

    clean = service.lint_fix(pptx_path=str(output), touched_slides=[0])
    assert clean.lint_issue_count_before == 0
    assert clean.fix_applied_count == 0
    assert clean.residue_count == 0


def test_deck_native_lint_ignores_thin_divider_thickness_near_peer_grid(
    tmp_path: Path,
) -> None:
    output = tmp_path / "thin-divider-thickness-grid.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(20)
    presentation.slide_height = Inches(11.25)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    for index, left in enumerate((1.0, 4.0, 7.0, 10.0, 13.0), start=1):
        peer = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left),
            Inches(4.42),
            Inches(2.0),
            Inches(1.0),
        )
        peer.name = f"bottom-grid-peer-{index}"

    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.25),
        Inches(5.42),
        Inches(17.5),
        Inches(1.98),
    )
    panel.name = "enclosing-panel"
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.25),
        Inches(5.42),
        Inches(17.5),
        Inches(0.10),
    )
    divider.name = "contained-divider"
    presentation.save(output)
    before = Presentation(output)
    before_divider = next(
        shape for shape in before.slides[0].shapes if shape.name == "contained-divider"
    )
    before_geometry = (
        before_divider.left,
        before_divider.top,
        before_divider.width,
        before_divider.height,
    )

    fixed = DeckNativeService().lint_fix(pptx_path=str(output), touched_slides=[0])

    assert fixed.success is True
    assert fixed.lint_issue_count_before == 0
    assert fixed.fix_applied_count == 0
    assert fixed.residue_count == 0
    repaired = Presentation(output)
    repaired_divider = next(
        shape for shape in repaired.slides[0].shapes if shape.name == "contained-divider"
    )
    assert (
        repaired_divider.left,
        repaired_divider.top,
        repaired_divider.width,
        repaired_divider.height,
    ) == before_geometry


def test_deck_native_lint_fix_translates_single_proven_text_edge(
    tmp_path: Path,
) -> None:
    output = tmp_path / "single-edge-alignment.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(20)
    presentation.slide_height = Inches(11.25)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    def add_text(name: str, left: float, top: float) -> None:
        shape = slide.shapes.add_textbox(
            Inches(left),
            Inches(top),
            Inches(7.91667),
            Inches(0.6),
        )
        shape.name = name
        shape.text = name
        shape.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    for index, top in enumerate((2.0, 3.0, 4.0), start=1):
        add_text(f"peer-{index}", 2.042, top)
    add_text("single-edge-target", 2.0, 1.0)
    presentation.save(output)

    service = DeckNativeService()
    fixed = service.lint_fix(pptx_path=str(output), touched_slides=[0])

    assert fixed.success is True
    assert fixed.fix_applied_count == 1
    assert fixed.issue_kinds == {"align-x": 1}
    assert fixed.residue_count == 0
    repaired = Presentation(output)
    target = next(shape for shape in repaired.slides[0].shapes if shape.name == "single-edge-target")
    assert target.left.inches == pytest.approx(2.042, abs=0.001)
    assert target.text_frame.paragraphs[0].runs[0].font.size.pt == 18

    clean = service.lint_fix(pptx_path=str(output), touched_slides=[0])
    assert clean.lint_issue_count_before == 0
    assert clean.fix_applied_count == 0
    assert clean.residue_count == 0


def test_deck_native_lint_fix_coordinates_shared_text_seam(
    tmp_path: Path,
) -> None:
    output = tmp_path / "shared-text-seam.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(20)
    presentation.slide_height = Inches(11.25)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    def add_text(name: str, left: float, top: float, width: float) -> None:
        shape = slide.shapes.add_textbox(
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(0.6),
        )
        shape.name = name
        shape.text = name
        shape.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    add_text("left-header", 7.0, 1.0, 5.76)
    add_text("seam-target", 12.76, 1.0, 6.01078)
    for index, top in enumerate((2.0, 3.0, 4.0, 5.0), start=1):
        add_text(f"row-peer-{index}", 12.669, top, 6.10178)
    presentation.save(output)
    original = Presentation(output)
    original_target = next(shape for shape in original.slides[0].shapes if shape.name == "seam-target")
    original_right = original_target.left.inches + original_target.width.inches

    service = DeckNativeService()
    fixed = service.lint_fix(pptx_path=str(output), touched_slides=[0])

    assert fixed.success is True
    assert fixed.fix_applied_count == 1
    assert fixed.issue_kinds == {"align-x-seam": 1}
    assert fixed.residue_count == 0
    repaired = Presentation(output)
    left = next(shape for shape in repaired.slides[0].shapes if shape.name == "left-header")
    target = next(shape for shape in repaired.slides[0].shapes if shape.name == "seam-target")
    assert left.left.inches + left.width.inches == pytest.approx(12.669, abs=0.001)
    assert target.left.inches == pytest.approx(12.669, abs=0.001)
    assert target.left.inches + target.width.inches == pytest.approx(original_right, abs=0.001)
    assert left.text_frame.paragraphs[0].runs[0].font.size.pt == 18
    assert target.text_frame.paragraphs[0].runs[0].font.size.pt == 18

    clean = service.lint_fix(pptx_path=str(output), touched_slides=[0])
    assert clean.lint_issue_count_before == 0
    assert clean.fix_applied_count == 0
    assert clean.residue_count == 0


def test_deck_native_lint_fix_rolls_back_unsafe_single_edge_translation(
    tmp_path: Path,
) -> None:
    output = tmp_path / "unsafe-single-edge.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(20)
    presentation.slide_height = Inches(11.25)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    def add_text(name: str, left: float, top: float, width: float, height: float = 0.6) -> None:
        shape = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.name = name
        shape.text = name
        shape.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    add_text("different-height-neighbor", 3.0, 1.0, 1.96, height=0.8)
    add_text("unsafe-target", 5.0, 1.0, 3.0)
    for index, top in enumerate((2.0, 3.0, 4.0), start=1):
        add_text(f"unsafe-peer-{index}", 4.9, top, 3.0)
    presentation.save(output)
    before = Presentation(output)
    before_target = next(shape for shape in before.slides[0].shapes if shape.name == "unsafe-target")
    before_left_emu = before_target.left
    before_width_emu = before_target.width

    fixed = DeckNativeService().lint_fix(pptx_path=str(output), touched_slides=[0])

    assert fixed.success is True
    assert fixed.fix_applied_count == 0
    assert fixed.residue_count == 1
    assert fixed.residue_kinds == {"misaligned": 1}
    repaired = Presentation(output)
    target = next(shape for shape in repaired.slides[0].shapes if shape.name == "unsafe-target")
    assert target.left == before_left_emu
    assert target.width == before_width_emu
    assert target.text_frame.paragraphs[0].runs[0].font.size.pt == 18


def test_deck_native_lint_fix_preserves_authored_gutter_during_left_snap(
    tmp_path: Path,
) -> None:
    output = tmp_path / "authored-gutter.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(20)
    presentation.slide_height = Inches(11.25)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    def add_text(name: str, left: float, top: float, width: float) -> None:
        shape = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(0.6)
        )
        shape.name = name
        shape.text = name
        shape.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    add_text("gutter-neighbor", 7.0, 1.0, 5.669)
    add_text("gutter-target", 12.76, 1.0, 6.01078)
    for index, top in enumerate((2.0, 3.0, 4.0), start=1):
        add_text(f"gutter-peer-{index}", 12.669, top, 6.10178)
    presentation.save(output)
    before = Presentation(output)
    before_target = next(shape for shape in before.slides[0].shapes if shape.name == "gutter-target")
    before_left = before_target.left

    fixed = DeckNativeService().lint_fix(pptx_path=str(output), touched_slides=[0])

    assert fixed.fix_applied_count == 0
    assert fixed.residue_kinds.get("misaligned", 0) >= 1
    repaired = Presentation(output)
    target = next(shape for shape in repaired.slides[0].shapes if shape.name == "gutter-target")
    assert target.left == before_left


def test_deck_native_lint_fix_preserves_original_text_container(
    tmp_path: Path,
) -> None:
    output = tmp_path / "contained-translation.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(20)
    presentation.slide_height = Inches(11.25)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2.0),
        Inches(1.0),
        Inches(7.91667),
        Inches(0.6),
    ).name = "target-container"

    def add_text(name: str, left: float, top: float) -> None:
        shape = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(7.91667), Inches(0.6)
        )
        shape.name = name
        shape.text = name
        shape.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    add_text("contained-target", 2.0, 1.0)
    for index, top in enumerate((2.0, 3.0, 4.0), start=1):
        add_text(f"contained-peer-{index}", 2.042, top)
    presentation.save(output)

    fixed = DeckNativeService().lint_fix(pptx_path=str(output), touched_slides=[0])

    assert fixed.fix_applied_count == 0
    assert fixed.residue_kinds.get("misaligned", 0) >= 1
    repaired = Presentation(output)
    target = next(shape for shape in repaired.slides[0].shapes if shape.name == "contained-target")
    assert target.left.inches == pytest.approx(2.0, abs=0.001)


def test_deck_native_lint_fix_does_not_compose_seam_with_grown_neighbor(
    tmp_path: Path,
) -> None:
    output = tmp_path / "grown-seam-neighbor.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(20)
    presentation.slide_height = Inches(11.25)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    neighbor = slide.shapes.add_textbox(
        Inches(7.0), Inches(1.0), Inches(5.76), Inches(0.33)
    )
    neighbor.name = "grown-neighbor"
    neighbor.text_frame.margin_left = 0
    neighbor.text_frame.margin_right = 0
    neighbor.text_frame.margin_top = Inches(0.05)
    neighbor.text_frame.margin_bottom = Inches(0.05)
    neighbor.text = "Caution requests explicit confirmation."
    neighbor.text_frame.paragraphs[0].line_spacing = Pt(21)
    neighbor.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    def add_text(name: str, left: float, top: float, width: float) -> None:
        shape = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(0.33)
        )
        shape.name = name
        shape.text = name
        shape.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    add_text("grown-neighbor-target", 12.76, 1.0, 6.01078)
    for index, top in enumerate((2.0, 3.0, 4.0), start=1):
        add_text(f"grown-neighbor-peer-{index}", 12.669, top, 6.10178)
    presentation.save(output)
    before = Presentation(output)
    before_target = next(shape for shape in before.slides[0].shapes if shape.name == "grown-neighbor-target")
    before_target_left = before_target.left

    fixed = DeckNativeService().lint_fix(pptx_path=str(output), touched_slides=[0])

    assert fixed.issue_kinds == {"grow": 1}
    assert fixed.residue_kinds == {"misaligned": 1}
    repaired = Presentation(output)
    target = next(shape for shape in repaired.slides[0].shapes if shape.name == "grown-neighbor-target")
    grown = next(shape for shape in repaired.slides[0].shapes if shape.name == "grown-neighbor")
    assert target.left == before_target_left
    assert grown.height.inches > 0.33


def test_deck_native_lint_fix_rolls_back_seam_that_would_wrap_text(
    tmp_path: Path,
    monkeypatch,
) -> None:
    home = tmp_path / "home"
    font_dir = home / "Library/Fonts"
    font_dir.mkdir(parents=True)
    embedded_font = ImageFont.load_default(size=12).path
    assert hasattr(embedded_font, "getvalue")
    (font_dir / "CanarySans.ttf").write_bytes(embedded_font.getvalue())
    monkeypatch.setenv("HOME", str(home))

    output = tmp_path / "wrapping-seam.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(20)
    presentation.slide_height = Inches(11.25)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    neighbor = slide.shapes.add_textbox(
        Inches(7.0), Inches(1.0), Inches(5.76), Inches(0.35)
    )
    neighbor.name = "wrap-neighbor"
    neighbor.text_frame.margin_left = 0
    neighbor.text_frame.margin_right = 0
    neighbor.text_frame.margin_top = 0
    neighbor.text_frame.margin_bottom = 0
    neighbor.text = "W" * 23 + " x"
    neighbor.text_frame.paragraphs[0].runs[0].font.name = "CanarySans"
    neighbor.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    def add_text(name: str, left: float, top: float, width: float) -> None:
        shape = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(0.35)
        )
        shape.name = name
        shape.text = name
        shape.text_frame.paragraphs[0].runs[0].font.name = "CanarySans"
        shape.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    add_text("wrap-target", 12.76, 1.0, 6.01078)
    for index, top in enumerate((2.0, 3.0, 4.0), start=1):
        add_text(f"wrap-peer-{index}", 12.669, top, 6.10178)
    presentation.save(output)
    before = Presentation(output)
    before_neighbor = next(shape for shape in before.slides[0].shapes if shape.name == "wrap-neighbor")
    before_target = next(shape for shape in before.slides[0].shapes if shape.name == "wrap-target")

    fixed = DeckNativeService().lint_fix(pptx_path=str(output), touched_slides=[0])

    assert fixed.fix_applied_count == 0
    assert fixed.residue_kinds == {"misaligned": 1}
    repaired = Presentation(output)
    repaired_neighbor = next(shape for shape in repaired.slides[0].shapes if shape.name == "wrap-neighbor")
    repaired_target = next(shape for shape in repaired.slides[0].shapes if shape.name == "wrap-target")
    assert repaired_neighbor.width == before_neighbor.width
    assert repaired_target.left == before_target.left
    assert repaired_neighbor.text_frame.paragraphs[0].runs[0].font.size.pt == 18


def test_deck_native_lint_fix_coordinates_production_shaped_closing_grid(
    tmp_path: Path,
) -> None:
    output = tmp_path / "coordinated-closing-grid.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(20)
    presentation.slide_height = Inches(11.25)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    def add_text(
        name: str,
        left: float,
        top: float,
        width: float,
        height: float,
        font_size: float,
        text: str,
    ) -> None:
        shape = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.name = name
        shape.text = text
        text_frame = shape.text_frame
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0
        text_frame.paragraphs[0].line_spacing = Pt(font_size)
        text_frame.paragraphs[0].runs[0].font.size = Pt(font_size)

    add_text("closing-title", 1.25, 0.73, 15.94, 1.25, 39, "Title")
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.25),
        Inches(2.04),
        Inches(1.0),
        Inches(0.06),
    )
    accent.name = "closing-accent"

    for index, (top, label_top) in enumerate(
        ((2.60, 2.80), (3.83, 4.03), (5.06, 5.26), (6.29, 6.49)),
        start=1,
    ):
        panel = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.0),
            Inches(top),
            Inches(17.77),
            Inches(1.02),
        )
        panel.name = f"question-panel-{index}"
        add_text(f"question-label-{index}", 1.30, label_top, 0.27, 0.29, 18, str(index))

    closing_panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.0),
        Inches(9.38),
        Inches(17.81),
        Inches(1.46),
    )
    closing_panel.name = "closing-panel"
    add_text("closing-thesis", 1.40, 9.73, 17.42, 0.19, 12, "Close")
    presentation.save(output)

    original = Presentation(output)
    original_text = {
        shape.name: shape.text
        for shape in original.slides[0].shapes
        if shape.has_text_frame
    }

    fixed = DeckNativeService().lint_fix(pptx_path=str(output), touched_slides=[0])

    assert fixed.success is True
    assert fixed.lint_issue_count_before == 4
    assert fixed.fix_applied_count == 4, fixed.residue
    assert fixed.issue_kinds == {"align-x": 3, "align-x-boundary": 1}
    assert fixed.residue_count == 0
    repaired = Presentation(output)
    by_name = {shape.name: shape for shape in repaired.slides[0].shapes}
    assert by_name["closing-title"].left.inches == pytest.approx(1.30, abs=0.001)
    assert by_name["closing-accent"].left.inches == pytest.approx(1.30, abs=0.001)
    assert by_name["closing-panel"].left.inches == pytest.approx(1.0, abs=0.001)
    assert by_name["closing-panel"].width.inches == pytest.approx(17.77, abs=0.001)
    assert by_name["closing-thesis"].left.inches == pytest.approx(1.30, abs=0.001)
    assert {
        shape.name: shape.text
        for shape in repaired.slides[0].shapes
        if shape.has_text_frame
    } == original_text
    assert by_name["closing-title"].text_frame.paragraphs[0].runs[0].font.size.pt == 39
    assert by_name["closing-thesis"].text_frame.paragraphs[0].runs[0].font.size.pt == 12

    clean = DeckNativeService().lint_fix(pptx_path=str(output), touched_slides=[0])
    assert clean.lint_issue_count_before == 0
    assert clean.fix_applied_count == 0
    assert clean.residue_count == 0


def test_deck_native_lint_fix_center_snaps_container_with_carried_content(
    tmp_path: Path,
) -> None:
    output = tmp_path / "centered-container-grid.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(20.0)
    presentation.slide_height = Inches(11.25)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    peer_geometries = (
        (2.19, 1.92),
        (4.23, 1.92),
        (6.27, 1.92),
        (8.31, 1.96),
    )
    for index, (top, height) in enumerate(peer_geometries, start=1):
        peer = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.77),
            Inches(top),
            Inches(17.50),
            Inches(height),
        )
        peer.name = f"center-peer-{index}"
        if index == 4:
            continue
        peer_text = slide.shapes.add_textbox(
            Inches(2.10),
            Inches(top + 0.30),
            Inches(17.17),
            Inches(0.36),
        )
        peer_text.name = f"center-peer-text-{index}"
        peer_text.text = f"Peer {index}"
        peer_text.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    lower_rule = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(1.77),
        Inches(10.23),
        Inches(19.27),
        Inches(10.23),
    )
    lower_rule.name = "center-container-lower-rule"
    heading = slide.shapes.add_textbox(
        Inches(2.10),
        Inches(8.92),
        Inches(17.17),
        Inches(0.36),
    )
    heading.name = "center-container-heading"
    heading.text = "Closing assertion"
    heading.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    target = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.25),
        Inches(8.54),
        Inches(18.33),
        Inches(2.44),
    )
    target.name = "center-container"
    vertical_rule = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(1.33),
        Inches(8.54),
        Inches(1.33),
        Inches(10.98),
    )
    vertical_rule.name = "center-container-vertical-rule"
    body = slide.shapes.add_textbox(
        Inches(1.75),
        Inches(9.44),
        Inches(17.52),
        Inches(0.62),
    )
    body.name = "center-container-body"
    body.text = "Carried content"
    body.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    presentation.save(output)

    before = Presentation(output)
    before_by_name = {shape.name: shape for shape in before.slides[0].shapes}
    stationary_names = (
        "center-container-lower-rule",
        "center-container-heading",
        "center-container-vertical-rule",
        "center-container-body",
    )
    original_geometries = {
        name: (
            before_by_name[name].left,
            before_by_name[name].top,
            before_by_name[name].width,
            before_by_name[name].height,
        )
        for name in stationary_names
    }
    original_text = {
        name: before_by_name[name].text
        for name in ("center-container-heading", "center-container-body")
    }

    fixed = DeckNativeService().lint_fix(pptx_path=str(output), touched_slides=[0])

    assert fixed.success is True
    assert fixed.lint_issue_count_before == 1
    assert fixed.fix_applied_count == 1, fixed.residue
    assert fixed.issue_kinds == {"align-x-container": 1}
    assert fixed.residue_count == 0
    repaired = Presentation(output)
    by_name = {shape.name: shape for shape in repaired.slides[0].shapes}
    expected_left = 10.52 - (18.33 / 2.0)
    assert by_name["center-container"].left.inches == pytest.approx(
        expected_left,
        abs=0.001,
    )
    assert {
        name: (
            by_name[name].left,
            by_name[name].top,
            by_name[name].width,
            by_name[name].height,
        )
        for name in stationary_names
    } == original_geometries
    assert {
        name: by_name[name].text
        for name in ("center-container-heading", "center-container-body")
    } == original_text

    clean = DeckNativeService().lint_fix(pptx_path=str(output), touched_slides=[0])
    assert clean.lint_issue_count_before == 0
    assert clean.fix_applied_count == 0
    assert clean.residue_count == 0


def test_deck_native_lint_fix_refuses_center_snap_with_rotated_carried_shape(
    tmp_path: Path,
) -> None:
    output = tmp_path / "rotated-center-container.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(20.0)
    presentation.slide_height = Inches(11.25)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    for index, top in enumerate((1.0, 2.2, 3.4, 4.6), start=1):
        peer = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.77),
            Inches(top),
            Inches(17.50),
            Inches(0.90),
        )
        peer.name = f"rotated-center-peer-{index}"

    target = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.25),
        Inches(7.00),
        Inches(18.33),
        Inches(2.44),
    )
    target.name = "rotated-center-container"
    rotated = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.30),
        Inches(7.50),
        Inches(0.40),
        Inches(0.40),
    )
    rotated.name = "rotated-carried-shape"
    rotated.rotation = 45
    presentation.save(output)

    before = Presentation(output)
    before_by_name = {shape.name: shape for shape in before.slides[0].shapes}
    original = {
        name: (
            before_by_name[name].left,
            before_by_name[name].top,
            before_by_name[name].width,
            before_by_name[name].height,
            before_by_name[name].rotation,
        )
        for name in ("rotated-center-container", "rotated-carried-shape")
    }

    fixed = DeckNativeService().lint_fix(pptx_path=str(output), touched_slides=[0])

    assert fixed.success is True
    assert fixed.lint_issue_count_before == 1
    assert fixed.fix_applied_count == 0
    assert fixed.residue_kinds == {"misaligned": 1}
    repaired = Presentation(output)
    by_name = {shape.name: shape for shape in repaired.slides[0].shapes}
    assert {
        name: (
            by_name[name].left,
            by_name[name].top,
            by_name[name].width,
            by_name[name].height,
            by_name[name].rotation,
        )
        for name in ("rotated-center-container", "rotated-carried-shape")
    } == original


def test_deck_native_lint_fix_rolls_back_container_boundary_without_child_snap(
    tmp_path: Path,
) -> None:
    output = tmp_path / "unsafe-container-boundary.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(20)
    presentation.slide_height = Inches(11.25)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    for index, top in enumerate((2.0, 3.2, 4.4, 5.6), start=1):
        panel = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.0),
            Inches(top),
            Inches(17.77),
            Inches(0.9),
        )
        panel.name = f"peer-panel-{index}"

    target = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.0),
        Inches(7.0),
        Inches(17.81),
        Inches(1.4),
    )
    target.name = "unsafe-container"
    child = slide.shapes.add_textbox(
        Inches(2.0),
        Inches(7.3),
        Inches(16.82),
        Inches(0.5),
    )
    child.name = "unmoved-child"
    child.text = "Child"
    child.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    presentation.save(output)
    before = Presentation(output)
    before_target = next(shape for shape in before.slides[0].shapes if shape.name == "unsafe-container")
    before_left, before_width = before_target.left, before_target.width
    before_child = next(
        shape for shape in before.slides[0].shapes if shape.name == "unmoved-child"
    )
    before_child_geometry = (
        before_child.left,
        before_child.top,
        before_child.width,
        before_child.height,
    )

    fixed = DeckNativeService().lint_fix(pptx_path=str(output), touched_slides=[0])

    assert fixed.success is True
    assert fixed.fix_applied_count == 0
    assert fixed.residue_kinds == {"misaligned": 1}
    assert "original containment" in fixed.residue[0]["issue"]
    repaired = Presentation(output)
    repaired_target = next(
        shape for shape in repaired.slides[0].shapes if shape.name == "unsafe-container"
    )
    assert repaired_target.left == before_left
    assert repaired_target.width == before_width
    repaired_child = next(
        shape for shape in repaired.slides[0].shapes if shape.name == "unmoved-child"
    )
    assert (
        repaired_child.left,
        repaired_child.top,
        repaired_child.width,
        repaired_child.height,
    ) == before_child_geometry
    assert repaired_child.text == child.text


def test_deck_native_lint_fix_rolls_back_multi_edge_container_resize(
    tmp_path: Path,
) -> None:
    output = tmp_path / "unsafe-multi-edge-container.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(20)
    presentation.slide_height = Inches(11.25)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    for index, top in enumerate((1.0, 2.5, 4.0, 5.5), start=1):
        peer = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.0),
            Inches(top),
            Inches(17.77),
            Inches(0.9),
        )
        peer.name = f"multi-edge-peer-{index}"

    target = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.0),
        Inches(7.0),
        Inches(17.87),
        Inches(1.4),
    )
    target.name = "multi-edge-container"
    child = slide.shapes.add_textbox(
        Inches(2.0),
        Inches(7.3),
        Inches(16.85),
        Inches(0.5),
    )
    child.name = "multi-edge-child"
    child.text = "Child"
    child.text_frame.margin_left = 0
    child.text_frame.margin_right = 0
    child.text_frame.margin_top = 0
    child.text_frame.margin_bottom = 0
    child.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
    presentation.save(output)

    before = Presentation(output)
    before_target = next(
        shape for shape in before.slides[0].shapes if shape.name == "multi-edge-container"
    )
    before_geometry = (
        before_target.left,
        before_target.top,
        before_target.width,
        before_target.height,
    )

    fixed = DeckNativeService().lint_fix(pptx_path=str(output), touched_slides=[0])

    assert fixed.success is True
    assert fixed.fix_applied_count == 0
    assert fixed.residue_kinds == {"misaligned": 1}
    assert "original containment" in fixed.residue[0]["issue"]
    repaired = Presentation(output)
    repaired_target = next(
        shape for shape in repaired.slides[0].shapes if shape.name == "multi-edge-container"
    )
    assert (
        repaired_target.left,
        repaired_target.top,
        repaired_target.width,
        repaired_target.height,
    ) == before_geometry


def test_deck_native_lint_fix_rolls_back_autoshape_external_text_occlusion(
    tmp_path: Path,
) -> None:
    output = tmp_path / "unsafe-autoshape-text-occlusion.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(20)
    presentation.slide_height = Inches(11.25)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    for index, (top, width) in enumerate(
        ((2.0, 2.0), (3.0, 2.5), (4.0, 3.0)),
        start=1,
    ):
        peer = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(2.0),
            Inches(top),
            Inches(width),
            Inches(0.4),
        )
        peer.name = f"occlusion-peer-{index}"

    external_text = slide.shapes.add_textbox(
        Inches(2.95),
        Inches(1.0),
        Inches(1.0),
        Inches(0.4),
    )
    external_text.name = "external-text"
    external_text.text = "External"
    external_text.text_frame.margin_left = 0
    external_text.text_frame.margin_right = 0
    external_text.text_frame.margin_top = 0
    external_text.text_frame.margin_bottom = 0
    external_text.text_frame.paragraphs[0].runs[0].font.size = Pt(12)

    target = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.9),
        Inches(1.0),
        Inches(1.0),
        Inches(0.4),
    )
    target.name = "occluding-target"
    presentation.save(output)

    before = Presentation(output)
    before_target = next(
        shape for shape in before.slides[0].shapes if shape.name == "occluding-target"
    )
    before_geometry = (
        before_target.left,
        before_target.top,
        before_target.width,
        before_target.height,
    )

    fixed = DeckNativeService().lint_fix(pptx_path=str(output), touched_slides=[0])

    assert fixed.success is True
    assert fixed.fix_applied_count == 0
    assert fixed.residue_kinds == {"misaligned": 1}
    assert "AUTO_SHAPE/external-text overlap" in fixed.residue[0]["issue"]
    repaired = Presentation(output)
    repaired_target = next(
        shape for shape in repaired.slides[0].shapes if shape.name == "occluding-target"
    )
    assert (
        repaired_target.left,
        repaired_target.top,
        repaired_target.width,
        repaired_target.height,
    ) == before_geometry


def test_deck_native_lint_fix_preserves_top_level_group_inside_container(
    tmp_path: Path,
) -> None:
    output = tmp_path / "grouped-container-content.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(20)
    presentation.slide_height = Inches(11.25)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    for index, (left, width, top) in enumerate(
        ((1.0, 17.77, 1.0), (2.0, 16.77, 2.5), (3.0, 15.77, 4.0)),
        start=1,
    ):
        peer = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(0.9),
        )
        peer.name = f"group-container-peer-{index}"

    target = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.0),
        Inches(7.0),
        Inches(17.81),
        Inches(1.4),
    )
    target.name = "group-container-target"
    group = slide.shapes.add_group_shape()
    group.name = "carried-group"
    grouped_content = group.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2.0),
        Inches(7.3),
        Inches(16.82),
        Inches(0.5),
    )
    grouped_content.name = "grouped-content"
    presentation.save(output)

    before = Presentation(output)
    before_target = next(
        shape for shape in before.slides[0].shapes if shape.name == "group-container-target"
    )
    before_group = next(
        shape for shape in before.slides[0].shapes if shape.name == "carried-group"
    )
    before_target_geometry = (
        before_target.left,
        before_target.top,
        before_target.width,
        before_target.height,
    )
    before_group_geometry = (
        before_group.left,
        before_group.top,
        before_group.width,
        before_group.height,
    )

    fixed = DeckNativeService().lint_fix(pptx_path=str(output), touched_slides=[0])

    assert fixed.success is True
    assert fixed.fix_applied_count == 0
    assert fixed.residue_kinds == {"misaligned": 1}
    assert "original containment" in fixed.residue[0]["issue"]
    repaired = Presentation(output)
    repaired_target = next(
        shape for shape in repaired.slides[0].shapes if shape.name == "group-container-target"
    )
    repaired_group = next(
        shape for shape in repaired.slides[0].shapes if shape.name == "carried-group"
    )
    assert (
        repaired_target.left,
        repaired_target.top,
        repaired_target.width,
        repaired_target.height,
    ) == before_target_geometry
    assert (
        repaired_group.left,
        repaired_group.top,
        repaired_group.width,
        repaired_group.height,
    ) == before_group_geometry


def test_deck_native_lint_fix_grows_small_overflow_inside_containing_panel(
    tmp_path: Path,
) -> None:
    output = tmp_path / "production-contained-overflow.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(20)
    presentation.slide_height = Inches(11.25)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.25),
        Inches(8.54),
        Inches(18.04),
        Inches(1.95),
    )
    panel.name = "production-container"
    text = slide.shapes.add_textbox(
        Inches(2.49),
        Inches(8.9),
        Inches(9.56),
        Inches(0.33),
    )
    text.name = "production-small-overflow"
    frame = text.text_frame
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = Inches(0.05)
    frame.margin_bottom = Inches(0.05)
    paragraph = frame.paragraphs[0]
    paragraph.text = "Caution requests explicit confirmation."
    paragraph.line_spacing = Pt(21)
    paragraph.runs[0].font.size = Pt(18)
    presentation.save(output)

    service = DeckNativeService()
    fixed = service.lint_fix(pptx_path=str(output), touched_slides=[0])

    assert fixed.success is True
    assert fixed.lint_issue_count_before == 1
    assert fixed.fix_applied_count == 1
    assert fixed.issue_kinds == {"grow": 1}
    assert fixed.residue_count == 0
    repaired = Presentation(output)
    repaired_text = next(shape for shape in repaired.slides[0].shapes if shape.name == "production-small-overflow")
    assert repaired_text.height.inches == pytest.approx(0.44, abs=0.01)
    assert repaired_text.text_frame.paragraphs[0].runs[0].font.size.pt == 18

    clean = service.lint_fix(pptx_path=str(output), touched_slides=[0])
    assert clean.lint_issue_count_before == 0
    assert clean.fix_applied_count == 0
    assert clean.residue_count == 0


def test_deck_native_lint_fix_repairs_canary_headline_and_kpi_overflow(
    tmp_path: Path,
    monkeypatch,
) -> None:
    """Exercise the production lint/fix subprocess with the failed canary geometry."""

    # Give both macOS and Linux subprocesses the same font bytes so this
    # regression tests geometry rather than whichever fonts the host installs.
    home = tmp_path / "home"
    embedded_font = ImageFont.load_default(size=12).path
    assert hasattr(embedded_font, "getvalue")
    for relative_dir in (Path("Library/Fonts"), Path(".fonts")):
        font_dir = home / relative_dir
        font_dir.mkdir(parents=True, exist_ok=True)
        (font_dir / "CanarySerif-Bold.ttf").write_bytes(embedded_font.getvalue())
    monkeypatch.setenv("HOME", str(home))

    output = tmp_path / "canary-overflow.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(20)
    presentation.slide_height = Inches(11.25)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    def add_textbox(
        *,
        name: str,
        at: tuple[float, float],
        size: tuple[float, float],
        text: str,
        font_size: float,
        line_spacing: float,
        bold: bool = False,
    ):
        shape = slide.shapes.add_textbox(
            Inches(at[0]),
            Inches(at[1]),
            Inches(size[0]),
            Inches(size[1]),
        )
        shape.name = name
        frame = shape.text_frame
        frame.margin_left = 0
        frame.margin_right = 0
        frame.margin_top = 0
        frame.margin_bottom = 0
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        paragraph.line_spacing = Pt(line_spacing)
        run = paragraph.runs[0]
        run.font.name = "CanarySerif"
        run.font.size = Pt(font_size)
        run.font.bold = bold
        return shape

    add_textbox(
        name="canary-headline",
        at=(1.25, 0.85),
        size=(7.75, 0.60),
        text="Habitat is fragmented, not absent",
        font_size=37.5,
        line_spacing=43.1,
        bold=True,
    )
    add_textbox(
        name="canary-narrative",
        at=(1.25, 1.77),
        size=(8.07, 1.56),
        text="Bees and butterflies forage between existing habitat patches.",
        font_size=18.75,
        line_spacing=25.6,
    )
    add_textbox(
        name="canary-kpi",
        at=(11.22, 6.87),
        # The production box was 1.47in. Keep the same tight composition with
        # a 0.02in cushion so the embedded cross-platform test font rewraps.
        size=(1.45, 0.61),
        text="0.9 mi",
        font_size=39,
        line_spacing=46.8,
        bold=True,
    )
    add_textbox(
        name="canary-kpi-label",
        at=(11.22, 7.59),
        size=(3.19, 0.69),
        text="average gap between usable forage patches",
        font_size=14.25,
        line_spacing=17.1,
    )
    mixed = add_textbox(
        name="canary-mixed-run",
        at=(1.25, 10.40),
        size=(4.5, 0.60),
        text="Signal ",
        font_size=12,
        line_spacing=43.1,
    )
    emphasis = mixed.text_frame.paragraphs[0].add_run()
    emphasis.text = "critical corridor failure"
    emphasis.font.name = "CanarySerif"
    emphasis.font.size = Pt(40)
    emphasis.font.bold = True
    presentation.save(output)

    service = DeckNativeService()
    fixed = service.lint_fix(pptx_path=str(output), touched_slides=[0])

    assert fixed.success is True
    assert fixed.lint_issue_count_before >= 2
    assert fixed.fix_applied_count >= 2
    assert fixed.residue_count == 0
    # The producer counts unresolved `remaining_issue_shapes` in the first
    # number, so equality proves none survived this repair pass.
    assert fixed.lint_issue_count_before == fixed.fix_applied_count

    clean = service.lint_fix(pptx_path=str(output), touched_slides=[0])
    assert clean.success is True
    assert clean.lint_issue_count_before == 0
    assert clean.fix_applied_count == 0
    assert clean.residue_count == 0

    repaired = Presentation(output)
    repaired_sizes = {
        shape.name: min(
            run.font.size.pt
            for paragraph in shape.text_frame.paragraphs
            for run in paragraph.runs
            if run.font.size is not None
        )
        for shape in repaired.slides[0].shapes
        if shape.name in {"canary-headline", "canary-kpi"}
    }
    assert set(repaired_sizes) == {"canary-headline", "canary-kpi"}
    assert 31 <= repaired_sizes["canary-headline"] < 37.5
    assert 31 <= repaired_sizes["canary-kpi"] < 39
    repaired_mixed = next(shape for shape in repaired.slides[0].shapes if shape.name == "canary-mixed-run")
    mixed_sizes = [
        run.font.size.pt
        for run in repaired_mixed.text_frame.paragraphs[0].runs
        if run.font.size is not None
    ]
    assert len(mixed_sizes) == 2
    assert mixed_sizes[0] >= 10
    assert mixed_sizes[1] >= mixed_sizes[0] * 2
    assert mixed_sizes[1] < 40


def test_deck_native_inspect_marks_full_slide_picture_inventory(tmp_path: Path) -> None:
    base = tmp_path / "base.pptx"
    patch = tmp_path / "full-slide-picture.patch.json"
    output = tmp_path / "full-slide-picture.pptx"
    image = tmp_path / "background.png"
    _wide_base_deck(base)
    Image.new("RGB", (1920, 1080), "blue").save(image)
    patch.write_text(
        json.dumps(
            {
                "ops": [
                    {"op": "add-slide", "layout": "Blank"},
                    {
                        "op": "add-picture",
                        "slide": 0,
                        "image": str(image),
                        "at": [0, 0],
                        "size": [20, 11.25],
                        "name": "hero-background",
                    },
                ]
            }
        ),
        encoding="utf-8",
    )
    service = DeckNativeService()

    applied = service.apply_patch(base_deck_path=str(base), patch_path=str(patch), output_path=str(output), fix=False)
    inspected = service.inspect(str(output))

    assert applied.success is True
    assert inspected.success is True
    assert inspected.full_slide_picture_count == 1
    inventory = json.loads(Path(inspected.shape_inventory_path or "").read_text(encoding="utf-8"))
    slide = inventory["slides"]["slide:1"]
    assert slide["full_slide_picture_count"] == 1
    assert slide["shapes"][0]["full_slide"] is True
    assert slide["shapes"][0]["name"] == "hero-background"


def test_native_mechanical_report_is_compact() -> None:
    from deerflow.sophia.deck_native.models import (
        NativeDeckInspectResult,
        NativeDeckLintFixResult,
        NativeDeckRenderResult,
    )

    report = native_mechanical_report(
        inspect=NativeDeckInspectResult(True, 2, 5, 4, 1, 1, 0.9, "inventory.json", "raw.json", []),
        lint_fix=NativeDeckLintFixResult(True, 1, 1, 0, 2, [], []),
        render=NativeDeckRenderResult(True, "rendered", 2, []),
        diff={"success": True, "changed": True, "diff_text": "x" * 2000},
    )

    assert report == {
        "inspect_success": True,
        "native_editability_score": 0.9,
        "native_text_shape_count": 4,
        "picture_shape_count": 1,
        "full_slide_picture_count": 1,
        "lint_fix_success": True,
        "lint_issue_count_before": 1,
        "lint_fix_applied_count": 1,
        "lint_residue_count": 0,
        "lint_residue_kinds": {},
        "lint_residue": [],
        "render_success": True,
        "rendered_slide_count": 2,
        "diff_success": True,
        "diff_changed": True,
    }


def test_deck_native_html_to_patch_compiles_when_playwright_available(tmp_path: Path) -> None:
    if importlib.util.find_spec("playwright") is None:
        pytest.skip("Python Playwright is not installed in this backend env")
    base = tmp_path / "base.pptx"
    html = tmp_path / "slide.html"
    patch = tmp_path / "html.patch.json"
    output = tmp_path / "html-native.pptx"
    _wide_base_deck(base)
    html.write_text(
        """<!doctype html><html><body style="margin:0;width:1920px;height:1080px">
        <h1 data-deck-id="title" data-deck-role="title" data-deck-required="true" style="position:absolute;left:96px;top:90px;font-size:72px;background:#fff;border:2px solid #111">Native HTML</h1>
        <p data-deck-id="narrative" data-deck-role="narrative" data-deck-required="true" style="position:absolute;left:96px;top:220px;font-size:36px">Compiled through html2patch.</p>
        </body></html>""",
        encoding="utf-8",
    )
    service = DeckNativeService()

    patched = service.html_to_patch(html_paths=[str(html)], base_deck_path=str(base), output_patch_path=str(patch))
    if not patched.success and any("browser" in error.lower() or "chromium" in error.lower() for error in patched.errors):
        pytest.skip("Python Playwright is installed but Chromium is unavailable")
    applied = service.apply_patch(base_deck_path=str(base), patch_path=str(patch), output_path=str(output), fix=True)
    inspected = service.inspect(str(output))

    assert patched.success is True
    assert patched.patch_op_count > 0
    source_map = json.loads(Path(patched.source_map_path or "").read_text(encoding="utf-8"))
    assert set(source_map["slides"]["slide:1"]["elements"]) == {"title", "narrative"}
    assert len(source_map["slides"]["slide:1"]["elements"]["title"]["shape_names"]) >= 2
    assert applied.success is True
    assert inspected.native_text_shape_count >= 2
    assert inspected.full_slide_picture_count == 0


def test_deck_native_html_to_patch_retains_transparent_required_wrapper(tmp_path: Path) -> None:
    if importlib.util.find_spec("playwright") is None:
        pytest.skip("Python Playwright is not installed in this backend env")
    base = tmp_path / "base.pptx"
    html = tmp_path / "transparent-wrapper.html"
    patch = tmp_path / "transparent-wrapper.patch.json"
    output = tmp_path / "transparent-wrapper.pptx"
    _wide_base_deck(base)
    html.write_text(
        """<!doctype html><html><body style="margin:0;width:1920px;height:1080px;background:#0A0E14;color:#EEF4FB">
        <section data-deck-id="system" data-deck-role="architecture" data-deck-required="true"
                 style="position:absolute;left:96px;top:90px;width:1500px;height:700px">
          <h1 style="position:absolute;left:0;top:0;width:1200px;font-size:72px">Transparent semantic wrapper</h1>
          <p style="position:absolute;left:0;top:180px;width:1100px;font-size:36px">Its native descendants retain the wrapper identity.</p>
        </section>
        </body></html>""",
        encoding="utf-8",
    )
    service = DeckNativeService()

    patched = service.html_to_patch(
        html_paths=[str(html)],
        base_deck_path=str(base),
        output_patch_path=str(patch),
    )
    if not patched.success and any(
        "browser" in error.lower() or "chromium" in error.lower()
        for error in patched.errors
    ):
        pytest.skip("Python Playwright is installed but Chromium is unavailable")
    applied = service.apply_patch(
        base_deck_path=str(base),
        patch_path=str(patch),
        output_path=str(output),
        fix=True,
    )
    inspected = service.inspect(str(output))

    assert patched.success is True
    source_map = json.loads(Path(patched.source_map_path or "").read_text(encoding="utf-8"))
    system = source_map["slides"]["slide:1"]["elements"]["system"]
    assert system["source_role"] == "architecture"
    assert system["source_required"] is True
    assert len(system["shape_names"]) == 2
    assert applied.success is True
    assert inspected.success is True
    inventory = json.loads(Path(inspected.shape_inventory_path or "").read_text(encoding="utf-8"))
    native_names = {
        shape["name"]
        for shape in inventory["slides"]["slide:1"]["shapes"]
        if shape.get("name")
    }
    assert set(system["shape_names"]).issubset(native_names)


def test_deck_native_html_to_patch_retains_nested_semantic_ids(tmp_path: Path) -> None:
    if importlib.util.find_spec("playwright") is None:
        pytest.skip("Python Playwright is not installed in this backend env")
    base = tmp_path / "base.pptx"
    html = tmp_path / "nested-semantic-ids.html"
    patch = tmp_path / "nested-semantic-ids.patch.json"
    output = tmp_path / "nested-semantic-ids.pptx"
    _wide_base_deck(base)
    html.write_text(
        """<!doctype html><html><body style="margin:0;width:1920px;height:1080px;background:#0A0E14;color:#EEF4FB">
        <section data-deck-id="system" data-deck-role="architecture" data-deck-required="true"
                 style="position:absolute;left:96px;top:90px;width:1500px;height:700px">
          <div data-deck-id="cluster" data-deck-role="diagram" data-deck-required="true"
               style="position:absolute;left:0;top:0;width:1300px;height:600px">
            <h1 data-deck-id="title" data-deck-role="title" data-deck-required="true"
                style="position:absolute;left:0;top:0;width:1200px;font-size:72px">Nested semantic identities</h1>
            <p data-deck-id="detail" data-deck-role="narrative" data-deck-required="true"
               style="position:absolute;left:0;top:180px;width:1100px;font-size:36px">One native shape satisfies its leaf and container identities.</p>
          </div>
        </section>
        </body></html>""",
        encoding="utf-8",
    )
    service = DeckNativeService()

    patched = service.html_to_patch(
        html_paths=[str(html)],
        base_deck_path=str(base),
        output_patch_path=str(patch),
    )
    if not patched.success and any(
        "browser" in error.lower() or "chromium" in error.lower()
        for error in patched.errors
    ):
        pytest.skip("Python Playwright is installed but Chromium is unavailable")
    applied = service.apply_patch(
        base_deck_path=str(base),
        patch_path=str(patch),
        output_path=str(output),
        fix=True,
    )
    inspected = service.inspect(str(output))

    assert patched.success is True
    source_map = json.loads(Path(patched.source_map_path or "").read_text(encoding="utf-8"))
    elements = source_map["slides"]["slide:1"]["elements"]
    assert set(elements) == {"system", "cluster", "title", "detail"}
    title_names = set(elements["title"]["shape_names"])
    detail_names = set(elements["detail"]["shape_names"])
    assert title_names
    assert detail_names
    assert title_names.isdisjoint(detail_names)
    assert set(elements["cluster"]["shape_names"]) == title_names | detail_names
    assert set(elements["system"]["shape_names"]) == title_names | detail_names
    assert applied.success is True
    assert inspected.success is True
    inventory = json.loads(Path(inspected.shape_inventory_path or "").read_text(encoding="utf-8"))
    native_names = {
        shape["name"]
        for shape in inventory["slides"]["slide:1"]["shapes"]
        if shape.get("name")
    }
    assert title_names | detail_names <= native_names


def test_deck_native_html_to_patch_still_rejects_duplicate_semantic_ids(tmp_path: Path) -> None:
    if importlib.util.find_spec("playwright") is None:
        pytest.skip("Python Playwright is not installed in this backend env")
    base = tmp_path / "base.pptx"
    html = tmp_path / "duplicate-semantic-ids.html"
    patch = tmp_path / "duplicate-semantic-ids.patch.json"
    _wide_base_deck(base)
    html.write_text(
        """<!doctype html><html><body style="margin:0;width:1920px;height:1080px;background:#0A0E14;color:#EEF4FB">
        <h1 data-deck-id="duplicate" style="position:absolute;left:96px;top:90px;font-size:72px">First</h1>
        <p data-deck-id="duplicate" style="position:absolute;left:96px;top:220px;font-size:36px">Second</p>
        </body></html>""",
        encoding="utf-8",
    )
    service = DeckNativeService()

    patched = service.html_to_patch(
        html_paths=[str(html)],
        base_deck_path=str(base),
        output_patch_path=str(patch),
    )
    if not patched.success and any(
        "browser" in error.lower() or "chromium" in error.lower()
        for error in patched.errors
    ):
        pytest.skip("Python Playwright is installed but Chromium is unavailable")

    assert patched.success is False
    assert patched.patch_path is None
    assert "duplicate data-deck-id: duplicate" in "\n".join(patched.errors)
