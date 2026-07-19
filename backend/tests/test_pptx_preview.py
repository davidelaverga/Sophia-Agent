"""Tests for the PPTX → PDF canvas-preview pipeline.

LibreOffice is not installed in the dev environment, so the soffice
subprocess is mocked. The tests exercise path naming, graceful
degradation, and the emit-time attachment in BuilderArtifactMiddleware.
"""

from __future__ import annotations

import os
import platform
import shutil
import subprocess
import tempfile
from pathlib import Path
from types import SimpleNamespace

import pytest
from pptx import Presentation

import deerflow.sophia.pptx_preview as pptx_preview
from deerflow.agents.sophia_agent.middlewares.builder_artifact import (
    BuilderArtifactMiddleware,
)
from deerflow.sophia.pptx_preview import maybe_render_pptx_preview, preview_path_for


def test_preview_path_never_collides_with_same_stem_pdf(tmp_path):
    deck = tmp_path / "deck.pptx"
    assert preview_path_for(deck) == tmp_path / "deck.preview.pdf"


def test_missing_soffice_skips_preview(tmp_path, monkeypatch):
    deck = tmp_path / "deck.pptx"
    deck.write_bytes(b"PK fake deck")
    monkeypatch.setattr(pptx_preview.shutil, "which", lambda _binary: None)

    assert maybe_render_pptx_preview(deck) is None


def test_missing_input_skips_preview(tmp_path):
    assert maybe_render_pptx_preview(tmp_path / "absent.pptx") is None


def test_symlink_input_is_not_staged_by_privileged_preview(
    tmp_path,
    monkeypatch,
):
    real_deck = tmp_path / "real.pptx"
    real_deck.write_bytes(b"PK fake deck")
    linked_deck = tmp_path / "linked.pptx"
    linked_deck.symlink_to(real_deck)
    monkeypatch.setattr(
        pptx_preview.shutil,
        "which",
        lambda _binary: "/fake/soffice",
    )
    monkeypatch.setattr(
        pptx_preview,
        "run_process_group",
        lambda *_args, **_kwargs: pytest.fail("renderer must not run"),
    )

    assert maybe_render_pptx_preview(linked_deck) is None


def test_successful_conversion_moves_preview_into_place(tmp_path, monkeypatch):
    deck = tmp_path / "deck.pptx"
    deck.write_bytes(b"PK fake deck")
    monkeypatch.setattr(pptx_preview.shutil, "which", lambda _binary: "/fake/soffice")
    captured: dict[str, object] = {}

    def fake_run(cmd, **kwargs):
        out_dir = Path(cmd[cmd.index("--outdir") + 1])
        staged_source = Path(cmd[-1])
        captured["out_dir"] = out_dir
        captured["staged_source"] = staged_source
        captured["staged_bytes"] = staged_source.read_bytes()
        captured.update(kwargs)
        (out_dir / (staged_source.stem + ".pdf")).write_bytes(b"%PDF-1.4 preview")
        return SimpleNamespace(returncode=0, stderr="", stdout="")

    monkeypatch.setattr(pptx_preview, "run_process_group", fake_run)

    preview = maybe_render_pptx_preview(deck)

    assert preview == tmp_path / "deck.preview.pdf"
    assert preview.read_bytes() == b"%PDF-1.4 preview"
    assert captured["staged_bytes"] == b"PK fake deck"
    assert captured["staged_source"] == Path(captured["out_dir"]) / "source.pptx"
    assert captured["private_read_dirs"] == [captured["staged_source"]]
    assert captured["writable_dirs"] == [str(captured["out_dir"])]
    assert captured["identity_paths"] == [deck]


def test_preview_default_timeout_supports_long_decks(tmp_path, monkeypatch):
    deck = tmp_path / "deck.pptx"
    deck.write_bytes(b"PK fake deck")
    monkeypatch.setattr(pptx_preview.shutil, "which", lambda _binary: "/fake/soffice")
    captured: dict = {}

    def fake_run(cmd, **kwargs):
        captured.update(kwargs)
        out_dir = Path(cmd[cmd.index("--outdir") + 1])
        staged_source = Path(cmd[-1])
        (out_dir / (staged_source.stem + ".pdf")).write_bytes(b"%PDF-1.4 preview")
        return SimpleNamespace(returncode=0, stderr="", stdout="")

    monkeypatch.setattr(pptx_preview, "run_process_group", fake_run)

    assert maybe_render_pptx_preview(deck) is not None
    assert captured["timeout"] == 300


def test_failed_conversion_returns_none(tmp_path, monkeypatch):
    deck = tmp_path / "deck.pptx"
    deck.write_bytes(b"PK fake deck")
    monkeypatch.setattr(pptx_preview.shutil, "which", lambda _binary: "/fake/soffice")
    monkeypatch.setattr(
        pptx_preview,
        "run_process_group",
        lambda cmd, **kwargs: SimpleNamespace(returncode=1, stderr="boom", stdout=""),
    )

    assert maybe_render_pptx_preview(deck) is None


def test_timeout_returns_none(tmp_path, monkeypatch):
    deck = tmp_path / "deck.pptx"
    deck.write_bytes(b"PK fake deck")
    monkeypatch.setattr(pptx_preview.shutil, "which", lambda _binary: "/fake/soffice")

    def fake_run(cmd, **kwargs):
        raise subprocess.TimeoutExpired(cmd=cmd, timeout=1)

    monkeypatch.setattr(pptx_preview, "run_process_group", fake_run)

    assert maybe_render_pptx_preview(deck) is None


@pytest.mark.skipif(
    platform.system() != "Linux"
    or not hasattr(os, "geteuid")
    or os.geteuid() != 0
    or shutil.which("soffice") is None
    or shutil.which("setpriv") is None,
    reason="requires the real root-Linux LibreOffice identity boundary",
)
def test_real_root_linux_preview_reads_source_below_private_ancestor() -> None:
    private_root = Path(tempfile.mkdtemp(prefix="dq1-private-pptx-source-"))
    private_root.chmod(0o700)
    try:
        deck = private_root / "private-deck.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_textbox(0, 0, 1_000_000, 500_000).text = "Private source"
        presentation.save(deck)

        preview = maybe_render_pptx_preview(deck)

        assert preview == private_root / "private-deck.preview.pdf"
        assert preview.read_bytes().startswith(b"%PDF")
    finally:
        shutil.rmtree(private_root, ignore_errors=True)


# ---- emit-time attachment ---------------------------------------------------


def _state_with_outputs(outputs: Path) -> dict:
    return {"thread_data": {"outputs_path": str(outputs)}}


def test_attach_preview_adds_filename_and_supporting_file(tmp_path, monkeypatch):
    outputs = tmp_path / "outputs"
    outputs.mkdir()
    deck = outputs / "deck.pptx"
    deck.write_bytes(b"PK fake deck")

    def fake_render(pptx_path: Path):
        preview = preview_path_for(pptx_path)
        preview.write_bytes(b"%PDF-1.4 preview")
        return preview

    monkeypatch.setattr(
        "deerflow.agents.sophia_agent.middlewares.builder_artifact.maybe_render_pptx_preview",
        fake_render,
    )

    args = BuilderArtifactMiddleware._attach_pptx_canvas_preview(
        {"artifact_path": "/mnt/user-data/outputs/deck.pptx"},
        _state_with_outputs(outputs),
    )

    assert args["artifact_preview_filename"] == "deck.preview.pdf"
    assert "/mnt/user-data/outputs/deck.preview.pdf" in args["supporting_files"]


def test_attach_preview_noop_for_non_pptx(tmp_path):
    outputs = tmp_path / "outputs"
    outputs.mkdir()
    (outputs / "report.pdf").write_bytes(b"%PDF-1.4")

    args = BuilderArtifactMiddleware._attach_pptx_canvas_preview(
        {"artifact_path": "/mnt/user-data/outputs/report.pdf"},
        _state_with_outputs(outputs),
    )

    assert "artifact_preview_filename" not in args


def test_attach_preview_noop_when_conversion_unavailable(tmp_path, monkeypatch):
    outputs = tmp_path / "outputs"
    outputs.mkdir()
    (outputs / "deck.pptx").write_bytes(b"PK fake deck")
    monkeypatch.setattr(
        "deerflow.agents.sophia_agent.middlewares.builder_artifact.maybe_render_pptx_preview",
        lambda _path: None,
    )

    args = BuilderArtifactMiddleware._attach_pptx_canvas_preview(
        {"artifact_path": "/mnt/user-data/outputs/deck.pptx"},
        _state_with_outputs(outputs),
    )

    assert "artifact_preview_filename" not in args
    assert "supporting_files" not in args
