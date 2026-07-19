from __future__ import annotations

import hashlib
import io
import json
import os
import platform
import shutil
import tempfile
from pathlib import Path
from types import SimpleNamespace

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter

import deerflow.sophia.deck_quality.snapshot as snapshot_module
from deerflow.sophia.deck_quality.snapshot import (
    SnapshotCompletionMetadata,
    SnapshotConflictError,
    SnapshotCoverageError,
    SnapshotMissingEvidenceError,
    SnapshotRunIdentity,
    SnapshotStaleError,
    SnapshotUploadError,
    ensure_committed_render_source,
    freeze_and_upload_evidence_snapshot,
    freeze_and_upload_pre_render_input_bundle,
    load_evidence_snapshot,
    load_pre_render_input_bundle,
    rasterize_preview_pdf,
)
from deerflow.sophia.pptx_preview import maybe_render_pptx_preview
from deerflow.sophia.storage.supabase_artifact_store import (
    immutable_builder_artifact_object_path,
)

QUALITY_RUN_ID = f"quality_{'a' * 64}"
BUILD_ID = "build_01SNAPSHOT"


class MemoryImmutableUploader:
    def __init__(self) -> None:
        self.objects: dict[str, bytes] = {}
        self.attempts: list[tuple[str, str]] = []
        self.reads: list[str] = []
        self.read_limits: list[tuple[str, int]] = []

    def create_if_absent(
        self,
        object_path: str,
        content: bytes,
        *,
        content_type: str,
    ) -> str:
        self.attempts.append((object_path, content_type))
        if object_path in self.objects:
            return "exists"
        self.objects[object_path] = bytes(content)
        return "created"

    def read(self, object_path: str) -> bytes | None:
        self.reads.append(object_path)
        value = self.objects.get(object_path)
        return bytes(value) if value is not None else None

    def read_bounded(self, object_path: str, *, max_bytes: int) -> bytes | None:
        self.read_limits.append((object_path, max_bytes))
        value = self.read(object_path)
        if value is not None and len(value) > max_bytes:
            raise RuntimeError("object exceeds read bound")
        return value


def _png_bytes(color: str, *, width: int = 320, height: int = 180) -> bytes:
    from io import BytesIO

    output = BytesIO()
    Image.new("RGB", (width, height), color).save(output, format="PNG")
    return output.getvalue()


def _write_json(path: Path, value: object) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(value, indent=2), encoding="utf-8")


def _write_deck(path: Path) -> None:
    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    first = presentation.slides.add_slide(blank)
    first.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1)).text = "Feedback loop"
    first.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1)).text = "Observe → appraise → act"
    second = presentation.slides.add_slide(blank)
    second.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1)).text = "Close the loop"
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def _pdf_bytes(*, pages: int = 2) -> bytes:
    output = io.BytesIO()
    writer = PdfWriter()
    for _ in range(pages):
        writer.add_blank_page(width=720, height=405)
    writer.write(output)
    return output.getvalue()


@pytest.fixture
def snapshot_inputs(tmp_path: Path) -> dict[str, object]:
    outputs = tmp_path / "outputs"
    artifact = outputs / "psi-deck.pptx"
    _write_deck(artifact)
    preview = outputs / "psi-deck.preview.pdf"
    preview.write_bytes(_pdf_bytes())
    _write_json(
        outputs / "deck_build" / "creative_plan.json",
        {
            "subject_materials": ["feedback", "appraisal"],
            "signature": "control loop",
        },
    )
    _write_json(
        outputs / "deck_build" / "design_plan.json",
        {"rhythm": "setup-mechanism-close", "palette": ["ink", "paper"]},
    )
    _write_json(
        outputs / "deck_build" / "build.json",
        {
            "build_id": BUILD_ID,
            "slides": [{"selector": "slide:1"}, {"selector": "slide:2"}],
        },
    )
    metadata = SnapshotCompletionMetadata(
        quality_run_id=QUALITY_RUN_ID,
        build_id=BUILD_ID,
        user_id="canary-user",
        thread_id="thread-01",
        task_id="task-01",
        builder_run_id="run-01",
        parent_builder_trace_id="trace-01",
        logical_artifact_id="artifact-01",
        artifact_version_id="artifact-version-01",
        manifest_revision=1,
        artifact_storage_object_path=immutable_builder_artifact_object_path(
            user_id="canary-user",
            thread_or_session_id="thread-01",
            logical_artifact_id="artifact-01",
            artifact_version_id="artifact-version-01",
            artifact_sha256=hashlib.sha256(artifact.read_bytes()).hexdigest(),
            filename="psi-deck.pptx",
        ),
    )
    return {
        "outputs": outputs,
        "artifact": artifact,
        "preview": preview,
        "metadata": metadata,
        "brief": {
            "request": ("Explain the PSI feedback loop.\n\nRelevant memories from this session:\n- private prior context"),
            "subject": "PSI motivation architecture",
            "audience": "AI product and engineering leaders",
            "goal": "Explain the control mechanism",
        },
        "mechanical": {
            "checks": {
                "authoritative_gate": True,
                "render_success": True,
                "native_editability": True,
            }
        },
    }


def _freeze(
    snapshot_inputs: dict[str, object],
    uploader: MemoryImmutableUploader,
    *,
    raster_pages: tuple[bytes, ...] | None = None,
):
    artifact = snapshot_inputs["artifact"]
    assert isinstance(artifact, Path)
    preview = snapshot_inputs["preview"]
    assert isinstance(preview, Path)
    pages = raster_pages or (_png_bytes("navy"), _png_bytes("white"))
    input_descriptor = _freeze_pre_render_inputs(snapshot_inputs, uploader)
    loaded_input = load_pre_render_input_bundle(
        descriptor=input_descriptor,
        expected_identity=_identity(snapshot_inputs, input_descriptor),
        reader=uploader,
        materialization_root=artifact.parent.parent / "dq1-test-inputs",
    )

    def render_from_fixture(_artifact: Path) -> Path:
        target = loaded_input.outputs_root / "accepted.rendered.pdf"
        target.write_bytes(preview.read_bytes())
        return target

    render_source = ensure_committed_render_source(
        loaded_input=loaded_input,
        uploader=uploader,
        renderer=render_from_fixture,
    )
    uploader.attempts.clear()
    return freeze_and_upload_evidence_snapshot(
        metadata=loaded_input.metadata,
        outputs_root=loaded_input.outputs_root,
        artifact_virtual_path=loaded_input.artifact_virtual_path,
        artifact_host_path=loaded_input.artifact_host_path,
        task_brief=loaded_input.brief,
        authoritative_mechanical=loaded_input.mechanical_record,
        uploader=uploader,
        render_source=render_source,
        pdf_rasterizer=lambda _preview: pages,
    )


def _freeze_pre_render_inputs(
    snapshot_inputs: dict[str, object],
    uploader: MemoryImmutableUploader,
):
    artifact = snapshot_inputs["artifact"]
    assert isinstance(artifact, Path)
    return freeze_and_upload_pre_render_input_bundle(
        metadata=snapshot_inputs["metadata"],  # type: ignore[arg-type]
        outputs_root=snapshot_inputs["outputs"],  # type: ignore[arg-type]
        artifact_virtual_path="/mnt/user-data/outputs/psi-deck.pptx",
        artifact_host_path=artifact,
        task_brief=snapshot_inputs["brief"],  # type: ignore[arg-type]
        authoritative_mechanical=snapshot_inputs["mechanical"],  # type: ignore[arg-type]
        uploader=uploader,
    )


def _identity(
    snapshot_inputs: dict[str, object],
    descriptor: object,
) -> SnapshotRunIdentity:
    metadata = snapshot_inputs["metadata"]
    assert isinstance(metadata, SnapshotCompletionMetadata)
    manifest_path = getattr(descriptor, "manifest_path")
    manifest_hash = getattr(descriptor, "manifest_hash")
    return SnapshotRunIdentity(
        quality_run_id=metadata.quality_run_id,
        user_id=metadata.user_id,
        thread_id=metadata.thread_id,
        task_id=metadata.task_id,
        build_id=metadata.build_id,
        builder_run_id=metadata.builder_run_id,
        parent_builder_trace_id=metadata.parent_builder_trace_id,
        logical_artifact_id=metadata.logical_artifact_id,
        artifact_version_id=metadata.artifact_version_id,
        manifest_revision=metadata.manifest_revision,
        input_manifest_object_path=manifest_path,
        input_manifest_hash=manifest_hash,
    )


def _identity_from_uploader(
    snapshot_inputs: dict[str, object],
    uploader: MemoryImmutableUploader,
) -> SnapshotRunIdentity:
    manifest_path = next(
        path for path in uploader.objects if path.endswith("/input_bundle/manifest.json")
    )
    descriptor = SimpleNamespace(
        manifest_path=manifest_path,
        manifest_hash=hashlib.sha256(uploader.objects[manifest_path]).hexdigest(),
    )
    return _identity(snapshot_inputs, descriptor)


def _source_bytes(outputs: Path) -> dict[str, bytes]:
    return {path.relative_to(outputs).as_posix(): path.read_bytes() for path in outputs.rglob("*") if path.is_file()}


def test_pre_render_input_bundle_round_trip_is_restart_ready_and_manifest_last(
    snapshot_inputs: dict[str, object],
    tmp_path: Path,
) -> None:
    uploader = MemoryImmutableUploader()
    source_metadata = snapshot_inputs["metadata"]
    source_artifact = snapshot_inputs["artifact"]
    source_preview = snapshot_inputs["preview"]
    assert isinstance(source_metadata, SnapshotCompletionMetadata)
    assert isinstance(source_artifact, Path)
    assert isinstance(source_preview, Path)
    descriptor = _freeze_pre_render_inputs(snapshot_inputs, uploader)

    assert descriptor.counts.model_dump() == {
        "content_object_count": 6,
        "native_json_count": 3,
        "control_record_count": 2,
        "total_object_count": 7,
    }
    assert len(uploader.attempts) == descriptor.counts.total_object_count
    assert uploader.attempts[-1] == (
        descriptor.manifest_path,
        "application/json",
    )
    serialized_descriptor = json.dumps(descriptor.model_dump(mode="json"))
    assert "private prior context" not in serialized_descriptor
    assert "Explain the PSI" not in serialized_descriptor

    loaded = load_pre_render_input_bundle(
        descriptor=descriptor,
        expected_identity=_identity(snapshot_inputs, descriptor),
        reader=uploader,
        materialization_root=tmp_path / "restart",
    )

    assert loaded.metadata.quality_run_id == QUALITY_RUN_ID
    assert loaded.metadata.artifact_storage_object_path.endswith(".pptx")
    assert loaded.metadata.artifact_storage_object_path == source_metadata.artifact_storage_object_path
    assert loaded.brief.request == "Explain the PSI feedback loop."
    assert loaded.mechanical_record == snapshot_inputs["mechanical"]
    assert loaded.artifact_host_path.read_bytes() == source_artifact.read_bytes()
    assert not (loaded.outputs_root / "psi-deck.preview.pdf").exists()
    assert (loaded.outputs_root / "deck_build" / "creative_plan.json").is_file()
    assert (loaded.outputs_root / "deck_build" / "design_plan.json").is_file()
    assert (loaded.outputs_root / "deck_build" / "build.json").is_file()

    def render_from_fixture(_artifact: Path) -> Path:
        target = loaded.outputs_root / "accepted.rendered.pdf"
        target.write_bytes(source_preview.read_bytes())
        return target

    render_source = ensure_committed_render_source(
        loaded_input=loaded,
        uploader=uploader,
        renderer=render_from_fixture,
    )
    committed_render_hash = render_source.reference.pdf.sha256
    source_preview.write_bytes(source_preview.read_bytes() + b"\n% later local mutation")
    evidence = freeze_and_upload_evidence_snapshot(
        metadata=loaded.metadata,
        outputs_root=loaded.outputs_root,
        artifact_virtual_path=loaded.artifact_virtual_path,
        artifact_host_path=loaded.artifact_host_path,
        task_brief=loaded.brief,
        authoritative_mechanical=loaded.mechanical_record,
        uploader=uploader,
        render_source=render_source,
        pdf_rasterizer=lambda _preview: (_png_bytes("navy"), _png_bytes("white")),
    )
    assert evidence.snapshot_path.endswith("/evidence_manifest.json")
    evidence_manifest = json.loads(uploader.objects[evidence.snapshot_path])
    assert evidence_manifest["source_hashes"]["render_source_pdf"] == committed_render_hash
    assert evidence_manifest["source_hashes"]["render_source_pdf"] != hashlib.sha256(
        source_preview.read_bytes()
    ).hexdigest()


@pytest.mark.skipif(
    shutil.which("soffice") is None or shutil.which("pdftoppm") is None,
    reason="real PPTX/PDF evidence integration requires LibreOffice and Poppler",
)
def test_post_row_render_source_is_committed_before_real_evidence_rendering(
    snapshot_inputs: dict[str, object],
    tmp_path: Path,
) -> None:
    artifact = snapshot_inputs["artifact"]
    preview = snapshot_inputs["preview"]
    assert isinstance(artifact, Path)
    assert isinstance(preview, Path)
    preview.unlink()

    uploader = MemoryImmutableUploader()
    input_descriptor = _freeze_pre_render_inputs(snapshot_inputs, uploader)
    assert not preview.exists()
    loaded_input = load_pre_render_input_bundle(
        descriptor=input_descriptor,
        expected_identity=_identity(snapshot_inputs, input_descriptor),
        reader=uploader,
        materialization_root=tmp_path / "real-restart",
    )
    render_source = ensure_committed_render_source(
        loaded_input=loaded_input,
        uploader=uploader,
        renderer=maybe_render_pptx_preview,
    )
    frozen_preview_hash = hashlib.sha256(render_source.pdf_host_path.read_bytes()).hexdigest()

    evidence_descriptor = freeze_and_upload_evidence_snapshot(
        metadata=loaded_input.metadata,
        outputs_root=loaded_input.outputs_root,
        artifact_virtual_path=loaded_input.artifact_virtual_path,
        artifact_host_path=loaded_input.artifact_host_path,
        task_brief=loaded_input.brief,
        authoritative_mechanical=loaded_input.mechanical_record,
        uploader=uploader,
        render_source=render_source,
    )
    loaded_evidence = load_evidence_snapshot(
        descriptor=evidence_descriptor,
        expected_identity=_identity(snapshot_inputs, input_descriptor),
        reader=uploader,
        materialization_root=tmp_path / "real-evidence",
    )

    assert evidence_descriptor.counts.slide_count == 2
    assert loaded_evidence.manifest.source_hashes.render_source_pdf == frozen_preview_hash
    assert tuple(str(selector) for selector in loaded_evidence.snapshot.renders.selectors) == (
        "slide:1",
        "slide:2",
    )
    assert all(
        max(image.width, image.height) == 2200
        for image in loaded_evidence.snapshot.renders.slides
    )
    assert max(
        loaded_evidence.snapshot.renders.contact_sheet.width,
        loaded_evidence.snapshot.renders.contact_sheet.height,
    ) <= 2048


def test_pre_render_input_bundle_replay_is_idempotent(
    snapshot_inputs: dict[str, object],
) -> None:
    uploader = MemoryImmutableUploader()
    first = _freeze_pre_render_inputs(snapshot_inputs, uploader)
    first_objects = dict(uploader.objects)
    uploader.attempts.clear()

    second = _freeze_pre_render_inputs(snapshot_inputs, uploader)

    assert second == first
    assert uploader.objects == first_objects
    assert len(uploader.attempts) == first.counts.total_object_count
    assert uploader.attempts[-1][0] == first.manifest_path


def test_pre_render_input_bundle_does_not_capture_or_require_a_preview(
    snapshot_inputs: dict[str, object],
) -> None:
    uploader = MemoryImmutableUploader()
    preview = snapshot_inputs["preview"]
    assert isinstance(preview, Path)
    preview.unlink()

    descriptor = _freeze_pre_render_inputs(snapshot_inputs, uploader)
    manifest = json.loads(uploader.objects[descriptor.manifest_path])

    assert {record["role"] for record in manifest["objects"]} == {
        "accepted_artifact",
        "creative_plan",
        "design_plan",
        "build_record",
        "blind_brief",
        "mechanical_record",
    }
    assert all(record["media_type"] != "application/pdf" for record in manifest["objects"])


def test_pre_render_row_identity_is_verified_before_referenced_reads(
    snapshot_inputs: dict[str, object],
    tmp_path: Path,
) -> None:
    uploader = MemoryImmutableUploader()
    descriptor = _freeze_pre_render_inputs(snapshot_inputs, uploader)
    expected = _identity(snapshot_inputs, descriptor).model_copy(
        update={"task_id": "different-task"}
    )
    uploader.reads.clear()

    with pytest.raises(SnapshotStaleError, match="durable row"):
        load_pre_render_input_bundle(
            descriptor=descriptor,
            expected_identity=expected,
            reader=uploader,
            materialization_root=tmp_path / "stale-input",
        )

    assert uploader.reads == [descriptor.manifest_path]
    assert not (tmp_path / "stale-input").exists()


def test_render_source_replay_uses_the_committed_manifest_without_rerendering(
    snapshot_inputs: dict[str, object],
    tmp_path: Path,
) -> None:
    uploader = MemoryImmutableUploader()
    descriptor = _freeze_pre_render_inputs(snapshot_inputs, uploader)
    loaded = load_pre_render_input_bundle(
        descriptor=descriptor,
        expected_identity=_identity(snapshot_inputs, descriptor),
        reader=uploader,
        materialization_root=tmp_path / "render-replay",
    )
    calls = 0

    def renderer(_artifact: Path) -> Path:
        nonlocal calls
        calls += 1
        target = loaded.outputs_root / "candidate.pdf"
        target.write_bytes(_pdf_bytes())
        return target

    first = ensure_committed_render_source(
        loaded_input=loaded,
        uploader=uploader,
        renderer=renderer,
    )
    second = ensure_committed_render_source(
        loaded_input=loaded,
        uploader=uploader,
        renderer=renderer,
    )

    assert calls == 1
    assert second.reference == first.reference
    assert second.pdf_host_path.read_bytes() == first.pdf_host_path.read_bytes()


def test_orphan_render_pdf_is_safe_after_manifest_commit_crash(
    snapshot_inputs: dict[str, object],
    tmp_path: Path,
) -> None:
    class FailRenderManifestOnce(MemoryImmutableUploader):
        failed = False

        def create_if_absent(
            self,
            object_path: str,
            content: bytes,
            *,
            content_type: str,
        ) -> str:
            if object_path.endswith("/render_source/manifest.json") and not self.failed:
                self.failed = True
                raise RuntimeError("simulated commit crash")
            return super().create_if_absent(
                object_path,
                content,
                content_type=content_type,
            )

    uploader = FailRenderManifestOnce()
    descriptor = _freeze_pre_render_inputs(snapshot_inputs, uploader)
    loaded = load_pre_render_input_bundle(
        descriptor=descriptor,
        expected_identity=_identity(snapshot_inputs, descriptor),
        reader=uploader,
        materialization_root=tmp_path / "render-crash",
    )

    def renderer(_artifact: Path) -> Path:
        target = loaded.outputs_root / "candidate.pdf"
        target.write_bytes(_pdf_bytes())
        return target

    with pytest.raises(SnapshotUploadError, match="manifest commit"):
        ensure_committed_render_source(
            loaded_input=loaded,
            uploader=uploader,
            renderer=renderer,
        )
    assert any(path.endswith(".pdf") for path in uploader.objects)
    assert not any(path.endswith("/render_source/manifest.json") for path in uploader.objects)

    committed = ensure_committed_render_source(
        loaded_input=loaded,
        uploader=uploader,
        renderer=renderer,
    )
    assert committed.reference.manifest_path in uploader.objects


def test_pre_render_input_bundle_rejects_conflicting_object_before_manifest(
    snapshot_inputs: dict[str, object],
) -> None:
    uploader = MemoryImmutableUploader()
    descriptor = _freeze_pre_render_inputs(snapshot_inputs, uploader)
    manifest = json.loads(uploader.objects[descriptor.manifest_path])
    design_path = next(item["object_path"] for item in manifest["objects"] if item["role"] == "design_plan")
    conflicting_bytes = b'{"different":true}'
    uploader.objects[design_path] = conflicting_bytes
    uploader.attempts.clear()

    with pytest.raises(SnapshotConflictError, match="design_plan"):
        _freeze_pre_render_inputs(snapshot_inputs, uploader)

    assert uploader.objects[design_path] == conflicting_bytes
    assert descriptor.manifest_path not in {path for path, _ in uploader.attempts}


def test_pre_render_input_bundle_rejects_partial_remote_without_local_writes(
    snapshot_inputs: dict[str, object],
    tmp_path: Path,
) -> None:
    uploader = MemoryImmutableUploader()
    descriptor = _freeze_pre_render_inputs(snapshot_inputs, uploader)
    manifest = json.loads(uploader.objects[descriptor.manifest_path])
    creative_path = next(item["object_path"] for item in manifest["objects"] if item["role"] == "creative_plan")
    del uploader.objects[creative_path]
    materialization_root = tmp_path / "partial-restart"

    with pytest.raises(SnapshotMissingEvidenceError, match="creative_plan"):
        load_pre_render_input_bundle(
            descriptor=descriptor,
            expected_identity=_identity(snapshot_inputs, descriptor),
            reader=uploader,
            materialization_root=materialization_root,
        )

    assert not materialization_root.exists()


def test_pre_render_bundle_survives_local_input_mutation_without_primary_duplication(
    snapshot_inputs: dict[str, object],
    tmp_path: Path,
) -> None:
    uploader = MemoryImmutableUploader()
    metadata = snapshot_inputs["metadata"]
    artifact = snapshot_inputs["artifact"]
    preview = snapshot_inputs["preview"]
    outputs = snapshot_inputs["outputs"]
    assert isinstance(metadata, SnapshotCompletionMetadata)
    assert isinstance(artifact, Path)
    assert isinstance(preview, Path)
    assert isinstance(outputs, Path)
    accepted_bytes = artifact.read_bytes()
    uploader.objects[metadata.artifact_storage_object_path] = accepted_bytes
    descriptor = _freeze_pre_render_inputs(snapshot_inputs, uploader)

    artifact.write_bytes(b"mutated local artifact")
    preview.write_bytes(b"%PDF-1.7\nmutated local preview\n%%EOF")
    (outputs / "deck_build" / "creative_plan.json").write_text(
        '{"subject_materials":["mutated"]}',
        encoding="utf-8",
    )
    loaded = load_pre_render_input_bundle(
        descriptor=descriptor,
        expected_identity=_identity(snapshot_inputs, descriptor),
        reader=uploader,
        materialization_root=tmp_path / "clean-restart",
    )

    assert loaded.artifact_host_path.read_bytes() == accepted_bytes
    assert not (loaded.outputs_root / "psi-deck.preview.pdf").exists()
    assert loaded.metadata.artifact_storage_object_path == metadata.artifact_storage_object_path
    assert uploader.objects[metadata.artifact_storage_object_path] == accepted_bytes
    reconstructed_creative = json.loads((loaded.outputs_root / "deck_build" / "creative_plan.json").read_bytes())
    assert reconstructed_creative["signature"] == "control loop"


def test_snapshot_reuses_exact_immutable_primary_for_evidence(
    snapshot_inputs: dict[str, object],
) -> None:
    uploader = MemoryImmutableUploader()
    outputs = snapshot_inputs["outputs"]
    artifact = snapshot_inputs["artifact"]
    metadata = snapshot_inputs["metadata"]
    assert isinstance(outputs, Path)
    assert isinstance(artifact, Path)
    assert isinstance(metadata, SnapshotCompletionMetadata)
    before = _source_bytes(outputs)

    descriptor = _freeze(snapshot_inputs, uploader)

    assert _source_bytes(outputs) == before
    assert descriptor.snapshot_id == QUALITY_RUN_ID
    assert descriptor.snapshot_path.endswith(f".builder/builds/{BUILD_ID}/quality/{QUALITY_RUN_ID}/evidence_manifest.json")
    assert descriptor.counts.model_dump() == {
        "slide_count": 2,
        "visible_text_slide_count": 2,
        "native_input_count": 3,
        "evidence_object_count": 8,
    }
    assert uploader.attempts[-1][0] == descriptor.snapshot_path
    immutable_paths = [path for path in uploader.objects if path.endswith(".pptx")]
    assert len(immutable_paths) == 1
    immutable_path = immutable_paths[0]
    artifact_hash = hashlib.sha256(artifact.read_bytes()).hexdigest()
    assert immutable_path == metadata.artifact_storage_object_path
    assert immutable_path == immutable_builder_artifact_object_path(
        user_id=metadata.user_id,
        thread_or_session_id=metadata.thread_id,
        logical_artifact_id=metadata.logical_artifact_id,
        artifact_version_id=metadata.artifact_version_id,
        artifact_sha256=artifact_hash,
        filename="psi-deck.pptx",
    )
    assert uploader.objects[immutable_path] == artifact.read_bytes()
    assert immutable_path in uploader.reads

    manifest = json.loads(uploader.objects[descriptor.snapshot_path])
    bundle_path = manifest["evidence_bundle_path"]
    bundle = json.loads(uploader.objects[bundle_path])
    frozen = bundle["snapshot"]
    assert frozen["brief"]["request"] == "Explain the PSI feedback loop."
    assert frozen["renders"]["selectors"] == ["slide:1", "slide:2"]
    assert [item["text"] for item in frozen["visible_text"]] == [
        "Feedback loop\nObserve → appraise → act",
        "Close the loop",
    ]
    assert set(manifest["source_hashes"]) == {
        "input_manifest",
        "artifact",
        "render_source_manifest",
        "render_source_pdf",
        "brief",
        "creative_plan",
        "design_plan",
        "build_record",
        "mechanical_record",
        "visible_text",
    }
    assert set(manifest["render_hashes"]) == {
        "slide:1",
        "slide:2",
        "contact-sheet",
    }
    assert bundle["artifact"]["storage_object_path"] == immutable_path
    assert bundle["artifact"]["sha256"] == artifact_hash


def test_snapshot_fails_closed_when_a_required_native_input_is_missing(
    snapshot_inputs: dict[str, object],
) -> None:
    outputs = snapshot_inputs["outputs"]
    assert isinstance(outputs, Path)
    (outputs / "deck_build" / "design_plan.json").unlink()
    uploader = MemoryImmutableUploader()

    with pytest.raises(SnapshotMissingEvidenceError, match="design_plan"):
        _freeze(snapshot_inputs, uploader)

    assert uploader.objects == {}


@pytest.mark.parametrize(
    "pages",
    [
        (_png_bytes("navy"),),
        (_png_bytes("navy"), b"not-a-png"),
    ],
)
def test_snapshot_rejects_incomplete_or_undecodable_png_coverage(
    snapshot_inputs: dict[str, object],
    pages: tuple[bytes, ...],
) -> None:
    uploader = MemoryImmutableUploader()

    with pytest.raises(SnapshotCoverageError):
        _freeze(snapshot_inputs, uploader, raster_pages=pages)

    assert not any(path.endswith("/evidence_manifest.json") for path in uploader.objects)


def test_identical_snapshot_replay_compares_hashes_and_is_idempotent(
    snapshot_inputs: dict[str, object],
) -> None:
    uploader = MemoryImmutableUploader()
    first = _freeze(snapshot_inputs, uploader)
    first_objects = dict(uploader.objects)
    uploader.attempts.clear()

    second = _freeze(snapshot_inputs, uploader)

    assert second == first
    assert uploader.objects == first_objects
    assert len(uploader.attempts) == first.counts.evidence_object_count - 2
    assert uploader.attempts[-1][0] == first.snapshot_path


def test_snapshot_create_only_artifact_replay_rejects_overwrite(
    snapshot_inputs: dict[str, object],
) -> None:
    uploader = MemoryImmutableUploader()
    descriptor = _freeze(snapshot_inputs, uploader)
    manifest = json.loads(uploader.objects[descriptor.snapshot_path])
    bundle = json.loads(uploader.objects[manifest["evidence_bundle_path"]])
    immutable_path = bundle["artifact"]["storage_object_path"]
    conflicting_bytes = b"pre-existing different immutable artifact"
    uploader.objects[immutable_path] = conflicting_bytes
    uploader.attempts.clear()

    with pytest.raises(SnapshotConflictError, match="accepted_artifact"):
        _freeze(snapshot_inputs, uploader)

    assert uploader.objects[immutable_path] == conflicting_bytes
    assert uploader.attempts == [
        (
            immutable_path,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    ]
    assert descriptor.snapshot_path not in {path for path, _ in uploader.attempts}


def test_later_immutable_primary_corruption_is_rejected_by_historical_snapshot(
    snapshot_inputs: dict[str, object],
    tmp_path: Path,
) -> None:
    uploader = MemoryImmutableUploader()
    metadata = snapshot_inputs["metadata"]
    artifact = snapshot_inputs["artifact"]
    assert isinstance(metadata, SnapshotCompletionMetadata)
    assert isinstance(artifact, Path)
    accepted_bytes = artifact.read_bytes()
    accepted_hash = hashlib.sha256(accepted_bytes).hexdigest()
    uploader.objects[metadata.artifact_storage_object_path] = accepted_bytes

    descriptor = _freeze(snapshot_inputs, uploader)
    manifest = json.loads(uploader.objects[descriptor.snapshot_path])
    bundle = json.loads(uploader.objects[manifest["evidence_bundle_path"]])
    immutable_path = bundle["artifact"]["storage_object_path"]
    assert immutable_path == metadata.artifact_storage_object_path

    uploader.objects[metadata.artifact_storage_object_path] = b"later primary overwrite"
    artifact.write_bytes(b"later local overwrite")
    with pytest.raises(
        SnapshotConflictError,
        match="immutable accepted artifact bytes do not match snapshot",
    ):
        load_evidence_snapshot(
            descriptor=descriptor,
            expected_identity=_identity_from_uploader(snapshot_inputs, uploader),
            reader=uploader,
            materialization_root=tmp_path / "historical-materialization",
        )

    assert accepted_hash != hashlib.sha256(uploader.objects[immutable_path]).hexdigest()
    assert uploader.objects[metadata.artifact_storage_object_path] == b"later primary overwrite"


def test_snapshot_replay_rejects_existing_content_mismatch_before_manifest(
    snapshot_inputs: dict[str, object],
) -> None:
    uploader = MemoryImmutableUploader()
    descriptor = _freeze(snapshot_inputs, uploader)
    manifest = json.loads(uploader.objects[descriptor.snapshot_path])
    bundle_path = manifest["evidence_bundle_path"]
    uploader.objects[bundle_path] = b'{"different":true}'
    uploader.attempts.clear()

    with pytest.raises(SnapshotConflictError, match="evidence_bundle"):
        _freeze(snapshot_inputs, uploader)

    attempted_paths = [path for path, _content_type in uploader.attempts]
    assert bundle_path in attempted_paths
    assert descriptor.snapshot_path not in attempted_paths


def test_snapshot_loader_verifies_and_materializes_only_lossless_renders(
    snapshot_inputs: dict[str, object],
    tmp_path: Path,
) -> None:
    uploader = MemoryImmutableUploader()
    descriptor = _freeze(snapshot_inputs, uploader)

    loaded = load_evidence_snapshot(
        descriptor=descriptor,
        expected_identity=_identity_from_uploader(snapshot_inputs, uploader),
        reader=uploader,
        materialization_root=tmp_path / "materialized",
    )

    assert loaded.descriptor == descriptor
    assert loaded.snapshot.build_id == BUILD_ID
    assert loaded.build_record["build_id"] == BUILD_ID
    local_paths = [
        Path(loaded.snapshot.renders.contact_sheet.path),
        *(Path(image.path) for image in loaded.snapshot.renders.slides),
    ]
    assert all(path.is_file() and path.suffix == ".png" for path in local_paths)
    assert all(QUALITY_RUN_ID in path.parts for path in local_paths)
    assert loaded.snapshot.artifact_path == "/mnt/user-data/outputs/psi-deck.pptx"
    assert not any(path.suffix == ".pptx" for path in local_paths)


def test_evidence_row_identity_is_verified_before_manifest_references(
    snapshot_inputs: dict[str, object],
    tmp_path: Path,
) -> None:
    uploader = MemoryImmutableUploader()
    descriptor = _freeze(snapshot_inputs, uploader)
    expected = _identity_from_uploader(snapshot_inputs, uploader).model_copy(
        update={"thread_id": "different-thread"}
    )
    uploader.reads.clear()

    with pytest.raises(SnapshotStaleError, match="durable row"):
        load_evidence_snapshot(
            descriptor=descriptor,
            expected_identity=expected,
            reader=uploader,
            materialization_root=tmp_path / "stale-evidence",
        )

    assert uploader.reads == [descriptor.snapshot_path]
    assert not (tmp_path / "stale-evidence").exists()


def test_snapshot_loader_rejects_remote_or_local_render_mismatch(
    snapshot_inputs: dict[str, object],
    tmp_path: Path,
) -> None:
    uploader = MemoryImmutableUploader()
    descriptor = _freeze(snapshot_inputs, uploader)
    loaded = load_evidence_snapshot(
        descriptor=descriptor,
        expected_identity=_identity_from_uploader(snapshot_inputs, uploader),
        reader=uploader,
        materialization_root=tmp_path / "materialized",
    )
    first_object_path = json.loads(uploader.objects[descriptor.snapshot_path])["objects"][0]["object_path"]
    original_remote = uploader.objects[first_object_path]
    uploader.objects[first_object_path] = b"corrupt"
    with pytest.raises(SnapshotConflictError, match="render object hash"):
        load_evidence_snapshot(
            descriptor=descriptor,
            expected_identity=_identity_from_uploader(snapshot_inputs, uploader),
            reader=uploader,
            materialization_root=tmp_path / "second-materialization",
        )

    uploader.objects[first_object_path] = original_remote
    first_local = Path(loaded.snapshot.renders.slides[0].path)
    first_local.write_bytes(_png_bytes("red"))
    with pytest.raises(SnapshotConflictError, match="local materialization"):
        load_evidence_snapshot(
            descriptor=descriptor,
            expected_identity=_identity_from_uploader(snapshot_inputs, uploader),
            reader=uploader,
            materialization_root=tmp_path / "materialized",
        )


def test_pdf_rasterizer_stages_source_in_fixed_resolved_bounded_poppler_command(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    preview = tmp_path / "deck.preview.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=720, height=405)
    with preview.open("wb") as stream:
        writer.write(stream)
    captured: dict[str, object] = {}
    png = _png_bytes("navy")
    monkeypatch.setattr(
        snapshot_module.shutil,
        "which",
        lambda binary: "/opt/poppler/bin/pdftoppm" if binary == "pdftoppm" else None,
    )

    def fake_run(
        command: list[str],
        *,
        timeout: int,
        private_read_dirs,
        writable_dirs,
        identity_paths,
    ):
        captured["command"] = command
        captured["timeout"] = timeout
        captured["private_read_dirs"] = private_read_dirs
        captured["writable_dirs"] = writable_dirs
        captured["identity_paths"] = identity_paths
        captured["staged_bytes"] = Path(command[-2]).read_bytes()
        prefix = Path(command[-1])
        prefix.with_name(prefix.name + "-1.png").write_bytes(png)
        return SimpleNamespace(returncode=0, stdout="", stderr="")

    monkeypatch.setattr(snapshot_module, "run_process_group", fake_run)

    rendered = rasterize_preview_pdf(preview)

    assert rendered == (png,)
    command = captured["command"]
    assert isinstance(command, list)
    staged_source = Path(command[-2])
    assert command == [
        "/opt/poppler/bin/pdftoppm",
        "-png",
        "-scale-to",
        "2200",
        str(staged_source),
        command[-1],
    ]
    assert staged_source == Path(command[-1]).parent / "source.pdf"
    assert captured["staged_bytes"] == preview.read_bytes()
    assert captured["timeout"] == 180
    assert captured["private_read_dirs"] == [staged_source]
    assert len(captured["writable_dirs"]) == 1  # type: ignore[arg-type]
    assert captured["writable_dirs"] == [str(staged_source.parent)]
    assert captured["identity_paths"] == [preview.resolve()]


def test_pdf_rasterizer_rejects_symlink_before_privileged_staging(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    real_preview = tmp_path / "real.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=720, height=405)
    with real_preview.open("wb") as stream:
        writer.write(stream)
    linked_preview = tmp_path / "linked.pdf"
    linked_preview.symlink_to(real_preview)
    monkeypatch.setattr(
        snapshot_module.shutil,
        "which",
        lambda binary: "/opt/poppler/bin/pdftoppm"
        if binary == "pdftoppm"
        else None,
    )
    monkeypatch.setattr(
        snapshot_module,
        "run_process_group",
        lambda *_args, **_kwargs: pytest.fail("renderer must not run"),
    )

    with pytest.raises(SnapshotMissingEvidenceError):
        rasterize_preview_pdf(linked_preview)


@pytest.mark.skipif(
    platform.system() != "Linux"
    or not hasattr(os, "geteuid")
    or os.geteuid() != 0
    or shutil.which("pdftoppm") is None
    or shutil.which("setpriv") is None,
    reason="requires the real root-Linux Poppler identity boundary",
)
def test_real_root_linux_pdf_rasterizer_reads_source_below_private_ancestor() -> None:
    private_root = Path(tempfile.mkdtemp(prefix="dq1-private-pdf-source-"))
    private_root.chmod(0o700)
    try:
        preview = private_root / "private.preview.pdf"
        writer = PdfWriter()
        writer.add_blank_page(width=720, height=405)
        with preview.open("wb") as stream:
            writer.write(stream)

        rendered = rasterize_preview_pdf(preview)

        assert len(rendered) == 1
        assert rendered[0].startswith(b"\x89PNG\r\n\x1a\n")
    finally:
        shutil.rmtree(private_root, ignore_errors=True)
