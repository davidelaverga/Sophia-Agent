from __future__ import annotations

from pathlib import Path

import pytest
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt
from test_deck_build_service import _creative_plan, _runtime, _slides

from deerflow.sophia.deck_build.creative_plan import normalize_creative_plan
from deerflow.sophia.deck_build.image_assets import apply_creative_asset_plan
from deerflow.sophia.deck_build.mechanical_gates import evaluate_mechanical_gates
from deerflow.sophia.deck_build.models import DeckBuild
from deerflow.sophia.deck_build.service import DeckBuildService


def _built_deck(tmp_path, *, repeated: bool = False, old_marker: bool = False):
    runtime = _runtime(tmp_path / "outputs")
    service = DeckBuildService()
    slides = service._build_slide_specs(
        _slides(),
        visual_policy="auto",
        runtime=runtime,
        style_profile={},
    )
    loaded = DeckBuild(
        build_id="deck-test",
        schema_version="sophia-deck-build/v1",
        user_id="user",
        thread_id="thread",
        parent_thread_id=None,
        run_id=None,
        task_id=None,
        requested_slide_count=len(slides),
        status="compiled",
        register="professional_technical",
        visual_policy="auto",
        style_profile={},
        deck_title="Technical Deck",
        output_path="/mnt/user-data/outputs/deck.pptx",
        slides=slides,
        expected_visual_count=0,
    )
    creative_plan = normalize_creative_plan(_creative_plan(), deck=loaded, request_context="")
    loaded.creative_plan = creative_plan
    loaded.design_plan = creative_plan.design_plan
    apply_creative_asset_plan(loaded, creative_plan)
    if repeated:
        for slide in loaded.slides:
            if slide.composition_plan is not None:
                slide.composition_plan.layout_name = "same_layout"
    if old_marker:
        loaded.slides[0].html_source = (loaded.slides[0].html_source or "") + "<!-- deck_build_templates_v1 -->"
    return loaded


def _render_dir(tmp_path: Path, *, light: bool = False, blank: bool = False) -> Path:
    render_dir = tmp_path / "rendered"
    render_dir.mkdir()
    for index in range(1, 4):
        image = Image.new("RGB", (320, 180), "#FFFFFF" if light else "#0A0E14")
        if not blank:
            draw = ImageDraw.Draw(image)
            draw.rectangle((30, 30, 260, 140), outline="#38BDF8", width=8)
            draw.rectangle((60, 70, 120, 110), fill="#38BDF8")
        image.save(render_dir / f"slide-{index}.jpg")
    return render_dir


def _native_pptx(tmp_path: Path, *, shape_name: str, text: str, font_size: float) -> Path:
    path = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    textbox.name = shape_name
    paragraph = textbox.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.runs[0].font.size = Pt(font_size)
    presentation.save(path)
    return path


def test_mechanical_gates_pass_for_nonblank_dark_native_deck(tmp_path) -> None:
    deck = _built_deck(tmp_path)

    result = evaluate_mechanical_gates(deck, rendered_dir=_render_dir(tmp_path))

    assert result.passed is True


def test_mechanical_gates_fail_sparse_rendered_slide(tmp_path) -> None:
    deck = _built_deck(tmp_path)

    result = evaluate_mechanical_gates(deck, rendered_dir=_render_dir(tmp_path, blank=True))

    assert result.passed is False
    assert any(issue.code == "sparse_rendered_slide" for issue in result.issues)


def test_mechanical_gates_fail_repeated_skeleton_and_old_renderer_marker(tmp_path) -> None:
    deck = _built_deck(tmp_path, repeated=True, old_marker=True)

    result = evaluate_mechanical_gates(deck, rendered_dir=_render_dir(tmp_path))

    assert result.passed is False
    assert {issue.code for issue in result.issues} >= {"repeated_slide_skeleton", "old_renderer_artifact"}


def test_single_legacy_class_name_is_not_misclassified_as_old_renderer(tmp_path) -> None:
    for marker in ("section-label", "system-diagram", "closing-synthesis"):
        deck = _built_deck(tmp_path / marker)
        deck.slides[0].html_source = f"<div class='{marker}'>Subject-specific content</div>"

        result = evaluate_mechanical_gates(deck, rendered_dir=_render_dir(tmp_path / marker))

        assert not any(issue.code == "old_renderer_artifact" for issue in result.issues)


def test_compound_legacy_class_signature_is_rejected(tmp_path) -> None:
    deck = _built_deck(tmp_path)
    deck.slides[0].html_source = """
    <div class='section-label'></div>
    <div class='system-diagram'></div>
    <div class='closing-synthesis'></div>
    """

    result = evaluate_mechanical_gates(deck, rendered_dir=_render_dir(tmp_path))

    assert any(issue.code == "old_renderer_artifact" for issue in result.issues)


def test_legacy_class_markers_split_across_slides_are_not_rejected(tmp_path) -> None:
    deck = _built_deck(tmp_path)
    for slide, marker in zip(
        deck.slides,
        ("section-label", "system-diagram", "closing-synthesis"),
        strict=False,
    ):
        slide.html_source = f"<div class='{marker}'>Subject-specific content</div>"

    result = evaluate_mechanical_gates(deck, rendered_dir=_render_dir(tmp_path))

    assert not any(issue.code == "old_renderer_artifact" for issue in result.issues)


def test_mechanical_gates_fail_dark_request_rendered_light(tmp_path) -> None:
    deck = _built_deck(tmp_path)

    result = evaluate_mechanical_gates(deck, rendered_dir=_render_dir(tmp_path, light=True))

    assert result.passed is False
    assert any(issue.code == "dark_request_rendered_light" for issue in result.issues)


def test_warm_ivory_plan_with_ink_black_text_is_not_treated_as_dark(tmp_path) -> None:
    deck = _built_deck(tmp_path)
    assert deck.design_plan is not None
    deck.style_profile = {"palette": "warm ivory, ink black, muted cobalt, ember"}
    deck.design_plan.style_lane = "executive_editorial"
    deck.design_plan.signature = "restrained editorial on warm ivory"
    deck.design_plan.requested_style_terms = [
        "restrained editorial",
        "warm ivory",
        "ink black",
        "muted cobalt",
        "ember accent",
    ]

    result = evaluate_mechanical_gates(deck, rendered_dir=_render_dir(tmp_path, light=True))

    assert not any(issue.code == "dark_request_rendered_light" for issue in result.issues)


def test_warm_ivory_plan_rendered_dark_fails_substrate_gate(tmp_path) -> None:
    deck = _built_deck(tmp_path)
    assert deck.design_plan is not None
    deck.style_profile = {"colors": {"background": "warm ivory", "text": "ink black"}}
    deck.design_plan.style_lane = "executive_editorial"
    deck.design_plan.signature = "restrained editorial on warm ivory"
    deck.design_plan.requested_style_terms = ["warm ivory", "ink black"]

    result = evaluate_mechanical_gates(deck, rendered_dir=_render_dir(tmp_path))

    assert any(issue.code == "light_request_rendered_dark" for issue in result.issues)


def test_requested_light_terms_override_stale_dark_style_profile(tmp_path) -> None:
    deck = _built_deck(tmp_path)
    assert deck.design_plan is not None
    deck.style_profile = {"background": "dark charcoal substrate"}
    deck.design_plan.style_lane = "executive_editorial"
    deck.design_plan.signature = "quiet editorial substrate"
    deck.design_plan.requested_style_terms = ["warm ivory", "ink black"]

    result = evaluate_mechanical_gates(deck, rendered_dir=_render_dir(tmp_path, light=True))

    assert not any(issue.code == "dark_request_rendered_light" for issue in result.issues)


def test_resolved_light_plan_overrides_stale_dark_profile_when_terms_are_neutral(tmp_path) -> None:
    deck = _built_deck(tmp_path)
    assert deck.design_plan is not None
    deck.style_profile = {"background": "dark charcoal substrate"}
    deck.design_plan.style_lane = "executive_editorial"
    deck.design_plan.signature = "warm ivory substrate with restrained editorial geometry"
    deck.design_plan.requested_style_terms = ["executive"]

    result = evaluate_mechanical_gates(deck, rendered_dir=_render_dir(tmp_path, light=True))

    assert not any(issue.code == "dark_request_rendered_light" for issue in result.issues)


def test_requested_dark_terms_override_stale_light_style_profile(tmp_path) -> None:
    deck = _built_deck(tmp_path)
    assert deck.design_plan is not None
    deck.style_profile = {"colors": {"background": "white", "text": "black"}}
    deck.design_plan.style_lane = "custom"
    deck.design_plan.signature = "restrained technical geometry"
    deck.design_plan.requested_style_terms = ["dark charcoal substrate", "warm ivory text"]

    result = evaluate_mechanical_gates(deck, rendered_dir=_render_dir(tmp_path, light=True))

    assert any(issue.code == "dark_request_rendered_light" for issue in result.issues)


@pytest.mark.parametrize(
    ("colors", "expected_code"),
    [
        ({"background": "white", "text": "black"}, None),
        ({"text": "black", "background": "white"}, None),
        ({"background": "black", "text": "white"}, "dark_request_rendered_light"),
        ({"text": "white", "background": "black"}, "dark_request_rendered_light"),
    ],
)
def test_structured_style_profile_preserves_background_text_roles(
    tmp_path,
    colors: dict[str, str],
    expected_code: str | None,
) -> None:
    deck = _built_deck(tmp_path)
    assert deck.design_plan is not None
    deck.style_profile = {"colors": colors}
    deck.design_plan.style_lane = "custom"
    deck.design_plan.signature = "restrained geometry"
    deck.design_plan.requested_style_terms = []

    result = evaluate_mechanical_gates(deck, rendered_dir=_render_dir(tmp_path, light=True))

    codes = {issue.code for issue in result.issues}
    if expected_code is None:
        assert "dark_request_rendered_light" not in codes
    else:
        assert expected_code in codes


def test_unknown_style_profile_key_cannot_trigger_substrate_gate(tmp_path) -> None:
    deck = _built_deck(tmp_path)
    assert deck.design_plan is not None
    deck.style_profile = {"custom_css": "body { background: black; }"}
    deck.design_plan.style_lane = "custom"
    deck.design_plan.signature = "restrained geometry"
    deck.design_plan.requested_style_terms = []

    result = evaluate_mechanical_gates(deck, rendered_dir=_render_dir(tmp_path, light=True))

    assert "dark_request_rendered_light" not in {issue.code for issue in result.issues}


def test_dark_charcoal_plan_with_warm_ivory_text_is_still_treated_as_dark(tmp_path) -> None:
    deck = _built_deck(tmp_path)
    assert deck.design_plan is not None
    deck.style_profile = {"palette": "dark charcoal substrate, warm ivory text"}
    deck.design_plan.style_lane = "technical_blueprint"
    deck.design_plan.signature = "dark charcoal substrate with warm ivory text"
    deck.design_plan.requested_style_terms = ["dark charcoal", "warm ivory text"]

    result = evaluate_mechanical_gates(deck, rendered_dir=_render_dir(tmp_path, light=True))

    assert any(issue.code == "dark_request_rendered_light" for issue in result.issues)


def test_non_text_overflow_requires_explicit_bleed_source_role(tmp_path) -> None:
    deck = _built_deck(tmp_path)
    deck.native_mechanical_report = {
        "lint_residue_count": 1,
        "lint_residue_kinds": {"slide_overflow_non_text": 1},
        "lint_residue": [
            {
                "slide": 0,
                "shape": "s1-background-box",
                "kind": "slide_overflow_non_text",
            }
        ],
    }

    failed = evaluate_mechanical_gates(deck, rendered_dir=_render_dir(tmp_path))
    assert any(issue.code == "native_lint_unapproved_bleed" for issue in failed.issues)

    deck.source_element_map = {
        "slides": {
            "slide:1": {
                "elements": {
                    "background": {
                        "source_role": "background",
                        "shape_names": ["s1-background-box"],
                    }
                }
            }
        }
    }
    allowed_root = tmp_path / "allowed"
    allowed_root.mkdir()
    allowed = evaluate_mechanical_gates(deck, rendered_dir=_render_dir(allowed_root))
    assert allowed.passed is True
    assert "native_lint_advisory:slide_overflow_non_text" in allowed.warnings


def test_non_text_overflow_resolves_native_inventory_id_to_direct_background_role(tmp_path) -> None:
    deck = _built_deck(tmp_path)
    deck.native_mechanical_report = {
        "lint_residue_count": 1,
        "lint_residue_kinds": {"slide_overflow_non_text": 1},
        "lint_residue": [{"slide": 0, "shape": "s10", "kind": "slide_overflow_non_text"}],
    }
    deck.native_shape_inventory = {
        "slide:1": {
            "shapes": [
                {"id": "s10", "name": "h2p-1-background-line-1-part-2"},
            ]
        }
    }
    deck.source_element_map = {
        "slides": {
            "slide:1": {
                "elements": {
                    "canvas": {
                        "source_role": "diagram",
                        "shape_names": ["h2p-1-background-line-1-part-2"],
                    },
                    "background": {
                        "source_role": "background",
                        "shape_names": ["h2p-1-background-line-1-part-2"],
                    }
                }
            }
        }
    }

    result = evaluate_mechanical_gates(deck, rendered_dir=_render_dir(tmp_path))

    assert result.passed is True
    assert "native_lint_advisory:slide_overflow_non_text" in result.warnings


def test_allowed_ancestor_role_does_not_hide_overflowing_semantic_child(tmp_path) -> None:
    deck = _built_deck(tmp_path)
    shape_name = "h2p-1-evidence-panel-box-1"
    deck.native_mechanical_report = {
        "lint_residue_count": 1,
        "lint_residue_kinds": {"slide_overflow_non_text": 1},
        "lint_residue": [{"slide": 0, "shape": "s10", "kind": "slide_overflow_non_text"}],
    }
    deck.native_shape_inventory = {
        "slide:1": {"shapes": [{"id": "s10", "name": shape_name}]}
    }
    deck.source_element_map = {
        "slides": {
            "slide:1": {
                "elements": {
                    "panel": {
                        "source_role": "background",
                        "shape_names": [shape_name],
                    },
                    "evidence-panel": {
                        "source_role": "evidence_panel",
                        "shape_names": [shape_name],
                    },
                }
            }
        }
    }

    result = evaluate_mechanical_gates(deck, rendered_dir=_render_dir(tmp_path))

    assert any(issue.code == "native_lint_unapproved_bleed" for issue in result.issues)


def test_post_fix_frame_overflow_and_misalignment_are_blocking_and_attributed(tmp_path) -> None:
    deck = _built_deck(tmp_path)
    deck.native_mechanical_report = {
        "lint_residue_count": 2,
        "lint_residue_kinds": {"frame_overflow": 1, "misaligned": 1},
        "lint_residue": [
            {
                "slide": 1,
                "shape": "s7",
                "kind": "frame_overflow",
                "overflow_bottom": 0.47,
            },
            {
                "slide": 2,
                "shape": "s11",
                "kind": "misaligned",
                "issue": 'vcenter edge 0.09" off gridline',
            },
        ],
    }

    result = evaluate_mechanical_gates(deck, rendered_dir=_render_dir(tmp_path))

    assert result.passed is False
    issues = {(issue.code, issue.selector) for issue in result.issues}
    assert ("native_lint_frame_overflow", "slide:2") in issues
    assert ("native_lint_misaligned", "slide:3") in issues
    assert not any(issue.selector == "deck" for issue in result.issues)


@pytest.mark.parametrize(
    ("font_size", "expected_passed"),
    [
        (17.9, False),
        (18.0, True),
    ],
)
def test_required_body_uses_computed_24px_floor(
    tmp_path: Path,
    font_size: float,
    expected_passed: bool,
) -> None:
    deck = _built_deck(tmp_path)
    deck.source_element_map = {
        "slides": {
            "slide:1": {
                "elements": {
                    "narrative-1": {
                        "source_role": "narrative",
                        "source_required": True,
                        "shape_names": ["s1-narrative-text"],
                    }
                }
            }
        }
    }
    pptx = _native_pptx(
        tmp_path,
        shape_name="s1-narrative-text",
        text="Rendered body copy",
        font_size=font_size,
    )

    result = evaluate_mechanical_gates(
        deck,
        rendered_dir=_render_dir(tmp_path),
        native_pptx_path=pptx,
    )

    assert result.passed is expected_passed
    codes = {issue.code for issue in result.issues}
    if expected_passed:
        assert "native_required_text_too_small" not in codes
    else:
        assert "native_required_text_too_small" in codes
        assert "24px" in result.failure_summary


@pytest.mark.parametrize(
    ("font_size", "expected_passed"),
    [
        (14.9, False),
        (15.0, True),
        (16.5, True),
    ],
)
def test_optional_compact_labels_use_computed_20px_floor(
    tmp_path: Path,
    font_size: float,
    expected_passed: bool,
) -> None:
    deck = _built_deck(tmp_path)
    deck.source_element_map = {
        "slides": {
            "slide:1": {
                "elements": {
                    "status-label": {
                        "source_role": "label",
                        "source_required": False,
                        "shape_names": ["s1-status-label-text"],
                    }
                }
            }
        }
    }
    pptx = _native_pptx(
        tmp_path,
        shape_name="s1-status-label-text",
        text="Success measure",
        font_size=font_size,
    )

    result = evaluate_mechanical_gates(
        deck,
        rendered_dir=_render_dir(tmp_path),
        native_pptx_path=pptx,
    )

    assert result.passed is expected_passed
    codes = {issue.code for issue in result.issues}
    if expected_passed:
        assert "native_compact_text_too_small" not in codes
    else:
        assert "native_compact_text_too_small" in codes
        assert "20px" in result.failure_summary


def test_compiled_typography_checks_only_emitted_text(tmp_path: Path) -> None:
    deck = _built_deck(tmp_path)
    deck.slides[0].html_source = (
        (deck.slides[0].html_source or "")
        + "<style>.unused-utility{font-size:12px}</style>"
    )
    pptx = _native_pptx(
        tmp_path,
        shape_name="s1-visible-label-text",
        text="Visible compact label",
        font_size=15.0,
    )

    result = evaluate_mechanical_gates(
        deck,
        rendered_dir=_render_dir(tmp_path),
        native_pptx_path=pptx,
    )

    assert result.passed is True
    assert not any("text_too_small" in issue.code for issue in result.issues)
