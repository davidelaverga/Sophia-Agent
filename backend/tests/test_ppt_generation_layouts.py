"""Tests for the v4.1 pure image-forward PPTX compiler wrapper."""

from __future__ import annotations

import importlib.util
import json
import os
import shutil
import subprocess
import sys
import zipfile
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation

_REPO_ROOT = Path(__file__).resolve().parents[2]
_SCRIPT_PATH = _REPO_ROOT / "skills" / "public" / "ppt-generation" / "scripts" / "generate.py"
_JS_COMPILER_PATH = (
    _REPO_ROOT / "backend" / "packages" / "harness" / "deerflow" / "sophia" / "js" / "compile_pptx.mjs"
)
_CODEX_NODE_BIN = Path.home() / ".cache/codex-runtimes/codex-primary-runtime/dependencies/node/bin/node"
_CODEX_NODE_MODULES = Path.home() / ".cache/codex-runtimes/codex-primary-runtime/dependencies/node/node_modules"

_REQUIRED_OFFICE_ENTRIES = {"[Content_Types].xml", "_rels/.rels", "ppt/presentation.xml"}


def _load_module():
    spec = importlib.util.spec_from_file_location("ppt_generation_layouts_script", _SCRIPT_PATH)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


gen = _load_module()


def _write_png(path: Path, size=(1280, 720), color=(40, 140, 180)) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", size, color=color).save(path)
    return path


def _slide_texts(slide) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.extend(paragraph.text for paragraph in shape.text_frame.paragraphs)
    return [text for text in texts if text]


def _pptx_media_count(path: Path) -> int:
    with zipfile.ZipFile(path) as archive:
        return sum(1 for name in archive.namelist() if name.startswith("ppt/media/"))


def _write_plan(path: Path, slides: list[dict], **extra) -> Path:
    path.write_text(json.dumps({"title": "Image Deck", "slides": slides, **extra}), encoding="utf-8")
    return path


def _pptxgenjs_runtime(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> Path:
    node_modules_candidates = [
        _JS_COMPILER_PATH.parent / "node_modules",
        _CODEX_NODE_MODULES,
    ]
    node_modules = next(
        (path for path in node_modules_candidates if (path / "pptxgenjs" / "package.json").is_file()),
        None,
    )
    if node_modules is None:
        pytest.skip("pptxgenjs is not installed for local JS compiler tests")
    runtime_dir = tmp_path / "pptxgenjs-runtime"
    runtime_dir.mkdir()
    shutil.copy2(_JS_COMPILER_PATH, runtime_dir / "compile_pptx.mjs")
    os.symlink(node_modules, runtime_dir / "node_modules", target_is_directory=True)
    if _CODEX_NODE_BIN.is_file():
        monkeypatch.setenv("PATH", f"{_CODEX_NODE_BIN.parent}{os.pathsep}{os.environ.get('PATH', '')}")
    monkeypatch.setenv("SOPHIA_PPTXGENJS", "1")
    monkeypatch.setenv("SOPHIA_ARTIFACT_JS_RUNTIME", str(runtime_dir))
    return runtime_dir


def test_generate_ppt_embeds_one_full_slide_image_per_slide(tmp_path: Path) -> None:
    slide_1 = _write_png(tmp_path / "slide-1.png", color=(20, 80, 160))
    slide_2 = _write_png(tmp_path / "slide-2.png", color=(160, 80, 20))
    plan = _write_plan(
        tmp_path / "plan.json",
        [
            {"slide_number": 1, "title": "Baked Title 1", "image_path": str(slide_1)},
            {"slide_number": 2, "title": "Baked Title 2", "image_path": str(slide_2)},
        ],
    )
    output = tmp_path / "deck.pptx"

    message = gen.generate_ppt(str(plan), [], str(output))

    assert message == "Successfully generated presentation with 2 slides (picture_count=2)"
    prs = Presentation(str(output))
    assert len(prs.slides) == 2
    assert _slide_texts(prs.slides[0]) == []
    assert _slide_texts(prs.slides[1]) == []
    with zipfile.ZipFile(output) as archive:
        assert _REQUIRED_OFFICE_ENTRIES.issubset(set(archive.namelist()))
    assert _pptx_media_count(output) >= 2


def test_generate_ppt_accepts_cli_slide_images_without_text_layout(tmp_path: Path) -> None:
    slide_1 = _write_png(tmp_path / "cli-1.png")
    slide_2 = _write_png(tmp_path / "cli-2.png", color=(10, 100, 90))
    plan = _write_plan(
        tmp_path / "plan.json",
        [
            {"slide_number": 1, "title": "Title in bitmap"},
            {"slide_number": 2, "title": "Also in bitmap", "key_points": ["Notes only"]},
        ],
    )
    output = tmp_path / "deck.pptx"

    message = gen.generate_ppt(str(plan), [str(slide_1), str(slide_2)], str(output))

    assert message == "Successfully generated presentation with 2 slides (picture_count=2)"
    prs = Presentation(str(output))
    assert all(_slide_texts(slide) == [] for slide in prs.slides)


def test_generate_ppt_rejects_plan_without_generated_images(tmp_path: Path) -> None:
    plan = _write_plan(
        tmp_path / "plan.json",
        [
            {"slide_number": 1, "title": "Native title must not render"},
            {"slide_number": 2, "title": "Native body must not render", "key_points": ["A", "B"]},
        ],
    )
    output = tmp_path / "deck.pptx"

    with pytest.raises(ValueError, match="Slide 1 is missing its generated slide image"):
        gen.generate_ppt(str(plan), [], str(output))

    assert not output.exists()


def test_generate_ppt_rejects_missing_image_instead_of_text_fallback(tmp_path: Path) -> None:
    plan = _write_plan(
        tmp_path / "plan.json",
        [{"slide_number": 1, "title": "Missing", "image_path": str(tmp_path / "missing.png")}],
    )
    output = tmp_path / "deck.pptx"

    with pytest.raises(FileNotFoundError, match="Slide 1 image not found"):
        gen.generate_ppt(str(plan), [], str(output))

    assert not output.exists()


def test_generate_ppt_rejects_cli_image_count_mismatch(tmp_path: Path) -> None:
    slide_1 = _write_png(tmp_path / "slide-1.png")
    plan = _write_plan(
        tmp_path / "plan.json",
        [
            {"slide_number": 1, "title": "One"},
            {"slide_number": 2, "title": "Two"},
        ],
    )

    with pytest.raises(ValueError, match="Slide image count does not match"):
        gen.generate_ppt(str(plan), [str(slide_1)], str(tmp_path / "deck.pptx"))


def test_cli_diagnostics_contract_for_image_only_deck(tmp_path: Path) -> None:
    slide_1 = _write_png(tmp_path / "slide-1.png")
    slide_2 = _write_png(tmp_path / "slide-2.png", color=(160, 80, 20))
    plan = _write_plan(
        tmp_path / "plan.json",
        [
            {"slide_number": 1, "title": "One", "image_path": str(slide_1)},
            {"slide_number": 2, "title": "Two", "image_path": str(slide_2)},
        ],
    )
    output = tmp_path / "deck.pptx"

    result = subprocess.run(
        [sys.executable, str(_SCRIPT_PATH), "--plan-file", str(plan), "--output-file", str(output)],
        capture_output=True,
        text=True,
        timeout=30,
    )

    assert result.returncode == 0, result.stderr
    assert "Successfully generated presentation with 2 slides (picture_count=2)" in result.stdout
    assert "PPT generation diagnostics:" in result.stderr
    assert "picture_count=2" in result.stderr
    assert "output_ext=.pptx" in result.stderr


def test_cli_failure_does_not_emit_text_layout_deck(tmp_path: Path) -> None:
    plan = _write_plan(tmp_path / "plan.json", [{"slide_number": 1, "title": "No image"}])
    output = tmp_path / "deck.pptx"

    result = subprocess.run(
        [sys.executable, str(_SCRIPT_PATH), "--plan-file", str(plan), "--output-file", str(output)],
        capture_output=True,
        text=True,
        timeout=30,
    )

    assert result.returncode == 1
    assert "missing its generated slide image" in result.stderr
    assert not output.exists()


def test_js_compiler_is_image_only() -> None:
    source = _JS_COMPILER_PATH.read_text(encoding="utf-8")

    assert "function renderImageForward" in source
    assert "addText(" not in source
    assert "function slideType" not in source
    assert "rendererForSlideType" not in source
    assert "native_title_overlay: false" in source
    assert "native_caption_overlay: false" in source


def test_js_compiler_delegation_stays_picture_only(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    _pptxgenjs_runtime(tmp_path, monkeypatch)
    slide = _write_png(tmp_path / "slide.png")
    plan = _write_plan(
        tmp_path / "plan.json",
        [{"slide_number": 1, "title": "Visible only inside bitmap", "image_path": str(slide)}],
    )
    output = tmp_path / "deck.pptx"

    message = gen.generate_ppt(str(plan), [], str(output))

    assert message == "Successfully generated presentation with PptxGenJS"
    prs = Presentation(str(output))
    assert len(prs.slides) == 1
    assert _slide_texts(prs.slides[0]) == []
