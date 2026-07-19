from __future__ import annotations

import hashlib
import hmac
import io
import json
import math
import os
import re
import shutil
import stat
import subprocess  # noqa: S404 - fixed arguments and a path-resolved pdftoppm binary
import tempfile
import zipfile
from collections.abc import Callable, Mapping
from pathlib import Path, PurePosixPath
from typing import Annotated, Any, Literal, Protocol

from PIL import Image
from pptx import Presentation
from pydantic import BaseModel, ConfigDict, Field, field_validator, model_validator
from pypdf import PdfReader

from deerflow.sophia.deck_quality.brief import sanitize_current_request
from deerflow.sophia.deck_quality.canonical import canonical_json_bytes, canonical_sha256
from deerflow.sophia.deck_quality.contact_sheet import create_contact_sheet
from deerflow.sophia.deck_quality.messages import (
    DIRECT_EVIDENCE_MAX_SLIDE_DIMENSION,
)
from deerflow.sophia.deck_quality.schemas import (
    BlindBrief,
    ImageEvidence,
    QualityEvidenceSnapshot,
    RenderEvidence,
    Sha256,
    StableSlideSelector,
    VisibleTextSlide,
)
from deerflow.sophia.deck_quality.visible_text import visible_text_sidecar
from deerflow.sophia.pptx_preview import maybe_render_pptx_preview
from deerflow.sophia.process_group import run_process_group
from deerflow.sophia.storage.supabase_artifact_store import (
    immutable_builder_artifact_object_path,
    normalize_object_path,
    safe_object_path_segment,
)

_QUALITY_RUN_RE = re.compile(r"^quality_[0-9a-f]{64}$")
_SAFE_BUILD_ID_RE = re.compile(r"^[A-Za-z0-9][A-Za-z0-9._-]{0,127}$")
_PDFTOPPM_MAX_DIMENSION = DIRECT_EVIDENCE_MAX_SLIDE_DIMENSION
_PDFTOPPM_TIMEOUT_SECONDS = 180
_MAX_ACCEPTED_PPTX_BYTES = 32 * 1024 * 1024
_MAX_ACCEPTED_PREVIEW_PDF_BYTES = 32 * 1024 * 1024
_MAX_NATIVE_JSON_BYTES = 2 * 1024 * 1024
_MAX_EVIDENCE_BUNDLE_BYTES = 4 * 1024 * 1024
_MAX_RENDER_PNG_BYTES = 1024 * 1024
_MAX_RENDER_TOTAL_BYTES = 3 * 1024 * 1024
_MAX_PPTX_MEMBER_BYTES = 64 * 1024 * 1024
_MAX_PPTX_UNCOMPRESSED_BYTES = 128 * 1024 * 1024
_MAX_PPTX_MEMBER_COUNT = 4_096
_MAX_RENDER_PAGES = 5
_MAX_RENDER_PAGE_PIXELS = _PDFTOPPM_MAX_DIMENSION**2
_MAX_RENDER_TOTAL_PIXELS = _MAX_RENDER_PAGE_PIXELS * _MAX_RENDER_PAGES
PreRenderInputRole = Literal[
    "accepted_artifact",
    "creative_plan",
    "design_plan",
    "build_record",
    "blind_brief",
    "mechanical_record",
]
_PRE_RENDER_INPUT_ROLES: tuple[PreRenderInputRole, ...] = (
    "accepted_artifact",
    "creative_plan",
    "design_plan",
    "build_record",
    "blind_brief",
    "mechanical_record",
)
_RENDER_SOURCE_PROFILE_VERSION = "libreoffice-impress-pdf-v1"
_RENDER_SOURCE_PROFILE_HASH = canonical_sha256(
    {
        "profile_version": _RENDER_SOURCE_PROFILE_VERSION,
        "artifact_format": "pptx",
        "render_format": "pdf",
        "maximum_pages": _MAX_RENDER_PAGES,
    }
)


class SnapshotError(RuntimeError):
    """A safe, content-free evidence snapshot failure."""

    code = "quality_persistence_error"


class SnapshotMissingEvidenceError(SnapshotError):
    code = "snapshot_evidence_missing"


class SnapshotCoverageError(SnapshotError):
    code = "coverage_error"


class SnapshotStaleError(SnapshotError):
    code = "artifact_snapshot_stale"


class SnapshotConflictError(SnapshotError):
    code = "immutable_snapshot_conflict"


class SnapshotUploadError(SnapshotError):
    code = "quality_persistence_error"


class _StrictFrozenModel(BaseModel):
    model_config = ConfigDict(extra="forbid", frozen=True)


def _strip_required_identity(value: str) -> str:
    stripped = value.strip()
    if not stripped:
        raise ValueError("identity cannot be blank")
    return stripped


class SnapshotCompletionMetadata(_StrictFrozenModel):
    """Only completion-safe correlation and durable artifact metadata."""

    campaign_id: Literal["DQ-1"] = "DQ-1"
    quality_run_id: str = Field(pattern=r"^quality_[0-9a-f]{64}$")
    build_id: str = Field(pattern=r"^[A-Za-z0-9][A-Za-z0-9._-]{0,127}$")
    user_id: str = Field(min_length=1, max_length=512)
    thread_id: str = Field(min_length=1, max_length=512)
    task_id: str = Field(min_length=1, max_length=512)
    builder_run_id: str = Field(min_length=1, max_length=512)
    parent_builder_trace_id: str = Field(min_length=1, max_length=512)
    logical_artifact_id: str = Field(min_length=1, max_length=512)
    artifact_version_id: str = Field(min_length=1, max_length=512)
    manifest_revision: int = Field(ge=1)
    artifact_storage_object_path: str = Field(min_length=1, max_length=4_096)

    @field_validator(
        "user_id",
        "thread_id",
        "task_id",
        "builder_run_id",
        "parent_builder_trace_id",
        "logical_artifact_id",
        "artifact_version_id",
    )
    @classmethod
    def normalize_identities(cls, value: str) -> str:
        return _strip_required_identity(value)

    @field_validator("artifact_storage_object_path")
    @classmethod
    def normalize_storage_path(cls, value: str) -> str:
        if "://" in value:
            raise ValueError("artifact storage reference must be an object path, not a URL")
        return normalize_object_path(value)


class SnapshotRunIdentity(_StrictFrozenModel):
    """Durable-row identity required before any manifest reference is read."""

    campaign_id: Literal["DQ-1"] = "DQ-1"
    quality_run_id: str = Field(pattern=r"^quality_[0-9a-f]{64}$")
    user_id: str = Field(min_length=1, max_length=512)
    thread_id: str = Field(min_length=1, max_length=512)
    task_id: str = Field(min_length=1, max_length=512)
    build_id: str = Field(pattern=r"^[A-Za-z0-9][A-Za-z0-9._-]{0,127}$")
    builder_run_id: str = Field(min_length=1, max_length=512)
    parent_builder_trace_id: str = Field(min_length=1, max_length=512)
    logical_artifact_id: str = Field(min_length=1, max_length=512)
    artifact_version_id: str = Field(min_length=1, max_length=512)
    manifest_revision: int = Field(ge=1)
    input_manifest_object_path: str = Field(min_length=1, max_length=4_096)
    input_manifest_hash: Sha256

    @field_validator(
        "user_id",
        "thread_id",
        "task_id",
        "builder_run_id",
        "parent_builder_trace_id",
        "logical_artifact_id",
        "artifact_version_id",
    )
    @classmethod
    def normalize_identities(cls, value: str) -> str:
        return _strip_required_identity(value)

    @field_validator("input_manifest_object_path")
    @classmethod
    def normalize_input_manifest_path(cls, value: str) -> str:
        return normalize_object_path(value)


class SnapshotArtifactReference(_StrictFrozenModel):
    virtual_path: str = Field(min_length=1, max_length=4_096)
    storage_object_path: str = Field(min_length=1, max_length=4_096)
    sha256: Sha256
    size_bytes: int = Field(gt=0)

    @field_validator("storage_object_path")
    @classmethod
    def normalize_storage_path(cls, value: str) -> str:
        return normalize_object_path(value)


class RenderSourcePdfReference(_StrictFrozenModel):
    object_path: str = Field(min_length=1, max_length=4_096)
    sha256: Sha256
    size_bytes: int = Field(gt=0)
    page_count: int = Field(ge=1, le=_MAX_RENDER_PAGES)
    media_type: Literal["application/pdf"] = "application/pdf"

    @field_validator("object_path")
    @classmethod
    def normalize_storage_path(cls, value: str) -> str:
        return normalize_object_path(value)


class RenderSourceManifest(_StrictFrozenModel):
    """Fixed create-only commit marker for the canonical PDF render source."""

    schema_version: Literal["deck-quality-render-source-manifest/v1"] = "deck-quality-render-source-manifest/v1"
    campaign_id: Literal["DQ-1"] = "DQ-1"
    quality_run_id: str = Field(pattern=r"^quality_[0-9a-f]{64}$")
    build_id: str = Field(pattern=r"^[A-Za-z0-9][A-Za-z0-9._-]{0,127}$")
    user_id: str = Field(min_length=1, max_length=512)
    thread_id: str = Field(min_length=1, max_length=512)
    task_id: str = Field(min_length=1, max_length=512)
    builder_run_id: str = Field(min_length=1, max_length=512)
    parent_builder_trace_id: str = Field(min_length=1, max_length=512)
    logical_artifact_id: str = Field(min_length=1, max_length=512)
    artifact_version_id: str = Field(min_length=1, max_length=512)
    artifact_manifest_revision: int = Field(ge=1)
    input_manifest_path: str = Field(min_length=1, max_length=4_096)
    input_manifest_hash: Sha256
    source_artifact: SnapshotArtifactReference
    pdf: RenderSourcePdfReference
    renderer_profile_version: Literal["libreoffice-impress-pdf-v1"] = _RENDER_SOURCE_PROFILE_VERSION
    renderer_profile_hash: Sha256 = _RENDER_SOURCE_PROFILE_HASH

    @field_validator("input_manifest_path")
    @classmethod
    def normalize_input_manifest_path(cls, value: str) -> str:
        return normalize_object_path(value)


class RenderSourceReference(_StrictFrozenModel):
    manifest_path: str = Field(min_length=1, max_length=4_096)
    manifest_hash: Sha256
    pdf: RenderSourcePdfReference
    renderer_profile_version: Literal["libreoffice-impress-pdf-v1"] = _RENDER_SOURCE_PROFILE_VERSION
    renderer_profile_hash: Sha256 = _RENDER_SOURCE_PROFILE_HASH

    @field_validator("manifest_path")
    @classmethod
    def normalize_manifest_path(cls, value: str) -> str:
        return normalize_object_path(value)


class LoadedRenderSource(_StrictFrozenModel):
    manifest: RenderSourceManifest
    reference: RenderSourceReference
    pdf_host_path: Path


class SnapshotObjectRecord(_StrictFrozenModel):
    role: Literal["render", "contact_sheet", "evidence_bundle"]
    object_path: str = Field(min_length=1, max_length=4_096)
    sha256: Sha256
    size_bytes: int = Field(gt=0)
    media_type: Literal["image/png", "application/json"]


class SnapshotSourceHashes(_StrictFrozenModel):
    input_manifest: Sha256
    artifact: Sha256
    render_source_manifest: Sha256
    render_source_pdf: Sha256
    brief: Sha256
    creative_plan: Sha256
    design_plan: Sha256
    build_record: Sha256
    mechanical_record: Sha256
    visible_text: Sha256


class SnapshotEvidenceBundle(_StrictFrozenModel):
    schema_version: Literal["deck-quality-evidence-bundle/v1"] = "deck-quality-evidence-bundle/v1"
    quality_run_id: str = Field(pattern=r"^quality_[0-9a-f]{64}$")
    thread_id: str = Field(min_length=1, max_length=512)
    artifact: SnapshotArtifactReference
    build_record: dict[str, Any]
    snapshot: QualityEvidenceSnapshot

    @model_validator(mode="after")
    def align_artifact_reference(self) -> SnapshotEvidenceBundle:
        if self.artifact.virtual_path != self.snapshot.artifact_path:
            raise ValueError("bundle artifact path does not match snapshot")
        if self.artifact.sha256 != self.snapshot.artifact_hash:
            raise ValueError("bundle artifact hash does not match snapshot")
        recorded_build_id = self.build_record.get("build_id")
        if recorded_build_id is not None and str(recorded_build_id) != self.snapshot.build_id:
            raise ValueError("bundle build record does not match snapshot")
        return self


class SnapshotEvidenceManifest(_StrictFrozenModel):
    schema_version: Literal["deck-quality-evidence-manifest/v2"] = "deck-quality-evidence-manifest/v2"
    campaign_id: Literal["DQ-1"] = "DQ-1"
    quality_run_id: str = Field(pattern=r"^quality_[0-9a-f]{64}$")
    snapshot_id: str = Field(pattern=r"^quality_[0-9a-f]{64}$")
    revision: Literal[2] = 2
    build_id: str
    user_id: str = Field(min_length=1, max_length=512)
    thread_id: str = Field(min_length=1, max_length=512)
    task_id: str = Field(min_length=1, max_length=512)
    builder_run_id: str = Field(min_length=1, max_length=512)
    parent_builder_trace_id: str = Field(min_length=1, max_length=512)
    logical_artifact_id: str = Field(min_length=1, max_length=512)
    artifact_version_id: str = Field(min_length=1, max_length=512)
    artifact_manifest_revision: int = Field(ge=1)
    input_manifest_path: str = Field(min_length=1, max_length=4_096)
    input_manifest_hash: Sha256
    artifact: SnapshotArtifactReference
    render_source: RenderSourceReference
    selectors: tuple[StableSlideSelector, ...]
    source_hashes: SnapshotSourceHashes
    render_hashes: dict[str, Sha256]
    objects: tuple[SnapshotObjectRecord, ...]
    evidence_bundle_path: str = Field(min_length=1, max_length=4_096)
    evidence_bundle_hash: Sha256

    @field_validator("input_manifest_path", "evidence_bundle_path")
    @classmethod
    def normalize_storage_path(cls, value: str) -> str:
        return normalize_object_path(value)

    @model_validator(mode="after")
    def validate_manifest_inventory(self) -> SnapshotEvidenceManifest:
        if not self.selectors or len(self.selectors) != len(set(self.selectors)):
            raise ValueError("manifest selectors must be nonempty and unique")
        expected_hash_keys = {str(selector) for selector in self.selectors} | {"contact-sheet"}
        if set(self.render_hashes) != expected_hash_keys:
            raise ValueError("manifest render hash coverage is incomplete")
        paths = tuple(record.object_path for record in self.objects)
        if len(paths) != len(set(paths)):
            raise ValueError("manifest object paths must be unique")
        if sum(record.role == "evidence_bundle" for record in self.objects) != 1:
            raise ValueError("manifest must contain exactly one evidence bundle object")
        if sum(record.role == "contact_sheet" for record in self.objects) != 1:
            raise ValueError("manifest must contain exactly one contact sheet object")
        if sum(record.role == "render" for record in self.objects) != len(self.selectors):
            raise ValueError("manifest render object coverage is incomplete")
        return self


class SnapshotCounts(_StrictFrozenModel):
    slide_count: int = Field(ge=1)
    visible_text_slide_count: int = Field(ge=1)
    native_input_count: Literal[3] = 3
    evidence_object_count: int = Field(ge=4)


class SnapshotDescriptor(_StrictFrozenModel):
    """Strict, content-free handoff returned to the durable dispatcher."""

    schema_version: Literal["deck-quality-snapshot-descriptor/v1"] = "deck-quality-snapshot-descriptor/v1"
    snapshot_id: str = Field(pattern=r"^quality_[0-9a-f]{64}$")
    snapshot_path: str = Field(min_length=1, max_length=4_096)
    snapshot_hash: Sha256
    revision: Literal[1] = 1
    counts: SnapshotCounts

    @field_validator("snapshot_path")
    @classmethod
    def normalize_snapshot_path(cls, value: str) -> str:
        return normalize_object_path(value)

    @model_validator(mode="after")
    def align_snapshot_path(self) -> SnapshotDescriptor:
        suffix = f"/quality/{self.snapshot_id}/evidence_manifest.json"
        if not self.snapshot_path.endswith(suffix):
            raise ValueError("snapshot descriptor path does not match snapshot identity")
        return self


class PreRenderInputObjectRecord(_StrictFrozenModel):
    role: PreRenderInputRole
    object_path: str = Field(min_length=1, max_length=4_096)
    sha256: Sha256
    size_bytes: int = Field(gt=0)
    media_type: Literal[
        "application/json",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ]

    @field_validator("object_path")
    @classmethod
    def normalize_storage_path(cls, value: str) -> str:
        return normalize_object_path(value)


ExactPreRenderInputObjects = Annotated[
    tuple[PreRenderInputObjectRecord, ...],
    Field(
        min_length=len(_PRE_RENDER_INPUT_ROLES),
        max_length=len(_PRE_RENDER_INPUT_ROLES),
    ),
]


class PreRenderInputBundleManifest(_StrictFrozenModel):
    schema_version: Literal["deck-quality-pre-render-input-manifest/v2"] = "deck-quality-pre-render-input-manifest/v2"
    campaign_id: Literal["DQ-1"] = "DQ-1"
    quality_run_id: str = Field(pattern=r"^quality_[0-9a-f]{64}$")
    revision: Literal[2] = 2
    build_id: str = Field(pattern=r"^[A-Za-z0-9][A-Za-z0-9._-]{0,127}$")
    user_id: str = Field(min_length=1, max_length=512)
    thread_id: str = Field(min_length=1, max_length=512)
    task_id: str = Field(min_length=1, max_length=512)
    builder_run_id: str = Field(min_length=1, max_length=512)
    parent_builder_trace_id: str = Field(min_length=1, max_length=512)
    logical_artifact_id: str = Field(min_length=1, max_length=512)
    artifact_version_id: str = Field(min_length=1, max_length=512)
    artifact_manifest_revision: int = Field(ge=1)
    artifact_virtual_path: str = Field(min_length=1, max_length=4_096)
    objects: ExactPreRenderInputObjects

    @model_validator(mode="after")
    def require_exact_ordered_inputs(self) -> PreRenderInputBundleManifest:
        roles = tuple(record.role for record in self.objects)
        if roles != _PRE_RENDER_INPUT_ROLES:
            raise ValueError("pre-render manifest requires all ordered input roles")
        paths = tuple(record.object_path for record in self.objects)
        if len(paths) != len(set(paths)):
            raise ValueError("pre-render manifest object paths must be unique")
        virtual = self.artifact_virtual_path.replace("\\", "/")
        pure = PurePosixPath(virtual)
        if not virtual.startswith("/mnt/user-data/outputs/") or ".." in pure.parts or pure.suffix.casefold() != ".pptx":
            raise ValueError("pre-render artifact virtual path is invalid")
        return self


class PreRenderInputBundleCounts(_StrictFrozenModel):
    content_object_count: Literal[6] = 6
    native_json_count: Literal[3] = 3
    control_record_count: Literal[2] = 2
    total_object_count: Literal[7] = 7


class PreRenderInputBundleDescriptor(_StrictFrozenModel):
    """Content-free durable handoff for restart-safe pre-render inputs."""

    schema_version: Literal["deck-quality-pre-render-input-descriptor/v2"] = "deck-quality-pre-render-input-descriptor/v2"
    bundle_id: str = Field(pattern=r"^quality_[0-9a-f]{64}$")
    manifest_path: str = Field(min_length=1, max_length=4_096)
    manifest_hash: Sha256
    revision: Literal[2] = 2
    counts: PreRenderInputBundleCounts

    @field_validator("manifest_path")
    @classmethod
    def normalize_manifest_path(cls, value: str) -> str:
        return normalize_object_path(value)

    @model_validator(mode="after")
    def align_manifest_path(self) -> PreRenderInputBundleDescriptor:
        suffix = f"/quality/{self.bundle_id}/input_bundle/manifest.json"
        if not self.manifest_path.endswith(suffix):
            raise ValueError("pre-render descriptor path does not match bundle identity")
        return self


class LoadedPreRenderInputBundle(_StrictFrozenModel):
    """Verified inputs materialized into an isolated restart workspace."""

    descriptor: PreRenderInputBundleDescriptor
    manifest: PreRenderInputBundleManifest
    metadata: SnapshotCompletionMetadata
    brief: BlindBrief
    mechanical_record: dict[str, Any]
    outputs_root: Path
    artifact_virtual_path: str
    artifact_host_path: Path


class LoadedEvidenceSnapshot(_StrictFrozenModel):
    descriptor: SnapshotDescriptor
    manifest: SnapshotEvidenceManifest
    build_record: dict[str, Any]
    snapshot: QualityEvidenceSnapshot


class ImmutableObjectReader(Protocol):
    def read(self, object_path: str) -> bytes | None: ...


class ImmutableObjectUploader(ImmutableObjectReader, Protocol):
    """Create-only object protocol; implementations must use upsert=false."""

    def create_if_absent(
        self,
        object_path: str,
        content: bytes,
        *,
        content_type: str,
    ) -> Literal["created", "exists"]: ...


PreviewResolver = Callable[[Path], Path | None]
PdfRasterizer = Callable[[Path], tuple[bytes, ...]]


def _sha256_bytes(content: bytes) -> str:
    return hashlib.sha256(content).hexdigest()


def _strict_json_object(raw: bytes, *, role: str) -> dict[str, Any]:
    def reject_constant(_value: str) -> None:
        raise ValueError

    try:
        value = json.loads(raw.decode("utf-8"), parse_constant=reject_constant)
    except (UnicodeDecodeError, ValueError, json.JSONDecodeError) as exc:
        raise SnapshotMissingEvidenceError(f"{role} is not valid JSON") from exc
    if not isinstance(value, dict):
        raise SnapshotMissingEvidenceError(f"{role} must contain a JSON object")
    return value


def _canonical_mapping(value: Mapping[str, Any], *, role: str) -> tuple[dict[str, Any], bytes]:
    try:
        encoded = canonical_json_bytes(dict(value))
        normalized = _strict_json_object(encoded, role=role)
    except SnapshotError:
        raise
    except (TypeError, ValueError) as exc:
        raise SnapshotMissingEvidenceError(f"{role} is not canonical JSON") from exc
    return normalized, encoded


def _canonical_blind_brief(
    task_brief: BlindBrief | Mapping[str, Any],
) -> tuple[BlindBrief, bytes]:
    try:
        raw_brief = task_brief.model_dump(mode="python") if isinstance(task_brief, BlindBrief) else dict(task_brief)
        brief = BlindBrief.model_validate(raw_brief)
        brief = brief.model_copy(update={"request": sanitize_current_request(brief.request)})
        encoded = canonical_json_bytes(brief)
    except (TypeError, ValueError) as exc:
        raise SnapshotMissingEvidenceError("task brief is incomplete or invalid") from exc
    return brief, encoded


def _read_required_file(
    path: Path,
    *,
    role: str,
    scope_root: Path | None = None,
    max_bytes: int | None = None,
) -> bytes:
    descriptor: int | None = None
    try:
        resolved = path.resolve(strict=True)
        if scope_root is not None:
            root = scope_root.resolve(strict=True)
            if root != resolved and root not in resolved.parents:
                raise SnapshotMissingEvidenceError(f"required {role} evidence is outside the declared outputs root")
        source_lstat = path.lstat()
        if (
            not stat.S_ISREG(source_lstat.st_mode)
            or source_lstat.st_nlink != 1
        ):
            raise SnapshotMissingEvidenceError(
                f"required {role} evidence is not a single-link regular file"
            )
        flags = os.O_RDONLY | getattr(os, "O_CLOEXEC", 0)
        flags |= getattr(os, "O_NOFOLLOW", 0)
        descriptor = os.open(path, flags)
        before = os.fstat(descriptor)
        if (
            not stat.S_ISREG(before.st_mode)
            or before.st_nlink != 1
            or (before.st_dev, before.st_ino)
            != (source_lstat.st_dev, source_lstat.st_ino)
        ):
            raise SnapshotStaleError(
                f"required {role} evidence changed before it was read"
            )
        size = before.st_size
        if size <= 0:
            raise SnapshotMissingEvidenceError(f"required {role} evidence is empty")
        if max_bytes is not None and size > max_bytes:
            raise SnapshotCoverageError(f"required {role} evidence exceeds its byte budget")
        with os.fdopen(descriptor, "rb", closefd=False) as stream:
            content = stream.read()
        after = os.fstat(descriptor)
    except OSError as exc:
        raise SnapshotMissingEvidenceError(f"required {role} evidence is missing") from exc
    finally:
        if descriptor is not None:
            os.close(descriptor)
    stable_fields = ("st_dev", "st_ino", "st_size", "st_mtime_ns", "st_ctime_ns")
    if (
        len(content) != size
        or any(
            getattr(before, field) != getattr(after, field)
            for field in stable_fields
        )
    ):
        raise SnapshotStaleError(f"required {role} evidence changed while it was read")
    return content


def _validate_artifact_paths(
    *,
    outputs_root: Path,
    artifact_virtual_path: str,
    artifact_host_path: Path,
) -> tuple[Path, str]:
    try:
        root = outputs_root.resolve(strict=True)
        host = artifact_host_path.resolve(strict=True)
    except OSError as exc:
        raise SnapshotMissingEvidenceError("accepted artifact path is missing") from exc
    if artifact_host_path.absolute() != host:
        raise SnapshotMissingEvidenceError("artifact host path cannot use a symlink")
    if root != host and root not in host.parents:
        raise SnapshotMissingEvidenceError("artifact host path is outside the declared outputs root")
    virtual = artifact_virtual_path.strip().replace("\\", "/")
    virtual_path = PurePosixPath(virtual)
    if not virtual.startswith("/mnt/user-data/outputs/") or ".." in virtual_path.parts or virtual_path.suffix.casefold() != ".pptx":
        raise SnapshotMissingEvidenceError("artifact virtual path is not an accepted PPTX output")
    if host.suffix.casefold() != ".pptx":
        raise SnapshotMissingEvidenceError("artifact host path is not a PPTX")
    return host, virtual


def _numbered_png_paths(root: Path) -> tuple[Path, ...]:
    candidates: list[tuple[int, Path]] = []
    for path in root.glob("page-*.png"):
        match = re.search(r"-([0-9]+)\.png$", path.name)
        if match:
            candidates.append((int(match.group(1)), path))
    candidates.sort(key=lambda item: item[0])
    page_numbers = tuple(item[0] for item in candidates)
    if page_numbers and page_numbers != tuple(range(1, len(page_numbers) + 1)):
        raise SnapshotCoverageError("PDF raster output pages are not contiguous")
    return tuple(item[1] for item in candidates)


def rasterize_preview_pdf(pdf_path: Path) -> tuple[bytes, ...]:
    """Render every PDF page as PNG with a fixed maximum dimension.

    A DPI lock is not a cost lock because PPTX page sizes can vary in physical
    inches. Bounding the long side makes the native ``detail=original`` token
    envelope independent of the authored page size.
    """

    executable = shutil.which("pdftoppm")
    if executable is None:
        raise SnapshotMissingEvidenceError("pdftoppm is unavailable for lossless render evidence")
    binary = str(Path(executable).resolve())
    source = pdf_path.absolute()
    source_bytes = _read_required_file(
        pdf_path,
        role="preview PDF",
        scope_root=pdf_path.parent,
        max_bytes=_MAX_ACCEPTED_PREVIEW_PDF_BYTES,
    )
    try:
        document = PdfReader(io.BytesIO(source_bytes), strict=True)
        if document.is_encrypted or not 1 <= len(document.pages) <= _MAX_RENDER_PAGES:
            raise SnapshotCoverageError("preview PDF page coverage exceeds the direct render budget")
        total_pixels = 0
        for page in document.pages:
            width_points = float(page.mediabox.width)
            height_points = float(page.mediabox.height)
            if width_points <= 0 or height_points <= 0:
                raise SnapshotCoverageError("preview PDF page exceeds the direct render pixel budget")
            scale = _PDFTOPPM_MAX_DIMENSION / max(width_points, height_points)
            scaled_width = max(1, math.ceil(width_points * scale))
            scaled_height = max(1, math.ceil(height_points * scale))
            pixels = scaled_width * scaled_height
            if pixels > _MAX_RENDER_PAGE_PIXELS:
                raise SnapshotCoverageError("preview PDF page exceeds the direct render pixel budget")
            total_pixels += pixels
        if total_pixels > _MAX_RENDER_TOTAL_PIXELS:
            raise SnapshotCoverageError("preview PDF exceeds the aggregate render pixel budget")
    except SnapshotError:
        raise
    except Exception as exc:
        raise SnapshotCoverageError("preview PDF structure is invalid") from exc
    try:
        with tempfile.TemporaryDirectory(prefix="dq1-pdf-raster-") as directory:
            staged_source = Path(directory) / "source.pdf"
            staged_source.write_bytes(source_bytes)
            prefix = Path(directory) / "page"
            completed = run_process_group(
                [
                    binary,
                    "-png",
                    "-scale-to",
                    str(_PDFTOPPM_MAX_DIMENSION),
                    str(staged_source),
                    str(prefix),
                ],
                timeout=_PDFTOPPM_TIMEOUT_SECONDS,
                private_read_dirs=[staged_source],
                writable_dirs=[directory],
                identity_paths=[source],
            )
            if completed.returncode != 0:
                raise SnapshotCoverageError("preview PDF rasterization failed")
            paths = _numbered_png_paths(Path(directory))
            if not paths:
                raise SnapshotCoverageError("preview PDF rasterization produced no pages")
            return tuple(path.read_bytes() for path in paths)
    except subprocess.TimeoutExpired as exc:
        raise SnapshotCoverageError("preview PDF rasterization timed out") from exc
    except OSError as exc:
        raise SnapshotCoverageError("preview PDF rasterization could not complete") from exc


def _shape_text_fragments(shape: Any) -> tuple[str, ...]:
    children = getattr(shape, "shapes", None)
    if children is not None:
        return tuple(fragment for child in children for fragment in _shape_text_fragments(child))
    if bool(getattr(shape, "has_table", False)):
        return tuple(paragraph.text for row in shape.table.rows for cell in row.cells for paragraph in cell.text_frame.paragraphs if paragraph.text.strip())
    if bool(getattr(shape, "has_text_frame", False)):
        return tuple(paragraph.text for paragraph in shape.text_frame.paragraphs if paragraph.text.strip())
    return ()


def _pptx_visible_text(artifact_bytes: bytes) -> tuple[int, tuple[VisibleTextSlide, ...]]:
    if not 0 < len(artifact_bytes) <= _MAX_ACCEPTED_PPTX_BYTES:
        raise SnapshotCoverageError("accepted PPTX exceeds its byte budget")
    try:
        with zipfile.ZipFile(io.BytesIO(artifact_bytes)) as archive:
            members = archive.infolist()
            if (
                not members
                or len(members) > _MAX_PPTX_MEMBER_COUNT
                or any(member.file_size > _MAX_PPTX_MEMBER_BYTES or member.flag_bits & 0x1 for member in members)
                or sum(member.file_size for member in members) > _MAX_PPTX_UNCOMPRESSED_BYTES
            ):
                raise SnapshotCoverageError("accepted PPTX package exceeds its expansion budget")
        presentation = Presentation(io.BytesIO(artifact_bytes))
    except SnapshotError:
        raise
    except Exception as exc:  # noqa: BLE001 - parser errors are reduced to a safe code
        raise SnapshotMissingEvidenceError("accepted PPTX cannot be decoded") from exc
    slide_count = len(presentation.slides)
    if slide_count < 1:
        raise SnapshotCoverageError("accepted PPTX contains no slides")
    slides = tuple(
        (
            f"slide:{index}",
            tuple(fragment for shape in slide.shapes for fragment in _shape_text_fragments(shape)),
        )
        for index, slide in enumerate(presentation.slides, start=1)
    )
    return slide_count, visible_text_sidecar(slides)


def _decode_png(content: bytes, *, role: str) -> tuple[int, int]:
    try:
        with Image.open(io.BytesIO(content)) as image:
            if image.format != "PNG":
                raise ValueError
            image.verify()
        with Image.open(io.BytesIO(content)) as image:
            width, height = image.size
    except Exception as exc:  # noqa: BLE001 - image parser errors are content-free externally
        raise SnapshotCoverageError(f"{role} is not a decodable PNG") from exc
    if width < 1 or height < 1:
        raise SnapshotCoverageError(f"{role} has invalid dimensions")
    return width, height


def _render_evidence(
    *,
    png_pages: tuple[bytes, ...],
    expected_slide_count: int,
    object_root: str,
    work_root: Path,
) -> tuple[RenderEvidence, bytes, tuple[tuple[str, bytes], ...]]:
    if len(png_pages) != expected_slide_count:
        raise SnapshotCoverageError("render page count does not match accepted PPTX slide count")
    selectors = tuple(f"slide:{index}" for index in range(1, expected_slide_count + 1))
    slide_records: list[ImageEvidence] = []
    local_paths: list[Path] = []
    objects: list[tuple[str, bytes]] = []
    for index, (selector, content) in enumerate(zip(selectors, png_pages, strict=True), start=1):
        width, height = _decode_png(content, role="individual render")
        local_path = work_root / f"slide-{index:04d}.png"
        local_path.write_bytes(content)
        local_paths.append(local_path)
        object_path = f"{object_root}/renders/slide-{index:04d}.png"
        objects.append((object_path, content))
        slide_records.append(
            ImageEvidence(
                selector=selector,  # type: ignore[arg-type]
                path=object_path,
                sha256=_sha256_bytes(content),
                width=width,
                height=height,
                decodes=True,
            )
        )
    contact_path = work_root / "contact-sheet.png"
    create_contact_sheet(tuple(local_paths), contact_path)
    contact_bytes = contact_path.read_bytes()
    contact_width, contact_height = _decode_png(contact_bytes, role="contact sheet")
    contact_object_path = f"{object_root}/renders/contact-sheet.png"
    renders = RenderEvidence(
        expected_slide_count=expected_slide_count,
        contact_sheet=ImageEvidence(
            selector="contact-sheet",
            path=contact_object_path,
            sha256=_sha256_bytes(contact_bytes),
            width=contact_width,
            height=contact_height,
            decodes=True,
        ),
        slides=tuple(slide_records),
        selectors=selectors,
    )
    return renders, contact_bytes, tuple(objects)


def _build_object_root(metadata: SnapshotCompletionMetadata) -> str:
    build_id = metadata.build_id
    if not _SAFE_BUILD_ID_RE.fullmatch(build_id):
        raise SnapshotMissingEvidenceError("build identity cannot address durable snapshot storage")
    return _build_object_root_for_identity(
        user_id=metadata.user_id,
        thread_id=metadata.thread_id,
        build_id=build_id,
    )


def _build_object_root_for_identity(
    *,
    user_id: str,
    thread_id: str,
    build_id: str,
) -> str:
    return normalize_object_path(f"artifacts/{safe_object_path_segment(user_id, default='user')}/{safe_object_path_segment(thread_id, default='thread')}/foundation/.builder/builds/{build_id}")


def _quality_object_root(metadata: SnapshotCompletionMetadata) -> str:
    return normalize_object_path(f"{_build_object_root(metadata)}/quality/{metadata.quality_run_id}")


def _quality_object_root_for_identity(
    *,
    user_id: str,
    thread_id: str,
    build_id: str,
    quality_run_id: str,
) -> str:
    return normalize_object_path(f"{_build_object_root_for_identity(user_id=user_id, thread_id=thread_id, build_id=build_id)}/quality/{quality_run_id}")


def _pre_render_input_root(metadata: SnapshotCompletionMetadata) -> str:
    return normalize_object_path(f"{_quality_object_root(metadata)}/input_bundle")


def _pre_render_json_object_path(
    *,
    input_root: str,
    role: PreRenderInputRole,
    content_hash: str,
) -> str:
    return normalize_object_path(f"{input_root}/objects/{role}/{content_hash}.json")


def _render_source_root(metadata: SnapshotCompletionMetadata) -> str:
    return normalize_object_path(f"{_quality_object_root(metadata)}/render_source")


def _render_source_manifest_path(metadata: SnapshotCompletionMetadata) -> str:
    return normalize_object_path(f"{_render_source_root(metadata)}/manifest.json")


def _render_source_pdf_path(
    metadata: SnapshotCompletionMetadata,
    *,
    content_hash: str,
) -> str:
    return normalize_object_path(
        f"{_render_source_root(metadata)}/objects/{content_hash}.pdf"
    )


def _immutable_artifact_object_path(
    *,
    user_id: str,
    thread_id: str,
    build_id: str,
    logical_artifact_id: str,
    artifact_version_id: str,
    artifact_hash: str,
    artifact_virtual_path: str,
) -> str:
    if not build_id.strip():
        raise SnapshotMissingEvidenceError("build identity is invalid")
    return immutable_builder_artifact_object_path(
        user_id=user_id,
        thread_or_session_id=thread_id,
        logical_artifact_id=logical_artifact_id,
        artifact_version_id=artifact_version_id,
        artifact_sha256=artifact_hash,
        filename=PurePosixPath(
            artifact_virtual_path.replace("\\", "/")
        ).name,
    )


def _read_native_inputs(
    outputs_root: Path,
) -> tuple[
    dict[str, Any],
    dict[str, Any],
    dict[str, Any],
    dict[str, tuple[Path, str]],
]:
    deck_build = outputs_root / "deck_build"
    loaded: dict[str, dict[str, Any]] = {}
    source_files: dict[str, tuple[Path, str]] = {}
    for role, filename in (
        ("creative_plan", "creative_plan.json"),
        ("design_plan", "design_plan.json"),
        ("build_record", "build.json"),
    ):
        path = deck_build / filename
        raw = _read_required_file(
            path,
            role=role,
            scope_root=outputs_root,
            max_bytes=_MAX_NATIVE_JSON_BYTES,
        )
        loaded[role] = _strict_json_object(raw, role=role)
        source_files[role] = (path, _sha256_bytes(raw))
    return (
        loaded["creative_plan"],
        loaded["design_plan"],
        loaded["build_record"],
        source_files,
    )


def _validate_artifact_storage_scope(metadata: SnapshotCompletionMetadata) -> None:
    expected_prefix = normalize_object_path(f"artifacts/{safe_object_path_segment(metadata.user_id, default='user')}/{safe_object_path_segment(metadata.thread_id, default='thread')}")
    if not metadata.artifact_storage_object_path.startswith(f"{expected_prefix}/"):
        raise SnapshotStaleError("durable artifact reference is outside the completed build scope")


def _validate_build_record(
    build_record: Mapping[str, Any],
    *,
    metadata: SnapshotCompletionMetadata,
    expected_slide_count: int,
) -> None:
    recorded_build_id = build_record.get("build_id")
    if recorded_build_id is not None and str(recorded_build_id) != metadata.build_id:
        raise SnapshotStaleError("native build record does not match the requested build")
    slides = build_record.get("slides")
    if isinstance(slides, list) and len(slides) != expected_slide_count:
        raise SnapshotCoverageError("native build record slide count does not match accepted PPTX")


def _verify_sources_unchanged(sources: Mapping[str, tuple[Path, str]]) -> None:
    for role, (path, expected_hash) in sources.items():
        try:
            actual_hash = _sha256_bytes(path.read_bytes())
        except OSError as exc:
            raise SnapshotStaleError(f"{role} changed while the evidence snapshot was prepared") from exc
        if actual_hash != expected_hash:
            raise SnapshotStaleError(f"{role} changed while the evidence snapshot was prepared")


def _upload_immutable(
    *,
    uploader: ImmutableObjectUploader,
    object_path: str,
    content: bytes,
    content_type: str,
    role: str,
    verify_after_create: bool = False,
) -> None:
    normalized = normalize_object_path(object_path)
    try:
        outcome = uploader.create_if_absent(
            normalized,
            content,
            content_type=content_type,
        )
    except Exception as exc:  # noqa: BLE001 - storage details must not escape or be logged
        raise SnapshotUploadError(f"immutable upload failed for {role}") from exc
    if outcome == "created" and not verify_after_create:
        return
    if outcome not in {"created", "exists"}:
        raise SnapshotUploadError(f"immutable uploader returned an invalid result for {role}")
    try:
        read_bounded = getattr(uploader, "read_bounded", None)
        if callable(read_bounded):
            existing = read_bounded(normalized, max_bytes=len(content))
        else:
            existing = uploader.read(normalized)
    except Exception as exc:  # noqa: BLE001 - storage details must not escape or be logged
        raise SnapshotUploadError(f"immutable replay verification failed for {role}") from exc
    if (
        existing is None
        or len(existing) != len(content)
        or not hmac.compare_digest(
            _sha256_bytes(existing),
            _sha256_bytes(content),
        )
    ):
        raise SnapshotConflictError(f"immutable evidence differs for {role}")


def freeze_and_upload_pre_render_input_bundle(
    *,
    metadata: SnapshotCompletionMetadata,
    outputs_root: Path,
    artifact_virtual_path: str,
    artifact_host_path: Path,
    task_brief: BlindBrief | Mapping[str, Any],
    authoritative_mechanical: Mapping[str, Any],
    uploader: ImmutableObjectUploader,
) -> PreRenderInputBundleDescriptor:
    """Commit every render input to immutable storage before rendering.

    The returned descriptor contains only correlation, path, hash, revision,
    and counts. The canonical manifest is uploaded after all six content
    objects and acts as the sole commit marker.
    """

    if not _QUALITY_RUN_RE.fullmatch(metadata.quality_run_id):
        raise SnapshotMissingEvidenceError("quality run identity is invalid")
    root = outputs_root.resolve()
    artifact_path, virtual_path = _validate_artifact_paths(
        outputs_root=root,
        artifact_virtual_path=artifact_virtual_path,
        artifact_host_path=artifact_host_path,
    )
    artifact_bytes = _read_required_file(
        artifact_path,
        role="accepted artifact",
        scope_root=root,
        max_bytes=_MAX_ACCEPTED_PPTX_BYTES,
    )
    artifact_hash = _sha256_bytes(artifact_bytes)
    creative_plan, design_plan, build_record, native_sources = _read_native_inputs(root)
    brief, brief_bytes = _canonical_blind_brief(task_brief)
    mechanical_record, mechanical_bytes = _canonical_mapping(
        authoritative_mechanical,
        role="authoritative mechanical record",
    )
    slide_count, _visible_text = _pptx_visible_text(artifact_bytes)
    _validate_build_record(
        build_record,
        metadata=metadata,
        expected_slide_count=slide_count,
    )
    _validate_artifact_storage_scope(metadata)

    creative_bytes = canonical_json_bytes(creative_plan)
    design_bytes = canonical_json_bytes(design_plan)
    build_bytes = canonical_json_bytes(build_record)
    input_root = _pre_render_input_root(metadata)
    immutable_artifact_path = _immutable_artifact_object_path(
        user_id=metadata.user_id,
        thread_id=metadata.thread_id,
        build_id=metadata.build_id,
        logical_artifact_id=metadata.logical_artifact_id,
        artifact_version_id=metadata.artifact_version_id,
        artifact_hash=artifact_hash,
        artifact_virtual_path=virtual_path,
    )
    content_by_role: tuple[
        tuple[PreRenderInputRole, str, bytes, str],
        ...,
    ] = (
        (
            "accepted_artifact",
            immutable_artifact_path,
            artifact_bytes,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ),
        (
            "creative_plan",
            _pre_render_json_object_path(
                input_root=input_root,
                role="creative_plan",
                content_hash=_sha256_bytes(creative_bytes),
            ),
            creative_bytes,
            "application/json",
        ),
        (
            "design_plan",
            _pre_render_json_object_path(
                input_root=input_root,
                role="design_plan",
                content_hash=_sha256_bytes(design_bytes),
            ),
            design_bytes,
            "application/json",
        ),
        (
            "build_record",
            _pre_render_json_object_path(
                input_root=input_root,
                role="build_record",
                content_hash=_sha256_bytes(build_bytes),
            ),
            build_bytes,
            "application/json",
        ),
        (
            "blind_brief",
            _pre_render_json_object_path(
                input_root=input_root,
                role="blind_brief",
                content_hash=_sha256_bytes(brief_bytes),
            ),
            brief_bytes,
            "application/json",
        ),
        (
            "mechanical_record",
            _pre_render_json_object_path(
                input_root=input_root,
                role="mechanical_record",
                content_hash=_sha256_bytes(mechanical_bytes),
            ),
            mechanical_bytes,
            "application/json",
        ),
    )
    records = tuple(
        PreRenderInputObjectRecord(
            role=role,
            object_path=object_path,
            sha256=_sha256_bytes(content),
            size_bytes=len(content),
            media_type=media_type,
        )
        for role, object_path, content, media_type in content_by_role
    )
    manifest = PreRenderInputBundleManifest(
        quality_run_id=metadata.quality_run_id,
        build_id=metadata.build_id,
        user_id=metadata.user_id,
        thread_id=metadata.thread_id,
        task_id=metadata.task_id,
        builder_run_id=metadata.builder_run_id,
        parent_builder_trace_id=metadata.parent_builder_trace_id,
        logical_artifact_id=metadata.logical_artifact_id,
        artifact_version_id=metadata.artifact_version_id,
        artifact_manifest_revision=metadata.manifest_revision,
        artifact_virtual_path=virtual_path,
        objects=records,
    )
    manifest_bytes = canonical_json_bytes(manifest)
    manifest_path = normalize_object_path(f"{input_root}/manifest.json")

    source_files = {
        "accepted artifact": (artifact_path, artifact_hash),
        **native_sources,
    }
    _verify_sources_unchanged(source_files)
    for role, object_path, content, media_type in content_by_role:
        _upload_immutable(
            uploader=uploader,
            object_path=object_path,
            content=content,
            content_type=media_type,
            role=role,
            verify_after_create=True,
        )
    # This canonical manifest is the commit marker and must remain last.
    _upload_immutable(
        uploader=uploader,
        object_path=manifest_path,
        content=manifest_bytes,
        content_type="application/json",
        role="pre_render_input_manifest",
        verify_after_create=True,
    )
    return PreRenderInputBundleDescriptor(
        bundle_id=metadata.quality_run_id,
        manifest_path=manifest_path,
        manifest_hash=_sha256_bytes(manifest_bytes),
        counts=PreRenderInputBundleCounts(),
    )


def _run_identity_matches_manifest(
    manifest: PreRenderInputBundleManifest | SnapshotEvidenceManifest,
    expected: SnapshotRunIdentity,
) -> bool:
    return all(
        getattr(manifest, field) == getattr(expected, expected_field)
        for field, expected_field in (
            ("campaign_id", "campaign_id"),
            ("quality_run_id", "quality_run_id"),
            ("user_id", "user_id"),
            ("thread_id", "thread_id"),
            ("task_id", "task_id"),
            ("build_id", "build_id"),
            ("builder_run_id", "builder_run_id"),
            ("parent_builder_trace_id", "parent_builder_trace_id"),
            ("logical_artifact_id", "logical_artifact_id"),
            ("artifact_version_id", "artifact_version_id"),
            ("artifact_manifest_revision", "manifest_revision"),
        )
    )


def verify_evidence_manifest_identity(
    manifest: SnapshotEvidenceManifest,
    expected: SnapshotRunIdentity,
) -> None:
    """Reject a row/manifest mismatch before any manifest reference is read."""

    if (
        not _run_identity_matches_manifest(manifest, expected)
        or manifest.input_manifest_path != expected.input_manifest_object_path
        or manifest.input_manifest_hash != expected.input_manifest_hash
    ):
        raise SnapshotStaleError("evidence manifest does not match the durable row")


def _source_artifact_reference(
    loaded: LoadedPreRenderInputBundle,
) -> SnapshotArtifactReference:
    accepted = next(
        record
        for record in loaded.manifest.objects
        if record.role == "accepted_artifact"
    )
    return SnapshotArtifactReference(
        virtual_path=loaded.artifact_virtual_path,
        storage_object_path=accepted.object_path,
        sha256=accepted.sha256,
        size_bytes=accepted.size_bytes,
    )


def _validate_render_source_pdf(
    content: bytes,
    *,
    expected_page_count: int,
) -> int:
    if len(content) > _MAX_ACCEPTED_PREVIEW_PDF_BYTES or not content.startswith(b"%PDF-"):
        raise SnapshotCoverageError("render-source PDF is invalid or oversized")
    try:
        reader = PdfReader(io.BytesIO(content))
        if reader.is_encrypted:
            raise ValueError
        page_count = len(reader.pages)
    except Exception as exc:
        raise SnapshotCoverageError("render-source PDF cannot be decoded") from exc
    if page_count != expected_page_count or not 1 <= page_count <= _MAX_RENDER_PAGES:
        raise SnapshotCoverageError("render-source PDF page coverage is invalid")
    return page_count


def _verify_render_source_manifest(
    loaded: LoadedPreRenderInputBundle,
    manifest: RenderSourceManifest,
) -> None:
    metadata = loaded.metadata
    expected_artifact = _source_artifact_reference(loaded)
    if any(
        actual != expected
        for actual, expected in (
            (manifest.quality_run_id, metadata.quality_run_id),
            (manifest.build_id, metadata.build_id),
            (manifest.user_id, metadata.user_id),
            (manifest.thread_id, metadata.thread_id),
            (manifest.task_id, metadata.task_id),
            (manifest.builder_run_id, metadata.builder_run_id),
            (manifest.parent_builder_trace_id, metadata.parent_builder_trace_id),
            (manifest.logical_artifact_id, metadata.logical_artifact_id),
            (manifest.artifact_version_id, metadata.artifact_version_id),
            (manifest.artifact_manifest_revision, metadata.manifest_revision),
            (manifest.input_manifest_path, loaded.descriptor.manifest_path),
            (manifest.input_manifest_hash, loaded.descriptor.manifest_hash),
            (manifest.source_artifact, expected_artifact),
            (manifest.renderer_profile_hash, _RENDER_SOURCE_PROFILE_HASH),
        )
    ):
        raise SnapshotStaleError("render-source manifest does not match immutable inputs")
    if manifest.pdf.object_path != _render_source_pdf_path(
        metadata,
        content_hash=manifest.pdf.sha256,
    ):
        raise SnapshotConflictError("render-source PDF path is not content addressed")


def _read_optional_snapshot_object(
    reader: ImmutableObjectReader,
    *,
    object_path: str,
    role: str,
    max_bytes: int,
) -> bytes | None:
    try:
        normalized = normalize_object_path(object_path)
        read_bounded = getattr(reader, "read_bounded", None)
        if callable(read_bounded):
            content = read_bounded(normalized, max_bytes=max_bytes)
        else:
            content = reader.read(normalized)
    except Exception as exc:  # noqa: BLE001 - storage details must remain content-free
        raise SnapshotUploadError(f"immutable read failed for {role}") from exc
    if content is not None and len(content) > max_bytes:
        raise SnapshotCoverageError(f"immutable {role} exceeds its byte budget")
    return content


def _load_committed_render_source(
    *,
    loaded_input: LoadedPreRenderInputBundle,
    reader: ImmutableObjectReader,
    manifest_bytes: bytes,
) -> LoadedRenderSource:
    manifest = _canonical_model_from_bytes(
        manifest_bytes,
        model_type=RenderSourceManifest,
        role="render_source_manifest",
    )
    _verify_render_source_manifest(loaded_input, manifest)
    pdf_bytes = _read_snapshot_object(
        reader,
        object_path=manifest.pdf.object_path,
        role="render_source_pdf",
        max_bytes=_MAX_ACCEPTED_PREVIEW_PDF_BYTES,
    )
    if (
        len(pdf_bytes) != manifest.pdf.size_bytes
        or not hmac.compare_digest(_sha256_bytes(pdf_bytes), manifest.pdf.sha256)
    ):
        raise SnapshotConflictError("render-source PDF does not match its manifest")
    expected_pages, _visible_text = _pptx_visible_text(
        loaded_input.artifact_host_path.read_bytes()
    )
    if _validate_render_source_pdf(
        pdf_bytes,
        expected_page_count=expected_pages,
    ) != manifest.pdf.page_count:
        raise SnapshotConflictError("render-source PDF page count changed")
    local_path = loaded_input.outputs_root / ".dq1" / "render-source.pdf"
    _write_materialized_input(local_path, pdf_bytes, role="render_source_pdf")
    manifest_path = _render_source_manifest_path(loaded_input.metadata)
    return LoadedRenderSource(
        manifest=manifest,
        reference=RenderSourceReference(
            manifest_path=manifest_path,
            manifest_hash=_sha256_bytes(manifest_bytes),
            pdf=manifest.pdf,
        ),
        pdf_host_path=local_path,
    )


def ensure_committed_render_source(
    *,
    loaded_input: LoadedPreRenderInputBundle,
    uploader: ImmutableObjectUploader,
    renderer: PreviewResolver = maybe_render_pptx_preview,
) -> LoadedRenderSource:
    """Commit one canonical post-row PDF source; first valid manifest wins."""

    manifest_path = _render_source_manifest_path(loaded_input.metadata)
    existing = _read_optional_snapshot_object(
        uploader,
        object_path=manifest_path,
        role="render_source_manifest",
        max_bytes=_MAX_NATIVE_JSON_BYTES,
    )
    if existing is not None:
        return _load_committed_render_source(
            loaded_input=loaded_input,
            reader=uploader,
            manifest_bytes=existing,
        )

    try:
        rendered_path = renderer(loaded_input.artifact_host_path)
    except Exception as exc:  # noqa: BLE001 - renderer details must remain private
        raise SnapshotCoverageError("accepted PPTX could not produce a render source") from exc
    if rendered_path is None:
        raise SnapshotMissingEvidenceError("accepted PPTX render source is unavailable")
    pdf_bytes = _read_required_file(
        rendered_path,
        role="render-source PDF",
        scope_root=loaded_input.outputs_root,
        max_bytes=_MAX_ACCEPTED_PREVIEW_PDF_BYTES,
    )
    expected_pages, _visible_text = _pptx_visible_text(
        loaded_input.artifact_host_path.read_bytes()
    )
    page_count = _validate_render_source_pdf(
        pdf_bytes,
        expected_page_count=expected_pages,
    )
    pdf_hash = _sha256_bytes(pdf_bytes)
    pdf_path = _render_source_pdf_path(
        loaded_input.metadata,
        content_hash=pdf_hash,
    )
    _upload_immutable(
        uploader=uploader,
        object_path=pdf_path,
        content=pdf_bytes,
        content_type="application/pdf",
        role="render_source_pdf",
        verify_after_create=True,
    )
    manifest = RenderSourceManifest(
        quality_run_id=loaded_input.metadata.quality_run_id,
        build_id=loaded_input.metadata.build_id,
        user_id=loaded_input.metadata.user_id,
        thread_id=loaded_input.metadata.thread_id,
        task_id=loaded_input.metadata.task_id,
        builder_run_id=loaded_input.metadata.builder_run_id,
        parent_builder_trace_id=loaded_input.metadata.parent_builder_trace_id,
        logical_artifact_id=loaded_input.metadata.logical_artifact_id,
        artifact_version_id=loaded_input.metadata.artifact_version_id,
        artifact_manifest_revision=loaded_input.metadata.manifest_revision,
        input_manifest_path=loaded_input.descriptor.manifest_path,
        input_manifest_hash=loaded_input.descriptor.manifest_hash,
        source_artifact=_source_artifact_reference(loaded_input),
        pdf=RenderSourcePdfReference(
            object_path=pdf_path,
            sha256=pdf_hash,
            size_bytes=len(pdf_bytes),
            page_count=page_count,
        ),
    )
    candidate_manifest = canonical_json_bytes(manifest)
    try:
        outcome = uploader.create_if_absent(
            manifest_path,
            candidate_manifest,
            content_type="application/json",
        )
    except Exception as exc:  # noqa: BLE001 - object details must not escape
        raise SnapshotUploadError("render-source manifest commit failed") from exc
    if outcome not in {"created", "exists"}:
        raise SnapshotUploadError("render-source manifest commit returned an invalid result")
    winner = _read_snapshot_object(
        uploader,
        object_path=manifest_path,
        role="render_source_manifest",
        max_bytes=_MAX_NATIVE_JSON_BYTES,
    )
    return _load_committed_render_source(
        loaded_input=loaded_input,
        reader=uploader,
        manifest_bytes=winner,
    )


def freeze_and_upload_evidence_snapshot(
    *,
    metadata: SnapshotCompletionMetadata,
    outputs_root: Path,
    artifact_virtual_path: str,
    artifact_host_path: Path,
    task_brief: BlindBrief | Mapping[str, Any],
    authoritative_mechanical: Mapping[str, Any],
    uploader: ImmutableObjectUploader,
    render_source: LoadedRenderSource,
    pdf_rasterizer: PdfRasterizer = rasterize_preview_pdf,
) -> SnapshotDescriptor:
    """Freeze and immutably upload one post-acceptance DQ-1 evidence bundle.

    Artifact bytes are copied to a build/version/hash-keyed create-only object
    without changing the ordinary delivery object. Source files are read
    without modification, all rendered evidence is lossless PNG derived only
    from the fixed committed render-source manifest, and the evidence manifest
    is always the final create-only object.
    """

    if not _QUALITY_RUN_RE.fullmatch(metadata.quality_run_id):
        raise SnapshotMissingEvidenceError("quality run identity is invalid")
    root = outputs_root.resolve()
    artifact_path, virtual_path = _validate_artifact_paths(
        outputs_root=root,
        artifact_virtual_path=artifact_virtual_path,
        artifact_host_path=artifact_host_path,
    )
    artifact_bytes = _read_required_file(
        artifact_path,
        role="accepted artifact",
        scope_root=root,
        max_bytes=_MAX_ACCEPTED_PPTX_BYTES,
    )
    artifact_hash = _sha256_bytes(artifact_bytes)
    source_files: dict[str, tuple[Path, str]] = {
        "accepted artifact": (artifact_path, artifact_hash),
    }
    creative_plan, design_plan, build_record, native_sources = _read_native_inputs(root)
    source_files.update(native_sources)
    brief, _brief_bytes = _canonical_blind_brief(task_brief)
    mechanical_record, _mechanical_bytes = _canonical_mapping(
        authoritative_mechanical,
        role="authoritative mechanical record",
    )
    slide_count, visible_text = _pptx_visible_text(artifact_bytes)
    _validate_build_record(
        build_record,
        metadata=metadata,
        expected_slide_count=slide_count,
    )
    expected_artifact = SnapshotArtifactReference(
        virtual_path=virtual_path,
        storage_object_path=_immutable_artifact_object_path(
            user_id=metadata.user_id,
            thread_id=metadata.thread_id,
            build_id=metadata.build_id,
            logical_artifact_id=metadata.logical_artifact_id,
            artifact_version_id=metadata.artifact_version_id,
            artifact_hash=artifact_hash,
            artifact_virtual_path=virtual_path,
        ),
        sha256=artifact_hash,
        size_bytes=len(artifact_bytes),
    )
    if any(
        actual != expected
        for actual, expected in (
            (render_source.manifest.quality_run_id, metadata.quality_run_id),
            (render_source.manifest.build_id, metadata.build_id),
            (render_source.manifest.user_id, metadata.user_id),
            (render_source.manifest.thread_id, metadata.thread_id),
            (render_source.manifest.task_id, metadata.task_id),
            (render_source.manifest.builder_run_id, metadata.builder_run_id),
            (render_source.manifest.parent_builder_trace_id, metadata.parent_builder_trace_id),
            (render_source.manifest.logical_artifact_id, metadata.logical_artifact_id),
            (render_source.manifest.artifact_version_id, metadata.artifact_version_id),
            (render_source.manifest.artifact_manifest_revision, metadata.manifest_revision),
            (render_source.manifest.source_artifact, expected_artifact),
            (render_source.manifest.input_manifest_path, f"{_pre_render_input_root(metadata)}/manifest.json"),
            (render_source.reference.manifest_path, _render_source_manifest_path(metadata)),
            (render_source.reference.manifest_hash, _sha256_bytes(canonical_json_bytes(render_source.manifest))),
            (render_source.reference.pdf, render_source.manifest.pdf),
            (render_source.reference.renderer_profile_version, _RENDER_SOURCE_PROFILE_VERSION),
            (render_source.reference.renderer_profile_hash, _RENDER_SOURCE_PROFILE_HASH),
        )
    ):
        raise SnapshotStaleError("render source does not match evidence inputs")
    preview_bytes = _read_required_file(
        render_source.pdf_host_path,
        role="render-source PDF",
        scope_root=root,
        max_bytes=_MAX_ACCEPTED_PREVIEW_PDF_BYTES,
    )
    preview_hash = _sha256_bytes(preview_bytes)
    if (
        preview_hash != render_source.reference.pdf.sha256
        or len(preview_bytes) != render_source.reference.pdf.size_bytes
        or _validate_render_source_pdf(
            preview_bytes,
            expected_page_count=slide_count,
        )
        != render_source.reference.pdf.page_count
    ):
        raise SnapshotConflictError("materialized render source differs from its commit")
    source_files["render-source PDF"] = (
        render_source.pdf_host_path,
        preview_hash,
    )
    _validate_artifact_storage_scope(metadata)
    quality_root = _quality_object_root(metadata)
    immutable_artifact_path = expected_artifact.storage_object_path

    with tempfile.TemporaryDirectory(prefix="dq1-evidence-snapshot-") as directory:
        work_root = Path(directory)
        frozen_preview = work_root / "accepted.preview.pdf"
        frozen_preview.write_bytes(preview_bytes)
        png_pages = pdf_rasterizer(frozen_preview)
        renders, contact_bytes, render_objects = _render_evidence(
            png_pages=png_pages,
            expected_slide_count=slide_count,
            object_root=quality_root,
            work_root=work_root,
        )

    artifact_reference = expected_artifact
    snapshot = QualityEvidenceSnapshot(
        campaign_id=metadata.campaign_id,
        build_id=metadata.build_id,
        user_id=metadata.user_id,
        task_id=metadata.task_id,
        builder_run_id=metadata.builder_run_id,
        parent_builder_trace_id=metadata.parent_builder_trace_id,
        logical_artifact_id=metadata.logical_artifact_id,
        artifact_version_id=metadata.artifact_version_id,
        manifest_revision=metadata.manifest_revision,
        artifact_path=virtual_path,
        artifact_hash=artifact_hash,
        brief_hash=canonical_sha256(brief),
        creative_plan_hash=canonical_sha256(creative_plan),
        design_plan_hash=canonical_sha256(design_plan),
        brief=brief,
        renders=renders,
        visible_text=visible_text,
        creative_plan=creative_plan,
        design_plan=design_plan,
        mechanical_record=mechanical_record,
        mechanical_record_hash=canonical_sha256(mechanical_record),
    )
    bundle = SnapshotEvidenceBundle(
        quality_run_id=metadata.quality_run_id,
        thread_id=metadata.thread_id,
        artifact=artifact_reference,
        build_record=build_record,
        snapshot=snapshot,
    )
    bundle_bytes = canonical_json_bytes(bundle)
    if len(bundle_bytes) > _MAX_EVIDENCE_BUNDLE_BYTES:
        raise SnapshotCoverageError("evidence bundle exceeds its byte budget")
    bundle_path = f"{quality_root}/evidence_bundle.json"
    object_records = [
        SnapshotObjectRecord(
            role="render",
            object_path=object_path,
            sha256=_sha256_bytes(content),
            size_bytes=len(content),
            media_type="image/png",
        )
        for object_path, content in render_objects
    ]
    contact_path = renders.contact_sheet.path
    object_records.append(
        SnapshotObjectRecord(
            role="contact_sheet",
            object_path=contact_path,
            sha256=_sha256_bytes(contact_bytes),
            size_bytes=len(contact_bytes),
            media_type="image/png",
        )
    )
    image_records = tuple(
        record for record in object_records if record.media_type == "image/png"
    )
    if any(record.size_bytes > _MAX_RENDER_PNG_BYTES for record in image_records):
        raise SnapshotCoverageError("render evidence exceeds its per-object byte budget")
    if sum(record.size_bytes for record in image_records) > _MAX_RENDER_TOTAL_BYTES:
        raise SnapshotCoverageError("render evidence exceeds its total byte budget")
    object_records.append(
        SnapshotObjectRecord(
            role="evidence_bundle",
            object_path=bundle_path,
            sha256=_sha256_bytes(bundle_bytes),
            size_bytes=len(bundle_bytes),
            media_type="application/json",
        )
    )
    render_hashes = {str(image.selector): image.sha256 for image in renders.slides}
    render_hashes["contact-sheet"] = renders.contact_sheet.sha256
    manifest = SnapshotEvidenceManifest(
        quality_run_id=metadata.quality_run_id,
        snapshot_id=metadata.quality_run_id,
        build_id=metadata.build_id,
        user_id=metadata.user_id,
        thread_id=metadata.thread_id,
        task_id=metadata.task_id,
        builder_run_id=metadata.builder_run_id,
        parent_builder_trace_id=metadata.parent_builder_trace_id,
        logical_artifact_id=metadata.logical_artifact_id,
        artifact_version_id=metadata.artifact_version_id,
        artifact_manifest_revision=metadata.manifest_revision,
        input_manifest_path=render_source.manifest.input_manifest_path,
        input_manifest_hash=render_source.manifest.input_manifest_hash,
        artifact=artifact_reference,
        render_source=render_source.reference,
        selectors=tuple(str(selector) for selector in renders.selectors),
        source_hashes=SnapshotSourceHashes(
            input_manifest=render_source.manifest.input_manifest_hash,
            artifact=artifact_hash,
            render_source_manifest=render_source.reference.manifest_hash,
            render_source_pdf=preview_hash,
            brief=snapshot.brief_hash,
            creative_plan=snapshot.creative_plan_hash,
            design_plan=snapshot.design_plan_hash,
            build_record=canonical_sha256(build_record),
            mechanical_record=snapshot.mechanical_record_hash,
            visible_text=canonical_sha256(visible_text),
        ),
        render_hashes=render_hashes,
        objects=tuple(object_records),
        evidence_bundle_path=bundle_path,
        evidence_bundle_hash=_sha256_bytes(bundle_bytes),
    )
    manifest_bytes = canonical_json_bytes(manifest)
    if len(manifest_bytes) > _MAX_NATIVE_JSON_BYTES:
        raise SnapshotCoverageError("evidence manifest exceeds its byte budget")
    manifest_path = f"{quality_root}/evidence_manifest.json"

    _verify_sources_unchanged(source_files)
    _upload_immutable(
        uploader=uploader,
        object_path=immutable_artifact_path,
        content=artifact_bytes,
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        role="accepted_artifact",
        verify_after_create=True,
    )
    for record, (object_path, content) in zip(
        object_records[:-2],
        render_objects,
        strict=True,
    ):
        _upload_immutable(
            uploader=uploader,
            object_path=object_path,
            content=content,
            content_type=record.media_type,
            role="render",
        )
    _upload_immutable(
        uploader=uploader,
        object_path=contact_path,
        content=contact_bytes,
        content_type="image/png",
        role="contact_sheet",
    )
    _upload_immutable(
        uploader=uploader,
        object_path=bundle_path,
        content=bundle_bytes,
        content_type="application/json",
        role="evidence_bundle",
    )
    # The manifest is the commit marker. It must remain the final upload.
    _upload_immutable(
        uploader=uploader,
        object_path=manifest_path,
        content=manifest_bytes,
        content_type="application/json",
        role="evidence_manifest",
    )
    return SnapshotDescriptor(
        snapshot_id=metadata.quality_run_id,
        snapshot_path=manifest_path,
        snapshot_hash=_sha256_bytes(manifest_bytes),
        counts=SnapshotCounts(
            slide_count=slide_count,
            visible_text_slide_count=len(visible_text),
            # Immutable artifact, render-source manifest/PDF, and this
            # evidence manifest are referenced outside the object inventory.
            evidence_object_count=len(object_records) + 4,
        ),
    )


def _read_snapshot_object(
    reader: ImmutableObjectReader,
    *,
    object_path: str,
    role: str,
    max_bytes: int,
) -> bytes:
    try:
        normalized = normalize_object_path(object_path)
        read_bounded = getattr(reader, "read_bounded", None)
        if callable(read_bounded):
            content = read_bounded(normalized, max_bytes=max_bytes)
        else:
            content = reader.read(normalized)
    except Exception as exc:  # noqa: BLE001 - storage details must not escape or be logged
        raise SnapshotUploadError(f"immutable read failed for {role}") from exc
    if not content:
        raise SnapshotMissingEvidenceError(f"immutable {role} object is missing")
    if len(content) > max_bytes:
        raise SnapshotCoverageError(f"immutable {role} exceeds its byte budget")
    return content


def _canonical_model_from_bytes[ModelT: BaseModel](
    content: bytes,
    *,
    model_type: type[ModelT],
    role: str,
) -> ModelT:
    payload = _strict_json_object(content, role=role)
    try:
        model = model_type.model_validate(payload)
    except ValueError as exc:
        raise SnapshotConflictError(f"immutable {role} object violates its schema") from exc
    if canonical_json_bytes(model) != content:
        raise SnapshotConflictError(f"immutable {role} object is not canonical")
    return model


def _verify_pre_render_input_manifest(
    *,
    descriptor: PreRenderInputBundleDescriptor,
    manifest: PreRenderInputBundleManifest,
) -> None:
    if manifest.quality_run_id != descriptor.bundle_id:
        raise SnapshotConflictError("pre-render input identities do not match")
    expected_input_root = normalize_object_path(f"{_quality_object_root_for_identity(user_id=manifest.user_id, thread_id=manifest.thread_id, build_id=manifest.build_id, quality_run_id=manifest.quality_run_id)}/input_bundle")
    if descriptor.manifest_path != f"{expected_input_root}/manifest.json":
        raise SnapshotConflictError("pre-render manifest path does not match build identity")
    for record in manifest.objects:
        if record.role == "accepted_artifact":
            expected_path = _immutable_artifact_object_path(
                user_id=manifest.user_id,
                thread_id=manifest.thread_id,
                build_id=manifest.build_id,
                logical_artifact_id=manifest.logical_artifact_id,
                artifact_version_id=manifest.artifact_version_id,
                artifact_hash=record.sha256,
                artifact_virtual_path=manifest.artifact_virtual_path,
            )
            expected_media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        else:
            expected_path = _pre_render_json_object_path(
                input_root=expected_input_root,
                role=record.role,
                content_hash=record.sha256,
            )
            expected_media_type = "application/json"
        if record.object_path != expected_path:
            raise SnapshotConflictError("pre-render input object path does not match immutable identity")
        if record.media_type != expected_media_type:
            raise SnapshotConflictError("pre-render input media type is invalid")


def _canonical_pre_render_mapping(content: bytes, *, role: str) -> dict[str, Any]:
    value = _strict_json_object(content, role=role)
    if canonical_json_bytes(value) != content:
        raise SnapshotConflictError(f"immutable {role} object is not canonical")
    return value


def _write_materialized_input(path: Path, content: bytes, *, role: str) -> None:
    if path.exists():
        try:
            existing = path.read_bytes()
        except OSError as exc:
            raise SnapshotUploadError(f"pre-render materialization failed for {role}") from exc
        if len(existing) != len(content) or not hmac.compare_digest(
            _sha256_bytes(existing),
            _sha256_bytes(content),
        ):
            raise SnapshotConflictError(f"pre-render materialization differs for {role}")
        return
    try:
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_bytes(content)
        written = path.read_bytes()
    except OSError as exc:
        raise SnapshotUploadError(f"pre-render materialization failed for {role}") from exc
    if len(written) != len(content) or not hmac.compare_digest(
        _sha256_bytes(written),
        _sha256_bytes(content),
    ):
        raise SnapshotUploadError(f"pre-render materialization verification failed for {role}")


def load_pre_render_input_bundle(
    *,
    descriptor: PreRenderInputBundleDescriptor,
    expected_identity: SnapshotRunIdentity,
    reader: ImmutableObjectReader,
    materialization_root: Path,
) -> LoadedPreRenderInputBundle:
    """Verify and reconstruct every input needed by evidence rendering.

    The fixed manifest is the only object read before the complete durable-row
    identity and input-manifest binding are verified.
    """

    manifest_bytes = _read_snapshot_object(
        reader,
        object_path=descriptor.manifest_path,
        role="pre_render_input_manifest",
        max_bytes=_MAX_NATIVE_JSON_BYTES,
    )
    if not hmac.compare_digest(
        _sha256_bytes(manifest_bytes),
        descriptor.manifest_hash,
    ):
        raise SnapshotConflictError("pre-render descriptor hash does not match manifest")
    manifest = _canonical_model_from_bytes(
        manifest_bytes,
        model_type=PreRenderInputBundleManifest,
        role="pre_render_input_manifest",
    )
    _verify_pre_render_input_manifest(
        descriptor=descriptor,
        manifest=manifest,
    )
    if (
        descriptor.manifest_path != expected_identity.input_manifest_object_path
        or descriptor.manifest_hash != expected_identity.input_manifest_hash
        or not _run_identity_matches_manifest(manifest, expected_identity)
    ):
        raise SnapshotStaleError("pre-render input manifest does not match the durable row")

    contents: dict[PreRenderInputRole, bytes] = {}
    for record in manifest.objects:
        max_bytes = (
            _MAX_ACCEPTED_PPTX_BYTES
            if record.role == "accepted_artifact"
            else _MAX_NATIVE_JSON_BYTES
        )
        if record.size_bytes > max_bytes:
            raise SnapshotCoverageError(f"immutable {record.role} exceeds its byte budget")
        content = _read_snapshot_object(
            reader,
            object_path=record.object_path,
            role=record.role,
            max_bytes=max_bytes,
        )
        if len(content) != record.size_bytes or not hmac.compare_digest(
            _sha256_bytes(content),
            record.sha256,
        ):
            raise SnapshotConflictError(f"immutable {record.role} bytes do not match pre-render manifest")
        contents[record.role] = content

    creative_plan = _canonical_pre_render_mapping(
        contents["creative_plan"],
        role="creative_plan",
    )
    design_plan = _canonical_pre_render_mapping(
        contents["design_plan"],
        role="design_plan",
    )
    build_record = _canonical_pre_render_mapping(
        contents["build_record"],
        role="build_record",
    )
    brief = _canonical_model_from_bytes(
        contents["blind_brief"],
        model_type=BlindBrief,
        role="blind_brief",
    )
    if sanitize_current_request(brief.request) != brief.request:
        raise SnapshotConflictError("immutable blind brief is not sanitized")
    mechanical_record = _canonical_pre_render_mapping(
        contents["mechanical_record"],
        role="mechanical_record",
    )
    artifact_bytes = contents["accepted_artifact"]
    slide_count, _visible_text = _pptx_visible_text(artifact_bytes)
    metadata = SnapshotCompletionMetadata(
        quality_run_id=manifest.quality_run_id,
        build_id=manifest.build_id,
        user_id=manifest.user_id,
        thread_id=manifest.thread_id,
        task_id=manifest.task_id,
        builder_run_id=manifest.builder_run_id,
        parent_builder_trace_id=manifest.parent_builder_trace_id,
        logical_artifact_id=manifest.logical_artifact_id,
        artifact_version_id=manifest.artifact_version_id,
        manifest_revision=manifest.artifact_manifest_revision,
        artifact_storage_object_path=manifest.objects[0].object_path,
    )
    _validate_build_record(
        build_record,
        metadata=metadata,
        expected_slide_count=slide_count,
    )

    outputs_root = materialization_root.resolve() / manifest.quality_run_id / "outputs"
    relative_artifact = manifest.artifact_virtual_path.removeprefix("/mnt/user-data/outputs/")
    artifact_host_path = (outputs_root / relative_artifact).resolve()
    resolved_outputs = outputs_root.resolve()
    if resolved_outputs != artifact_host_path and resolved_outputs not in artifact_host_path.parents:
        raise SnapshotConflictError("pre-render artifact materialization is out of scope")

    # No local writes occur until the complete remote bundle has verified.
    _write_materialized_input(
        artifact_host_path,
        artifact_bytes,
        role="accepted_artifact",
    )
    deck_build_root = resolved_outputs / "deck_build"
    for role, filename, content in (
        ("creative_plan", "creative_plan.json", canonical_json_bytes(creative_plan)),
        ("design_plan", "design_plan.json", canonical_json_bytes(design_plan)),
        ("build_record", "build.json", canonical_json_bytes(build_record)),
    ):
        _write_materialized_input(
            deck_build_root / filename,
            content,
            role=role,
        )
    _validate_artifact_paths(
        outputs_root=resolved_outputs,
        artifact_virtual_path=manifest.artifact_virtual_path,
        artifact_host_path=artifact_host_path,
    )
    return LoadedPreRenderInputBundle(
        descriptor=descriptor,
        manifest=manifest,
        metadata=metadata,
        brief=brief,
        mechanical_record=mechanical_record,
        outputs_root=resolved_outputs,
        artifact_virtual_path=manifest.artifact_virtual_path,
        artifact_host_path=artifact_host_path,
    )


def _verify_loaded_manifest(
    *,
    descriptor: SnapshotDescriptor,
    manifest: SnapshotEvidenceManifest,
    bundle: SnapshotEvidenceBundle,
) -> None:
    if manifest.quality_run_id != descriptor.snapshot_id or manifest.snapshot_id != descriptor.snapshot_id:
        raise SnapshotConflictError("snapshot identities do not match")
    expected_suffix = f"/.builder/builds/{manifest.build_id}/quality/{manifest.quality_run_id}/evidence_manifest.json"
    if not descriptor.snapshot_path.endswith(expected_suffix):
        raise SnapshotConflictError("snapshot manifest path does not match build identity")
    snapshot = bundle.snapshot
    if (
        bundle.quality_run_id != manifest.quality_run_id
        or bundle.thread_id != manifest.thread_id
        or snapshot.campaign_id != manifest.campaign_id
        or snapshot.build_id != manifest.build_id
        or snapshot.user_id != manifest.user_id
        or snapshot.task_id != manifest.task_id
        or snapshot.builder_run_id != manifest.builder_run_id
        or snapshot.parent_builder_trace_id != manifest.parent_builder_trace_id
        or snapshot.logical_artifact_id != manifest.logical_artifact_id
        or snapshot.artifact_version_id != manifest.artifact_version_id
        or snapshot.manifest_revision != manifest.artifact_manifest_revision
        or bundle.artifact != manifest.artifact
    ):
        raise SnapshotConflictError("snapshot bundle correlation does not match manifest")
    if tuple(str(selector) for selector in snapshot.renders.selectors) != tuple(str(selector) for selector in manifest.selectors):
        raise SnapshotCoverageError("snapshot selectors do not match manifest")
    source_hashes = manifest.source_hashes
    expected_source_hashes = {
        "input_manifest": manifest.input_manifest_hash,
        "artifact": snapshot.artifact_hash,
        "render_source_manifest": manifest.render_source.manifest_hash,
        "render_source_pdf": manifest.render_source.pdf.sha256,
        "brief": snapshot.brief_hash,
        "creative_plan": snapshot.creative_plan_hash,
        "design_plan": snapshot.design_plan_hash,
        "build_record": canonical_sha256(bundle.build_record),
        "mechanical_record": snapshot.mechanical_record_hash,
        "visible_text": canonical_sha256(snapshot.visible_text),
    }
    actual_source_hashes = source_hashes.model_dump(mode="python")
    for key, expected_hash in expected_source_hashes.items():
        if actual_source_hashes[key] != expected_hash:
            raise SnapshotConflictError(f"snapshot {key} hash does not match manifest")
    expected_render_hashes = {str(image.selector): image.sha256 for image in snapshot.renders.slides}
    expected_render_hashes["contact-sheet"] = snapshot.renders.contact_sheet.sha256
    if manifest.render_hashes != expected_render_hashes:
        raise SnapshotConflictError("snapshot render hashes do not match manifest")
    if descriptor.counts.slide_count != snapshot.renders.expected_slide_count:
        raise SnapshotCoverageError("snapshot descriptor slide count does not match bundle")
    if descriptor.counts.visible_text_slide_count != len(snapshot.visible_text):
        raise SnapshotCoverageError("snapshot descriptor text coverage does not match bundle")
    if descriptor.counts.evidence_object_count != len(manifest.objects) + 4:
        raise SnapshotConflictError("snapshot descriptor object count does not match manifest")
    expected_artifact_path = _immutable_artifact_object_path(
        user_id=snapshot.user_id,
        thread_id=bundle.thread_id,
        build_id=snapshot.build_id,
        logical_artifact_id=snapshot.logical_artifact_id,
        artifact_version_id=snapshot.artifact_version_id,
        artifact_hash=snapshot.artifact_hash,
        artifact_virtual_path=snapshot.artifact_path,
    )
    if bundle.artifact.storage_object_path != expected_artifact_path:
        raise SnapshotConflictError("snapshot artifact object path does not match immutable identity")


def _verify_evidence_render_source_manifest(
    *,
    evidence_manifest: SnapshotEvidenceManifest,
    render_manifest: RenderSourceManifest,
) -> None:
    quality_root = _quality_object_root_for_identity(
        user_id=evidence_manifest.user_id,
        thread_id=evidence_manifest.thread_id,
        build_id=evidence_manifest.build_id,
        quality_run_id=evidence_manifest.quality_run_id,
    )
    expected_manifest_path = normalize_object_path(
        f"{quality_root}/render_source/manifest.json"
    )
    expected_pdf_path = normalize_object_path(
        f"{quality_root}/render_source/objects/{render_manifest.pdf.sha256}.pdf"
    )
    reference = evidence_manifest.render_source
    if any(
        actual != expected
        for actual, expected in (
            (reference.manifest_path, expected_manifest_path),
            (reference.pdf, render_manifest.pdf),
            (reference.renderer_profile_version, render_manifest.renderer_profile_version),
            (reference.renderer_profile_hash, render_manifest.renderer_profile_hash),
            (render_manifest.campaign_id, evidence_manifest.campaign_id),
            (render_manifest.quality_run_id, evidence_manifest.quality_run_id),
            (render_manifest.build_id, evidence_manifest.build_id),
            (render_manifest.user_id, evidence_manifest.user_id),
            (render_manifest.thread_id, evidence_manifest.thread_id),
            (render_manifest.task_id, evidence_manifest.task_id),
            (render_manifest.builder_run_id, evidence_manifest.builder_run_id),
            (render_manifest.parent_builder_trace_id, evidence_manifest.parent_builder_trace_id),
            (render_manifest.logical_artifact_id, evidence_manifest.logical_artifact_id),
            (render_manifest.artifact_version_id, evidence_manifest.artifact_version_id),
            (render_manifest.artifact_manifest_revision, evidence_manifest.artifact_manifest_revision),
            (render_manifest.input_manifest_path, evidence_manifest.input_manifest_path),
            (render_manifest.input_manifest_hash, evidence_manifest.input_manifest_hash),
            (render_manifest.source_artifact, evidence_manifest.artifact),
            (render_manifest.pdf.object_path, expected_pdf_path),
            (render_manifest.renderer_profile_hash, _RENDER_SOURCE_PROFILE_HASH),
        )
    ):
        raise SnapshotStaleError("render-source commit does not match the evidence manifest")


def _write_materialized_png(path: Path, content: bytes, *, role: str) -> None:
    if path.exists():
        try:
            existing = path.read_bytes()
        except OSError as exc:
            raise SnapshotUploadError(f"local materialization failed for {role}") from exc
        if not hmac.compare_digest(_sha256_bytes(existing), _sha256_bytes(content)):
            raise SnapshotConflictError(f"local materialization differs for {role}")
        return
    try:
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_bytes(content)
    except OSError as exc:
        raise SnapshotUploadError(f"local materialization failed for {role}") from exc
    if _sha256_bytes(path.read_bytes()) != _sha256_bytes(content):
        raise SnapshotUploadError(f"local materialization verification failed for {role}")


def load_evidence_snapshot(
    *,
    descriptor: SnapshotDescriptor,
    expected_identity: SnapshotRunIdentity,
    reader: ImmutableObjectReader,
    materialization_root: Path,
) -> LoadedEvidenceSnapshot:
    """Verify an immutable snapshot and materialize only its lossless PNGs.

    No manifest reference is followed until the evidence manifest has matched
    the complete durable-row identity and its fixed input-manifest binding.
    """

    manifest_bytes = _read_snapshot_object(
        reader,
        object_path=descriptor.snapshot_path,
        role="evidence_manifest",
        max_bytes=_MAX_NATIVE_JSON_BYTES,
    )
    if not hmac.compare_digest(
        _sha256_bytes(manifest_bytes),
        descriptor.snapshot_hash,
    ):
        raise SnapshotConflictError("snapshot descriptor hash does not match manifest")
    manifest = _canonical_model_from_bytes(
        manifest_bytes,
        model_type=SnapshotEvidenceManifest,
        role="evidence_manifest",
    )
    verify_evidence_manifest_identity(manifest, expected_identity)
    if manifest.snapshot_id != descriptor.snapshot_id:
        raise SnapshotStaleError("evidence manifest snapshot identity is stale")

    render_manifest_bytes = _read_snapshot_object(
        reader,
        object_path=manifest.render_source.manifest_path,
        role="render_source_manifest",
        max_bytes=_MAX_NATIVE_JSON_BYTES,
    )
    if not hmac.compare_digest(
        _sha256_bytes(render_manifest_bytes),
        manifest.render_source.manifest_hash,
    ):
        raise SnapshotConflictError("render-source manifest hash does not match evidence")
    render_manifest = _canonical_model_from_bytes(
        render_manifest_bytes,
        model_type=RenderSourceManifest,
        role="render_source_manifest",
    )
    _verify_evidence_render_source_manifest(
        evidence_manifest=manifest,
        render_manifest=render_manifest,
    )
    if render_manifest.pdf.size_bytes > _MAX_ACCEPTED_PREVIEW_PDF_BYTES:
        raise SnapshotCoverageError("render-source PDF exceeds its byte budget")
    render_pdf_bytes = _read_snapshot_object(
        reader,
        object_path=render_manifest.pdf.object_path,
        role="render_source_pdf",
        max_bytes=_MAX_ACCEPTED_PREVIEW_PDF_BYTES,
    )
    if (
        len(render_pdf_bytes) != render_manifest.pdf.size_bytes
        or not hmac.compare_digest(
            _sha256_bytes(render_pdf_bytes),
            render_manifest.pdf.sha256,
        )
        or _validate_render_source_pdf(
            render_pdf_bytes,
            expected_page_count=len(manifest.selectors),
        )
        != render_manifest.pdf.page_count
    ):
        raise SnapshotConflictError("render-source PDF does not match its commit")

    bundle_bytes = _read_snapshot_object(
        reader,
        object_path=manifest.evidence_bundle_path,
        role="evidence_bundle",
        max_bytes=_MAX_EVIDENCE_BUNDLE_BYTES,
    )
    if not hmac.compare_digest(
        _sha256_bytes(bundle_bytes),
        manifest.evidence_bundle_hash,
    ):
        raise SnapshotConflictError("evidence bundle hash does not match manifest")
    bundle = _canonical_model_from_bytes(
        bundle_bytes,
        model_type=SnapshotEvidenceBundle,
        role="evidence_bundle",
    )
    _verify_loaded_manifest(
        descriptor=descriptor,
        manifest=manifest,
        bundle=bundle,
    )

    if bundle.artifact.size_bytes > _MAX_ACCEPTED_PPTX_BYTES:
        raise SnapshotCoverageError("immutable accepted artifact exceeds its byte budget")
    artifact_bytes = _read_snapshot_object(
        reader,
        object_path=bundle.artifact.storage_object_path,
        role="accepted_artifact",
        max_bytes=_MAX_ACCEPTED_PPTX_BYTES,
    )
    if len(artifact_bytes) != bundle.artifact.size_bytes or not hmac.compare_digest(
        _sha256_bytes(artifact_bytes),
        bundle.artifact.sha256,
    ):
        raise SnapshotConflictError("immutable accepted artifact bytes do not match snapshot")

    records_by_path = {record.object_path: record for record in manifest.objects}
    expected_paths = {image.path for image in bundle.snapshot.renders.slides} | {
        bundle.snapshot.renders.contact_sheet.path,
        manifest.evidence_bundle_path,
    }
    if set(records_by_path) != expected_paths:
        raise SnapshotConflictError("manifest object inventory does not match bundle")
    bundle_record = records_by_path[manifest.evidence_bundle_path]
    if (
        bundle_record.role != "evidence_bundle"
        or bundle_record.media_type != "application/json"
        or bundle_record.sha256 != manifest.evidence_bundle_hash
        or bundle_record.size_bytes != len(bundle_bytes)
    ):
        raise SnapshotConflictError("evidence bundle record does not match manifest")
    render_records = tuple(
        record for record in manifest.objects if record.role != "evidence_bundle"
    )
    if any(record.size_bytes > _MAX_RENDER_PNG_BYTES for record in render_records):
        raise SnapshotCoverageError("render object exceeds its byte budget")
    if sum(record.size_bytes for record in render_records) > _MAX_RENDER_TOTAL_BYTES:
        raise SnapshotCoverageError("render object inventory exceeds its total byte budget")
    render_root = materialization_root.resolve() / descriptor.snapshot_id / "renders"
    materialized_slides: list[ImageEvidence] = []
    for index, image in enumerate(bundle.snapshot.renders.slides, start=1):
        record = records_by_path[image.path]
        if record.role != "render" or record.media_type != "image/png" or record.sha256 != image.sha256:
            raise SnapshotConflictError("render object record does not match bundle")
        content = _read_snapshot_object(
            reader,
            object_path=image.path,
            role="render",
            max_bytes=_MAX_RENDER_PNG_BYTES,
        )
        if (
            len(content) != record.size_bytes
            or not hmac.compare_digest(_sha256_bytes(content), image.sha256)
        ):
            raise SnapshotConflictError("render object hash does not match bundle")
        width, height = _decode_png(content, role="render object")
        if (width, height) != (image.width, image.height):
            raise SnapshotCoverageError("render dimensions do not match bundle")
        local_path = render_root / f"slide-{index:04d}.png"
        _write_materialized_png(local_path, content, role="render")
        materialized_slides.append(image.model_copy(update={"path": local_path.as_posix()}))
    contact = bundle.snapshot.renders.contact_sheet
    contact_record = records_by_path[contact.path]
    if contact_record.role != "contact_sheet" or contact_record.media_type != "image/png" or contact_record.sha256 != contact.sha256:
        raise SnapshotConflictError("contact sheet object record does not match bundle")
    contact_content = _read_snapshot_object(
        reader,
        object_path=contact.path,
        role="contact_sheet",
        max_bytes=_MAX_RENDER_PNG_BYTES,
    )
    if (
        len(contact_content) != contact_record.size_bytes
        or not hmac.compare_digest(_sha256_bytes(contact_content), contact.sha256)
    ):
        raise SnapshotConflictError("contact sheet hash does not match bundle")
    contact_width, contact_height = _decode_png(contact_content, role="contact sheet object")
    if (contact_width, contact_height) != (contact.width, contact.height):
        raise SnapshotCoverageError("contact sheet dimensions do not match bundle")
    local_contact_path = render_root / "contact-sheet.png"
    _write_materialized_png(local_contact_path, contact_content, role="contact_sheet")
    materialized_contact = contact.model_copy(update={"path": local_contact_path.as_posix()})
    materialized_renders = RenderEvidence(
        expected_slide_count=bundle.snapshot.renders.expected_slide_count,
        contact_sheet=materialized_contact,
        slides=tuple(materialized_slides),
        selectors=bundle.snapshot.renders.selectors,
    )
    snapshot_payload = bundle.snapshot.model_dump(mode="python")
    snapshot_payload["renders"] = materialized_renders.model_dump(mode="python")
    materialized_snapshot = QualityEvidenceSnapshot.model_validate(snapshot_payload)
    return LoadedEvidenceSnapshot(
        descriptor=descriptor,
        manifest=manifest,
        build_record=bundle.build_record,
        snapshot=materialized_snapshot,
    )
