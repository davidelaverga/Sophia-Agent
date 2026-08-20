from __future__ import annotations

import contextlib
import hashlib
import json
import logging
import math
import os
import re
import subprocess  # noqa: S404 - fixed Python script path with sanitized args.
import sys
import time
from collections.abc import Callable, Iterator
from datetime import UTC, datetime
from html.parser import HTMLParser
from pathlib import Path
from typing import Any

import tinycss2
from bs4 import BeautifulSoup, Tag
from langchain.tools import ToolRuntime

from deerflow.sandbox.tools import get_thread_data, replace_virtual_path
from deerflow.sophia.build_runtime.identity import new_build_id
from deerflow.sophia.deck_build.asset_policy import (
    generated_asset_slides,
    normalize_visual_policy,
    write_asset_policy,
)
from deerflow.sophia.deck_build.creative_plan import (
    CreativePlanValidationError,
    normalize_creative_plan,
    write_creative_plan,
)
from deerflow.sophia.deck_build.design_plan import write_design_plan
from deerflow.sophia.deck_build.evaluator import DeckEvaluator
from deerflow.sophia.deck_build.foundation import BuildFoundationPersistenceError, materialize_deck_foundation_safely
from deerflow.sophia.deck_build.html_sanitizer import (
    assemble_compact_slide_html,
    validate_and_sanitize_slide_html,
    validation_summary,
)
from deerflow.sophia.deck_build.image_assets import (
    apply_creative_asset_plan,
    planned_asset_ref_basenames,
)
from deerflow.sophia.deck_build.image_prompting import deck_asset_prompt_payload
from deerflow.sophia.deck_build.ir_repair import deck_mechanical_repair_instruction_from_reports
from deerflow.sophia.deck_build.mechanical_gates import evaluate_mechanical_gates
from deerflow.sophia.deck_build.models import DeckBuild, DeckBuildResult, DeckEvaluation, DeckSlideSpec
from deerflow.sophia.deck_build.native_contrast import evaluate_native_contrast
from deerflow.sophia.deck_build.source_retention import (
    evaluate_source_retention,
    retention_summary,
)
from deerflow.sophia.deck_build.storage import save_deck_build
from deerflow.sophia.deck_build.tool_contract import (
    COMPACT_V2_MAX_SLIDE_HTML_BODY_BYTES,
    COMPACT_V2_TARGET_SLIDE_HTML_BODY_BYTES,
)
from deerflow.sophia.deck_build.tracing import (
    DEFAULT_ARTIFACT_TARGET_EXT,
    DEFAULT_DECK_COMPILE_MODE,
    DEFAULT_DECK_ROUTE,
    FORBIDDEN_SCREENSHOT_COMPILE_MODES,
    HTML_SCREENSHOT_DEBUG_COMPILE_MODE,
    NATIVE_DECK_COMPILE_MODE,
    NATIVE_UNAVAILABLE_DECK_COMPILE_MODE,
    basename,
    deck_span,
    finish_span,
    safe_excerpt,
    stable_hash,
)
from deerflow.sophia.deck_native import DeckNativeService, native_mechanical_report
from deerflow.sophia.deck_native.errors import DeckNativePathError
from deerflow.sophia.deck_native.models import (
    NativeDeckInspectResult,
    NativeDeckLintFixResult,
    NativeDeckPatchResult,
    NativeDeckPreflight,
    NativeDeckRenderResult,
)
from deerflow.sophia.deck_native.policy import classify_native_deck_substrate
from deerflow.sophia.image_subprocess import (
    ImageThreadRoots,
    TrustedImageRequest,
    run_trusted_image_request,
)
from deerflow.sophia.subprocess_env import trusted_subprocess_env
from deerflow.sophia.tools.prepare_pptx_image_manifest import create_pptx_image_manifest_core
from deerflow.sophia.tools.render_markdown_to_pdf import _ensure_relative_to_outputs

_OUTPUTS = "/mnt/user-data/outputs/"
_ASSETS = f"{_OUTPUTS}assets"
_PROMPTS = f"{_ASSETS}/prompts"
_SLIDES = f"{_OUTPUTS}slides"
_DECK_BUILD = f"{_OUTPUTS}deck_build"
_DESIGN_PLAN = f"{_DECK_BUILD}/design_plan.json"
_CREATIVE_PLAN = f"{_DECK_BUILD}/creative_plan.json"
_ASSET_POLICY = f"{_DECK_BUILD}/asset_policy.json"
_MANIFEST = f"{_ASSETS}/slide-visuals.manifest.json"
_NATIVE = f"{_OUTPUTS}.builder/deck_native"
_NATIVE_BASE = f"{_NATIVE}/base.pptx"
_NATIVE_PATCH = f"{_NATIVE}/deck.patch.json"
_NATIVE_RENDER_DIR = f"{_NATIVE}/rendered"
_SCHEMA = "sophia-deck-build/v1"
_SAFE_PPTX_FONTS = {"arial", "calibri", "cambria"}
_GENERIC_CSS_FONTS = {"cursive", "fantasy", "monospace", "sans-serif", "serif", "system-ui"}
_FONT_SIZE_KEYWORDS = {
    "large",
    "larger",
    "medium",
    "small",
    "smaller",
    "x-large",
    "x-small",
    "xx-large",
    "xx-small",
    "xxx-large",
}
_BASE_FONT_SELECTORS = {"*", ".slide-root", "body", "html", "main", "main.slide-root"}
_COMPACT_SOURCE_ID_RE = re.compile(r"^[a-z][a-z0-9_-]{0,31}$")
_COMPACT_SOURCE_GEOMETRY_PROPERTIES = ("left", "top", "width", "height")
_COMPACT_SOURCE_ANCHOR_INVARIANT_PROPERTIES = (
    "position",
    "box-sizing",
    "margin",
)
_COMPACT_SOURCE_PHYSICAL_MARGIN_PROPERTIES = frozenset(
    {
        "margin",
        "margin-bottom",
        "margin-left",
        "margin-right",
        "margin-top",
    }
)
_COMPACT_SOURCE_LOGICAL_MARGIN_PROPERTIES = frozenset(
    {
        "margin-block",
        "margin-block-end",
        "margin-block-start",
        "margin-inline",
        "margin-inline-end",
        "margin-inline-start",
    }
)
_COMPACT_SOURCE_VENDOR_MARGIN_PROPERTIES = frozenset(
    {
        "-moz-margin-end",
        "-moz-margin-left-value",
        "-moz-margin-right-value",
        "-moz-margin-start",
        "-webkit-margin-after",
        "-webkit-margin-before",
        "-webkit-margin-bottom-collapse",
        "-webkit-margin-collapse",
        "-webkit-margin-end",
        "-webkit-margin-start",
        "-webkit-margin-top-collapse",
    }
)
_COMPACT_SOURCE_INLINE_DUPLICATE_PROPERTIES = frozenset(
    {
        "position",
        "box-sizing",
        "margin",
        *_COMPACT_SOURCE_GEOMETRY_PROPERTIES,
    }
)
_BANNED_STYLE_RE = re.compile(r"\b(chalkboard|blackboard|whiteboard|handwritten|sketch|cyberpunk|neon)\b", re.I)
_BANNED_TEXT_RE = re.compile(r"THE TEXT READS|title reads|large readable text|paragraph text|axis labels?|formula", re.I)
_NEGATED_BANNED_TERM_RE = re.compile(r"(?:\bno\b|\bnot\b|\bwithout\b|\bavoid\b|\bnever\b|\bdo\s+not\b)\W*$", re.I)
_TEXT_ONLY_REQUEST_RE = re.compile(
    r"\b(?:plain\s+text[-\s]?only|text[-\s]?only|no[-\s]?image|no\s+images?"
    r"|no\s+visuals?|without\s+(?:images?|visuals?)|with\s+no\s+(?:images?|visuals?))\b",
    re.I,
)
_REQUEST_CONTEXT_TEXT_KEYS = (
    "user_request",
    "request",
    "prompt",
    "task",
    "task_brief",
    "title",
    "description",
    "artifact_title",
    "task_title",
)
_SERIAL_REPAIR_MAX_ATTEMPTS = 2
_NON_REPAIRABLE_IMAGE_ERROR_CLASSES = {
    "auth_invalid",
    "content_blocked",
    "deck_deadline_exceeded",
    "egress_blocked",
    "file_not_found",
    "image_script_not_found",
    "import_error",
    "invalid_reference_image",
    "invalid_size",
    "manifest_empty",
    "manifest_prompt_missing",
    "manifest_unreadable",
    "missing_api_key",
    "missing_prompt_or_output",
    "org_not_verified",
    "permission_denied",
    "process_exit",
    "python_not_found",
    "quota_exceeded",
    "sandbox_path_rejected",
    "shell_error",
}

logger = logging.getLogger(__name__)


DeckCompiler = Callable[[ToolRuntime, str, str, str], dict[str, Any]]
ImageSingleRunner = Callable[[DeckSlideSpec, ToolRuntime, int], dict[str, Any]]


class DeckBuildService:
    def __init__(
        self,
        *,
        image_batch_runner: Callable[[str, ToolRuntime], dict[str, Any]] | None = None,
        image_single_runner: ImageSingleRunner | None = None,
        deck_compiler: DeckCompiler | None = None,
        native_service: DeckNativeService | None = None,
        evaluator: DeckEvaluator | None = None,
    ) -> None:
        self._image_batch_runner = image_batch_runner or self._run_image_batch_subprocess
        self._image_single_runner = image_single_runner or self._run_image_single_subprocess
        self._deck_compiler = deck_compiler or _compile_with_build_deck_from_slides
        self._native_service = native_service or DeckNativeService()
        self._evaluator = evaluator or DeckEvaluator()

    def prepare_and_build(
        self,
        *,
        runtime: ToolRuntime,
        deck_title: str,
        slides: list[dict[str, Any]],
        output_path: str,
        register: str = "professional_technical",
        visual_policy: str = "auto",
        deck_stylesheet: str | None = None,
        authoring_contract: str | None = None,
        style_profile: dict[str, Any] | None = None,
        design_plan: dict[str, Any] | None = None,
        creative_plan: dict[str, Any] | None = None,
        native_lint_slide_indices: tuple[int, ...] | None = None,
    ) -> DeckBuildResult:
        service_started = time.perf_counter()
        # Fresh production builds receive this at dispatch so identity survives
        # retries, resume, revision, and thread handoff. The fallback is kept
        # only for direct service callers and legacy queued payloads.
        build_id = str(_state_value(runtime, "builder_build_id") or _state_value(runtime, "build_id") or new_build_id())
        now = _now()
        resolved_visual_policy = normalize_visual_policy(visual_policy)
        deck = DeckBuild(
            build_id=build_id,
            schema_version=_SCHEMA,
            user_id=_runtime_identity_value(runtime, "user_id"),
            thread_id=str(_runtime_identity_value(runtime, "thread_id") or ""),
            parent_thread_id=_runtime_identity_value(runtime, "parent_thread_id"),
            run_id=_runtime_identity_value(runtime, "run_id"),
            task_id=_runtime_identity_value(runtime, "task_id"),
            requested_slide_count=len(slides),
            status="planned",
            register=register,
            visual_policy=resolved_visual_policy,
            style_profile=style_profile or {},
            design_plan=None,
            deck_title=deck_title,
            output_path=output_path,
            slides=[],
            expected_visual_count=0,
            deck_stylesheet=(deck_stylesheet or "").strip() or None,
            deck_authoring_contract=(
                str(authoring_contract or "compact_model_html_v1")
                if (deck_stylesheet or "").strip()
                else "legacy_full_html_v1"
            ),
            deck_stylesheet_hash=(hashlib.sha256((deck_stylesheet or "").encode("utf-8")).hexdigest() if (deck_stylesheet or "").strip() else None),
            deck_route=DEFAULT_DECK_ROUTE,
            deck_compile_mode=DEFAULT_DECK_COMPILE_MODE,
            native_required=True,
            legacy_screenshot_debug=False,
            native_editability_score=0.0,
            created_at=now,
            updated_at=now,
        )
        try:
            if (
                deck.deck_authoring_contract == "compact_model_html_v2"
                and _state_value(runtime, "deck_candidate_compile") is not True
            ):
                slides, font_normalization = _normalize_compact_v2_inline_font_fallbacks(
                    slides,
                )
                if font_normalization["normalized_declaration_count"] > 0:
                    logger.info(
                        "[DeckIRNormalization] compact-v2 inline font fallbacks normalized "
                        "normalized_slides=%d normalized_attributes=%d normalized_declarations=%d "
                        "rawContentExcluded=true",
                        font_normalization["normalized_slide_count"],
                        font_normalization["normalized_attribute_count"],
                        font_normalization["normalized_declaration_count"],
                    )
                    with deck_span(
                        "deck.ir.normalize",
                        runtime=runtime,
                        build_id=build_id,
                        visual_policy=deck.visual_policy,
                        status=deck.status,
                        slide_count=len(slides),
                        run_type="tool",
                        inputs={
                            "authoring_contract": deck.deck_authoring_contract,
                            "normalization_kind": "inline_font_fallbacks",
                        },
                    ) as run:
                        finish_span(run, font_normalization)
                normalized_stylesheet, invariant_normalization = (
                    _normalize_compact_v2_anchor_invariant_contract(
                        deck.deck_stylesheet or "",
                        slides,
                    )
                )
                if invariant_normalization["normalized_anchor_rule_count"] > 0:
                    deck.deck_stylesheet = normalized_stylesheet
                    deck.deck_stylesheet_hash = hashlib.sha256(
                        normalized_stylesheet.encode("utf-8")
                    ).hexdigest()
                    logger.info(
                        "[DeckIRNormalization] compact-v2 anchor invariants completed "
                        "normalized_rules=%d injected_declarations=%d rawContentExcluded=true",
                        invariant_normalization["normalized_anchor_rule_count"],
                        invariant_normalization["injected_declaration_count"],
                    )
                    with deck_span(
                        "deck.ir.normalize",
                        runtime=runtime,
                        build_id=build_id,
                        visual_policy=deck.visual_policy,
                        status=deck.status,
                        slide_count=len(slides),
                        run_type="tool",
                        inputs={
                            "authoring_contract": deck.deck_authoring_contract,
                            "normalization_kind": "anchor_invariant_contract_completion",
                        },
                    ) as run:
                        finish_span(run, invariant_normalization)
                slides, normalization = _normalize_compact_v2_anchor_inline_geometry(
                    deck.deck_stylesheet or "",
                    slides,
                )
                if normalization["normalized_anchor_count"] > 0:
                    logger.info(
                        "[DeckIRNormalization] compact-v2 duplicate anchor geometry removed "
                        "normalized_slides=%d normalized_anchors=%d removed_declarations=%d "
                        "rawContentExcluded=true",
                        normalization["normalized_slide_count"],
                        normalization["normalized_anchor_count"],
                        normalization["removed_declaration_count"],
                    )
                    with deck_span(
                        "deck.ir.normalize",
                        runtime=runtime,
                        build_id=build_id,
                        visual_policy=deck.visual_policy,
                        status=deck.status,
                        slide_count=len(slides),
                        run_type="tool",
                        inputs={
                            "authoring_contract": deck.deck_authoring_contract,
                            "normalization_kind": "duplicate_anchor_inline_geometry",
                        },
                    ) as run:
                        finish_span(run, normalization)
            deadline_setter = getattr(self._native_service, "set_deadline_epoch_ms", None)
            if callable(deadline_setter):
                deadline_setter(_service_deadline_epoch_ms(runtime))
            _assert_deck_deadline(runtime, stage="input_validation")
            self._validate_inputs(deck, slides, output_path, runtime)
            deck.slides = self._build_slide_specs(
                slides,
                visual_policy=deck.visual_policy,
                runtime=runtime,
                style_profile=deck.style_profile,
            )
            with deck_span(
                "deck.ir.validate",
                runtime=runtime,
                build_id=build_id,
                visual_policy=deck.visual_policy,
                status=deck.status,
                slide_count=len(slides),
                inputs={
                    "slide_count": len(slides),
                    "register": register,
                    "visual_policy": deck.visual_policy,
                    "layout_kinds": [slide.layout_kind for slide in deck.slides],
                    "slide_roles": [slide.role for slide in deck.slides],
                },
            ) as run:
                finish_span(run, {"valid": True, "failure_code": None, "issue_count": 0, "hard_issue_count": 0})
            self._resolve_creative_plan_and_asset_policy(
                deck,
                runtime,
                source_design_plan=design_plan,
                source_creative_plan=creative_plan,
            )
            _assert_deck_deadline(runtime, stage="slide_html")
            self._render_slide_html(deck, runtime)
            if deck.expected_visual_count > 0:
                self._write_prompt_files(deck, runtime)
                self._prepare_manifest(deck, runtime)
                _assert_deck_deadline(runtime, stage="image_batch")
                with _deck_trace_runtime_context(runtime, deck):
                    summary = self._run_visual_batch(deck, runtime)
                    self._apply_batch_summary(deck, runtime, summary)
                    self._repair_visuals_after_batch(deck, runtime, summary)
                self._verify_visuals(deck, runtime)
            else:
                for slide in deck.slides:
                    slide.visual_required = False
                    slide.visual_status = "not_required"
                deck.image_generation_status = "not_required"
                deck.primary_image_batch_status = "not_required"
            _assert_deck_deadline(runtime, stage="source_quality_evaluation")
            source_evaluation = self._evaluate_source_quality(deck, runtime)
            _assert_deck_deadline(runtime, stage="native_compile")
            try:
                # Compile even when source quality has gaps so the single repair
                # receives both static and native/mechanical targets together.
                self._compile_pptx(
                    deck,
                    runtime,
                    native_lint_slide_indices=native_lint_slide_indices,
                )
            except DeckBuildFailure as exc:
                if exc.code == "deck_mechanical_gate_failed" and not source_evaluation.passed:
                    raise DeckBuildFailure(
                        exc.code,
                        _combined_source_and_mechanical_summary(source_evaluation, exc.summary),
                        retryable=exc.retryable,
                    ) from exc
                raise
            if not source_evaluation.passed:
                raise DeckBuildFailure(
                    "deck_source_quality_failed",
                    _quality_failure_summary(source_evaluation),
                    retryable=True,
                )
            _assert_deck_deadline(runtime, stage="deck_evaluation")
            self._evaluate(deck, runtime)
            deck.status = "evaluated"
            _finalize_image_generation_status(deck, success=True)
            self._assert_deck_success_allowed(deck, runtime)
            try:
                materialize_deck_foundation_safely(deck, runtime)
            except BuildFoundationPersistenceError as exc:
                raise DeckBuildFailure("build_manifest_persistence_failed", str(exc), retryable=False) from exc
            deck.updated_at = _now()
            deck_path = save_deck_build(deck, runtime)
            return self._success_result(
                deck,
                deck_path,
                runtime,
                success_allowed_checked=True,
                service_elapsed_ms=int((time.perf_counter() - service_started) * 1000),
            )
        except DeckBuildFailure as exc:
            deck.status = "failed_terminal"
            deck.failure_code = exc.code
            deck.failure_summary = exc.summary
            if deck.expected_visual_count > 0 and deck.slides:
                self._refresh_visual_counts(deck, runtime)
            _finalize_image_generation_status(deck, success=False)
            deck.updated_at = _now()
            deck_path = save_deck_build(deck, runtime)
            source_report = deck.source_quality_report if isinstance(deck.source_quality_report, dict) else {}
            mechanical_report = deck.mechanical_gate_results if isinstance(deck.mechanical_gate_results, dict) else {}
            native_report = deck.native_mechanical_report if isinstance(deck.native_mechanical_report, dict) else {}
            logger.warning(
                "[DeckBuildFailure] build_id=%s code=%s retryable=%s slides=%d "
                "source_passed=%s source_hard_failures=%d source_issues=%d "
                "mechanical_passed=%s mechanical_issues=%d native_passed=%s "
                "deck_build_path=%s elapsed_ms=%d summary=%s",
                deck.build_id,
                exc.code,
                exc.retryable,
                len(deck.slides),
                source_report.get("passed"),
                _report_count(source_report, "hard_failures", "hard_issues"),
                _report_count(source_report, "issues", "findings"),
                mechanical_report.get("passed"),
                _report_count(mechanical_report, "issues", "failures"),
                native_report.get("passed"),
                deck_path,
                int((time.perf_counter() - service_started) * 1000),
                safe_excerpt(exc.summary, limit=800),
            )
            self._trace_terminal(deck, runtime, success=False, deck_path=deck_path, retryable=exc.retryable)
            return self._failure_result(
                deck,
                deck_path,
                exc,
                runtime,
                service_elapsed_ms=int((time.perf_counter() - service_started) * 1000),
            )

    def _validate_inputs(self, deck: DeckBuild, slides: list[dict[str, Any]], output_path: str, runtime: ToolRuntime) -> None:
        if not 1 <= len(slides) <= 30:
            raise DeckBuildFailure("invalid_deck_ir", "Deck slide count must be between 1 and 30.", retryable=True)
        requested_slide_count = _requested_slide_count_from_state(runtime)
        if requested_slide_count and requested_slide_count != len(slides):
            raise DeckBuildFailure(
                "invalid_deck_ir",
                f"Deck IR has {len(slides)} slides, but the request targeted {requested_slide_count}.",
                retryable=True,
            )
        path_error = _ensure_relative_to_outputs("output_path", output_path)
        if path_error or not output_path.lower().endswith(".pptx"):
            raise DeckBuildFailure("invalid_deck_ir", path_error or "output_path must end with .pptx", retryable=True)
        if deck.visual_policy not in {"auto", "auto_with_images_allowed", "text_only"}:
            raise DeckBuildFailure(
                "invalid_deck_ir",
                "visual_policy must be auto, required, or text_only.",
                retryable=True,
            )
        if deck.visual_policy == "text_only" and not _explicit_text_only_requested(runtime):
            raise DeckBuildFailure(
                "invalid_deck_ir",
                "visual_policy='text_only' requires an explicit plain/text-only/no-visual request.",
                retryable=True,
            )
        _validate_authoring_inputs(
            deck,
            slides,
            allow_repair_overlay=_state_value(runtime, "deck_candidate_compile") is True,
        )

    def _build_slide_specs(
        self,
        slides: list[dict[str, Any]],
        *,
        visual_policy: str,
        runtime: ToolRuntime,
        style_profile: dict[str, Any],
    ) -> list[DeckSlideSpec]:
        specs: list[DeckSlideSpec] = []
        for index, raw in enumerate(slides, start=1):
            title = _clean_text(raw.get("title"))
            narrative = _clean_text(raw.get("narrative"))
            visual_prompt = _clean_text(raw.get("visual_prompt")) or None
            if not title or len(title) > 90:
                raise DeckBuildFailure("invalid_deck_ir", f"Slide {index} title is required and must be <= 90 chars.", retryable=True)
            if not narrative or len(narrative) > 280:
                raise DeckBuildFailure("invalid_deck_ir", f"Slide {index} narrative is required and must be <= 280 chars.", retryable=True)
            prompt_sanitized = False
            prompt_warning: str | None = None
            if visual_prompt:
                unrequested_style_terms = _unrequested_banned_style_terms(visual_prompt, runtime, style_profile)
                if _contains_unnegated_match(_BANNED_TEXT_RE, visual_prompt) or unrequested_style_terms:
                    visual_prompt = None
                    prompt_sanitized = True
                    prompt_warning = "visual_prompt_dropped_for_asset_policy_guardrail"
            authoring_sources = _slide_authoring_sources(raw)
            specs.append(
                DeckSlideSpec(
                    selector=f"slide:{index}",
                    index=index,
                    role=_clean_text(raw.get("role")) or ("cover" if index == 1 else "context"),
                    layout_kind=_clean_text(raw.get("layout_kind")) or ("cover_hero" if index == 1 else "single_visual_focus"),
                    title=title,
                    narrative=narrative,
                    claim=_clean_text(raw.get("claim")) or None,
                    visual_prompt=visual_prompt,
                    html_body=authoring_sources["html_body"],
                    slide_css=authoring_sources["slide_css"],
                    html_source=authoring_sources["html_source"],
                    speaker_notes=_clean_text(raw.get("speaker_notes")) or None,
                    visual_required=False,
                    gate_results={
                        "visual_prompt_sanitized": prompt_sanitized,
                        "style_warning": prompt_warning,
                    }
                    if prompt_sanitized
                    else {},
                )
            )
        return specs

    def _resolve_creative_plan_and_asset_policy(
        self,
        deck: DeckBuild,
        runtime: ToolRuntime,
        *,
        source_design_plan: dict[str, Any] | None,
        source_creative_plan: dict[str, Any] | None,
    ) -> None:
        request_context = _request_context_text(runtime)
        with deck_span(
            "deck.creative_plan.validate",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            run_type="tool",
            inputs={
                "register": deck.register,
                "style_profile_keys": sorted(str(key) for key in deck.style_profile.keys()),
                "source_design_plan_present": bool(source_design_plan),
                "creative_plan_present": bool(source_creative_plan),
            },
        ) as run:
            try:
                creative = normalize_creative_plan(
                    source_creative_plan,
                    deck=deck,
                    request_context=request_context,
                    source_design_plan=source_design_plan,
                )
            except CreativePlanValidationError as exc:
                finish_span(
                    run,
                    {
                        "valid": False,
                        "failure_code": exc.code,
                        "failure_summary": safe_excerpt(exc.summary),
                    },
                )
                raise DeckBuildFailure(exc.code, exc.summary, retryable=True) from exc
            deck.creative_plan = creative
            deck.design_plan = creative.design_plan
            creative_path = _host_path(_CREATIVE_PLAN, runtime)
            write_creative_plan(creative, creative_path)
            deck.creative_plan_path = _CREATIVE_PLAN
            design_path = _host_path(_DESIGN_PLAN, runtime)
            write_design_plan(deck.design_plan, design_path)
            deck.design_plan_path = _DESIGN_PLAN
            finish_span(
                run,
                {
                    "valid": True,
                    "creative_plan_file": basename(deck.creative_plan_path),
                    "image_strategy": creative.image_strategy,
                    "planned_image_asset_count": len(creative.image_assets),
                    "slide_composition_count": len(creative.slide_compositions),
                    "style_lane": deck.design_plan.style_lane,
                    "palette": [token.name for token in deck.design_plan.palette],
                    "slide_width_px": deck.design_plan.grid.slide_width_px,
                    "slide_height_px": deck.design_plan.grid.slide_height_px,
                    "design_plan_file": basename(deck.design_plan_path),
                },
            )
        with deck_span(
            "deck.asset_policy.resolve",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            run_type="tool",
            inputs={"slide_count": len(deck.slides)},
        ) as run:
            try:
                apply_creative_asset_plan(deck, deck.creative_plan)
            except CreativePlanValidationError as exc:
                finish_span(
                    run,
                    {
                        "success": False,
                        "failure_code": exc.code,
                        "failure_summary": safe_excerpt(exc.summary),
                    },
                )
                raise DeckBuildFailure(exc.code, exc.summary, retryable=True) from exc
            asset_policy_path = _host_path(_ASSET_POLICY, runtime)
            write_asset_policy(deck, asset_policy_path)
            deck.asset_policy_path = _ASSET_POLICY
            prompt_warnings = [str(slide.gate_results.get("style_warning")) for slide in deck.slides if slide.gate_results.get("style_warning")]
            deck.style_warnings = sorted(set([*deck.style_warnings, *prompt_warnings]))
            finish_span(
                run,
                {
                    "success": True,
                    "expected_visual_count": deck.expected_visual_count,
                    "generated_asset_count": deck.generated_asset_count,
                    "native_html_slide_count": deck.native_html_slide_count,
                    "hybrid_slide_count": deck.hybrid_slide_count,
                    "text_only_slide_count": deck.text_only_slide_count,
                    "asset_policy_file": basename(deck.asset_policy_path),
                    "visual_modes": [slide.asset_plan.visual_mode if slide.asset_plan else None for slide in deck.slides],
                    "layout_families": [getattr(slide.composition_plan, "layout_name", None) for slide in deck.slides],
                    "style_warning_count": len(deck.style_warnings),
                },
            )

    def _write_prompt_files(self, deck: DeckBuild, runtime: ToolRuntime) -> None:
        for slide in deck.slides:
            slide.visual_status = "not_required"
            slide.visual_required = bool(slide.asset_plan and slide.asset_plan.image_gen_required)
        for slide in generated_asset_slides(deck):
            prompt_path = f"{_PROMPTS}/slide-{slide.index:02d}.json"
            host = _host_path(prompt_path, runtime)
            host.parent.mkdir(parents=True, exist_ok=True)
            payload = deck_asset_prompt_payload(slide, deck)
            host.write_text(json.dumps(payload, indent=2), encoding="utf-8")
            slide.visual_prompt_path = prompt_path
            slide.visual_asset_path = f"{_ASSETS}/slide-{slide.index:02d}.png"
            slide.visual_status = "pending"
        deck.status = "visual_specs_ready"

    def _prepare_manifest(self, deck: DeckBuild, runtime: ToolRuntime) -> None:
        generated_slides = generated_asset_slides(deck)
        prompt_files = [slide.visual_prompt_path or "" for slide in generated_slides]
        with deck_span(
            "deck.image_manifest.prepare",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            run_type="tool",
            inputs={"prompt_count": len(prompt_files), "manifest_file": basename(_MANIFEST)},
        ) as run:
            result = create_pptx_image_manifest_core(
                thread_data=get_thread_data(runtime),
                prompt_files=prompt_files,
                manifest_path=_MANIFEST,
                manifest_author="DeckBuildService",
                trace=False,
            )
            raw_items = result.get("items") if isinstance(result.get("items"), list) else []
            items = [item for item in raw_items if isinstance(item, dict)]
            finish_span(
                run,
                {
                    "success": bool(result.get("success")),
                    "manifest_author": "DeckBuildService",
                    "item_count": result.get("expected_count", 0),
                    "prompt_count": len(prompt_files),
                    "generated_asset_slide_indices": [slide.index for slide in generated_slides],
                    "prompt_basenames": [basename(item.get("prompt_file")) for item in items],
                    "prompt_hashes": [item.get("prompt_hash") for item in items if item.get("prompt_hash")],
                    "prompt_chars_total": sum(int(item.get("prompt_chars") or 0) for item in items),
                    "output_basenames": [basename(item.get("output_path")) for item in items],
                    "schema_version": result.get("schema_version"),
                },
            )
            if not result.get("success"):
                raise DeckBuildFailure(str(result.get("error_type") or "invalid_deck_ir"), str(result.get("error") or "manifest failed"), retryable=True)
            for slide, item in zip(generated_slides, result.get("items", []), strict=True):
                slide.visual_asset_path = item.get("output_path")
            self._clear_expected_visual_outputs(deck, runtime)

    def _run_visual_batch(self, deck: DeckBuild, runtime: ToolRuntime) -> dict[str, Any]:
        deck.status = "visual_batch_running"
        started = time.perf_counter()
        with deck_span(
            "deck.image_batch.run",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            run_type="tool",
            inputs={
                "requested": deck.expected_visual_count,
                "concurrency": os.getenv("SOPHIA_IMAGE_GEN_CONCURRENCY", "2"),
                "sdk_timeout_seconds": os.getenv("SOPHIA_IMAGE_GEN_TIMEOUT", "120"),
                "sdk_max_retries": os.getenv("SOPHIA_IMAGE_GEN_MAX_RETRIES", "1"),
                "manifest_file": basename(_MANIFEST),
            },
        ) as run:
            summary = self._image_batch_runner(_MANIFEST, runtime)
            summary.setdefault("duration_ms", int((time.perf_counter() - started) * 1000))
            if not summary.get("summary_present", True):
                summary = self._reconcile_missing_batch_summary(deck, runtime, summary)
            self._record_primary_batch_state(deck, summary)
            finish_span(run, _safe_batch_summary(summary))
            if str(summary.get("error_class") or "") == "deck_deadline_exceeded":
                raise DeckBuildFailure(
                    "deck_deadline_exceeded",
                    "The shared builder deadline expired during image generation.",
                    retryable=False,
                )
            if not summary.get("summary_present", True) and not summary.get("batch_attempted"):
                raise DeckBuildFailure("deck_visual_batch_startup_failed", "Image batch did not emit IMAGEGEN_BATCH.", retryable=False)
            if not summary.get("complete"):
                if summary.get("requested") is None:
                    summary["requested"] = deck.expected_visual_count
                return summary
            return summary

    def _clear_expected_visual_outputs(self, deck: DeckBuild, runtime: ToolRuntime) -> None:
        for slide in generated_asset_slides(deck):
            if not slide.visual_asset_path:
                continue
            try:
                host = _host_path(slide.visual_asset_path, runtime)
                if host.is_file():
                    host.unlink()
            except OSError:
                continue

    def _record_primary_batch_state(self, deck: DeckBuild, summary: dict[str, Any]) -> None:
        error_class = summary.get("error_class")
        if not error_class:
            histogram = summary.get("error_class_histogram")
            if isinstance(histogram, dict) and histogram:
                error_class = next(iter(histogram))
        if error_class:
            deck.primary_image_batch_error_class = str(error_class)
            deck.image_generation_reason = str(error_class)
        if str(error_class or "") == "timeout":
            deck.batch_timeout_count += 1
        if summary.get("complete"):
            deck.primary_image_batch_status = "success"
        elif summary.get("batch_attempted") or summary.get("summary_present", True):
            deck.primary_image_batch_status = "partial" if int(summary.get("images_generated") or 0) > 0 else "failed"
        else:
            deck.primary_image_batch_status = "failed"
        if summary.get("partial_batch_salvaged"):
            deck.partial_batch_salvaged = True

    def _reconcile_missing_batch_summary(
        self,
        deck: DeckBuild,
        runtime: ToolRuntime,
        summary: dict[str, Any],
    ) -> dict[str, Any]:
        progress_items = summary.get("items") if isinstance(summary.get("items"), list) else []
        error_class = str(summary.get("error_class") or "batch_summary_missing")
        batch_attempted = error_class == "timeout" or bool(progress_items)
        if not batch_attempted:
            return {**summary, "batch_attempted": False}
        progress_by_output = {str(item.get("output_file")): item for item in progress_items if isinstance(item, dict) and item.get("output_file")}
        items: list[dict[str, Any]] = []
        succeeded = 0
        histogram: dict[str, int] = {}
        with deck_span(
            "deck.image_batch.timeout_reconcile",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            run_type="tool",
            inputs={
                "summary_error_class": error_class,
                "progress_item_count": len(progress_items),
                "expected_visual_count": deck.expected_visual_count,
            },
        ) as run:
            for slide in generated_asset_slides(deck):
                output_file = str(slide.visual_asset_path or "")
                progress = progress_by_output.get(output_file, {})
                output_bytes = _file_size(_host_path(output_file, runtime)) if output_file else None
                exists = bool(output_bytes and output_bytes > 0)
                success = exists or (bool(progress.get("success")) and exists)
                item_error_class = None if success else str(progress.get("error_class") or error_class or "timeout")
                if success:
                    succeeded += 1
                else:
                    histogram[item_error_class or "unknown"] = histogram.get(item_error_class or "unknown", 0) + 1
                item = {
                    "item_index": slide.index,
                    "output_file": output_file,
                    "success": success,
                    "bytes": output_bytes or 0,
                    "duration_ms": progress.get("duration_ms"),
                    "error_class": item_error_class,
                }
                raw_error = progress.get("raw_error_excerpt") or summary.get("raw_error_excerpt")
                if raw_error and not success:
                    item["raw_error_excerpt"] = safe_excerpt(raw_error)
                items.append({key: value for key, value in item.items() if value not in (None, "")})
            reconciled = {
                **summary,
                "summary_present": False,
                "batch_attempted": True,
                "requested": deck.expected_visual_count,
                "images_generated": succeeded,
                "failed": max(0, deck.expected_visual_count - succeeded),
                "complete": succeeded == deck.expected_visual_count,
                "items": items,
                "error_class": error_class,
                "error_class_histogram": histogram,
                "partial_batch_salvaged": succeeded > 0,
            }
            finish_span(
                run,
                {
                    "batch_attempted": True,
                    "partial_batch_salvaged": succeeded > 0,
                    "successful_visual_count": succeeded,
                    "missing_visual_count": max(0, deck.expected_visual_count - succeeded),
                    "error_class": error_class,
                },
            )
            return reconciled

    def _apply_batch_summary(self, deck: DeckBuild, runtime: ToolRuntime, summary: dict[str, Any]) -> None:
        items = summary.get("items") if isinstance(summary.get("items"), list) else []
        by_output = {str(item.get("output_file")): item for item in items if isinstance(item, dict)}
        for slide in generated_asset_slides(deck):
            item = by_output.get(str(slide.visual_asset_path or ""))
            if item and item.get("success"):
                slide.visual_status = "generated"
            else:
                slide.visual_status = "failed"
                slide.visual_error_class = str((item or {}).get("error_class") or summary.get("error_class") or "image_generation_failed")
            with deck_span(
                "deck.image_batch.item",
                runtime=runtime,
                build_id=deck.build_id,
                visual_policy=deck.visual_policy,
                status=deck.status,
                slide_count=len(deck.slides),
                run_type="tool",
                inputs={
                    "item_index": slide.index,
                    "selector": slide.selector,
                    "prompt_file": basename(slide.visual_prompt_path),
                    "prompt_hash": _file_hash(_host_path(slide.visual_prompt_path or "", runtime)),
                    "output_file": basename(slide.visual_asset_path),
                    "deck_asset": True,
                    "slide_visual": False,
                },
            ) as run:
                finish_span(
                    run,
                    {
                        "success": slide.visual_status == "generated",
                        "error_class": slide.visual_error_class,
                        "output_bytes": _file_size(_host_path(slide.visual_asset_path or "", runtime)),
                    },
                )

    def _repair_visuals_after_batch(self, deck: DeckBuild, runtime: ToolRuntime, summary: dict[str, Any]) -> None:
        missing = self._slides_needing_repair(deck, runtime)
        if not missing:
            return
        if not _batch_summary_allows_serial_repair(summary):
            deck.image_generation_reason = _primary_error_class(summary) or deck.image_generation_reason
            return
        attempts = 0
        repaired = 0
        stopped_error_class: str | None = None
        with deck_span(
            "deck.image_repair.run",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            run_type="tool",
            inputs={
                "missing_visual_count": len(missing),
                "max_attempts_per_visual": _SERIAL_REPAIR_MAX_ATTEMPTS,
                "batch_error_class": _primary_error_class(summary),
            },
        ) as run:
            for slide in missing:
                for attempt_no in range(1, _SERIAL_REPAIR_MAX_ATTEMPTS + 1):
                    _assert_deck_deadline(
                        runtime,
                        stage="serial_image_repair",
                    )
                    attempts += 1
                    result = self._image_single_runner(slide, runtime, attempt_no)
                    result_error_class = str(result.get("error_class") or "")
                    if result_error_class == "deck_deadline_exceeded":
                        raise DeckBuildFailure(
                            "deck_deadline_exceeded",
                            "The shared builder deadline expired during serial image repair.",
                            retryable=False,
                        )
                    output_exists = bool(slide.visual_asset_path and (_file_size(_host_path(slide.visual_asset_path, runtime)) or 0) > 0)
                    success = bool(result.get("success")) and output_exists
                    deck.serial_repair_count += 1
                    with deck_span(
                        "deck.image_repair.item",
                        runtime=runtime,
                        build_id=deck.build_id,
                        visual_policy=deck.visual_policy,
                        status=deck.status,
                        slide_count=len(deck.slides),
                        run_type="tool",
                        inputs={
                            "mode": "serial_repair",
                            "repair_attempt_no": attempt_no,
                            "item_index": slide.index,
                            "selector": slide.selector,
                            "prompt_file": basename(slide.visual_prompt_path),
                            "prompt_hash": _file_hash(_host_path(slide.visual_prompt_path or "", runtime)),
                            "output_file": basename(slide.visual_asset_path),
                            "allowed_manifest_output": True,
                        },
                    ) as item_run:
                        finish_span(
                            item_run,
                            {
                                "success": success,
                                "error_class": None if success else result_error_class or "image_generation_failed",
                                "exit_code": result.get("exit_code"),
                                "duration_ms": result.get("duration_ms"),
                                "output_bytes": (_file_size(_host_path(slide.visual_asset_path, runtime)) if slide.visual_asset_path else None),
                                "raw_error_excerpt": safe_excerpt(result.get("raw_error_excerpt")),
                            },
                        )
                    if success:
                        slide.visual_status = "generated"
                        slide.visual_error_class = None
                        repaired += 1
                        break
                    slide.visual_status = "failed"
                    slide.visual_error_class = result_error_class or "image_generation_failed"
                    deck.image_generation_reason = slide.visual_error_class
                    if slide.visual_error_class in _NON_REPAIRABLE_IMAGE_ERROR_CLASSES:
                        stopped_error_class = slide.visual_error_class
                        break
                if stopped_error_class:
                    break
            successful, remaining_missing = self._refresh_visual_counts(deck, runtime)
            if attempts and remaining_missing == 0:
                deck.primary_image_batch_status = "repaired"
            finish_span(
                run,
                {
                    "attempt_count": attempts,
                    "repaired_visual_count": repaired,
                    "successful_visual_count": successful,
                    "remaining_missing_visual_count": remaining_missing,
                    "stopped_error_class": stopped_error_class,
                },
            )

    def _slides_needing_repair(self, deck: DeckBuild, runtime: ToolRuntime) -> list[DeckSlideSpec]:
        slides: list[DeckSlideSpec] = []
        for slide in generated_asset_slides(deck):
            output_exists = bool(slide.visual_asset_path and (_file_size(_host_path(slide.visual_asset_path, runtime)) or 0) > 0)
            if slide.visual_status != "generated" or not output_exists:
                slides.append(slide)
        return slides

    def _refresh_visual_counts(self, deck: DeckBuild, runtime: ToolRuntime) -> tuple[int, int]:
        successful = 0
        missing = 0
        for slide in generated_asset_slides(deck):
            exists = bool(slide.visual_asset_path and (_file_size(_host_path(slide.visual_asset_path, runtime)) or 0) > 0)
            if slide.visual_status == "generated" and exists:
                successful += 1
            else:
                missing += 1
        deck.successful_visual_count = successful
        deck.missing_visual_count = missing
        return successful, missing

    def _verify_visuals(self, deck: DeckBuild, runtime: ToolRuntime) -> None:
        successful, missing = self._refresh_visual_counts(deck, runtime)
        with deck_span(
            "deck.visuals.verify",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            inputs={},
        ) as run:
            complete = deck.expected_visual_count == successful and missing == 0
            finish_span(
                run,
                {
                    "visual_policy": deck.visual_policy,
                    "expected_visual_count": deck.expected_visual_count,
                    "successful_visual_count": successful,
                    "referenced_visual_count": deck.referenced_visual_count,
                    "missing_visual_count": missing,
                    "complete": complete,
                },
            )
        if deck.expected_visual_count != successful or missing:
            raise DeckBuildFailure("deck_visuals_incomplete", "Required slide visuals are incomplete.", retryable=False)
        deck.status = "visuals_complete"

    def _render_slide_html(self, deck: DeckBuild, runtime: ToolRuntime) -> None:
        with deck_span(
            "deck.slide_html.validate",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            run_type="tool",
            inputs={
                "slide_count": len(deck.slides),
                "planned_asset_count": deck.expected_visual_count,
                "canvas_width_px": 1920,
                "canvas_height_px": 1080,
            },
        ) as run:
            _clear_slide_html_directory(_host_path(_SLIDES, runtime))
            validations = []
            allowed_asset_refs = planned_asset_ref_basenames(deck)
            for slide in deck.slides:
                if not slide.html_source and slide.html_body and deck.deck_stylesheet:
                    slide.html_source = assemble_compact_slide_html(
                        deck_stylesheet=deck.deck_stylesheet,
                        html_body=slide.html_body,
                        slide_css=slide.slide_css,
                    )
                    deck.deck_html_fragment_count += 1
                if not slide.html_source:
                    summary = f"Slide {slide.index} is missing html_body."
                    finish_span(run, {"valid": False, "failure_code": "deck_slide_html_missing", "failure_summary": summary})
                    raise DeckBuildFailure("deck_slide_html_missing", summary, retryable=True)
                virtual = _slide_html_virtual_path(slide)
                host = _host_path(virtual, runtime)
                sanitized, validation = validate_and_sanitize_slide_html(
                    slide,
                    allowed_asset_refs=allowed_asset_refs,
                )
                validations.append(validation)
                if not validation.valid:
                    summary = "; ".join(validation.errors[:3]) or f"Slide {slide.index} html_source is invalid."
                    deck.html_source_validation = validation_summary(validations)
                    finish_span(
                        run,
                        {
                            **deck.html_source_validation,
                            "failure_code": "deck_slide_html_invalid",
                            "failure_summary": safe_excerpt(summary),
                        },
                    )
                    raise DeckBuildFailure("deck_slide_html_invalid", summary, retryable=True)
                host.parent.mkdir(parents=True, exist_ok=True)
                host.write_text(sanitized, encoding="utf-8")
                slide.html_source = sanitized
                deck.deck_assembled_html_bytes += len(sanitized.encode("utf-8"))
                slide.html_source_path = virtual
                slide.gate_results["chrome_detected"] = False
                slide.gate_results["html_source_validation"] = validation.to_dict()
            deck.html_source_validation = validation_summary(validations)
            deck.referenced_visual_count = _referenced_planned_asset_count(deck, validations)
            if deck.expected_visual_count > deck.referenced_visual_count:
                missing = deck.expected_visual_count - deck.referenced_visual_count
                summary = f"Deck creative plan declares {deck.expected_visual_count} generated image asset(s), but slide HTML references only {deck.referenced_visual_count} planned asset(s)."
                finish_span(
                    run,
                    {
                        **deck.html_source_validation,
                        "failure_code": "deck_slide_html_invalid",
                        "failure_summary": safe_excerpt(summary),
                        "expected_visual_count": deck.expected_visual_count,
                        "referenced_visual_count": deck.referenced_visual_count,
                        "missing_expected_visual_count": missing,
                    },
                )
                raise DeckBuildFailure("deck_slide_html_invalid", summary, retryable=True)
            deck.status = "slides_rendered"
            finish_span(
                run,
                {
                    **deck.html_source_validation,
                    "html_basenames": [basename(slide.html_source_path) for slide in deck.slides],
                    "selectors": [slide.selector for slide in deck.slides],
                    "style_lane": deck.design_plan.style_lane if hasattr(deck.design_plan, "style_lane") else None,
                    "layout_families": [getattr(slide.composition_plan, "layout_name", None) for slide in deck.slides],
                    "generated_asset_count": deck.generated_asset_count,
                    "native_html_slide_count": deck.native_html_slide_count,
                    "hybrid_slide_count": deck.hybrid_slide_count,
                    "text_only_slide_count": deck.text_only_slide_count,
                },
            )

    def _compile_pptx(
        self,
        deck: DeckBuild,
        runtime: ToolRuntime,
        *,
        native_lint_slide_indices: tuple[int, ...] | None = None,
    ) -> None:
        preflight = self._trace_native_requirement(deck, runtime)
        if not preflight.success:
            deck.deck_compile_mode = NATIVE_UNAVAILABLE_DECK_COMPILE_MODE
            raise DeckBuildFailure(
                "deck_native_unavailable",
                _native_error_summary(preflight.errors, "Native deck service is unavailable."),
                retryable=False,
            )
        try:
            self._compile_native_pptx(
                deck,
                runtime,
                native_lint_slide_indices=native_lint_slide_indices,
            )
        except DeckNativePathError as exc:
            deck.deck_compile_mode = NATIVE_UNAVAILABLE_DECK_COMPILE_MODE
            raise DeckBuildFailure(
                "deck_native_unavailable",
                safe_excerpt(exc) or "Native deck service is unavailable.",
                retryable=False,
            ) from exc

    def _trace_native_requirement(self, deck: DeckBuild, runtime: ToolRuntime):
        deck.native_required = True
        deck.legacy_screenshot_debug = _legacy_screenshot_debug_allowed(runtime)
        preflight = _native_preflight(self._native_service)
        with deck_span(
            "deck.native.requirement",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            run_type="tool",
            deck_compile_mode=deck.deck_compile_mode,
            inputs={},
        ) as run:
            finish_span(
                run,
                {
                    "native_required": True,
                    "legacy_screenshot_debug_allowed": deck.legacy_screenshot_debug,
                    "preflight_success": preflight.success,
                    "scripts_dir_exists": preflight.scripts_dir_exists,
                    "deck_py_exists": preflight.deck_py_exists,
                    "html2patch_py_exists": preflight.html2patch_py_exists,
                    "error_count": len(preflight.errors),
                    "deck_compile_mode_before": deck.deck_compile_mode,
                },
            )
        return preflight

    def _compile_screenshot_debug_pptx(self, deck: DeckBuild, runtime: ToolRuntime) -> None:
        deck.deck_compile_mode = HTML_SCREENSHOT_DEBUG_COMPILE_MODE
        deck.legacy_screenshot_debug = True
        deck.native_editability_score = 0.0
        deck.native_text_shape_count = 0
        deck.picture_shape_count = 0
        deck.full_slide_picture_count = len(deck.slides)
        deck.quality_warning = _merge_warning(deck.quality_warning, "screenshot_deck_debug_only")
        with deck_span(
            "deck.pptx.compile",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            run_type="tool",
            inputs={"slide_count": len(deck.slides), "output_file": basename(deck.output_path)},
        ) as run:
            result = self._deck_compiler(
                runtime,
                deck.output_path,
                deck.deck_title,
                _SLIDES,
            )
            finish_span(
                run,
                {
                    "success": bool(result.get("success")),
                    "deck_route": deck.deck_route,
                    "deck_compile_mode": deck.deck_compile_mode,
                    "artifact_target_ext": DEFAULT_ARTIFACT_TARGET_EXT,
                    "pptx_file": basename(result.get("pptx_path")),
                    "size_bytes": result.get("size_bytes"),
                    "engine": result.get("engine"),
                    "overflow_slide_count": len(result.get("overflow_slides") or []),
                },
            )
            if not result.get("success"):
                raise DeckBuildFailure(str(result.get("error_type") or "deck_compile_failed"), str(result.get("error") or "PPTX compile failed."), retryable=False)
            deck.pptx_path = result.get("pptx_path")
            deck.compile_overflow_slides = [entry for entry in (result.get("overflow_slides") or []) if isinstance(entry, dict)]
            deck.status = "compiled"

    def _compile_native_pptx(
        self,
        deck: DeckBuild,
        runtime: ToolRuntime,
        *,
        native_lint_slide_indices: tuple[int, ...] | None = None,
    ) -> None:
        deck.deck_compile_mode = NATIVE_DECK_COMPILE_MODE
        base_host = _host_path(_NATIVE_BASE, runtime)
        patch_host = _host_path(_NATIVE_PATCH, runtime)
        output_host = _host_path(deck.output_path, runtime)
        render_host = _host_path(_NATIVE_RENDER_DIR, runtime)
        html_hosts = [_host_path(slide.html_source_path or "", runtime) for slide in deck.slides]
        _write_native_base_deck(base_host)
        _assert_deck_deadline(runtime, stage="native_html2patch")
        with deck_span(
            "deck.native.html2patch",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            run_type="tool",
            deck_compile_mode=deck.deck_compile_mode,
            inputs={"html_count": len(html_hosts), "base_file": basename(str(base_host)), "patch_file": basename(str(patch_host))},
        ) as run:
            html2patch = self._native_service.html_to_patch(
                html_paths=[str(path) for path in html_hosts],
                base_deck_path=str(base_host),
                output_patch_path=str(patch_host),
            )
            finish_span(run, _native_patch_span_outputs(html2patch))
        if not html2patch.success:
            if _native_deadline_error(html2patch.errors):
                raise DeckBuildFailure("deck_deadline_exceeded", _native_error_summary(html2patch.errors, "Deck deadline exceeded."), retryable=False)
            code = "deck_native_startup_failed" if _native_startup_error(html2patch.errors) else "deck_native_html2patch_failed"
            raise DeckBuildFailure(code, _native_error_summary(html2patch.errors, "Native html2patch failed."), retryable=False)
        deck.source_element_map = _load_json_dict(html2patch.source_map_path)
        _assert_deck_deadline(runtime, stage="native_patch_apply")
        with deck_span(
            "deck.native.patch_apply",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            run_type="tool",
            deck_compile_mode=deck.deck_compile_mode,
            inputs={"patch_file": basename(html2patch.patch_path), "output_file": basename(deck.output_path), "fix": False},
        ) as run:
            applied = self._native_service.apply_patch(
                base_deck_path=str(base_host),
                patch_path=str(patch_host),
                output_path=str(output_host),
                # Lint/fix runs as the next explicit stage. Keeping apply
                # mutation-free makes its issue/fix telemetry honest and
                # prevents the first repair pass from disappearing before the
                # mechanical report is assembled.
                fix=False,
            )
            finish_span(run, _native_patch_span_outputs(applied))
        if not applied.success:
            if _native_deadline_error(applied.errors):
                raise DeckBuildFailure("deck_deadline_exceeded", _native_error_summary(applied.errors, "Deck deadline exceeded."), retryable=False)
            code = "deck_native_patch_validation_failed" if applied.validation_error_count else "deck_native_patch_apply_failed"
            raise DeckBuildFailure(code, _native_error_summary(applied.errors, "Native patch apply failed."), retryable=False)
        _assert_deck_deadline(runtime, stage="native_inspect")
        with deck_span(
            "deck.native.inspect",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            run_type="tool",
            deck_compile_mode=deck.deck_compile_mode,
            inputs={"pptx_file": basename(deck.output_path)},
        ) as run:
            inspected = self._native_service.inspect(str(output_host))
            finish_span(run, _native_inspect_span_outputs(inspected))
        if not inspected.success:
            if _native_deadline_error(inspected.errors):
                raise DeckBuildFailure("deck_deadline_exceeded", _native_error_summary(inspected.errors, "Deck deadline exceeded."), retryable=False)
            raise DeckBuildFailure("deck_native_inspect_failed", _native_error_summary(inspected.errors, "Native inspect failed."), retryable=False)
        self._record_native_inspect(deck, inspected)
        if (deck.native_editability_score or 0.0) < 0.60:
            raise DeckBuildFailure(
                "deck_native_editability_failed",
                f"Native editability score {deck.native_editability_score:.2f} is below the 0.60 D1 gate.",
                retryable=False,
            )
        all_slide_indices = [slide.index - 1 for slide in deck.slides]
        if native_lint_slide_indices is None:
            lint_slide_indices = all_slide_indices
        else:
            lint_slide_indices = list(native_lint_slide_indices)
            if (
                not lint_slide_indices
                or len(set(lint_slide_indices)) != len(lint_slide_indices)
                or any(
                    isinstance(index, bool)
                    or not isinstance(index, int)
                    or index not in all_slide_indices
                    for index in lint_slide_indices
                )
            ):
                raise DeckBuildFailure(
                    "invalid_deck_ir",
                    "Native lint slide scope is invalid.",
                    retryable=False,
                )
            lint_slide_indices.sort()
        _assert_deck_deadline(runtime, stage="native_lint_fix")
        with deck_span(
            "deck.native.lint_fix",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            run_type="tool",
            deck_compile_mode=deck.deck_compile_mode,
            inputs={"touched_slide_count": len(lint_slide_indices)},
        ) as run:
            lint_fix = self._native_service.lint_fix(
                pptx_path=str(output_host),
                touched_slides=lint_slide_indices,
            )
            finish_span(run, _native_lint_fix_span_outputs(lint_fix))
        if not lint_fix.success:
            if _native_deadline_error(lint_fix.errors):
                raise DeckBuildFailure("deck_deadline_exceeded", _native_error_summary(lint_fix.errors, "Deck deadline exceeded."), retryable=False)
            raise DeckBuildFailure("deck_native_lint_fix_failed", _native_error_summary(lint_fix.errors, "Native lint/fix failed."), retryable=False)
        if lint_fix.residue_count:
            deck.quality_warning = _merge_warning(deck.quality_warning, "native_lint_residue")
        for slide in deck.slides:
            slide.gate_results["native_editability_score"] = deck.native_editability_score
            slide.gate_results["lint_residue_count"] = lint_fix.residue_count
        final_inspected = self._inspect_final_native_pptx(
            deck=deck,
            runtime=runtime,
            output_host=output_host,
            fix_applied_count=lint_fix.fix_applied_count,
        )
        with deck_span(
            "deck.source_retention.evaluate",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            run_type="tool",
            deck_compile_mode=deck.deck_compile_mode,
            inputs={"source_map_slide_count": len(deck.source_element_map.get("slides") or {})},
        ) as run:
            retention_reports = evaluate_source_retention(
                slides=deck.slides,
                native_shape_inventory=deck.native_shape_inventory,
                source_element_map=deck.source_element_map,
            )
            deck.source_retention_report = retention_summary(retention_reports)
            finish_span(
                run,
                {
                    "passed": deck.source_retention_report.get("passed"),
                    "missing_required_count": deck.source_retention_report.get("missing_required_count"),
                    "duplicate_source_id_count": deck.source_retention_report.get("duplicate_source_id_count"),
                    "low_retention_count": len(deck.source_retention_report.get("low_retention") or []),
                },
            )
        self._evaluate_final_native_contrast(deck=deck, runtime=runtime, output_host=output_host)
        _assert_deck_deadline(runtime, stage="native_render")
        with deck_span(
            "deck.native.render",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            run_type="tool",
            deck_compile_mode=deck.deck_compile_mode,
            inputs={
                "pptx_file": basename(deck.output_path),
                "slide_count": len(all_slide_indices),
            },
        ) as run:
            rendered = self._native_service.render(
                pptx_path=str(output_host),
                output_dir=str(render_host),
                slides=all_slide_indices,
            )
            finish_span(run, _native_render_span_outputs(rendered))
        if not rendered.success:
            if _native_deadline_error(rendered.errors):
                raise DeckBuildFailure("deck_deadline_exceeded", _native_error_summary(rendered.errors, "Deck deadline exceeded."), retryable=False)
            raise DeckBuildFailure("deck_native_render_failed", _native_error_summary(rendered.errors, "Native render failed."), retryable=False)
        _assert_deck_deadline(runtime, stage="native_diff")
        with deck_span(
            "deck.native.diff",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            run_type="tool",
            deck_compile_mode=deck.deck_compile_mode,
            inputs={"before_file": basename(str(base_host)), "after_file": basename(deck.output_path)},
        ) as run:
            diff = self._native_service.diff(before_path=str(base_host), after_path=str(output_host))
            finish_span(run, {"success": bool(diff.get("success")), "changed": bool(diff.get("changed")), "error_count": len(diff.get("errors") or [])})
        if not diff.get("success"):
            if _native_deadline_error([str(item) for item in diff.get("errors") or []]):
                raise DeckBuildFailure("deck_deadline_exceeded", "The shared builder deadline expired during native diff.", retryable=False)
            deck.quality_warning = _merge_warning(deck.quality_warning, "native_diff_unavailable")
        deck.native_mechanical_report = native_mechanical_report(
            inspect=final_inspected,
            lint_fix=lint_fix,
            render=rendered,
            diff=diff,
        )
        deck.native_mechanical_report["source_retention"] = deck.source_retention_report
        deck.native_mechanical_report["contrast"] = deck.native_contrast_report
        with deck_span(
            "deck.native.mechanical_report",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            run_type="tool",
            deck_compile_mode=deck.deck_compile_mode,
            inputs={},
        ) as run:
            finish_span(run, deck.native_mechanical_report)
        with deck_span(
            "deck.mechanical_gates.evaluate",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            run_type="tool",
            deck_compile_mode=deck.deck_compile_mode,
            inputs={
                "render_dir": basename(str(render_host)),
                "native_text_shape_count": deck.native_text_shape_count,
                "native_editability_score": deck.native_editability_score,
            },
        ) as run:
            gate_result = evaluate_mechanical_gates(
                deck,
                rendered_dir=render_host,
                native_pptx_path=output_host,
            )
            deck.mechanical_gate_results = gate_result.to_dict()
            finish_span(run, deck.mechanical_gate_results)
        if not deck.mechanical_gate_results.get("passed"):
            raise DeckBuildFailure(
                str(deck.mechanical_gate_results.get("failure_code") or "deck_mechanical_gate_failed"),
                str(deck.mechanical_gate_results.get("failure_summary") or "Deck mechanical gates failed."),
                retryable=True,
            )
        deck.pptx_path = deck.output_path
        deck.compile_overflow_slides = []
        deck.status = "compiled"

    def _inspect_final_native_pptx(
        self,
        *,
        deck: DeckBuild,
        runtime: ToolRuntime,
        output_host: Path,
        fix_applied_count: int,
    ) -> NativeDeckInspectResult:
        _assert_deck_deadline(runtime, stage="native_final_inspect")
        with deck_span(
            "deck.native.inspect_final",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            run_type="tool",
            deck_compile_mode=deck.deck_compile_mode,
            inputs={"pptx_file": basename(deck.output_path), "fix_applied_count": fix_applied_count},
        ) as run:
            inspected = self._native_service.inspect(str(output_host))
            finish_span(run, _native_inspect_span_outputs(inspected))
        if not inspected.success:
            raise DeckBuildFailure(
                "deck_native_inspect_failed",
                _native_error_summary(inspected.errors, "Final native inspect failed."),
                retryable=False,
            )
        self._record_native_inspect(deck, inspected)
        return inspected

    def _evaluate_final_native_contrast(
        self,
        *,
        deck: DeckBuild,
        runtime: ToolRuntime,
        output_host: Path,
    ) -> None:
        with deck_span(
            "deck.native.contrast",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            run_type="tool",
            deck_compile_mode=deck.deck_compile_mode,
            inputs={"pptx_file": basename(deck.output_path)},
        ) as run:
            try:
                deck.native_contrast_report = evaluate_native_contrast(
                    pptx_path=output_host,
                    source_element_map=deck.source_element_map,
                )
            except Exception as exc:  # noqa: BLE001 - analyzer failures become clean terminal results.
                finish_span(
                    run,
                    {
                        "passed": False,
                        "failure_code": "deck_native_contrast_analysis_failed",
                        "error_class": type(exc).__name__,
                    },
                )
                raise DeckBuildFailure(
                    "deck_native_contrast_analysis_failed",
                    "The final native PPTX could not be analyzed for deterministic text contrast.",
                    retryable=False,
                ) from exc
            finish_span(
                run,
                {
                    "passed": deck.native_contrast_report.get("passed"),
                    "checked_run_count": deck.native_contrast_report.get("checked_run_count"),
                    "required_issue_count": deck.native_contrast_report.get("required_issue_count"),
                    "indeterminate_required_count": deck.native_contrast_report.get("indeterminate_required_count"),
                },
            )

    def _record_native_inspect(self, deck: DeckBuild, inspected: NativeDeckInspectResult) -> None:
        deck.native_editability_score = inspected.native_editability_score
        deck.native_text_shape_count = inspected.native_text_shape_count
        deck.picture_shape_count = inspected.picture_shape_count
        deck.full_slide_picture_count = inspected.full_slide_picture_count
        inventory = _load_native_shape_inventory(inspected.shape_inventory_path)
        deck.native_shape_inventory = inventory
        for slide in deck.slides:
            slide_inventory = inventory.get(slide.selector) if isinstance(inventory, dict) else None
            if isinstance(slide_inventory, dict):
                slide.gate_results["native_shape_inventory"] = {
                    "native_slide_index": slide_inventory.get("native_slide_index"),
                    "title": slide_inventory.get("title"),
                    "body": slide_inventory.get("body"),
                    "visual": slide_inventory.get("visual"),
                }

    def _evaluate_source_quality(self, deck: DeckBuild, runtime: ToolRuntime) -> DeckEvaluation:
        output_host = _host_path(deck.output_path, runtime)
        with deck_span(
            "deck.evaluate.source",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            inputs={},
        ) as run:
            evaluation = self._evaluator.evaluate(
                deck,
                output_host_path=output_host,
                allowed_style_terms=_requested_banned_style_terms(runtime, deck.style_profile),
                require_compiled_output=False,
            )
            deck.source_quality_report = evaluation.to_dict()
            finish_span(
                run,
                {
                    "passed": evaluation.passed,
                    "hard_failure_count": len(evaluation.hard_failures),
                    "soft_warning_count": len(evaluation.soft_warnings),
                    "checks": sorted({issue.check for issue in [*evaluation.hard_failures, *evaluation.soft_warnings]}),
                    "affected_selectors": sorted({issue.selector for issue in evaluation.hard_failures}),
                    "quality_warning": evaluation.quality_warning,
                },
            )
        return evaluation

    def _evaluate(self, deck: DeckBuild, runtime: ToolRuntime) -> None:
        output_host = _host_path(deck.output_path, runtime)
        with deck_span(
            "deck.evaluate",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            inputs={},
        ) as run:
            evaluation = self._evaluator.evaluate(
                deck,
                output_host_path=output_host,
                allowed_style_terms=_requested_banned_style_terms(runtime, deck.style_profile),
            )
            finish_span(
                run,
                {
                    "passed": evaluation.passed,
                    "hard_failure_count": len(evaluation.hard_failures),
                    "soft_warning_count": len(evaluation.soft_warnings),
                    "checks": list({issue.check for issue in [*evaluation.hard_failures, *evaluation.soft_warnings]}),
                    "quality_warning": evaluation.quality_warning,
                },
            )
        if not evaluation.passed:
            raise DeckBuildFailure(
                "deck_quality_failed",
                "; ".join(issue.detail for issue in evaluation.hard_failures[:3]) or "Deck quality gate failed.",
                retryable=False,
            )
        deck.quality_warning = _merge_warning(deck.quality_warning, evaluation.quality_warning)

    def _trace_terminal(self, deck: DeckBuild, runtime: ToolRuntime, *, success: bool, deck_path: str, retryable: bool) -> None:
        with deck_span(
            "deck.terminal",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            inputs={},
        ) as run:
            finish_span(
                run,
                {
                    "terminal": True,
                    "status": "success" if success else "failed",
                    "failure_code": deck.failure_code,
                    "retryable": retryable,
                    "artifact_file": basename(deck.pptx_path) if success else None,
                    "deck_build_file": basename(deck_path),
                    "deck_route": deck.deck_route,
                    "deck_compile_mode": deck.deck_compile_mode,
                    "native_required": deck.native_required,
                    "legacy_screenshot_debug": deck.legacy_screenshot_debug,
                    "native_editability_score": deck.native_editability_score,
                    "native_text_shape_count": deck.native_text_shape_count,
                    "picture_shape_count": deck.picture_shape_count,
                    "full_slide_picture_count": deck.full_slide_picture_count,
                    "deck_authoring_contract": deck.deck_authoring_contract,
                    "deck_stylesheet_hash": deck.deck_stylesheet_hash,
                    "deck_html_fragment_count": deck.deck_html_fragment_count,
                    "deck_assembled_html_bytes": deck.deck_assembled_html_bytes,
                    "quality_warning": deck.quality_warning,
                    "failure_summary": safe_excerpt(deck.failure_summary),
                    "image_generation_status": deck.image_generation_status,
                    "image_generation_reason": deck.image_generation_reason,
                    "primary_image_batch_status": deck.primary_image_batch_status,
                    "primary_image_batch_error_class": deck.primary_image_batch_error_class,
                    "serial_repair_count": deck.serial_repair_count,
                    "batch_timeout_count": deck.batch_timeout_count,
                    "partial_batch_salvaged": deck.partial_batch_salvaged,
                    "native_mechanical_report": deck.native_mechanical_report,
                    "mechanical_gate_results": deck.mechanical_gate_results,
                    "html_source_validation": deck.html_source_validation,
                    "source_quality_report": deck.source_quality_report,
                    "creative_plan_file": basename(deck.creative_plan_path),
                    "design_plan_file": basename(deck.design_plan_path),
                    "asset_policy_file": basename(deck.asset_policy_path),
                    "style_warnings": deck.style_warnings,
                    "generated_asset_count": deck.generated_asset_count,
                    "native_html_slide_count": deck.native_html_slide_count,
                    "hybrid_slide_count": deck.hybrid_slide_count,
                    "text_only_slide_count": deck.text_only_slide_count,
                },
            )

    def _trace_emit_decision(self, deck: DeckBuild, runtime: ToolRuntime, *, success: bool, retryable: bool) -> None:
        with deck_span(
            "deck.emit.decision",
            runtime=runtime,
            build_id=deck.build_id,
            visual_policy=deck.visual_policy,
            status=deck.status,
            slide_count=len(deck.slides),
            inputs={},
        ) as run:
            finish_span(
                run,
                {
                    "emit_allowed": success,
                    "artifact_file": basename(deck.pptx_path) if success else None,
                    "deck_route": deck.deck_route,
                    "deck_compile_mode": deck.deck_compile_mode,
                    "native_required": deck.native_required,
                    "legacy_screenshot_debug": deck.legacy_screenshot_debug,
                    "native_editability_score": deck.native_editability_score,
                    "native_text_shape_count": deck.native_text_shape_count,
                    "picture_shape_count": deck.picture_shape_count,
                    "full_slide_picture_count": deck.full_slide_picture_count,
                    "deck_authoring_contract": deck.deck_authoring_contract,
                    "deck_html_fragment_count": deck.deck_html_fragment_count,
                    "deck_assembled_html_bytes": deck.deck_assembled_html_bytes,
                    "failure_code": deck.failure_code,
                    "failure_summary": safe_excerpt(deck.failure_summary),
                    "retryable": retryable,
                    "expected_visual_count": deck.expected_visual_count,
                    "successful_visual_count": deck.successful_visual_count,
                    "referenced_visual_count": deck.referenced_visual_count,
                    "missing_visual_count": deck.missing_visual_count,
                    "quality_warning": deck.quality_warning,
                    "image_generation_status": deck.image_generation_status,
                    "image_generation_reason": deck.image_generation_reason,
                    "primary_image_batch_status": deck.primary_image_batch_status,
                    "primary_image_batch_error_class": deck.primary_image_batch_error_class,
                    "serial_repair_count": deck.serial_repair_count,
                    "batch_timeout_count": deck.batch_timeout_count,
                    "partial_batch_salvaged": deck.partial_batch_salvaged,
                    "native_mechanical_report": deck.native_mechanical_report,
                    "mechanical_gate_results": deck.mechanical_gate_results,
                    "html_source_validation": deck.html_source_validation,
                    "creative_plan_file": basename(deck.creative_plan_path),
                    "design_plan_file": basename(deck.design_plan_path),
                    "asset_policy_file": basename(deck.asset_policy_path),
                    "style_warnings": deck.style_warnings,
                    "generated_asset_count": deck.generated_asset_count,
                    "native_html_slide_count": deck.native_html_slide_count,
                    "hybrid_slide_count": deck.hybrid_slide_count,
                    "text_only_slide_count": deck.text_only_slide_count,
                },
            )

    def _success_result(
        self,
        deck: DeckBuild,
        deck_path: str,
        runtime: ToolRuntime,
        *,
        success_allowed_checked: bool = False,
        service_elapsed_ms: int = 0,
    ) -> DeckBuildResult:
        if not success_allowed_checked:
            self._assert_deck_success_allowed(deck, runtime)
        _finalize_image_generation_status(deck, success=True)
        self._trace_terminal(deck, runtime, success=True, deck_path=deck_path, retryable=False)
        self._trace_emit_decision(deck, runtime, success=True, retryable=False)
        return DeckBuildResult(
            success=True,
            build_id=deck.build_id,
            deck_build_path=deck_path,
            pptx_path=deck.pptx_path,
            deck_route=deck.deck_route,
            deck_compile_mode=deck.deck_compile_mode,
            native_required=deck.native_required,
            legacy_screenshot_debug=deck.legacy_screenshot_debug,
            native_editability_score=deck.native_editability_score,
            native_text_shape_count=deck.native_text_shape_count,
            picture_shape_count=deck.picture_shape_count,
            full_slide_picture_count=deck.full_slide_picture_count,
            slide_count=len(deck.slides),
            expected_visual_count=deck.expected_visual_count,
            successful_visual_count=deck.successful_visual_count,
            referenced_visual_count=deck.referenced_visual_count,
            missing_visual_count=deck.missing_visual_count,
            deck_authoring_contract=deck.deck_authoring_contract,
            deck_stylesheet_hash=deck.deck_stylesheet_hash,
            deck_html_fragment_count=deck.deck_html_fragment_count,
            deck_assembled_html_bytes=deck.deck_assembled_html_bytes,
            creative_plan_path=deck.creative_plan_path,
            design_plan_path=deck.design_plan_path,
            asset_policy_path=deck.asset_policy_path,
            html_source_validation=deck.html_source_validation,
            source_quality_report=deck.source_quality_report,
            mechanical_gate_results=deck.mechanical_gate_results,
            style_warnings=deck.style_warnings,
            generated_asset_count=deck.generated_asset_count,
            native_html_slide_count=deck.native_html_slide_count,
            hybrid_slide_count=deck.hybrid_slide_count,
            text_only_slide_count=deck.text_only_slide_count,
            quality_status="warning" if deck.quality_warning else "passed",
            quality_warning=deck.quality_warning,
            warnings=[deck.quality_warning] if deck.quality_warning else [],
            image_generation_status=deck.image_generation_status,
            image_generation_reason=deck.image_generation_reason,
            primary_image_batch_status=deck.primary_image_batch_status,
            primary_image_batch_error_class=deck.primary_image_batch_error_class,
            serial_repair_count=deck.serial_repair_count,
            batch_timeout_count=deck.batch_timeout_count,
            partial_batch_salvaged=deck.partial_batch_salvaged,
            native_mechanical_report=deck.native_mechanical_report,
            source_retention_report=deck.source_retention_report,
            native_contrast_report=deck.native_contrast_report,
            repair_instruction=None,
            source_bundle_path=deck.source_bundle_path,
            manifest_path=deck.manifest_path,
            manifest_revision=deck.manifest_revision,
            logical_artifact_id=deck.logical_artifact_id,
            current_artifact_version_id=deck.current_artifact_version_id,
            foundation_status=deck.foundation_status,
            service_elapsed_ms=max(0, service_elapsed_ms),
        )

    def _assert_deck_success_allowed(self, deck: DeckBuild, runtime: ToolRuntime | None = None) -> None:
        if deck.deck_compile_mode in FORBIDDEN_SCREENSHOT_COMPILE_MODES:
            raise DeckBuildFailure(
                "deck_screenshot_compile_forbidden",
                "Screenshot-backed PPTX is not an allowed fresh-deck output.",
                retryable=False,
            )
        if deck.deck_compile_mode != NATIVE_DECK_COMPILE_MODE:
            raise DeckBuildFailure(
                "deck_native_compile_required",
                "Fresh PPTX decks must compile through the native PowerPoint substrate.",
                retryable=False,
            )
        verdict = classify_native_deck_substrate(
            slide_count=len(deck.slides),
            native_editability_score=deck.native_editability_score,
            native_text_shape_count=deck.native_text_shape_count,
            picture_shape_count=deck.picture_shape_count,
            full_slide_picture_count=deck.full_slide_picture_count,
            native_shape_inventory=deck.native_shape_inventory,
        )
        if runtime is not None:
            with deck_span(
                "deck.native.substrate_classify",
                runtime=runtime,
                build_id=deck.build_id,
                visual_policy=deck.visual_policy,
                status=deck.status,
                slide_count=len(deck.slides),
                run_type="tool",
                deck_compile_mode=deck.deck_compile_mode,
                inputs={
                    "slide_count": len(deck.slides),
                    "native_editability_score": deck.native_editability_score,
                    "native_text_shape_count": deck.native_text_shape_count,
                    "picture_shape_count": deck.picture_shape_count,
                    "full_slide_picture_count": deck.full_slide_picture_count,
                },
            ) as run:
                finish_span(
                    run,
                    {
                        "passed": verdict.passed,
                        "verdict": verdict.verdict,
                        "hard_failure_code": verdict.hard_failure_code,
                        "warnings": verdict.warnings,
                    },
                )
        if verdict.warnings:
            deck.quality_warning = _merge_warning(deck.quality_warning, *verdict.warnings)
        for slide in deck.slides:
            slide.gate_results["native_substrate_verdict"] = verdict.verdict
            slide.gate_results["native_substrate_warnings"] = list(verdict.warnings)
        if not verdict.passed:
            raise DeckBuildFailure(
                verdict.hard_failure_code or "deck_native_substrate_failed",
                verdict.hard_failure_summary or "Native deck substrate gate failed.",
                retryable=False,
            )
        if deck.quality_warning and "screenshot_deck" in deck.quality_warning:
            raise DeckBuildFailure(
                "deck_screenshot_compile_forbidden",
                "Screenshot-backed PPTX is not an allowed fresh-deck output.",
                retryable=False,
            )

    def _failure_result(
        self,
        deck: DeckBuild,
        deck_path: str,
        exc: DeckBuildFailure,
        runtime: ToolRuntime,
        *,
        service_elapsed_ms: int = 0,
    ) -> DeckBuildResult:
        self._trace_emit_decision(deck, runtime, success=False, retryable=exc.retryable)
        return DeckBuildResult(
            success=False,
            build_id=deck.build_id,
            deck_build_path=deck_path,
            deck_route=deck.deck_route,
            deck_compile_mode=deck.deck_compile_mode,
            native_required=deck.native_required,
            legacy_screenshot_debug=deck.legacy_screenshot_debug,
            native_editability_score=deck.native_editability_score,
            native_text_shape_count=deck.native_text_shape_count,
            picture_shape_count=deck.picture_shape_count,
            full_slide_picture_count=deck.full_slide_picture_count,
            slide_count=len(deck.slides),
            expected_visual_count=deck.expected_visual_count,
            successful_visual_count=deck.successful_visual_count,
            referenced_visual_count=deck.referenced_visual_count,
            missing_visual_count=deck.missing_visual_count,
            deck_authoring_contract=deck.deck_authoring_contract,
            deck_stylesheet_hash=deck.deck_stylesheet_hash,
            deck_html_fragment_count=deck.deck_html_fragment_count,
            deck_assembled_html_bytes=deck.deck_assembled_html_bytes,
            creative_plan_path=deck.creative_plan_path,
            design_plan_path=deck.design_plan_path,
            asset_policy_path=deck.asset_policy_path,
            html_source_validation=deck.html_source_validation,
            source_quality_report=deck.source_quality_report,
            mechanical_gate_results=deck.mechanical_gate_results,
            style_warnings=deck.style_warnings,
            generated_asset_count=deck.generated_asset_count,
            native_html_slide_count=deck.native_html_slide_count,
            hybrid_slide_count=deck.hybrid_slide_count,
            text_only_slide_count=deck.text_only_slide_count,
            failure_code=exc.code,
            failure_summary=exc.summary,
            retryable=exc.retryable,
            image_generation_status=deck.image_generation_status,
            image_generation_reason=deck.image_generation_reason,
            primary_image_batch_status=deck.primary_image_batch_status,
            primary_image_batch_error_class=deck.primary_image_batch_error_class,
            serial_repair_count=deck.serial_repair_count,
            batch_timeout_count=deck.batch_timeout_count,
            partial_batch_salvaged=deck.partial_batch_salvaged,
            quality_status="failed",
            quality_warning=deck.quality_warning,
            warnings=[deck.quality_warning] if deck.quality_warning else [],
            native_mechanical_report=deck.native_mechanical_report,
            source_retention_report=deck.source_retention_report,
            native_contrast_report=deck.native_contrast_report,
            root_failure_code=exc.code,
            root_failure_summary=exc.summary,
            repair_instruction=_repair_instruction_for_failure(exc, deck=deck),
            source_bundle_path=deck.source_bundle_path,
            manifest_path=deck.manifest_path,
            manifest_revision=deck.manifest_revision,
            logical_artifact_id=deck.logical_artifact_id,
            current_artifact_version_id=deck.current_artifact_version_id,
            foundation_status=deck.foundation_status,
            service_elapsed_ms=max(0, service_elapsed_ms),
        )

    def _run_image_batch_subprocess(self, manifest_path: str, runtime: ToolRuntime) -> dict[str, Any]:
        script = _image_script_path()
        if script is None:
            return {"summary_present": False, "complete": False, "error_class": "image_script_not_found", "raw_error_excerpt": "image generation script not found"}
        env = _image_subprocess_env(runtime)
        timeout = _deck_image_batch_timeout_seconds(manifest_path, runtime)
        try:
            completed = run_trusted_image_request(
                TrustedImageRequest(
                    python_executable=sys.executable,
                    script=script,
                    roots=_image_thread_roots(runtime),
                    mode="manifest",
                    manifest_file=_host_path(manifest_path, runtime),
                ),
                timeout=timeout,
                env=env,
            )
        except subprocess.TimeoutExpired as exc:
            stdout = _timeout_stream_text(getattr(exc, "stdout", None) or getattr(exc, "output", None))
            stderr = _timeout_stream_text(getattr(exc, "stderr", None))
            remaining = _remaining_deadline_seconds(runtime)
            deadline_near = remaining is not None and remaining <= 90
            return {
                "summary_present": False,
                "complete": False,
                "exit_code": 124,
                "error_class": "deck_deadline_exceeded" if deadline_near else "timeout",
                "items": _parse_batch_item_progress(stdout),
                "raw_error_excerpt": safe_excerpt(f"image batch subprocess timed out after {timeout}s; stdout_chars={len(stdout)} stderr_chars={len(stderr)} {(stderr or stdout).strip()}"),
            }
        summary = _parse_batch_summary(completed.stdout)
        if summary is None:
            return {
                "summary_present": False,
                "complete": False,
                "exit_code": completed.returncode,
                "error_class": "batch_summary_missing",
                "items": _parse_batch_item_progress(completed.stdout),
                "raw_error_excerpt": safe_excerpt((completed.stderr or completed.stdout or "").strip()),
            }
        summary["summary_present"] = True
        summary["exit_code"] = completed.returncode
        return summary

    def _run_image_single_subprocess(self, slide: DeckSlideSpec, runtime: ToolRuntime, attempt_no: int) -> dict[str, Any]:
        script = _image_script_path()
        if script is None:
            return {"success": False, "error_class": "image_script_not_found", "raw_error_excerpt": "image generation script not found"}
        if not slide.visual_prompt_path or not slide.visual_asset_path:
            return {"success": False, "error_class": "missing_prompt_or_output"}
        started = time.perf_counter()
        remaining = _remaining_deadline_seconds(runtime)
        if remaining is not None and remaining <= 0:
            return {
                "success": False,
                "error_class": "deck_deadline_exceeded",
                "raw_error_excerpt": "shared builder deadline exhausted before serial repair",
            }
        timeout = min(
            _image_single_timeout_seconds(),
            max(1, remaining) if remaining is not None else _image_single_timeout_seconds(),
        )
        aspect_ratio = (slide.asset_plan.aspect_ratio if slide.asset_plan else None) or "16:9"
        try:
            completed = run_trusted_image_request(
                TrustedImageRequest(
                    python_executable=sys.executable,
                    script=script,
                    roots=_image_thread_roots(runtime),
                    mode="single",
                    prompt_file=_host_path(slide.visual_prompt_path, runtime),
                    output_file=_host_path(slide.visual_asset_path, runtime),
                    aspect_ratio=aspect_ratio,
                ),
                timeout=timeout,
                env=_image_subprocess_env(runtime),
            )
        except subprocess.TimeoutExpired as exc:
            stdout = _timeout_stream_text(getattr(exc, "stdout", None) or getattr(exc, "output", None))
            stderr = _timeout_stream_text(getattr(exc, "stderr", None))
            return {
                "success": False,
                "exit_code": 124,
                "error_class": "timeout",
                "duration_ms": int((time.perf_counter() - started) * 1000),
                "raw_error_excerpt": safe_excerpt(f"serial image repair attempt {attempt_no} timed out after {timeout}s; stdout_chars={len(stdout)} stderr_chars={len(stderr)} {(stderr or stdout).strip()}"),
            }
        output_bytes = _file_size(_host_path(slide.visual_asset_path, runtime))
        success = completed.returncode == 0 and bool(output_bytes and output_bytes > 0)
        return {
            "success": success,
            "exit_code": completed.returncode,
            "error_class": None if success else _image_generation_error_class_from_output(completed.stderr, completed.stdout, completed.returncode),
            "duration_ms": int((time.perf_counter() - started) * 1000),
            "bytes": output_bytes or 0,
            "raw_error_excerpt": safe_excerpt((completed.stderr or completed.stdout or "").strip()),
        }


class DeckBuildFailure(Exception):
    def __init__(self, code: str, summary: str, *, retryable: bool) -> None:
        super().__init__(summary)
        self.code = code
        self.summary = summary
        self.retryable = retryable


def _report_count(report: dict[str, Any], *keys: str) -> int:
    for key in keys:
        value = report.get(key)
        if isinstance(value, (list, tuple, set, dict)):
            return len(value)
    return 0


def _quality_failure_summary(evaluation: DeckEvaluation) -> str:
    grouped: dict[tuple[str, str], list[str]] = {}
    for issue in evaluation.hard_failures:
        key = (str(issue.check or "quality"), str(issue.detail or "Deck source quality failed."))
        selectors = grouped.setdefault(key, [])
        if issue.selector not in selectors:
            selectors.append(issue.selector)
    parts = [
        f"{check} on {', '.join(selectors)}: {detail}"
        for (check, detail), selectors in grouped.items()
    ]
    return safe_excerpt("; ".join(parts) or "Deck source quality gate failed.", limit=1600)


def _combined_source_and_mechanical_summary(evaluation: DeckEvaluation, mechanical_summary: str) -> str:
    return safe_excerpt(
        f"{mechanical_summary} Source quality also failed: {_quality_failure_summary(evaluation)}",
        limit=1800,
    )


def _state_value(runtime: Any, key: str) -> Any:
    state = getattr(runtime, "state", None)
    return state.get(key) if isinstance(state, dict) else None


def _runtime_identity_value(runtime: Any, key: str) -> Any:
    state = getattr(runtime, "state", None)
    if isinstance(state, dict):
        for source in (
            state,
            state.get("builder_task"),
            state.get("delegation_context"),
        ):
            if isinstance(source, dict) and source.get(key) not in (None, ""):
                return source[key]

    execution_info = getattr(runtime, "execution_info", None)
    value = getattr(execution_info, key, None) if execution_info is not None else None
    if value not in (None, ""):
        return value

    context = getattr(runtime, "context", None)
    if isinstance(context, dict) and context.get(key) not in (None, ""):
        return context[key]

    config = getattr(runtime, "config", None)
    if isinstance(config, dict):
        for source_name in ("configurable", "metadata"):
            source = config.get(source_name)
            if isinstance(source, dict) and source.get(key) not in (None, ""):
                return source[key]
    return None


def _builder_deadline_epoch_ms(runtime: Any) -> int | None:
    raw = _state_value(runtime, "builder_deadline_epoch_ms")
    if isinstance(raw, (int, float)) and raw > 0:
        return int(raw)
    kickoff = _state_value(runtime, "builder_task_kickoff_ms")
    timeout = _state_value(runtime, "builder_timeout_seconds")
    if isinstance(kickoff, (int, float)) and kickoff > 0 and isinstance(timeout, (int, float)) and timeout > 0:
        return int(kickoff) + int(timeout) * 1000
    return None


def _terminal_reserve_seconds(runtime: Any) -> int:
    budget = _state_value(runtime, "builder_budget")
    if not isinstance(budget, dict):
        return 0
    try:
        return max(0, int(budget.get("terminal_reserve_seconds", 0) or 0))
    except (TypeError, ValueError):
        return 0


def _service_deadline_epoch_ms(runtime: Any) -> int | None:
    deadline = _builder_deadline_epoch_ms(runtime)
    if deadline is None:
        return None
    return max(1, deadline - _terminal_reserve_seconds(runtime) * 1_000)


def _remaining_deadline_seconds(
    runtime: Any,
    *,
    reserve_seconds: int | None = None,
) -> int | None:
    deadline = _builder_deadline_epoch_ms(runtime)
    if deadline is None:
        return None
    if reserve_seconds is None:
        reserve_seconds = _terminal_reserve_seconds(runtime)
    return max(
        0,
        int((deadline - int(time.time() * 1000)) / 1000) - max(0, reserve_seconds),
    )


def _assert_deck_deadline(
    runtime: Any,
    *,
    stage: str,
    reserve_seconds: int | None = None,
) -> None:
    remaining = _remaining_deadline_seconds(runtime, reserve_seconds=reserve_seconds)
    if remaining is not None and remaining <= 0:
        raise DeckBuildFailure(
            "deck_deadline_exceeded",
            f"The shared builder deadline was exhausted before {stage}.",
            retryable=False,
        )


def _truthy_env(name: str) -> bool:
    return (os.getenv(name, "") or "").strip().lower() in {"1", "true", "yes", "on"}


def _native_preflight(native_service: Any) -> NativeDeckPreflight:
    preflight = getattr(native_service, "preflight", None)
    if callable(preflight):
        result = preflight()
        if isinstance(result, NativeDeckPreflight):
            return result
        if isinstance(result, dict):
            return NativeDeckPreflight(
                success=bool(result.get("success")),
                scripts_dir_exists=bool(result.get("scripts_dir_exists")),
                deck_py_exists=bool(result.get("deck_py_exists")),
                html2patch_py_exists=bool(result.get("html2patch_py_exists")),
                errors=[str(item) for item in result.get("errors", []) if item],
            )
    return NativeDeckPreflight(
        success=True,
        scripts_dir_exists=True,
        deck_py_exists=True,
        html2patch_py_exists=True,
        errors=[],
    )


def _legacy_screenshot_debug_allowed(runtime: ToolRuntime) -> bool:
    if not _truthy_env("SOPHIA_DECK_LEGACY_SCREENSHOT_DEBUG"):
        return False
    return not _is_production_runtime(runtime)


def _is_production_runtime(runtime: ToolRuntime) -> bool:
    if any(os.getenv(name) for name in ("RENDER", "RENDER_SERVICE_ID", "RENDER_SERVICE_NAME")):
        return True
    for key in ("SOPHIA_ENV", "SOPHIA_RUNTIME_ENV", "APP_ENV", "ENV", "ENVIRONMENT"):
        if (os.getenv(key, "") or "").strip().lower() == "production":
            return True
    value = _state_value(runtime, "environment") or _state_value(runtime, "runtime_env")
    return str(value or "").strip().lower() == "production"


def _native_startup_error(errors: list[str]) -> bool:
    text = "\n".join(errors).lower()
    return any(
        marker in text
        for marker in (
            "playwright is required",
            "hands-on-deck script not found",
            "no module named",
            "python-pptx is required",
            "chromium",
            "browser",
        )
    )


def _native_deadline_error(errors: list[str]) -> bool:
    return "deck deadline exceeded" in "\n".join(errors).lower()


def _native_error_summary(errors: list[str], fallback: str) -> str:
    if not errors:
        return fallback
    return safe_excerpt(errors[0], limit=600) or fallback


def _merge_warning(*warnings: str | None) -> str | None:
    parts: list[str] = []
    for warning in warnings:
        if not warning:
            continue
        for part in str(warning).split(";"):
            clean = part.strip()
            if clean and clean not in parts:
                parts.append(clean)
    return "; ".join(parts) if parts else None


def _write_native_base_deck(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = Inches(20)
    presentation.slide_height = Inches(11.25)
    presentation.save(path)


def _native_patch_span_outputs(result: NativeDeckPatchResult) -> dict[str, Any]:
    return {
        "success": result.success,
        "patch_file": basename(result.patch_path),
        "output_file": basename(result.output_pptx_path),
        "patch_op_count": result.patch_op_count,
        "source_map_file": basename(result.source_map_path),
        "validation_error_count": result.validation_error_count,
        "error_count": len(result.errors),
        "error_excerpt": safe_excerpt(result.errors[0]) if result.errors else None,
    }


def _native_inspect_span_outputs(result: NativeDeckInspectResult) -> dict[str, Any]:
    return {
        "success": result.success,
        "slide_count": result.slide_count,
        "shape_count": result.shape_count,
        "native_text_shape_count": result.native_text_shape_count,
        "picture_shape_count": result.picture_shape_count,
        "full_slide_picture_count": result.full_slide_picture_count,
        "native_editability_score": result.native_editability_score,
        "shape_inventory_file": basename(result.shape_inventory_path),
        "raw_json_file": basename(result.raw_json_path),
        "error_count": len(result.errors),
        "error_excerpt": safe_excerpt(result.errors[0]) if result.errors else None,
    }


def _native_lint_fix_span_outputs(result: NativeDeckLintFixResult) -> dict[str, Any]:
    return {
        "success": result.success,
        "lint_issue_count_before": result.lint_issue_count_before,
        "fix_applied_count": result.fix_applied_count,
        "residue_count": result.residue_count,
        "touched_slide_count": result.touched_slide_count,
        "issue_kinds": result.issue_kinds,
        "residue_kinds": result.residue_kinds,
        "error_count": len(result.errors),
        "error_excerpt": safe_excerpt(result.errors[0]) if result.errors else None,
    }


def _native_render_span_outputs(result: NativeDeckRenderResult) -> dict[str, Any]:
    return {
        "success": result.success,
        "render_dir": basename(result.render_dir),
        "rendered_slide_count": result.rendered_slide_count,
        "error_count": len(result.errors),
        "error_excerpt": safe_excerpt(result.errors[0]) if result.errors else None,
    }


def _load_native_shape_inventory(path: str | None) -> dict[str, Any]:
    if not path:
        return {}
    try:
        payload = json.loads(Path(path).read_text(encoding="utf-8"))
    except (OSError, ValueError):
        return {}
    slides = payload.get("slides") if isinstance(payload, dict) else None
    return slides if isinstance(slides, dict) else {}


def _load_json_dict(path: str | None) -> dict[str, Any]:
    if not path:
        return {}
    try:
        payload = json.loads(Path(path).read_text(encoding="utf-8"))
    except (OSError, ValueError):
        return {}
    return payload if isinstance(payload, dict) else {}


def _slide_html_virtual_path(slide: DeckSlideSpec) -> str:
    role = re.sub(r"[^a-z0-9]+", "-", str(slide.role or "slide").lower()).strip("-") or "slide"
    return f"{_SLIDES}/{slide.index:02d}-{role}.html"


def _referenced_planned_asset_count(deck: DeckBuild, validations: list[Any]) -> int:
    refs_by_selector = {validation.selector: set(validation.image_refs) for validation in validations}
    count = 0
    for slide in generated_asset_slides(deck):
        expected = f"../assets/slide-{slide.index:02d}.png"
        if expected in refs_by_selector.get(slide.selector, set()):
            count += 1
    return count


def _repair_instruction_for_failure(
    exc: DeckBuildFailure,
    *,
    deck: DeckBuild | None = None,
) -> dict[str, Any] | None:
    if not exc.retryable:
        return None
    if exc.code not in {
        "deck_creative_plan_required",
        "deck_creative_plan_invalid",
        "deck_slide_html_missing",
        "deck_slide_html_invalid",
        "deck_image_asset_plan_invalid",
        "deck_mechanical_gate_failed",
        "deck_source_quality_failed",
        "invalid_deck_ir",
    }:
        return None
    instruction = {
        "should_retry": True,
        "max_retry_count": 1,
        "failure_scope": "creative_deck_plan" if exc.code.startswith("deck_creative") else "slide_html_or_mechanical_gate",
        "message": ("Create or repair the subject-derived DeckCreativePlan, shared deck_stylesheet, and every slide.html_body, using planned generated images only as assets, then call prepare_deck_build exactly once more."),
        "repair_message": (
            "Repair the D2.1 deck input and call prepare_deck_build exactly once more. "
            "Include authoring_contract=compact_model_html_v2, creative_plan with design_plan, image_assets, "
            "slide_compositions, one concise shared deck_stylesheet, and one html_body per slide. Put all authored CSS "
            "in deck_stylesheet and keep slide_css empty. Keep the slide canvas 1920x1080 and keep every shape inside it. If an intentional decorative bleed is unavoidable, mark its owning HTML element with data-deck-role=\"decorative\", \"background\", or "bleed" so the native gate can identify it. Use an opaque background, no scripts/external URLs, "
            "and reference only planned assets as ../assets/slide-XX.png. "
            f"Previous failure: {exc.code}: {safe_excerpt(exc.summary, limit=400)}"
        ),
    }
    if exc.code in {"deck_mechanical_gate_failed", "deck_source_quality_failed"} and deck is not None:
        targeted = deck_mechanical_repair_instruction_from_reports(
            native_contrast_report=deck.native_contrast_report,
            source_element_map=deck.source_element_map,
            native_mechanical_report=deck.native_mechanical_report,
            mechanical_gate_results=deck.mechanical_gate_results,
            native_shape_inventory=deck.native_shape_inventory,
            source_quality_report=getattr(deck, "source_quality_report", {}),
        )
        if targeted is not None:
            instruction.update(targeted)
    return instruction


@contextlib.contextmanager
def _deck_trace_runtime_context(runtime: ToolRuntime, deck: DeckBuild) -> Iterator[None]:
    state = getattr(runtime, "state", None)
    if not isinstance(state, dict):
        yield
        return
    sentinel = object()
    patch = {
        "current_deck_build_id": deck.build_id,
        "current_deck_route": deck.deck_route,
        "current_deck_compile_mode": deck.deck_compile_mode,
        "current_deck_artifact_target_ext": DEFAULT_ARTIFACT_TARGET_EXT,
    }
    previous = {key: state.get(key, sentinel) for key in patch}
    state.update(patch)
    try:
        yield
    finally:
        for key, old_value in previous.items():
            if old_value is sentinel:
                state.pop(key, None)
            else:
                state[key] = old_value


def _compile_with_build_deck_from_slides(
    runtime: ToolRuntime,
    output_path: str,
    title: str,
    slides_dir: str,
) -> dict[str, Any]:
    from deerflow.sophia.tools.build_deck_from_slides import build_deck_from_slides

    raw = build_deck_from_slides.func(
        runtime=runtime,
        output_path=output_path,
        title=title,
        slides_dir=slides_dir,
    )
    try:
        payload = json.loads(raw)
    except (TypeError, ValueError):
        return {"success": False, "error_type": "deck_compile_failed", "error": "Compiler returned non-JSON output."}
    return payload if isinstance(payload, dict) else {"success": False, "error_type": "deck_compile_failed", "error": "Compiler returned invalid output."}


def _contains_unnegated_match(pattern: re.Pattern[str], value: str) -> bool:
    return bool(_unnegated_matches(pattern, value))


def _unnegated_matches(pattern: re.Pattern[str], value: str) -> set[str]:
    matches: set[str] = set()
    for match in pattern.finditer(value):
        prefix = value[max(0, match.start() - 32) : match.start()]
        if _NEGATED_BANNED_TERM_RE.search(prefix):
            continue
        matches.add(match.group(0).lower())
    return matches


def _unrequested_banned_style_terms(value: str, runtime: ToolRuntime, style_profile: dict[str, Any]) -> set[str]:
    style_terms = _unnegated_matches(_BANNED_STYLE_RE, value)
    if not style_terms:
        return set()
    requested_style_terms = _requested_banned_style_terms(runtime, style_profile)
    return style_terms - requested_style_terms


def _requested_banned_style_terms(runtime: ToolRuntime, style_profile: dict[str, Any]) -> set[str]:
    requested_style_text = _normalize_style_text(
        "\n".join(
            [
                _request_context_text(runtime),
                "\n".join(_string_values(style_profile)),
            ]
        )
    )
    return _unnegated_matches(_BANNED_STYLE_RE, requested_style_text)


def _normalize_style_text(value: str) -> str:
    return re.sub(r"[_-]+", " ", value.lower())


def _request_context_text(runtime: ToolRuntime) -> str:
    state = getattr(runtime, "state", None)
    if not isinstance(state, dict):
        return ""
    haystack_parts: list[str] = []
    for source in (state, state.get("delegation_context"), state.get("builder_task"), state.get("artifact_request")):
        if not isinstance(source, dict):
            continue
        for key in _REQUEST_CONTEXT_TEXT_KEYS:
            value = source.get(key)
            if isinstance(value, str):
                haystack_parts.append(value)
    return "\n".join(haystack_parts)


def _string_values(value: Any) -> list[str]:
    if isinstance(value, str):
        return [value]
    if isinstance(value, dict):
        values: list[str] = []
        for item in value.values():
            values.extend(_string_values(item))
        return values
    if isinstance(value, (list, tuple, set)):
        values = []
        for item in value:
            values.extend(_string_values(item))
        return values
    return []


def _clear_slide_html_directory(slides_dir: Path) -> None:
    if not slides_dir.exists():
        return
    for path in slides_dir.glob("*.html"):
        if path.is_file():
            path.unlink()


def _requested_slide_count_from_state(runtime: ToolRuntime) -> int:
    state = getattr(runtime, "state", None)
    if not isinstance(state, dict):
        return 0
    sources: list[Any] = [state, state.get("delegation_context"), state.get("builder_task")]
    for source in sources:
        if not isinstance(source, dict):
            continue
        for key in ("builder_pptx_requested_slide_count", "target_slide_count", "requested_slide_count"):
            try:
                value = int(source.get(key) or 0)
            except (TypeError, ValueError):
                value = 0
            if value > 0:
                return value
    return 0


def _explicit_text_only_requested(runtime: ToolRuntime) -> bool:
    return bool(_TEXT_ONLY_REQUEST_RE.search(_request_context_text(runtime)))


def _finalize_image_generation_status(deck: DeckBuild, *, success: bool) -> None:
    if deck.expected_visual_count <= 0:
        deck.image_generation_status = deck.image_generation_status or "not_required"
        deck.primary_image_batch_status = deck.primary_image_batch_status or "not_required"
        return
    if success:
        deck.image_generation_status = "success_after_repair" if deck.serial_repair_count > 0 else "success"
        deck.image_generation_reason = None
        if deck.serial_repair_count > 0:
            deck.primary_image_batch_status = "repaired"
        else:
            deck.primary_image_batch_status = deck.primary_image_batch_status or "success"
        return
    if deck.successful_visual_count == deck.expected_visual_count and deck.missing_visual_count == 0:
        deck.image_generation_status = "success_after_repair" if deck.serial_repair_count > 0 else "success"
        deck.image_generation_reason = None
        deck.primary_image_batch_status = deck.primary_image_batch_status or ("repaired" if deck.serial_repair_count > 0 else "success")
        return
    if deck.successful_visual_count > 0:
        deck.image_generation_status = "partial"
    else:
        deck.image_generation_status = "failed"
    deck.image_generation_reason = deck.image_generation_reason or deck.primary_image_batch_error_class or deck.failure_code or "image_generation_failed"
    deck.primary_image_batch_status = deck.primary_image_batch_status or "failed"
    deck.primary_image_batch_error_class = deck.primary_image_batch_error_class or deck.image_generation_reason


def _current_image_trace_env(runtime: ToolRuntime) -> dict[str, str]:
    env: dict[str, str] = {}
    try:
        from langsmith.run_helpers import get_current_run_tree

        run_tree = get_current_run_tree()
        dotted_order = getattr(run_tree, "dotted_order", None)
        trace_id = getattr(run_tree, "trace_id", None)
        run_id = getattr(run_tree, "id", None)
        if dotted_order:
            env["SOPHIA_PARENT_DOTTED_ORDER"] = str(dotted_order)
        if trace_id:
            env["SOPHIA_PARENT_TRACE_ID"] = str(trace_id)
        if run_id:
            env["SOPHIA_PARENT_RUN_ID"] = str(run_id)
    except Exception:  # noqa: BLE001 - tracing env must never block generation.
        pass
    thread_id = _runtime_identity_value(runtime, "thread_id")
    if thread_id:
        env["SOPHIA_THREAD_ID"] = str(thread_id)
    session_id = (
        _runtime_identity_value(runtime, "session_id")
        or _runtime_identity_value(runtime, "parent_thread_id")
        or _runtime_identity_value(runtime, "companion_session_id")
    )
    if session_id:
        env["SOPHIA_SESSION_ID"] = str(session_id)
    task_id = _runtime_identity_value(runtime, "task_id") or _runtime_identity_value(runtime, "builder_task_id")
    if task_id:
        env["SOPHIA_TASK_ID"] = str(task_id)
    run_id = _runtime_identity_value(runtime, "run_id") or _runtime_identity_value(runtime, "builder_run_id")
    if run_id:
        env["SOPHIA_RUN_ID"] = str(run_id)
    user_id_hash = stable_hash(_runtime_identity_value(runtime, "user_id"))
    if user_id_hash:
        env["SOPHIA_USER_ID_HASH"] = user_id_hash
    build_id = _state_value(runtime, "current_deck_build_id") or _state_value(runtime, "deck_build_id")
    if build_id:
        env["SOPHIA_BUILD_ID"] = str(build_id)
        env["SOPHIA_DECK_BUILD_ID"] = str(build_id)
    env["SOPHIA_DECK_ROUTE"] = str(_state_value(runtime, "current_deck_route") or _state_value(runtime, "deck_route") or DEFAULT_DECK_ROUTE)
    env["SOPHIA_DECK_COMPILE_MODE"] = str(_state_value(runtime, "current_deck_compile_mode") or _state_value(runtime, "deck_compile_mode") or DEFAULT_DECK_COMPILE_MODE)
    env["SOPHIA_ARTIFACT_TARGET_EXT"] = str(_state_value(runtime, "current_deck_artifact_target_ext") or _state_value(runtime, "artifact_target_ext") or DEFAULT_ARTIFACT_TARGET_EXT)
    return env


def _image_subprocess_env(runtime: ToolRuntime) -> dict[str, str]:
    env = trusted_subprocess_env(
        allow_openai=True,
        allow_langsmith=True,
    )
    thread_data = get_thread_data(runtime) or {}
    if thread_data.get("outputs_path"):
        env["SOPHIA_OUTPUTS_HOST_PATH"] = str(thread_data["outputs_path"])
    if thread_data.get("workspace_path"):
        env["SOPHIA_WORKSPACE_HOST_PATH"] = str(thread_data["workspace_path"])
    env.update(_current_image_trace_env(runtime))
    return env


def _int_env(name: str, default: int) -> int:
    try:
        value = int(os.getenv(name, "") or default)
    except (TypeError, ValueError):
        return default
    return value if value > 0 else default


def _nonnegative_int_env(name: str, default: int) -> int:
    try:
        value = int(os.getenv(name, "") or default)
    except (TypeError, ValueError):
        return default
    return value if value >= 0 else default


def _manifest_item_count_for_timeout(manifest_path: str, runtime: ToolRuntime) -> int:
    try:
        host_path = _host_path(manifest_path, runtime)
        payload = json.loads(host_path.read_text(encoding="utf-8"))
    except (OSError, ValueError):
        return 0
    items = payload.get("items") if isinstance(payload, dict) else None
    return len(items) if isinstance(items, list) else 0


def _image_thread_roots(runtime: ToolRuntime) -> ImageThreadRoots:
    """Return the one canonical thread boundary used by the image broker."""

    thread_data = get_thread_data(runtime) or {}
    values: dict[str, str] = {}
    for name, key in (
        ("workspace", "workspace_path"),
        ("outputs", "outputs_path"),
        ("uploads", "uploads_path"),
    ):
        value = thread_data.get(key)
        if not isinstance(value, str) or not value.strip():
            raise ValueError(f"image subprocess requires thread {key}")
        values[name] = value
    return ImageThreadRoots.create(**values)


def _deck_image_batch_timeout_seconds(manifest_path: str, runtime: ToolRuntime) -> int:
    override = os.getenv("SOPHIA_DECK_IMAGE_BATCH_TIMEOUT")
    if override and override.strip():
        calculated = _int_env("SOPHIA_DECK_IMAGE_BATCH_TIMEOUT", 1800)
    else:
        per_image_timeout = _int_env("SOPHIA_IMAGE_GEN_TIMEOUT", 120)
        max_retries = _nonnegative_int_env("SOPHIA_IMAGE_GEN_MAX_RETRIES", 1)
        concurrency = _int_env("SOPHIA_IMAGE_GEN_CONCURRENCY", 2)
        item_count = max(1, _manifest_item_count_for_timeout(manifest_path, runtime))
        waves = max(1, (item_count + concurrency - 1) // concurrency)
        calculated = max(
            _image_single_timeout_seconds(),
            waves * per_image_timeout * (max_retries + 1) + 60,
        )
    remaining = _remaining_deadline_seconds(runtime)
    return min(calculated, max(1, remaining)) if remaining is not None else calculated


def _image_single_timeout_seconds() -> int:
    per_image_timeout = _int_env("SOPHIA_IMAGE_GEN_TIMEOUT", 120)
    max_retries = _nonnegative_int_env("SOPHIA_IMAGE_GEN_MAX_RETRIES", 1)
    return per_image_timeout * (max_retries + 1) + 30


def _timeout_stream_text(value: object) -> str:
    if value is None:
        return ""
    if isinstance(value, bytes):
        return value.decode("utf-8", errors="replace")
    return str(value)


def _authoring_failure(summary: str) -> DeckBuildFailure:
    return DeckBuildFailure("invalid_deck_ir", summary, retryable=True)


def _normalize_compact_v2_anchor_invariant_contract(
    stylesheet: str,
    slides: list[dict[str, Any]],
) -> tuple[str, dict[str, Any]]:
    """Complete one mechanically unambiguous fresh-authoring anchor contract.

    Some model-authored compact-v2 decks contain a unique, unused class rule
    with exactly the three mandatory anchor invariants while every declared
    anchor already has a standalone, safe four-field ``#id`` geometry rule.
    That source is invalid and does not compile, but its intended absolute-box
    contract is mechanically unambiguous. In that one bounded case, copy only
    the mandatory constants into the existing ID rules. Every other source
    remains unchanged and proceeds to the normal fail-closed validator.
    """

    def report(
        *,
        normalized_rule_count: int = 0,
        injected_declaration_count: int = 0,
        carrier_selector_hash: str | None = None,
        stylesheet_after: str = stylesheet,
    ) -> dict[str, Any]:
        return {
            "normalization_applied": normalized_rule_count > 0,
            "normalized_anchor_rule_count": normalized_rule_count,
            "injected_declaration_count": injected_declaration_count,
            "carrier_selector_sha256": carrier_selector_hash,
            "stylesheet_bytes_before": len(stylesheet.encode("utf-8")),
            "stylesheet_bytes_after": len(stylesheet_after.encode("utf-8")),
            "html_body_changed": False,
            "strict_validator_bypassed": False,
            "candidate_compile_changed": False,
            "raw_content_excluded": True,
        }

    try:
        rules = tinycss2.parse_stylesheet(
            stylesheet,
            skip_comments=False,
            skip_whitespace=False,
        )
    except Exception:  # noqa: BLE001 - strict validation reports source errors.
        return stylesheet, report()
    if any(getattr(rule, "type", "") in {"at-rule", "error"} for rule in rules):
        return stylesheet, report()

    soups: list[BeautifulSoup] = []
    selector_soups: list[BeautifulSoup] = []
    anchor_markers_by_soup: list[set[int]] = []
    selector_anchor_markers_by_soup: list[set[int]] = []
    declared_ids_by_soup: list[set[str]] = []
    declared_ids: list[str] = []
    for raw in slides:
        declared = raw.get("repair_anchor_ids")
        if (
            not isinstance(declared, list)
            or len(declared) != 2
            or any(
                not isinstance(identifier, str)
                or _COMPACT_SOURCE_ID_RE.fullmatch(identifier) is None
                for identifier in declared
            )
            or len(set(declared)) != 2
        ):
            return stylesheet, report()
        body = raw.get("html_body")
        if not isinstance(body, str) or not body.strip():
            return stylesheet, report()
        try:
            soup = BeautifulSoup(body, "html.parser")
        except Exception:  # noqa: BLE001 - strict validation reports source errors.
            return stylesheet, report()
        anchors: list[Tag] = []
        anchor_data_ids: set[str] = set()
        for identifier in declared:
            matches = soup.find_all(id=identifier)
            if len(matches) != 1 or not isinstance(matches[0], Tag):
                return stylesheet, report()
            anchor = matches[0]
            if (
                not any(anchor is element for element in soup.contents)
                or not _compact_source_anchor_is_eligible(
                    anchor,
                    {identifier: (0.0, 0.0, 48.0, 24.0)},
                )
            ):
                return stylesheet, report()
            data_id = str(anchor.attrs.get("data-deck-id") or "").strip()
            if data_id in anchor_data_ids:
                return stylesheet, report()
            anchor_data_ids.add(data_id)
            anchors.append(anchor)
            declared_ids.append(identifier)
        soups.append(soup)
        anchor_markers_by_soup.append({id(anchor) for anchor in anchors})
        declared_ids_by_soup.append(set(declared))
        selector_soup = BeautifulSoup(
            assemble_compact_slide_html(
                deck_stylesheet="",
                html_body=body,
                slide_css="",
            ),
            "html.parser",
        )
        selector_anchors = [selector_soup.find(id=identifier) for identifier in declared]
        if any(not isinstance(anchor, Tag) for anchor in selector_anchors):
            return stylesheet, report()
        selector_soups.append(selector_soup)
        selector_anchor_markers_by_soup.append(
            {id(anchor) for anchor in selector_anchors}
        )

    unique_ids = tuple(dict.fromkeys(declared_ids))
    if not unique_ids:
        return stylesheet, report()
    for identifier in unique_ids:
        for soup, anchor_markers, slide_declared_ids in zip(
            soups,
            anchor_markers_by_soup,
            declared_ids_by_soup,
            strict=True,
        ):
            occurrences = soup.find_all(id=identifier)
            expected_count = 1 if identifier in slide_declared_ids else 0
            if len(occurrences) != expected_count or (
                expected_count == 1 and id(occurrences[0]) not in anchor_markers
            ):
                return stylesheet, report()
    qualified_rules = [
        rule for rule in rules if getattr(rule, "type", "") == "qualified-rule"
    ]
    target_rules: dict[str, Any] = {}
    target_rule_markers: set[int] = set()
    geometry_names = set(_COMPACT_SOURCE_GEOMETRY_PROPERTIES)
    for identifier in unique_ids:
        matches = [
            rule
            for rule in qualified_rules
            if tinycss2.serialize(rule.prelude).strip() == f"#{identifier}"
        ]
        if len(matches) != 1:
            return stylesheet, report()
        rule = matches[0]
        declarations = tinycss2.parse_declaration_list(
            rule.content,
            skip_comments=True,
            skip_whitespace=True,
        )
        if (
            any(getattr(item, "type", "") != "declaration" for item in declarations)
            or any(bool(getattr(item, "important", False)) for item in declarations)
        ):
            return stylesheet, report()
        final = {item.lower_name: item for item in declarations}
        if len(final) != len(declarations) or set(final) != geometry_names:
            return stylesheet, report()
        values = tuple(
            _css_absolute_px_value(final.get(property_name))
            for property_name in _COMPACT_SOURCE_GEOMETRY_PROPERTIES
        )
        if any(value is None for value in values):
            return stylesheet, report()
        left, top, width, height = (float(value) for value in values if value is not None)
        if (
            left < 0
            or top < 0
            or width < 48
            or height < 24
            or left + width > 1920
            or top + height > 1080
            or not (
                left >= 8
                or top >= 8
                or left + width <= 1912
                or top + height <= 1072
            )
        ):
            return stylesheet, report()
        target_rules[identifier] = rule
        target_rule_markers.add(id(rule))

    carrier_candidates: list[tuple[Any, str]] = []
    invariant_names = set(_COMPACT_SOURCE_ANCHOR_INVARIANT_PROPERTIES)
    for rule in qualified_rules:
        if id(rule) in target_rule_markers:
            continue
        selector = tinycss2.serialize(rule.prelude).strip()
        if (
            selector in _BASE_FONT_SELECTORS
            or re.fullmatch(r"\.[a-z][a-z0-9_-]{0,31}", selector) is None
        ):
            continue
        declarations = tinycss2.parse_declaration_list(
            rule.content,
            skip_comments=True,
            skip_whitespace=True,
        )
        if (
            any(getattr(item, "type", "") != "declaration" for item in declarations)
            or any(bool(getattr(item, "important", False)) for item in declarations)
        ):
            continue
        final = {item.lower_name: item for item in declarations}
        if len(final) != len(declarations) or set(final) != invariant_names:
            continue
        if (
            _css_single_identifier(final.get("position")) != "absolute"
            or _css_single_identifier(final.get("box-sizing")) != "border-box"
            or _css_zero_value(final.get("margin")) is not True
        ):
            continue
        try:
            if any(soup.select(selector) for soup in selector_soups):
                continue
        except Exception:
            continue
        carrier_candidates.append((rule, selector))
    if len(carrier_candidates) != 1:
        return stylesheet, report()
    carrier_rule, carrier_selector = carrier_candidates[0]

    protected_names = geometry_names | invariant_names
    for rule in qualified_rules:
        if id(rule) in target_rule_markers or rule is carrier_rule:
            continue
        declarations = tinycss2.parse_declaration_list(
            rule.content,
            skip_comments=True,
            skip_whitespace=True,
        )
        if not any(
            getattr(item, "type", "") == "declaration"
            and item.lower_name in protected_names
            for item in declarations
        ):
            continue
        selector = tinycss2.serialize(rule.prelude).strip()
        try:
            for soup, anchor_markers in zip(
                selector_soups,
                selector_anchor_markers_by_soup,
                strict=True,
            ):
                if any(id(match) in anchor_markers for match in soup.select(selector)):
                    return stylesheet, report()
        except Exception:
            return stylesheet, report()

    invariant_prefix = "position:absolute;box-sizing:border-box;margin:0;"
    for rule in target_rules.values():
        rule.content = tinycss2.parse_component_value_list(
            invariant_prefix + tinycss2.serialize(rule.content)
        )
    normalized_stylesheet = tinycss2.serialize(rules)
    if len(normalized_stylesheet.encode("utf-8")) > 8 * 1024:
        return stylesheet, report()
    if not set(unique_ids) <= set(_compact_shared_id_geometry(normalized_stylesheet)):
        return stylesheet, report()
    selector_hash = hashlib.sha256(carrier_selector.encode("utf-8")).hexdigest()
    return normalized_stylesheet, report(
        normalized_rule_count=len(target_rules),
        injected_declaration_count=(
            len(target_rules) * len(_COMPACT_SOURCE_ANCHOR_INVARIANT_PROPERTIES)
        ),
        carrier_selector_hash=selector_hash,
        stylesheet_after=normalized_stylesheet,
    )


def _normalize_compact_v2_anchor_inline_geometry(
    stylesheet: str,
    slides: list[dict[str, Any]],
) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    """Remove only redundant inline geometry backed by a strict shared rule.

    Compact-v2 deliberately keeps repair-anchor geometry in the shared
    stylesheet. Models occasionally duplicate that already-authoritative
    geometry inline. When the declared anchor has a complete, safe ``#id``
    geometry witness, removing the duplicate declarations is semantics
    preserving and lets the unchanged strict validator remain authoritative.
    Invalid or missing shared geometry, malformed inline CSS, nested anchors,
    and every non-geometry declaration continue through to the normal
    fail-closed validation path.
    """

    shared_geometry = _compact_shared_id_geometry(stylesheet)
    normalized_slides = slides
    normalized_slide_indexes: list[int] = []
    normalized_anchor_count = 0
    removed_declaration_count = 0
    removed_properties: set[str] = set()

    for index, raw in enumerate(slides):
        declared = raw.get("repair_anchor_ids")
        if (
            not isinstance(declared, list)
            or len(declared) != 2
            or any(not isinstance(identifier, str) for identifier in declared)
        ):
            continue
        body = raw.get("html_body")
        if not isinstance(body, str) or not body.strip():
            continue
        try:
            soup = BeautifulSoup(body, "html.parser")
        except Exception:  # noqa: BLE001 - strict validation reports the source error.
            continue

        slide_changed = False
        for element in soup.contents:
            if not isinstance(element, Tag):
                continue
            element_id = element.attrs.get("id")
            if (
                str(element.name).casefold() not in {"div", "section"}
                or not isinstance(element_id, str)
                or element_id not in declared
                or element_id not in shared_geometry
            ):
                continue
            inline_style = element.attrs.get("style")
            if not isinstance(inline_style, str) or not inline_style.strip():
                continue
            declarations = tinycss2.parse_declaration_list(
                inline_style,
                skip_comments=True,
                skip_whitespace=True,
            )
            if any(getattr(item, "type", "") != "declaration" for item in declarations):
                continue
            removed = [
                item
                for item in declarations
                if item.lower_name in _COMPACT_SOURCE_INLINE_DUPLICATE_PROPERTIES
                and (
                    item.lower_name != "margin"
                    or _compact_margin_declaration_is_literal_zero(item)
                )
            ]
            if not removed:
                continue
            retained = [
                item
                for item in declarations
                if item.lower_name not in _COMPACT_SOURCE_INLINE_DUPLICATE_PROPERTIES
            ]
            serialized = tinycss2.serialize(retained).strip()
            if serialized:
                element.attrs["style"] = serialized
            else:
                element.attrs.pop("style", None)
            slide_changed = True
            normalized_anchor_count += 1
            removed_declaration_count += len(removed)
            removed_properties.update(item.lower_name for item in removed)

        if not slide_changed:
            continue
        if normalized_slides is slides:
            normalized_slides = [dict(slide) for slide in slides]
        normalized_slides[index]["html_body"] = str(soup)
        normalized_slide_indexes.append(index)

    return normalized_slides, {
        "normalization_applied": normalized_anchor_count > 0,
        "normalized_slide_count": len(normalized_slide_indexes),
        "normalized_anchor_count": normalized_anchor_count,
        "removed_declaration_count": removed_declaration_count,
        "removed_property_names": sorted(removed_properties),
        "strict_validator_bypassed": False,
        "candidate_compile_changed": False,
        "raw_content_excluded": True,
    }


class _InlineFontFallbackNormalizer(HTMLParser):
    """Collect byte-local start-tag rewrites without serializing the DOM."""

    def __init__(self, source: str) -> None:
        super().__init__(convert_charrefs=False)
        self._source = source
        self._line_offsets = [0, *(match.end() for match in re.finditer(r"\n", source))]
        self.replacements: list[tuple[int, int, str]] = []
        self.normalized_attribute_count = 0
        self.normalized_declaration_count = 0

    def handle_starttag(
        self,
        _tag: str,
        _attrs: list[tuple[str, str | None]],
    ) -> None:
        self._collect_start_tag()

    def handle_startendtag(
        self,
        _tag: str,
        _attrs: list[tuple[str, str | None]],
    ) -> None:
        self._collect_start_tag()

    def _collect_start_tag(self) -> None:
        raw_tag = self.get_starttag_text()
        if not isinstance(raw_tag, str) or not raw_tag:
            return
        normalized_tag, attribute_count, declaration_count = (
            _normalize_start_tag_inline_font_fallbacks(raw_tag)
        )
        if not declaration_count:
            return
        line, column = self.getpos()
        if line < 1 or line > len(self._line_offsets):
            return
        start = self._line_offsets[line - 1] + column
        end = start + len(raw_tag)
        if self._source[start:end] != raw_tag:
            return
        self.replacements.append((start, end, normalized_tag))
        self.normalized_attribute_count += attribute_count
        self.normalized_declaration_count += declaration_count


def _normalize_start_tag_inline_font_fallbacks(
    raw_tag: str,
) -> tuple[str, int, int]:
    replacements: list[tuple[int, int, str]] = []
    normalized_declaration_count = 0
    for value_start, value_end in _quoted_style_attribute_value_spans(raw_tag):
        value = raw_tag[value_start:value_end]
        normalized, declaration_count = _normalize_inline_font_style_value(value)
        if not declaration_count:
            continue
        replacements.append((value_start, value_end, normalized))
        normalized_declaration_count += declaration_count
    if not replacements:
        return raw_tag, 0, 0
    normalized_tag = raw_tag
    for start, end, normalized in reversed(replacements):
        normalized_tag = normalized_tag[:start] + normalized + normalized_tag[end:]
    return normalized_tag, len(replacements), normalized_declaration_count


def _quoted_style_attribute_value_spans(
    raw_tag: str,
) -> tuple[tuple[int, int], ...]:
    """Locate real quoted style attributes while skipping other values."""

    length = len(raw_tag)
    index = 1 if raw_tag.startswith("<") else 0
    if index < length and raw_tag[index] == "/":
        index += 1
    while (
        index < length
        and not raw_tag[index].isspace()
        and raw_tag[index] not in {"/", ">"}
    ):
        index += 1
    spans: list[tuple[int, int]] = []
    while index < length:
        while index < length and raw_tag[index].isspace():
            index += 1
        if index >= length or raw_tag[index] == ">":
            break
        if raw_tag[index] == "/" and index + 1 < length and raw_tag[index + 1] == ">":
            break
        name_start = index
        while (
            index < length
            and not raw_tag[index].isspace()
            and raw_tag[index] not in {"=", "/", ">"}
        ):
            index += 1
        name = raw_tag[name_start:index].casefold()
        while index < length and raw_tag[index].isspace():
            index += 1
        if index >= length or raw_tag[index] != "=":
            if index == name_start:
                index += 1
            continue
        index += 1
        while index < length and raw_tag[index].isspace():
            index += 1
        if index >= length:
            break
        quote = raw_tag[index]
        if quote not in {'"', "'"}:
            while (
                index < length
                and not raw_tag[index].isspace()
                and raw_tag[index] not in {"/", ">"}
            ):
                index += 1
            continue
        value_start = index + 1
        value_end = raw_tag.find(quote, value_start)
        if value_end < 0:
            break
        if name == "style":
            spans.append((value_start, value_end))
        index = value_end + 1
    return tuple(spans)


def _normalize_inline_font_style_value(value: str) -> tuple[str, int]:
    try:
        declarations = tinycss2.parse_declaration_list(
            value,
            skip_comments=False,
            skip_whitespace=False,
        )
    except (RecursionError, TypeError, ValueError):
        return value, 0
    if any(
        getattr(item, "type", "") in {"at-rule", "error", "qualified-rule"}
        for item in declarations
    ):
        return value, 0
    normalized_count = 0
    for declaration in declarations:
        if (
            getattr(declaration, "type", "") != "declaration"
            or declaration.lower_name not in {"font", "font-family"}
        ):
            continue
        family_tokens = _font_family_tokens(declaration)
        if family_tokens is None:
            continue
        normalized_families = _portable_font_family_fallbacks(family_tokens)
        if normalized_families is None:
            continue
        family_start = len(declaration.value) - len(family_tokens)
        prefix = tinycss2.serialize(declaration.value[:family_start])
        declaration.value = tinycss2.parse_component_value_list(
            prefix + normalized_families
        )
        normalized_count += 1
    if not normalized_count:
        return value, 0
    normalized = tinycss2.serialize(declarations)
    if (
        not value.rstrip().endswith(";")
        and normalized.rstrip().endswith(";")
    ):
        trailing = normalized[len(normalized.rstrip()) :]
        normalized = normalized.rstrip()[:-1] + trailing
    return normalized, normalized_count


def _normalize_compact_v2_inline_font_fallbacks(
    slides: list[dict[str, Any]],
) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    """Remove unsupported secondary font families from fresh inline CSS."""

    normalized_slides = slides
    normalized_slide_count = 0
    normalized_attribute_count = 0
    normalized_declaration_count = 0
    for index, raw in enumerate(slides):
        body = raw.get("html_body")
        if not isinstance(body, str) or not body.strip():
            continue
        normalizer = _InlineFontFallbackNormalizer(body)
        try:
            normalizer.feed(body)
            normalizer.close()
        except (RecursionError, TypeError, ValueError):
            continue
        if not normalizer.replacements:
            continue
        normalized_body = body
        for start, end, replacement in reversed(normalizer.replacements):
            normalized_body = (
                normalized_body[:start] + replacement + normalized_body[end:]
            )
        if normalized_slides is slides:
            normalized_slides = [dict(slide) for slide in slides]
        normalized_slides[index]["html_body"] = normalized_body
        normalized_slide_count += 1
        normalized_attribute_count += normalizer.normalized_attribute_count
        normalized_declaration_count += normalizer.normalized_declaration_count
    return normalized_slides, {
        "normalization_applied": normalized_declaration_count > 0,
        "normalized_slide_count": normalized_slide_count,
        "normalized_attribute_count": normalized_attribute_count,
        "normalized_declaration_count": normalized_declaration_count,
        "unsafe_primary_accepted": False,
        "strict_validator_bypassed": False,
        "candidate_compile_changed": False,
        "raw_content_excluded": True,
    }


def _validate_authoring_inputs(
    deck: DeckBuild,
    slides: list[dict[str, Any]],
    *,
    allow_repair_overlay: bool = False,
) -> None:
    compact_mode = bool(deck.deck_stylesheet)
    stylesheet = deck.deck_stylesheet or ""
    if len(stylesheet.encode("utf-8")) > 24 * 1024:
        raise _authoring_failure("deck_stylesheet exceeds 24576 bytes.")
    if "</style" in stylesheet.lower():
        raise _authoring_failure("deck_stylesheet contains a forbidden closing style tag.")
    if deck.deck_authoring_contract == "compact_model_html_v2":
        if len(stylesheet.encode("utf-8")) > 8 * 1024:
            raise _authoring_failure("deck_stylesheet exceeds the compact-v2 8192-byte limit.")
        normalized_stylesheet = _normalize_compact_pptx_stylesheet_font_fallbacks(stylesheet)
        if normalized_stylesheet != stylesheet:
            deck.deck_stylesheet = normalized_stylesheet
        stylesheet = normalized_stylesheet
        deck.deck_stylesheet_hash = hashlib.sha256(stylesheet.encode("utf-8")).hexdigest() if stylesheet else None
    if deck.deck_authoring_contract == "compact_model_html_v2":
        _validate_compact_pptx_font_contract(stylesheet, slides)
        if not allow_repair_overlay:
            _validate_compact_source_addressability(stylesheet, slides)
    total_bytes = len(stylesheet.encode("utf-8"))
    for index, raw in enumerate(slides):
        total_bytes += _validate_slide_authoring_input(
            compact_mode=compact_mode,
            index=index,
            raw=raw,
        )
    if total_bytes > 128 * 1024:
        raise _authoring_failure("Deck authoring payload exceeds 131072 bytes.")
    if deck.deck_authoring_contract == "compact_model_html_v2":
        _validate_v2_authoring_sizes(deck, slides)


def _validate_compact_pptx_font_contract(stylesheet: str, slides: list[dict[str, Any]]) -> None:
    base_font_declared = _validate_css_font_declarations(
        stylesheet,
        label="deck_stylesheet",
        require_base=True,
    )
    if not base_font_declared:
        raise _authoring_failure(
            "deck_stylesheet must set Cambria, Calibri, or Arial on main, body, html, .slide-root, or the universal selector."
        )
    for index, raw in enumerate(slides):
        _validate_css_font_declarations(
            str(raw.get("slide_css") or ""),
            label=f"slides[{index}].slide_css",
            require_base=False,
        )
        html_body = str(raw.get("html_body") or "")
        for inline_style in _inline_style_values(html_body):
            _validate_font_declaration_list(
                inline_style,
                label=f"slides[{index}].html_body inline style",
                base_selector=False,
            )


def _compact_margin_declaration_is_literal_zero(declaration: Any) -> bool:
    tokens = [
        token
        for token in declaration.value
        if getattr(token, "type", "") not in {"comment", "whitespace"}
    ]
    maximum = 4 if declaration.lower_name == "margin" else 1
    if not 1 <= len(tokens) <= maximum:
        return False
    for token in tokens:
        if getattr(token, "type", "") not in {
            "dimension",
            "number",
            "percentage",
        }:
            return False
        try:
            value = float(token.value)
        except (TypeError, ValueError):
            return False
        if not math.isfinite(value) or value != 0:
            return False
    return True


def _compact_margin_issue_in_declarations(
    declarations: list[Any],
) -> str | None:
    for declaration in declarations:
        if getattr(declaration, "type", "") != "declaration":
            continue
        name = declaration.lower_name
        if name in _COMPACT_SOURCE_LOGICAL_MARGIN_PROPERTIES:
            return "logical"
        if name in _COMPACT_SOURCE_VENDOR_MARGIN_PROPERTIES:
            return "vendor"
        if (
            name in _COMPACT_SOURCE_PHYSICAL_MARGIN_PROPERTIES
            and not _compact_margin_declaration_is_literal_zero(declaration)
        ):
            return "physical_not_literal_zero"
    return None


def _compact_declared_anchor_margin_issue(
    stylesheet: str,
    soup: BeautifulSoup,
    anchors: tuple[Tag, ...],
) -> str | None:
    anchor_markers = {id(anchor) for anchor in anchors}
    rules = tinycss2.parse_stylesheet(
        stylesheet,
        skip_comments=True,
        skip_whitespace=True,
    )
    for rule in rules:
        if getattr(rule, "type", "") != "qualified-rule":
            continue
        selector = tinycss2.serialize(rule.prelude).strip()
        try:
            matches_anchor = any(
                id(match) in anchor_markers for match in soup.select(selector)
            )
        except Exception:
            continue
        if not matches_anchor:
            continue
        declarations = tinycss2.parse_declaration_list(
            rule.content,
            skip_comments=True,
            skip_whitespace=True,
        )
        issue = _compact_margin_issue_in_declarations(declarations)
        if issue is not None:
            return issue
    for anchor in anchors:
        inline_style = anchor.attrs.get("style")
        if not isinstance(inline_style, str) or not inline_style.strip():
            continue
        declarations = tinycss2.parse_declaration_list(
            inline_style,
            skip_comments=True,
            skip_whitespace=True,
        )
        issue = _compact_margin_issue_in_declarations(declarations)
        if issue is not None:
            return issue
    return None


def _validate_compact_source_addressability(
    stylesheet: str,
    slides: list[dict[str, Any]],
) -> None:
    """Require a source pair that the sealed DQ2 geometry witness can address.

    This intentionally accepts only a small, mechanically provable subset of
    compact-v2: two direct section/div semantic text owners with short
    slide-unique IDs and complete bounded geometry in standalone shared-
    stylesheet ID rules.
    The existing strict DQ2 witness remains the final proof of cascade safety.
    """

    shared_geometry = _compact_shared_id_geometry(stylesheet)
    for index, raw in enumerate(slides):
        body = str(raw.get("html_body") or "")
        slide_css = str(raw.get("slide_css") or "")
        if slide_css.strip():
            raise _authoring_failure(
                f"slides[{index}].slide_css must be empty for compact-v2 source addressability; "
                "put shared anchor geometry in deck_stylesheet."
            )
        declared_anchor_ids = raw.get("repair_anchor_ids")
        if (
            not isinstance(declared_anchor_ids, list)
            or len(declared_anchor_ids) != 2
            or any(
                not isinstance(identifier, str)
                or _COMPACT_SOURCE_ID_RE.fullmatch(identifier) is None
                for identifier in declared_anchor_ids
            )
            or len(set(declared_anchor_ids)) != 2
        ):
            raise _authoring_failure(
                f"slides[{index}].repair_anchor_ids must declare exactly two distinct short HTML ids matching "
                "[a-z][a-z0-9_-]{0,31}."
            )
        declared_anchor_id_set = set(declared_anchor_ids)
        try:
            soup = BeautifulSoup(body, "html.parser")
        except Exception as exc:
            raise _authoring_failure(
                f"slides[{index}].html_body cannot prove compact-v2 source addressability."
            ) from exc

        body_ids: list[str] = []
        for element in soup.find_all(True):
            element_id = element.attrs.get("id")
            if not isinstance(element_id, str) or not element_id:
                continue
            if element_id in body_ids:
                raise _authoring_failure(
                    f"slides[{index}].html_body contains duplicate HTML id '{safe_excerpt(element_id, limit=40)}'; "
                    "HTML IDs must be unique within one slide."
                )
            body_ids.append(element_id)
        assembled_soup = BeautifulSoup(
            assemble_compact_slide_html(
                deck_stylesheet="",
                html_body=body,
                slide_css="",
            ),
            "html.parser",
        )
        assembled_anchors = tuple(
            anchor
            for identifier in declared_anchor_ids
            for anchor in (assembled_soup.find(id=identifier),)
            if isinstance(anchor, Tag)
        )
        if (
            len(assembled_anchors) == 2
            and _compact_declared_anchor_margin_issue(
                stylesheet,
                assembled_soup,
                assembled_anchors,
            )
            is not None
        ):
            raise _authoring_failure(
                f"slides[{index}].html_body or deck_stylesheet has a margin declaration matching a declared "
                "repair anchor that is auto, nonzero, or otherwise not literal zero, or uses a logical/vendor "
                "margin property. Remove the unsafe declaration; do not override it with a later margin:0 reset."
            )
        anchors = tuple(
            element
            for element in soup.contents
            if isinstance(element, Tag)
            and _compact_source_anchor_is_eligible(element, shared_geometry)
        )
        eligible_anchor_ids = {str(anchor.attrs["id"]) for anchor in anchors}
        if not declared_anchor_id_set <= eligible_anchor_ids:
            raise _authoring_failure(
                f"slides[{index}].html_body must contain both repair anchors declared by "
                "repair_anchor_ids as independent visible text-bearing section/div direct children of main with "
                "short unique HTML IDs, nonempty data-deck-id/data-deck-role, data-deck-required='true', and "
                "complete safe absolute px geometry in deck_stylesheet."
            )
        # Import locally to avoid making ordinary deck-service module loading
        # depend on the campaign repair author. The witness itself remains
        # unchanged and authoritative.
        from deerflow.sophia.deck_design_lift.repair_author import (  # noqa: PLC0415
            _strict_geometry_effective_box,
            _strict_geometry_source_witness,
            _strict_geometry_translation_origin,
        )

        translation_blocked = False
        for anchor in assembled_anchors:
            effective_box = _strict_geometry_effective_box(
                anchor,
                assembled_soup,
                deck_css=stylesheet,
                baseline_slide_css=slide_css,
            )
            if (
                effective_box is not None
                and _strict_geometry_translation_origin(effective_box) is None
            ):
                translation_blocked = True
                break
        if translation_blocked:
            raise _authoring_failure(
                f"slides[{index}].deck_stylesheet effective repair-anchor geometry must remain wholly inside the "
                "1920x1080 canvas and leave at least 8px translation clearance on one horizontal or vertical side "
                "of each whole anchor. Adjust only its literal px left/top/width/height."
            )

        anchors = tuple(
            anchor
            for anchor in anchors
            if _compact_anchor_has_visible_text(
                anchor,
                soup,
                deck_css=stylesheet,
                baseline_slide_css=slide_css,
            )
        )
        visible_anchor_ids = {str(anchor.attrs["id"]) for anchor in anchors}
        if not declared_anchor_id_set <= visible_anchor_ids:
            raise _authoring_failure(
                f"slides[{index}].html_body must keep both declared compact-v2 repair anchors visible and "
                "text-bearing; hidden, inert, aria-hidden, display:none, visibility:hidden, and opacity below 1 "
                "do not qualify."
            )
        semantic_ids = [str(anchor.attrs["data-deck-id"]).strip() for anchor in anchors]
        if len(set(semantic_ids)) != len(semantic_ids):
            duplicate = next(value for value in semantic_ids if semantic_ids.count(value) > 1)
            raise _authoring_failure(
                f"slides[{index}].html_body reuses anchor data-deck-id "
                f"'{safe_excerpt(duplicate, limit=40)}'; eligible compact-v2 anchors must use distinct "
                "data-deck-id values within the slide."
            )

        shared_witness = _strict_geometry_source_witness(
            body=body,
            baseline_slide_css="",
            deck_css=stylesheet,
            minimum=2,
            target_element_ids=frozenset(declared_anchor_id_set),
        )
        effective_witness = shared_witness
        if slide_css.strip():
            effective_witness = _strict_geometry_source_witness(
                body=body,
                baseline_slide_css=slide_css,
                deck_css=stylesheet,
                minimum=2,
                target_element_ids=frozenset(declared_anchor_id_set),
            )
        shared_witness_ids = _compact_witness_anchor_ids(shared_witness)
        effective_witness_ids = _compact_witness_anchor_ids(effective_witness)
        if (
            shared_witness_ids is None
            or effective_witness_ids is None
            or shared_witness_ids != declared_anchor_id_set
            or effective_witness_ids != declared_anchor_id_set
        ):
            raise _authoring_failure(
                f"slides[{index}] cannot prove two compact-v2 source geometry anchors; "
                "keep their complete baseline geometry in deck_stylesheet and any authenticated repair overlay safe. "
                "Avoid any other matching nonzero, logical, or vendor margin rule."
            )


def _compact_witness_anchor_ids(witness: str | None) -> set[str] | None:
    if witness is None:
        return None
    rules = tinycss2.parse_stylesheet(witness, skip_comments=True, skip_whitespace=True)
    selectors: set[str] = set()
    for rule in rules:
        if getattr(rule, "type", "") != "qualified-rule":
            return None
        selector = tinycss2.serialize(rule.prelude).strip()
        match = re.fullmatch(r"#([a-z][a-z0-9_-]{0,31})", selector)
        if match is None:
            return None
        selectors.add(match.group(1))
    return selectors if len(rules) == 2 and len(selectors) == 2 else None


def _compact_anchor_has_visible_text(
    anchor: Tag,
    soup: BeautifulSoup,
    *,
    deck_css: str,
    baseline_slide_css: str,
) -> bool:
    from deerflow.sophia.deck_design_lift.repair_author import (  # noqa: PLC0415
        _authenticated_display_generates_box,
        _authenticated_element_is_fully_opaque,
        _authenticated_target_is_visible,
    )

    for text_node in anchor.find_all(string=True):
        if not str(text_node).strip() or not isinstance(text_node.parent, Tag):
            continue
        current: Tag | None = text_node.parent
        ancestry: list[Tag] = []
        while current is not None:
            ancestry.append(current)
            if current is anchor:
                break
            current = current.parent if isinstance(current.parent, Tag) else None
        if not ancestry or ancestry[-1] is not anchor:
            continue
        if any(
            "hidden" in element.attrs
            or "inert" in element.attrs
            or str(element.attrs.get("aria-hidden") or "").strip().casefold() == "true"
            or not _authenticated_display_generates_box(
                element,
                soup,
                deck_css=deck_css,
                baseline_slide_css=baseline_slide_css,
            )
            or not _authenticated_target_is_visible(
                element,
                soup,
                deck_css=deck_css,
                baseline_slide_css=baseline_slide_css,
            )
            or not _authenticated_element_is_fully_opaque(
                element,
                soup,
                deck_css=deck_css,
                baseline_slide_css=baseline_slide_css,
            )
            for element in ancestry
        ):
            continue
        return True
    return False


def _compact_source_anchor_is_eligible(
    element: Tag,
    shared_geometry: dict[str, tuple[float, float, float, float]],
) -> bool:
    element_id = element.attrs.get("id")
    if not isinstance(element_id, str) or not _COMPACT_SOURCE_ID_RE.fullmatch(element_id):
        return False
    if element_id not in shared_geometry:
        return False
    if str(element.name).casefold() not in {"div", "section"}:
        return False
    if not str(element.attrs.get("data-deck-id") or "").strip():
        return False
    if not str(element.attrs.get("data-deck-role") or "").strip():
        return False
    if str(element.attrs.get("data-deck-required") or "").strip().casefold() != "true":
        return False
    if not element.get_text(" ", strip=True):
        return False
    inline_style = element.attrs.get("style")
    if isinstance(inline_style, str) and inline_style.strip():
        declarations = tinycss2.parse_declaration_list(
            inline_style,
            skip_comments=True,
            skip_whitespace=True,
        )
        if any(
            getattr(item, "type", "") != "declaration"
            or item.lower_name in {"position", "box-sizing", "margin", *_COMPACT_SOURCE_GEOMETRY_PROPERTIES}
            for item in declarations
        ):
            return False
    return True


def _compact_shared_id_geometry(stylesheet: str) -> dict[str, tuple[float, float, float, float]]:
    geometry: dict[str, tuple[float, float, float, float]] = {}
    rules = tinycss2.parse_stylesheet(stylesheet, skip_comments=True, skip_whitespace=True)
    for rule in rules:
        if getattr(rule, "type", "") != "qualified-rule":
            continue
        selector = tinycss2.serialize(rule.prelude).strip()
        match = re.fullmatch(r"#([a-z][a-z0-9_-]{0,31})", selector)
        if match is None:
            continue
        declarations = tinycss2.parse_declaration_list(
            rule.content,
            skip_comments=True,
            skip_whitespace=True,
        )
        if any(getattr(item, "type", "") != "declaration" for item in declarations):
            continue
        if any(bool(getattr(item, "important", False)) for item in declarations):
            continue
        final = {item.lower_name: item for item in declarations}
        if _css_single_identifier(final.get("position")) != "absolute":
            continue
        if _css_single_identifier(final.get("box-sizing")) != "border-box":
            continue
        if _css_zero_value(final.get("margin")) is not True:
            continue
        values = tuple(
            _css_absolute_px_value(final.get(property_name))
            for property_name in _COMPACT_SOURCE_GEOMETRY_PROPERTIES
        )
        if any(value is None for value in values):
            continue
        left, top, width, height = (float(value) for value in values if value is not None)
        if (
            left < 0
            or top < 0
            or width < 48
            or height < 24
            or left + width > 1920
            or top + height > 1080
        ):
            continue
        geometry[match.group(1)] = (left, top, width, height)
    return geometry


def _css_single_identifier(declaration: Any | None) -> str | None:
    if declaration is None:
        return None
    tokens = [
        token
        for token in declaration.value
        if getattr(token, "type", "") not in {"comment", "whitespace"}
    ]
    if len(tokens) != 1 or getattr(tokens[0], "type", "") != "ident":
        return None
    return str(tokens[0].value).casefold()


def _css_absolute_px_value(declaration: Any | None) -> float | None:
    if declaration is None:
        return None
    tokens = [
        token
        for token in declaration.value
        if getattr(token, "type", "") not in {"comment", "whitespace"}
    ]
    if len(tokens) != 1:
        return None
    token = tokens[0]
    if getattr(token, "type", "") == "number" and float(token.value) == 0:
        return 0.0
    if getattr(token, "type", "") != "dimension" or str(getattr(token, "unit", "")).casefold() != "px":
        return None
    value = float(token.value)
    return value if value == value and abs(value) != float("inf") else None


def _css_zero_value(declaration: Any | None) -> bool | None:
    value = _css_absolute_px_value(declaration)
    return value == 0 if value is not None else None


def _validate_css_font_declarations(css: str, *, label: str, require_base: bool) -> bool:
    try:
        rules = tinycss2.parse_stylesheet(css, skip_comments=True, skip_whitespace=True)
    except (RecursionError, TypeError, ValueError) as exc:
        raise _authoring_failure(
            f"{label} contains excessively nested or malformed CSS, which is unsupported for compact-v2 authoring."
        ) from exc
    if any(getattr(rule, "type", "") == "error" for rule in rules):
        raise _authoring_failure(
            f"{label} contains malformed CSS, which is unsupported for compact PPTX font validation."
        )
    if any(getattr(rule, "type", "") == "at-rule" for rule in rules):
        raise _authoring_failure(
            f"{label} uses a nested at-rule CSS construct, which is unsupported for compact-v2 authoring."
        )
    try:
        return _validate_css_font_rule_nodes(
            rules,
            label=label,
            allow_base=require_base,
        )
    except DeckBuildFailure:
        raise
    except (RecursionError, TypeError, ValueError) as exc:
        raise _authoring_failure(
            f"{label} contains excessively nested or malformed CSS, which is unsupported for compact-v2 authoring."
        ) from exc


def _normalize_compact_pptx_stylesheet_font_fallbacks(css: str) -> str:
    """Normalize non-portable secondary families in the shared stylesheet.

    The first family still determines whether the declaration is admissible. A
    model-authored unsafe primary therefore continues to fail closed, while a
    harmless browser-only fallback cannot spend the single bounded repair.
    """

    try:
        rules = tinycss2.parse_stylesheet(css, skip_comments=False, skip_whitespace=False)
    except (RecursionError, TypeError, ValueError):
        return css
    if any(getattr(rule, "type", "") in {"at-rule", "error"} for rule in rules):
        return css
    try:
        changed = _normalize_css_font_fallback_rule_nodes(rules)
        return tinycss2.serialize(rules) if changed else css
    except (RecursionError, TypeError, ValueError):
        return css


def _normalize_css_font_fallback_rule_nodes(rules: list[Any]) -> bool:
    changed = False
    for rule in rules:
        rule_type = getattr(rule, "type", "")
        if rule_type != "qualified-rule":
            continue
        declarations = tinycss2.parse_declaration_list(
            rule.content,
            skip_comments=False,
            skip_whitespace=False,
        )
        if any(getattr(item, "type", "") in {"at-rule", "error", "qualified-rule"} for item in declarations):
            continue
        rule_changed = False
        for declaration in declarations:
            if getattr(declaration, "type", "") != "declaration" or declaration.lower_name not in {
                "font",
                "font-family",
            }:
                continue
            family_tokens = _font_family_tokens(declaration)
            if family_tokens is None:
                continue
            normalized_families = _portable_font_family_fallbacks(family_tokens)
            if normalized_families is None:
                continue
            family_start = len(declaration.value) - len(family_tokens)
            prefix = tinycss2.serialize(declaration.value[:family_start])
            declaration.value = tinycss2.parse_component_value_list(prefix + normalized_families)
            rule_changed = True
        if rule_changed:
            rule.content = tinycss2.parse_component_value_list(tinycss2.serialize(declarations))
            changed = True
    return changed


def _portable_font_family_fallbacks(tokens: list[Any]) -> str | None:
    groups: list[list[Any]] = [[]]
    for token in tokens:
        if token.type == "literal" and token.value == ",":
            groups.append([])
        else:
            groups[-1].append(token)
    serialized_groups = [tinycss2.serialize(group).strip() for group in groups]
    if not serialized_groups or any(not value for value in serialized_groups):
        return None
    names = _font_family_names(tokens)
    if len(names) != len(serialized_groups) or names[0] not in _SAFE_PPTX_FONTS:
        return None
    allowed = _SAFE_PPTX_FONTS | _GENERIC_CSS_FONTS
    retained = [value for value, name in zip(serialized_groups, names, strict=True) if name in allowed]
    if len(retained) == len(serialized_groups):
        return None
    return ", ".join(retained)


def _validate_css_font_rule_nodes(rules: list[Any], *, label: str, allow_base: bool) -> bool:
    base_font_declared = False
    for rule in rules:
        rule_type = getattr(rule, "type", "")
        if rule_type == "error":
            raise _authoring_failure(
                f"{label} contains malformed CSS, which is unsupported for compact PPTX font validation."
            )
        if rule_type == "at-rule":
            raise _authoring_failure(
                f"{label} uses a nested at-rule CSS construct, which is unsupported for compact-v2 authoring."
            )
        if rule_type != "qualified-rule":
            continue
        selector = tinycss2.serialize(rule.prelude).strip()
        is_base = bool(allow_base and _selector_declares_base_font(selector))
        if _validate_font_declaration_list(
            tinycss2.serialize(rule.content),
            label=f"{label} selector {safe_excerpt(selector, limit=80)}",
            base_selector=is_base,
        ):
            base_font_declared = True
    return base_font_declared


def _selector_declares_base_font(selector: str) -> bool:
    return any(arm.strip().lower() in _BASE_FONT_SELECTORS for arm in selector.split(","))


class _InlineStyleCollector(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.styles: list[str] = []

    def handle_starttag(self, _tag: str, attrs: list[tuple[str, str | None]]) -> None:
        for name, value in attrs:
            if name.lower() == "style" and value:
                self.styles.append(value)


def _inline_style_values(html: str) -> list[str]:
    collector = _InlineStyleCollector()
    try:
        collector.feed(html)
    except Exception:
        return []
    return collector.styles


def _validate_font_declaration_list(css: str, *, label: str, base_selector: bool) -> bool:
    safe_explicit_family = False
    try:
        declarations = tinycss2.parse_declaration_list(css, skip_comments=True, skip_whitespace=True)
    except (RecursionError, TypeError, ValueError) as exc:
        raise _authoring_failure(
            f"{label} contains excessively nested or malformed CSS, which is unsupported for compact-v2 authoring."
        ) from exc
    for declaration in declarations:
        if getattr(declaration, "type", "") in {"at-rule", "error", "qualified-rule"}:
            raise _authoring_failure(
                f"{label} uses nested or malformed CSS rules, which are unsupported for compact PPTX font validation."
            )
        if getattr(declaration, "type", "") == "declaration" and declaration.lower_name == "all":
            raise _authoring_failure(
                f"{label} uses the CSS all shorthand, which can reset the required Office-safe font family."
            )
        if getattr(declaration, "type", "") != "declaration" or declaration.lower_name not in {"font", "font-family"}:
            continue
        family_tokens = _font_family_tokens(declaration)
        try:
            serialized = tinycss2.serialize(declaration.value).strip()
        except (RecursionError, TypeError, ValueError) as exc:
            raise _authoring_failure(
                f"{label} contains excessively nested or malformed CSS, which is unsupported for compact-v2 authoring."
            ) from exc
        if family_tokens is None and serialized.lower() == "inherit":
            continue
        if family_tokens is not None and _font_family_list_has_empty_group(family_tokens):
            raise _authoring_failure(
                f"{label} uses malformed PPTX {declaration.lower_name} '{safe_excerpt(serialized, limit=80)}'. "
                "Font-family lists cannot contain empty entries."
            )
        families = _font_family_names(family_tokens or [])
        if not families or families[0] not in _SAFE_PPTX_FONTS:
            raise _authoring_failure(
                f"{label} uses unsupported PPTX {declaration.lower_name} '{safe_excerpt(serialized, limit=80)}'. "
                "Use Cambria for headings and Calibri or Arial for body and utility text."
            )
        if any(family not in _SAFE_PPTX_FONTS | _GENERIC_CSS_FONTS for family in families[1:]):
            raise _authoring_failure(
                f"{label} uses a non-portable PPTX font fallback in '{safe_excerpt(serialized, limit=80)}'. "
                "Use only Cambria, Calibri, Arial, and generic fallback families."
            )
        safe_explicit_family = True
    return bool(base_selector and safe_explicit_family)


def _font_family_tokens(declaration: Any) -> list[Any] | None:
    tokens = list(declaration.value)
    if declaration.lower_name == "font-family":
        return tokens
    significant = [index for index, token in enumerate(tokens) if token.type not in {"comment", "whitespace"}]
    size_position = next(
        (
            position
            for position, index in enumerate(significant)
            if _is_css_font_size_token(tokens[index])
        ),
        None,
    )
    if size_position is None:
        return None
    family_position = size_position + 1
    if family_position < len(significant) and getattr(tokens[significant[family_position]], "value", None) == "/":
        family_position += 2
    if family_position >= len(significant):
        return []
    return tokens[significant[family_position] :]


def _is_css_font_size_token(token: Any) -> bool:
    if token.type in {"dimension", "percentage"}:
        return True
    return token.type == "ident" and str(token.value).lower() in _FONT_SIZE_KEYWORDS


def _font_family_names(tokens: list[Any]) -> list[str]:
    groups: list[list[Any]] = [[]]
    for token in tokens:
        if token.type == "literal" and token.value == ",":
            groups.append([])
        else:
            groups[-1].append(token)
    names: list[str] = []
    for group in groups:
        significant = [token for token in group if token.type not in {"comment", "whitespace"}]
        if len(significant) == 1 and significant[0].type in {"ident", "string"}:
            value = str(significant[0].value)
        else:
            value = tinycss2.serialize(group).strip().strip("\"'")
        normalized = re.sub(r"\s+", " ", value).strip().lower()
        if normalized:
            names.append(normalized)
    return names


def _font_family_list_has_empty_group(tokens: list[Any]) -> bool:
    groups: list[list[Any]] = [[]]
    for token in tokens:
        if token.type == "literal" and token.value == ",":
            groups.append([])
        else:
            groups[-1].append(token)
    return any(not [token for token in group if token.type not in {"comment", "whitespace"}] for group in groups)


def _validate_v2_authoring_sizes(deck: DeckBuild, slides: list[dict[str, Any]]) -> None:
    stylesheet = deck.deck_stylesheet or ""
    if len(stylesheet.encode("utf-8")) > 8 * 1024:
        raise _authoring_failure("deck_stylesheet exceeds the compact-v2 8192-byte limit.")
    body_sizes: list[int] = []
    for index, raw in enumerate(slides):
        body = str(raw.get("html_body") or "").strip()
        slide_css = str(raw.get("slide_css") or "").strip()
        body_size = len(body.encode("utf-8"))
        body_sizes.append(body_size)
        if body_size > COMPACT_V2_MAX_SLIDE_HTML_BODY_BYTES:
            raise _authoring_failure(
                f"slides[{index}].html_body exceeds the compact-v2 hard "
                f"{COMPACT_V2_MAX_SLIDE_HTML_BODY_BYTES}-byte limit."
            )
        if len(slide_css.encode("utf-8")) > 1024:
            raise _authoring_failure(f"slides[{index}].slide_css exceeds the compact-v2 1024-byte limit.")
    body_total = sum(body_sizes)
    body_budget = len(body_sizes) * COMPACT_V2_TARGET_SLIDE_HTML_BODY_BYTES
    if body_total > body_budget:
        raise _authoring_failure(
            f"slides.html_body_total is {body_total} bytes; compact-v2 aggregate budget is "
            f"{body_budget} bytes ({len(body_sizes)} slides x "
            f"{COMPACT_V2_TARGET_SLIDE_HTML_BODY_BYTES} bytes)."
        )


def _slide_authoring_sources(raw: dict[str, Any]) -> dict[str, str | None]:
    return {key: str(raw.get(key) or "").strip() or None for key in ("html_body", "slide_css", "html_source")}


def _validate_slide_authoring_input(*, compact_mode: bool, index: int, raw: dict[str, Any]) -> int:
    body = str(raw.get("html_body") or "").strip()
    slide_css = str(raw.get("slide_css") or "").strip()
    source = str(raw.get("html_source") or "").strip()
    _validate_slide_authoring_mode(compact_mode, index, body, slide_css, source)
    _validate_slide_authoring_tags(compact_mode, index, body, slide_css)
    _validate_slide_authoring_sizes(index, body, slide_css)
    return sum(len(value.encode("utf-8")) for value in (body, slide_css, source))


def _validate_slide_authoring_mode(
    compact_mode: bool,
    index: int,
    body: str,
    slide_css: str,
    source: str,
) -> None:
    if compact_mode and (not body or source):
        raise _authoring_failure(f"slides[{index}] must use html_body only in compact authoring mode.")
    if not compact_mode and (body or slide_css):
        raise _authoring_failure(f"slides[{index}] must use legacy html_source only when deck_stylesheet is absent.")


def _validate_slide_authoring_tags(compact_mode: bool, index: int, body: str, slide_css: str) -> None:
    if compact_mode and _fragment_contains_document_tag(body):
        raise _authoring_failure(f"slides[{index}].html_body contains a document-level tag.")
    if "</style" in slide_css.lower():
        raise _authoring_failure(f"slides[{index}].slide_css contains a forbidden closing style tag.")


def _validate_slide_authoring_sizes(index: int, body: str, slide_css: str) -> None:
    if len(body.encode("utf-8")) > 16 * 1024:
        raise _authoring_failure(f"slides[{index}].html_body exceeds 16384 bytes.")
    if len(slide_css.encode("utf-8")) > 8 * 1024:
        raise _authoring_failure(f"slides[{index}].slide_css exceeds 8192 bytes.")


def _fragment_contains_document_tag(body: str) -> bool:
    lowered = body.lower()
    return any(tag in lowered for tag in ("<html", "</html", "<head", "</head", "<body", "</body", "<style", "</style"))


def _now() -> str:
    return datetime.now(UTC).isoformat()


def _clean_text(value: object) -> str:
    return str(value or "").strip()


def _host_path(virtual_path: str, runtime: ToolRuntime) -> Path:
    return Path(replace_virtual_path(virtual_path, get_thread_data(runtime)))


def _composition_for_layout(layout_kind: str) -> str:
    return {
        "cover_hero": "hero visual with one confident focal system metaphor",
        "visual_left_text_right": "visual weighted to the left with clean negative space",
        "text_left_visual_right": "visual weighted to the right with clean negative space",
        "comparison_two_column": "two balanced visual fields with mostly text-free contrast",
        "timeline_flow": "clear process flow with shapes and spatial progression",
        "closing_summary": "calm synthesis visual with restrained emphasis",
    }.get(layout_kind, "single technical visual focus with restrained professional composition")


def _image_script_path() -> Path | None:
    candidates = [
        Path("/mnt/skills/public/image-generation/scripts/generate.py"),
        Path("/app/skills/public/image-generation/scripts/generate.py"),
        Path(__file__).resolve().parents[6] / "skills/public/image-generation/scripts/generate.py",
    ]
    for candidate in candidates:
        if candidate.is_file():
            return candidate
    return None


def _parse_batch_summary(stdout: str) -> dict[str, Any] | None:
    for line in stdout.splitlines():
        if line.startswith("IMAGEGEN_BATCH "):
            try:
                return json.loads(line.removeprefix("IMAGEGEN_BATCH ").strip())
            except ValueError:
                return None
    return None


def _parse_batch_item_progress(stdout: str) -> list[dict[str, Any]]:
    items: list[dict[str, Any]] = []
    for line in stdout.splitlines():
        if not line.startswith("IMAGEGEN_BATCH_ITEM "):
            continue
        try:
            payload = json.loads(line.removeprefix("IMAGEGEN_BATCH_ITEM ").strip())
        except ValueError:
            continue
        if isinstance(payload, dict):
            if "raw_error_excerpt" in payload:
                payload["raw_error_excerpt"] = safe_excerpt(payload["raw_error_excerpt"])
            items.append(payload)
    return items


def _image_generation_error_class_from_output(stderr: str, stdout: str, returncode: int) -> str:
    text = f"{stderr}\n{stdout}"
    match = re.search(r"IMAGEGEN_FAIL\s+reason=([A-Za-z0-9_:-]+)", text)
    if match:
        return match.group(1)
    if returncode == 2:
        return "process_exit"
    if "timed out" in text.lower() or "timeout" in text.lower():
        return "timeout"
    if "no bytes" in text.lower() or "empty output" in text.lower():
        return "empty_output"
    return "api_error"


def _summary_error_classes(summary: dict[str, Any]) -> set[str]:
    classes: set[str] = set()
    error_class = summary.get("error_class")
    if error_class:
        classes.add(str(error_class))
    histogram = summary.get("error_class_histogram")
    if isinstance(histogram, dict):
        classes.update(str(key) for key in histogram if key)
    items = summary.get("items")
    if isinstance(items, list):
        for item in items:
            if isinstance(item, dict) and item.get("error_class"):
                classes.add(str(item.get("error_class")))
    return classes


def _primary_error_class(summary: dict[str, Any]) -> str | None:
    classes = _summary_error_classes(summary)
    if classes:
        return sorted(classes)[0]
    return None


def _batch_summary_allows_serial_repair(summary: dict[str, Any]) -> bool:
    if summary.get("complete"):
        return False
    classes = _summary_error_classes(summary)
    if classes.intersection(_NON_REPAIRABLE_IMAGE_ERROR_CLASSES):
        return False
    if summary.get("summary_present", True):
        return bool(int(summary.get("requested") or 0) > int(summary.get("images_generated") or 0))
    return bool(summary.get("batch_attempted"))


def _safe_batch_summary(summary: dict[str, Any]) -> dict[str, Any]:
    payload = dict(summary)
    if "raw_error_excerpt" in payload:
        payload["raw_error_excerpt"] = safe_excerpt(payload["raw_error_excerpt"])
    items = payload.get("items")
    if isinstance(items, list):
        safe_items: list[Any] = []
        for item in items:
            if not isinstance(item, dict):
                safe_items.append(item)
                continue
            safe_item = dict(item)
            if "raw_error_excerpt" in safe_item:
                safe_item["raw_error_excerpt"] = safe_excerpt(safe_item["raw_error_excerpt"])
            for path_key in ("output_file", "output_path", "prompt_file"):
                if path_key in safe_item:
                    safe_item[path_key] = basename(safe_item[path_key])
            safe_items.append(safe_item)
        payload["items"] = safe_items
    return payload


def _file_hash(path: Path) -> str | None:
    try:
        return hashlib.sha256(path.read_bytes()).hexdigest()[:16]
    except OSError:
        return None


def _file_size(path: Path) -> int | None:
    try:
        return path.stat().st_size
    except OSError:
        return None
