"""Deterministic contrast analysis for final native PPTX text shapes."""

from __future__ import annotations

from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation


@dataclass
class NativeContrastIssue:
    selector: str
    shape_name: str
    text_excerpt: str
    foreground: str | None
    background: str | None
    contrast_ratio: float | None
    required_ratio: float | None
    required_semantic: bool
    indeterminate: bool
    reason: str

    def to_dict(self) -> dict[str, Any]:
        return asdict(self)


def evaluate_native_contrast(
    *,
    pptx_path: str | Path,
    source_element_map: dict[str, Any] | None = None,
) -> dict[str, Any]:
    presentation = Presentation(str(pptx_path))
    required_names = _required_shape_names(source_element_map)
    issues: list[NativeContrastIssue] = []
    checked_runs = 0
    for slide_index, slide in enumerate(presentation.slides, start=1):
        selector = f"slide:{slide_index}"
        slide_background = _solid_fill_rgb(slide.background.fill)
        prior_shapes: list[Any] = []
        for shape in slide.shapes:
            shape_checked, shape_issues = _shape_contrast_findings(
                shape=shape,
                prior_shapes=prior_shapes,
                slide_background=slide_background,
                selector=selector,
                required_names=required_names.get(selector, set()),
            )
            checked_runs += shape_checked
            issues.extend(shape_issues)
            prior_shapes.append(shape)
    required_issues = [issue for issue in issues if issue.required_semantic]
    return {
        "passed": not required_issues,
        "checked_run_count": checked_runs,
        "issue_count": len(issues),
        "required_issue_count": len(required_issues),
        "indeterminate_required_count": sum(1 for issue in required_issues if issue.indeterminate),
        "issues": [issue.to_dict() for issue in issues],
    }


def _shape_contrast_findings(
    *,
    shape: Any,
    prior_shapes: list[Any],
    slide_background: str | None,
    selector: str,
    required_names: set[str],
) -> tuple[int, list[NativeContrastIssue]]:
    if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
        return 0, []
    shape_name = str(shape.name or "")
    required = shape_name in required_names
    background, background_reason = _effective_background(shape, prior_shapes, slide_background)
    issues: list[NativeContrastIssue] = []
    checked_runs = 0
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = str(run.text or "").strip()
            if not text:
                continue
            checked_runs += 1
            issue = _run_contrast_issue(
                run=run,
                text=text,
                selector=selector,
                shape_name=shape_name,
                background=background,
                background_reason=background_reason,
                required=required,
            )
            if issue is not None:
                issues.append(issue)
    return checked_runs, issues


def _run_contrast_issue(
    *,
    run: Any,
    text: str,
    selector: str,
    shape_name: str,
    background: str | None,
    background_reason: str,
    required: bool,
) -> NativeContrastIssue | None:
    foreground = _font_rgb(run.font)
    size_pt = float(run.font.size.pt) if run.font.size is not None else None
    required_ratio = _required_contrast_ratio(size_pt=size_pt, bold=bool(run.font.bold))
    if foreground is None or background is None or size_pt is None:
        if not required:
            return None
        missing = "foreground" if foreground is None else "font_size" if size_pt is None else background_reason
        return NativeContrastIssue(
            selector=selector,
            shape_name=shape_name,
            text_excerpt=text[:80],
            foreground=foreground,
            background=background,
            contrast_ratio=None,
            required_ratio=required_ratio,
            required_semantic=True,
            indeterminate=True,
            reason=f"contrast_indeterminate:{missing}",
        )
    ratio = _contrast_ratio(foreground, background)
    if ratio + 1e-6 >= required_ratio:
        return None
    return NativeContrastIssue(
        selector=selector,
        shape_name=shape_name,
        text_excerpt=text[:80],
        foreground=foreground,
        background=background,
        contrast_ratio=round(ratio, 3),
        required_ratio=required_ratio,
        required_semantic=required,
        indeterminate=False,
        reason="contrast_below_threshold",
    )


def _required_contrast_ratio(*, size_pt: float | None, bold: bool) -> float:
    if size_pt is not None and (size_pt >= 18 or (bold and size_pt >= 14)):
        return 3.0
    return 4.5


def _required_shape_names(source_element_map: dict[str, Any] | None) -> dict[str, set[str]]:
    result: dict[str, set[str]] = {}
    slides = source_element_map.get("slides") if isinstance(source_element_map, dict) else None
    for selector, slide in (slides or {}).items():
        elements = slide.get("elements") if isinstance(slide, dict) else None
        for record in (elements or {}).values():
            if not isinstance(record, dict) or not record.get("source_required"):
                continue
            result.setdefault(str(selector), set()).update(str(name) for name in record.get("shape_names") or [])
    return result


def _effective_background(shape: Any, prior_shapes: list[Any], slide_background: str | None) -> tuple[str | None, str]:
    own = _solid_fill_rgb(getattr(shape, "fill", None))
    if own is not None:
        return own, "own_fill"
    candidates: list[str] = []
    for candidate in prior_shapes:
        fill = _solid_fill_rgb(getattr(candidate, "fill", None))
        if fill is None:
            continue
        if _overlap_fraction(shape, candidate) >= 0.80:
            candidates.append(fill)
    if candidates:
        return candidates[-1], "containing_shape"
    if slide_background is not None:
        return slide_background, "slide_background"
    return None, "unknown_background"


def _solid_fill_rgb(fill: Any) -> str | None:
    if fill is None:
        return None
    try:
        if str(fill.type).upper() not in {"SOLID (1)", "SOLID", "1"} and int(fill.type) != 1:
            return None
    except (TypeError, ValueError, AttributeError):
        return None
    try:
        rgb = fill.fore_color.rgb
    except (AttributeError, TypeError, ValueError):
        return None
    return str(rgb).upper() if rgb is not None else None


def _font_rgb(font: Any) -> str | None:
    try:
        rgb = font.color.rgb
    except (AttributeError, TypeError, ValueError):
        return None
    return str(rgb).upper() if rgb is not None else None


def _overlap_fraction(inner: Any, outer: Any) -> float:
    il, it = int(inner.left), int(inner.top)
    ir, ib = il + int(inner.width), it + int(inner.height)
    ol, ot = int(outer.left), int(outer.top)
    oright, obottom = ol + int(outer.width), ot + int(outer.height)
    width = max(0, min(ir, oright) - max(il, ol))
    height = max(0, min(ib, obottom) - max(it, ot))
    area = max(1, int(inner.width) * int(inner.height))
    return (width * height) / area


def _contrast_ratio(foreground: str, background: str) -> float:
    light = _relative_luminance(foreground)
    dark = _relative_luminance(background)
    high, low = max(light, dark), min(light, dark)
    return (high + 0.05) / (low + 0.05)


def _relative_luminance(value: str) -> float:
    channels = [int(value[index : index + 2], 16) / 255 for index in (0, 2, 4)]
    linear = [channel / 12.92 if channel <= 0.04045 else ((channel + 0.055) / 1.055) ** 2.4 for channel in channels]
    return 0.2126 * linear[0] + 0.7152 * linear[1] + 0.0722 * linear[2]
