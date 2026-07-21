from __future__ import annotations

import re
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any

from pptx import Presentation

from deerflow.sophia.deck_build.design_plan import (
    classify_substrate_intent,
    normalize_deck_style_profile,
)
from deerflow.sophia.deck_build.models import DeckBuild

try:  # pragma: no cover - Pillow is present in the backend image, optional in tiny test envs.
    from PIL import Image, ImageStat
except Exception:  # noqa: BLE001
    Image = None  # type: ignore[assignment]
    ImageStat = None  # type: ignore[assignment]

_OLD_RENDERER_CLASS_MARKERS = (
    "section-label",
    "system-diagram",
    "closing-synthesis",
)
_OLD_RENDERER_STRONG_MARKER = "deck_build_templates_v1"
_HARD_RESIDUE_KINDS = {
    "frame_overflow",
    "misaligned",
    "slide_overflow_text",
    "covered_by_picture",
    "repair_still_failing",
}
_ADVISORY_RESIDUE_KINDS = {"overlap", "slide_overflow_non_text"}
_KNOWN_RESIDUE_KINDS = _HARD_RESIDUE_KINDS | _ADVISORY_RESIDUE_KINDS
_POINTS_PER_CSS_PX = 0.75
_REQUIRED_TEXT_MIN_CSS_PX = 24.0
_COMPACT_TEXT_MIN_CSS_PX = 20.0
_REQUIRED_TEXT_MIN_PT = _REQUIRED_TEXT_MIN_CSS_PX * _POINTS_PER_CSS_PX
_COMPACT_TEXT_MIN_PT = _COMPACT_TEXT_MIN_CSS_PX * _POINTS_PER_CSS_PX
_REQUIRED_TEXT_ROLES = {
    "body",
    "callout",
    "content",
    "description",
    "evidence",
    "heading",
    "headline",
    "narrative",
    "paragraph",
    "title",
}


@dataclass
class MechanicalGateIssue:
    code: str
    selector: str
    summary: str
    repair_hint: str

    def to_dict(self) -> dict[str, Any]:
        return asdict(self)


@dataclass
class MechanicalGateResult:
    passed: bool
    failure_code: str | None = None
    failure_summary: str | None = None
    issues: list[MechanicalGateIssue] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)
    slide_render_metrics: list[dict[str, Any]] = field(default_factory=list)

    def to_dict(self) -> dict[str, Any]:
        return {
            "passed": self.passed,
            "failure_code": self.failure_code,
            "failure_summary": self.failure_summary,
            "issues": [issue.to_dict() for issue in self.issues],
            "warnings": self.warnings,
            "slide_render_metrics": self.slide_render_metrics,
        }


def evaluate_mechanical_gates(
    deck: DeckBuild,
    *,
    rendered_dir: Path | None,
    native_pptx_path: Path | None = None,
) -> MechanicalGateResult:
    issues: list[MechanicalGateIssue] = []
    warnings: list[str] = []
    metrics = _render_metrics(rendered_dir)
    issues.extend(_old_renderer_issues(deck))
    issues.extend(_repeated_structure_issues(deck))
    residue_issues, residue_warnings = _native_residue_findings(deck)
    issues.extend(residue_issues)
    warnings.extend(residue_warnings)
    issues.extend(_source_retention_issues(deck))
    issues.extend(_native_contrast_issues(deck))
    issues.extend(
        _compiled_typography_issues(
            deck,
            native_pptx_path=native_pptx_path,
        )
    )
    issues.extend(_sparse_render_issues(metrics))
    issues.extend(_dark_request_light_render_issues(deck, metrics))
    issues.extend(_light_request_dark_render_issues(deck, metrics))
    passed = not issues
    return MechanicalGateResult(
        passed=passed,
        failure_code=None if passed else "deck_mechanical_gate_failed",
        failure_summary=None if passed else "; ".join(issue.summary for issue in issues[:3]),
        issues=issues,
        warnings=warnings,
        slide_render_metrics=metrics,
    )


def _compiled_typography_issues(
    deck: DeckBuild,
    *,
    native_pptx_path: Path | None,
) -> list[MechanicalGateIssue]:
    """Evaluate final native text, never declarations in repeated shared CSS."""

    if native_pptx_path is None or not native_pptx_path.is_file():
        return []
    try:
        presentation = Presentation(str(native_pptx_path))
    except (OSError, ValueError):
        return []

    issues: list[MechanicalGateIssue] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        selector = f"slide:{slide_index}"
        for shape in _native_shapes(slide.shapes):
            text, font_sizes = _native_shape_text_and_sizes(shape)
            if not text or not font_sizes:
                continue
            shape_name = str(getattr(shape, "name", "") or "").strip()
            source_ids, roles, source_required = _typography_source_context(
                deck,
                selector=selector,
                shape_name=shape_name,
            )
            required = bool(
                source_required
                or roles & _REQUIRED_TEXT_ROLES
                or _matches_required_slide_copy(deck, selector=selector, text=text)
            )
            minimum_pt = _REQUIRED_TEXT_MIN_PT if required else _COMPACT_TEXT_MIN_PT
            actual_pt = min(font_sizes)
            if actual_pt + 0.05 >= minimum_pt:
                continue
            actual_px = actual_pt / _POINTS_PER_CSS_PX
            minimum_px = minimum_pt / _POINTS_PER_CSS_PX
            source_label = ", ".join(source_ids[:3]) or shape_name
            code = "native_required_text_too_small" if required else "native_compact_text_too_small"
            kind = "Required/body" if required else "Compact"
            issues.append(
                MechanicalGateIssue(
                    code=code,
                    selector=selector,
                    summary=(
                        f"{kind} text '{source_label}' compiles at {actual_pt:g}pt "
                        f"({actual_px:g}px), below the {minimum_pt:g}pt ({minimum_px:g}px) floor."
                    ),
                    repair_hint=(
                        "Use at least 24px for required body/narrative text and at least 20px for optional "
                        "labels/captions; cut copy instead of shrinking type."
                    ),
                )
            )
    return issues


def _native_shapes(shapes: Any) -> list[Any]:
    flattened: list[Any] = []
    for shape in shapes:
        flattened.append(shape)
        children = getattr(shape, "shapes", None)
        if children is not None:
            flattened.extend(_native_shapes(children))
    return flattened


def _native_shape_text_and_sizes(shape: Any) -> tuple[str, list[float]]:
    frames: list[Any] = []
    if bool(getattr(shape, "has_text_frame", False)):
        frames.append(shape.text_frame)
    if bool(getattr(shape, "has_table", False)):
        frames.extend(cell.text_frame for row in shape.table.rows for cell in row.cells)
    text_parts: list[str] = []
    sizes: list[float] = []
    for frame in frames:
        for paragraph in frame.paragraphs:
            for run in paragraph.runs:
                text = str(run.text or "").strip()
                if not text:
                    continue
                text_parts.append(text)
                if run.font.size is not None:
                    sizes.append(float(run.font.size.pt))
    return _normalize_visible_text(" ".join(text_parts)), sizes


def _typography_source_context(
    deck: DeckBuild,
    *,
    selector: str,
    shape_name: str,
) -> tuple[list[str], set[str], bool]:
    slides = deck.source_element_map.get("slides") if isinstance(deck.source_element_map, dict) else None
    slide_map = slides.get(selector) if isinstance(slides, dict) else None
    elements = slide_map.get("elements") if isinstance(slide_map, dict) else None
    source_ids: list[str] = []
    roles: set[str] = set()
    required = False
    for source_id, record in (elements or {}).items():
        if not isinstance(record, dict) or shape_name not in _record_shape_names(record):
            continue
        source_ids.append(str(source_id))
        role = str(record.get("source_role") or "").strip().lower()
        if role:
            roles.add(role)
        required = bool(required or record.get("source_required"))
    return source_ids, roles, required


def _matches_required_slide_copy(deck: DeckBuild, *, selector: str, text: str) -> bool:
    if not text:
        return False
    for slide in deck.slides:
        if slide.selector != selector:
            continue
        required_copy = {
            _normalize_visible_text(slide.title),
            _normalize_visible_text(slide.narrative),
        }
        return text in required_copy
    return False


def _normalize_visible_text(value: Any) -> str:
    return re.sub(r"\s+", " ", str(value or "")).strip()


def _old_renderer_issues(deck: DeckBuild) -> list[MechanicalGateIssue]:
    sources = [slide.html_source or "" for slide in deck.slides]
    explicit_fingerprint = any(_OLD_RENDERER_STRONG_MARKER in source for source in sources)
    compound_fingerprint = any(
        all(marker in source for marker in _OLD_RENDERER_CLASS_MARKERS)
        for source in sources
    )
    if not explicit_fingerprint and not compound_fingerprint:
        return []
    return [
        MechanicalGateIssue(
            code="old_renderer_artifact",
            selector="deck",
            summary="Slide HTML appears to contain the retired deterministic renderer skeleton.",
            repair_hint="Author subject-specific slide HTML from the creative plan instead of using the old template structure.",
        )
    ]


def _repeated_structure_issues(deck: DeckBuild) -> list[MechanicalGateIssue]:
    layout_names = [
        str(getattr(slide.composition_plan, "layout_name", "") or "")
        for slide in deck.slides
        if slide.composition_plan is not None
    ]
    if len(layout_names) < 3:
        return []
    most_common = max((layout_names.count(name), name) for name in set(layout_names) if name)
    if most_common[0] < max(3, len(layout_names) - 1):
        return []
    return [
        MechanicalGateIssue(
            code="repeated_slide_skeleton",
            selector="deck",
            summary=f"Creative plan repeats layout '{most_common[1]}' across too many slides.",
            repair_hint="Vary slide structure, spatial rhythm, and diagram language across the deck.",
        )
    ]


def _native_residue_findings(deck: DeckBuild) -> tuple[list[MechanicalGateIssue], list[str]]:
    report = deck.native_mechanical_report if isinstance(deck.native_mechanical_report, dict) else {}
    residue_count = int(report.get("lint_residue_count") or 0)
    if residue_count <= 0:
        return [], []
    residue_kinds = report.get("lint_residue_kinds")
    kinds = set(str(key) for key in residue_kinds) if isinstance(residue_kinds, dict) else set()
    residue = report.get("lint_residue") if isinstance(report.get("lint_residue"), list) else []
    item_kinds = {
        str(item.get("kind") or "")
        for item in residue
        if isinstance(item, dict) and item.get("kind")
    }
    issues = _residue_kind_issues(kinds, item_kinds=item_kinds)
    warnings: list[str] = []

    for item in residue:
        if not isinstance(item, dict):
            continue
        item_issues, item_warnings = _residue_item_findings(deck, item)
        issues.extend(item_issues)
        warnings.extend(item_warnings)
    return issues, sorted(set(warnings))


def _residue_kind_issues(
    kinds: set[str],
    *,
    item_kinds: set[str] | None = None,
) -> list[MechanicalGateIssue]:
    unknown = sorted(kind for kind in kinds if kind not in _KNOWN_RESIDUE_KINDS)
    issues: list[MechanicalGateIssue] = []
    if unknown:
        issues.append(
            MechanicalGateIssue(
                code="unknown_native_lint_residue",
                selector="deck",
                summary=f"Native lint/fix left unknown residue kinds: {', '.join(unknown[:5])}.",
                repair_hint="Simplify or revise HTML geometry so hands-on-deck can produce clean native shapes.",
            )
        )
    if not kinds:
        issues.append(
            MechanicalGateIssue(
                code="unknown_native_lint_residue",
                selector="deck",
                summary="Native lint/fix reported residue without stable producer-side kinds.",
                repair_hint="Re-run through the pinned native compiler and preserve residue kind metadata.",
            )
        )
    issues.extend(
        [
            MechanicalGateIssue(
                code=f"native_lint_{kind}",
                selector="deck",
                summary=f"Native lint/fix left blocking residue: {kind}.",
                repair_hint="Repair the exact affected source element and re-run prepare_deck_build once.",
            )
            for kind in sorted((kinds & _HARD_RESIDUE_KINDS) - (item_kinds or set()))
        ]
    )
    return issues


def _residue_item_findings(
    deck: DeckBuild,
    item: dict[str, Any],
) -> tuple[list[MechanicalGateIssue], list[str]]:
    kind = str(item.get("kind") or "")
    selector = f"slide:{int(item.get('slide') or 0) + 1}"
    if kind == "frame_overflow":
        overflow = float(item.get("overflow_bottom") or 0.0)
        return [
            MechanicalGateIssue(
                code="native_lint_frame_overflow",
                selector=selector,
                summary=f'Native text frame still overflows by {overflow:g}" after repair.',
                repair_hint="Increase the matching source box height or shorten the copy without shrinking below the type floor.",
            )
        ], []
    if kind == "misaligned":
        detail = str(item.get("issue") or "Shape remains off its inferred alignment grid.")
        return [
            MechanicalGateIssue(
                code="native_lint_misaligned",
                selector=selector,
                summary=f"Native shape alignment remains inconsistent: {detail}",
                repair_hint="Align the matching source connector or shape to its intended peer edge or centerline.",
            )
        ], []
    if kind in {"slide_overflow_text", "covered_by_picture", "repair_still_failing"}:
        return [
            MechanicalGateIssue(
                code=f"native_lint_{kind}",
                selector=selector,
                summary=str(item.get("issue") or f"Native lint/fix left blocking residue: {kind}."),
                repair_hint=str(item.get("suggest") or "Repair the matching source element and re-run native compilation."),
            )
        ], []
    if kind == "overlap" and float(item.get("overlap_area") or 0.0) >= 0.08:
        return [
            MechanicalGateIssue(
                code="native_lint_severe_overlap",
                selector=selector,
                summary="Native lint/fix left a material shape overlap.",
                repair_hint="Separate or intentionally restack the named semantic elements.",
            )
        ], []
    if kind == "slide_overflow_non_text":
        if _residue_source_role(deck, item) in {"background", "bleed", "decorative"}:
            return [], [f"native_lint_advisory:{kind}"]
        return [
            MechanicalGateIssue(
                code="native_lint_unapproved_bleed",
                selector=selector,
                summary="A non-text shape extends off-slide without an explicit bleed/background role.",
                repair_hint="Keep the shape inside the slide or mark its source role as background, bleed, or decorative.",
            )
        ], []
    if kind in _ADVISORY_RESIDUE_KINDS:
        return [], [f"native_lint_advisory:{kind}"]
    return [], []


def _residue_source_role(deck: DeckBuild, item: dict[str, Any]) -> str | None:
    try:
        selector = f"slide:{int(item.get('slide') or 0) + 1}"
    except (TypeError, ValueError):
        return None
    shape_name = _native_shape_name(deck, selector=selector, shape_id=str(item.get("shape") or ""))
    matches = [
        (source_id, record)
        for source_id, record in _source_records_with_ids(deck, selector)
        if shape_name in _record_shape_names(record)
    ]
    direct_matches = [
        (source_id, record)
        for source_id, record in matches
        if _is_direct_compiler_shape_name(shape_name=shape_name, source_id=source_id)
    ]
    if direct_matches:
        _source_id, record = max(direct_matches, key=lambda item: len(item[0]))
        return str(record.get("source_role") or "").strip().lower() or None
    if len(matches) == 1:
        return str(matches[0][1].get("source_role") or "").strip().lower() or None
    return None


def _source_records(deck: DeckBuild, selector: str) -> list[dict[str, Any]]:
    return [record for _source_id, record in _source_records_with_ids(deck, selector)]


def _source_records_with_ids(deck: DeckBuild, selector: str) -> list[tuple[str, dict[str, Any]]]:
    slides = deck.source_element_map.get("slides") if isinstance(deck.source_element_map, dict) else None
    slide_map = slides.get(selector) if isinstance(slides, dict) else None
    elements = slide_map.get("elements") if isinstance(slide_map, dict) else None
    return [
        (str(source_id), record)
        for source_id, record in (elements or {}).items()
        if isinstance(record, dict)
    ]


def _native_shape_name(deck: DeckBuild, *, selector: str, shape_id: str) -> str:
    inventory = deck.native_shape_inventory if isinstance(deck.native_shape_inventory, dict) else {}
    wrapped_slides = inventory.get("slides")
    if isinstance(wrapped_slides, dict):
        inventory = wrapped_slides
    slide = inventory.get(selector) if isinstance(inventory, dict) else None
    shapes = slide.get("shapes") if isinstance(slide, dict) else None
    for record in shapes if isinstance(shapes, list) else []:
        if isinstance(record, dict) and str(record.get("id") or "") == shape_id:
            return str(record.get("name") or shape_id)
    return shape_id


def _is_direct_compiler_shape_name(*, shape_name: str, source_id: str) -> bool:
    suffix_re = re.compile(
        rf"-{re.escape(source_id)}-"
        r"(?:(?:box|text|image|table)(?:-\d+)?|line-\d+(?:-part-\d+)?)$",
        re.I,
    )
    return bool(suffix_re.search(shape_name))


def _record_shape_names(record: dict[str, Any]) -> set[str]:
    return {str(name) for name in record.get("shape_names") or []}


def _source_retention_issues(deck: DeckBuild) -> list[MechanicalGateIssue]:
    report = deck.source_retention_report if isinstance(deck.source_retention_report, dict) else {}
    issues: list[MechanicalGateIssue] = []
    for item in report.get("missing_required") or []:
        if not isinstance(item, dict):
            continue
        source_id = str(item.get("source_id") or "unknown")
        issues.append(
            MechanicalGateIssue(
                code="required_source_element_missing",
                selector=str(item.get("selector") or "deck"),
                summary=f"Required semantic element '{source_id}' is missing from the native PPTX.",
                repair_hint="Use supported HTML/CSS and preserve the same data-deck-id through native compilation.",
            )
        )
    for item in report.get("duplicates") or []:
        if not isinstance(item, dict):
            continue
        issues.append(
            MechanicalGateIssue(
                code="duplicate_source_element_id",
                selector=str(item.get("selector") or "deck"),
                summary=f"Duplicate semantic source ID: {item.get('source_id')}",
                repair_hint="Give every semantic element on the slide a unique data-deck-id.",
            )
        )
    return issues


def _native_contrast_issues(deck: DeckBuild) -> list[MechanicalGateIssue]:
    report = deck.native_contrast_report if isinstance(deck.native_contrast_report, dict) else {}
    issues: list[MechanicalGateIssue] = []
    for item in report.get("issues") or []:
        if not isinstance(item, dict) or not item.get("required_semantic"):
            continue
        indeterminate = bool(item.get("indeterminate"))
        issues.append(
            MechanicalGateIssue(
                code="native_text_contrast_indeterminate" if indeterminate else "native_text_contrast_failed",
                selector=str(item.get("selector") or "deck"),
                summary=(
                    f"Required text '{item.get('text_excerpt')}' has no deterministic opaque background."
                    if indeterminate
                    else f"Required text contrast {item.get('contrast_ratio')} is below {item.get('required_ratio')}."
                ),
                repair_hint="Place required text on an opaque compiler-supported fill with a compliant text color.",
            )
        )
    return issues


def _sparse_render_issues(metrics: list[dict[str, Any]]) -> list[MechanicalGateIssue]:
    issues: list[MechanicalGateIssue] = []
    for metric in metrics:
        if metric.get("unreadable"):
            continue
        if float(metric.get("non_background_ratio") or 0.0) < 0.008:
            issues.append(
                MechanicalGateIssue(
                    code="sparse_rendered_slide",
                    selector=str(metric.get("selector") or "slide"),
                    summary="Rendered slide is near-blank and likely mechanically corrupted.",
                    repair_hint="Restore the required semantic elements with supported native HTML/CSS.",
                )
            )
    return issues


def _dark_request_light_render_issues(deck: DeckBuild, metrics: list[dict[str, Any]]) -> list[MechanicalGateIssue]:
    if not _dark_requested(deck):
        return []
    readable = [metric for metric in metrics if not metric.get("unreadable")]
    if not readable:
        return []
    light_count = sum(1 for metric in readable if float(metric.get("mean_luminance") or 0.0) > 205)
    if light_count <= len(readable) / 2:
        return []
    return [
        MechanicalGateIssue(
            code="dark_request_rendered_light",
            selector="deck",
            summary="Dark/technical deck request rendered as majority-light slides.",
            repair_hint="Use an opaque dark slide substrate consistently across the deck.",
        )
    ]


def _light_request_dark_render_issues(deck: DeckBuild, metrics: list[dict[str, Any]]) -> list[MechanicalGateIssue]:
    if _substrate_intent(deck) != "light":
        return []
    readable = [metric for metric in metrics if not metric.get("unreadable")]
    if not readable:
        return []
    dark_count = sum(
        1
        for metric in readable
        if float(metric.get("mean_luminance") if metric.get("mean_luminance") is not None else 255.0) < 50
    )
    if dark_count <= len(readable) / 2:
        return []
    return [
        MechanicalGateIssue(
            code="light_request_rendered_dark",
            selector="deck",
            summary="Light deck request rendered as majority-dark slides.",
            repair_hint="Use an opaque light slide substrate consistently across the deck.",
        )
    ]


def _dark_requested(deck: DeckBuild) -> bool:
    return _substrate_intent(deck) == "dark"


def _substrate_intent(deck: DeckBuild) -> str | None:
    plan = deck.design_plan
    intent = classify_substrate_intent(
        [
            str(getattr(plan, "style_lane", "")),
            str(getattr(plan, "signature", "")),
        ]
    )
    if intent is None:
        intent = classify_substrate_intent(getattr(plan, "requested_style_terms", []))
    if intent is None:
        intent = classify_substrate_intent(normalize_deck_style_profile(deck.style_profile))
    return intent


def _render_metrics(rendered_dir: Path | None) -> list[dict[str, Any]]:
    if rendered_dir is None or Image is None or ImageStat is None or not rendered_dir.is_dir():
        return []
    metrics: list[dict[str, Any]] = []
    for index, path in enumerate(sorted(rendered_dir.glob("slide-*.*")), start=1):
        if path.suffix.lower() not in {".jpg", ".jpeg", ".png"}:
            continue
        selector = f"slide:{index}"
        try:
            with Image.open(path) as image:
                rgb = image.convert("RGB")
                stat = ImageStat.Stat(rgb)
                mean = sum(stat.mean) / 3.0
                sample = rgb.resize((64, 36))
                bg = sample.getpixel((0, 0))
                pixels = list(sample.get_flattened_data())
                changed = sum(1 for pixel in pixels if _distance(pixel, bg) > 18)
                metrics.append(
                    {
                        "selector": selector,
                        "file": path.name,
                        "mean_luminance": round(mean, 2),
                        "non_background_ratio": round(changed / max(1, len(pixels)), 4),
                    }
                )
        except Exception:  # noqa: BLE001
            metrics.append({"selector": selector, "file": path.name, "unreadable": True})
    return metrics


def _distance(a: tuple[int, int, int], b: tuple[int, int, int]) -> float:
    return ((a[0] - b[0]) ** 2 + (a[1] - b[1]) ** 2 + (a[2] - b[2]) ** 2) ** 0.5
