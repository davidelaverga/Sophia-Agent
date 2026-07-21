"""Production compiler for one manifest-addressed DQ-2 deck candidate.

The compiler deliberately has no authority to invent content.  It receives the
complete compact source graph from the durable materializer, reconstructs the
deck with the ordinary :class:`DeckBuildService`, and proves that the frozen
brief, plans, visible text, native substrate, and non-target renders survived.
Any visual asset required by the frozen creative plan is copied byte-for-byte
from an injected trusted DQ-1 baseline; image and language models are never
called here.

The inner deck build runs with build-foundation persistence and ambient
LangSmith tracing explicitly disabled.  DQ-2 owns the candidate manifest CAS,
while DQ-1 receives only the bounded, immutable producer bundle created after
the candidate has passed every deterministic gate in this module.
"""

from __future__ import annotations

import hashlib
import inspect
import io
import json
import logging
import math
import re
import tempfile
import unicodedata
import zipfile
from collections.abc import Awaitable, Callable, Mapping
from dataclasses import dataclass
from functools import partial
from pathlib import Path
from types import SimpleNamespace
from typing import Any, Literal, Protocol

import anyio
from PIL import Image, ImageChops, ImageStat
from pptx import Presentation
from pydantic import BaseModel, ConfigDict, Field, ValidationError, model_validator

from deerflow.config.build_foundation_config import BuildFoundationConfig
from deerflow.sandbox.tools import replace_virtual_path
from deerflow.sophia.build_manifest import DECK_STYLE_ROOT_SELECTOR
from deerflow.sophia.deck_build.models import DeckBuildResult
from deerflow.sophia.deck_build.service import DeckBuildService
from deerflow.sophia.deck_build.tracing import (
    FORBIDDEN_SCREENSHOT_COMPILE_MODES,
    NATIVE_DECK_COMPILE_MODE,
)
from deerflow.sophia.deck_design_lift.materializer import (
    DeckCandidateCompilation,
    DeckCandidateCompileRequest,
    DerivedDeckSource,
)
from deerflow.sophia.deck_design_lift.schemas import (
    ContentPreservationProof,
    LocalityProof,
)
from deerflow.sophia.deck_quality.canonical import canonical_json_bytes, canonical_sha256
from deerflow.sophia.deck_quality.instrument import DeckQualityRuntimeInstrument
from deerflow.sophia.deck_quality.publisher import (
    DeckQualityProducerBundleReceipt,
    DeckQualitySourcePack,
    PreparedDeckQualityPublication,
    capture_deck_quality_source_pack,
    deck_quality_immutable_artifact_snapshot_path,
    deck_quality_producer_archive_path,
    deck_quality_producer_bundle_path,
    persist_deck_quality_producer_bundle,
    safe_deck_quality_publication_error_code,
)
from deerflow.sophia.deck_quality.schemas import (
    MechanicalCheck,
    MechanicalProjection,
    VisibleTextSlide,
)
from deerflow.sophia.observability import langsmith_tracing_disabled
from deerflow.sophia.storage.supabase_artifact_store import (
    SupabaseImmutableObjectStore,
)

MAX_BASELINE_RENDER_BYTES = 8 * 1024 * 1024
MAX_BASELINE_ASSET_BYTES = 16 * 1024 * 1024
MAX_CANDIDATE_PPTX_BYTES = 32 * 1024 * 1024
MAX_CANDIDATE_PACKAGE_MEMBERS = 4_096
MAX_CANDIDATE_PACKAGE_MEMBER_BYTES = 64 * 1024 * 1024
MAX_CANDIDATE_PACKAGE_UNCOMPRESSED_BYTES = 128 * 1024 * 1024
MAX_NATIVE_RECORD_BYTES = 2 * 1024 * 1024
MAX_JSON_INPUT_BYTES = 4 * 1024 * 1024
_VISIBLE_TEXT_TOKEN_PATTERN = re.compile(r"\w+|[^\w\s]")
RENDER_COMPARE_WIDTH = 480
RENDER_COMPARE_HEIGHT = 270
MAX_UNCHANGED_RENDER_MEAN_DELTA = 8.0
MAX_NATIVE_EDITABILITY_DROP = 0.01

logger = logging.getLogger(__name__)

_CANONICAL_ZIP_DATETIME = (1980, 1, 1, 0, 0, 0)
_VOLATILE_BUILD_RECORD_KEYS = frozenset(
    {
        "created_at",
        "langsmith_trace_ids",
        "service_elapsed_ms",
        "updated_at",
    }
)

_SLIDE_SELECTOR_RE = re.compile(r"^slide:([1-9][0-9]*)$")
_MECHANICAL_CHECK_IDS: tuple[str, ...] = (
    "authoritative_gate",
    "source_retention",
    "native_editability",
    "contrast",
    "native_lint",
    "overflow_collision_clipping",
    "render_success",
    "visual_asset_completeness",
    "artifact_identity",
)

CandidateCompilationErrorCode = Literal[
    "baseline_unavailable",
    "baseline_invalid",
    "identity_mismatch",
    "instrument_mismatch",
    "source_graph_invalid",
    "source_hash_mismatch",
    "source_decode_failed",
    "plan_revision_forbidden",
    "baseline_asset_missing",
    "baseline_asset_invalid",
    "service_failed",
    "service_result_invalid",
    "candidate_artifact_invalid",
    "plan_changed",
    "derived_source_invalid",
    "mechanical_gate_failed",
    "native_inventory_changed",
    "render_collateral_changed",
    "content_changed",
    "publication_failed",
]


class DeckCandidateCompilationError(RuntimeError):
    """A content-free candidate compilation failure safe for durable state."""

    def __init__(self, code: CandidateCompilationErrorCode) -> None:
        self.code = code
        super().__init__(code)


class _StrictFrozenModel(BaseModel):
    model_config = ConfigDict(extra="forbid", frozen=True, arbitrary_types_allowed=True)


def _sha256(content: bytes) -> str:
    return hashlib.sha256(content).hexdigest()


class BaselineDeckRender(_StrictFrozenModel):
    """One exact DQ-1 slide render loaded from immutable evidence."""

    selector: str = Field(pattern=r"^slide:[1-9][0-9]*$")
    content: bytes = Field(min_length=1, max_length=MAX_BASELINE_RENDER_BYTES)
    sha256: str = Field(pattern=r"^[0-9a-f]{64}$")

    @model_validator(mode="after")
    def validate_hash(self) -> BaselineDeckRender:
        if _sha256(self.content) != self.sha256:
            raise ValueError("baseline render hash does not match")
        return self


class BaselineVisualAsset(_StrictFrozenModel):
    """A frozen generated asset that may only be copied, never regenerated."""

    asset_id: str = Field(pattern=r"^[a-zA-Z0-9][a-zA-Z0-9_.:-]*$")
    selector: str = Field(pattern=r"^slide:[1-9][0-9]*$")
    content: bytes = Field(min_length=1, max_length=MAX_BASELINE_ASSET_BYTES)
    sha256: str = Field(pattern=r"^[0-9a-f]{64}$")

    @model_validator(mode="after")
    def validate_hash(self) -> BaselineVisualAsset:
        if _sha256(self.content) != self.sha256:
            raise ValueError("baseline visual asset hash does not match")
        return self


class DeckCandidateBaseline(_StrictFrozenModel):
    """Trusted, immutable DQ-1 inputs needed to compile a candidate."""

    build_id: str = Field(min_length=1, max_length=512)
    user_id: str = Field(min_length=1, max_length=256)
    thread_id: str = Field(min_length=1, max_length=256)
    task_id: str = Field(min_length=1, max_length=256)
    builder_run_id: str = Field(min_length=1, max_length=256)
    parent_builder_trace_id: str = Field(min_length=1, max_length=256)
    initial_quality_run_id: str = Field(min_length=1, max_length=512)
    logical_artifact_id: str = Field(min_length=1, max_length=256)
    initial_artifact_version_id: str = Field(min_length=1, max_length=256)
    initial_manifest_revision: int = Field(ge=1)
    task_brief: str = Field(min_length=1, max_length=20_000)
    build_record: dict[str, Any]
    creative_plan_record: dict[str, Any]
    design_plan_record: dict[str, Any]
    instrument: DeckQualityRuntimeInstrument
    visible_text: tuple[VisibleTextSlide, ...]
    renders: tuple[BaselineDeckRender, ...]
    visual_assets: tuple[BaselineVisualAsset, ...] = ()

    @model_validator(mode="after")
    def validate_inventory(self) -> DeckCandidateBaseline:
        visible_selectors = tuple(item.selector for item in self.visible_text)
        render_selectors = tuple(item.selector for item in self.renders)
        if not visible_selectors or visible_selectors != render_selectors:
            raise ValueError("baseline visible-text/render coverage differs")
        if len(visible_selectors) != len(set(visible_selectors)):
            raise ValueError("baseline selectors are duplicated")
        asset_ids = tuple(item.asset_id for item in self.visual_assets)
        asset_selectors = tuple(item.selector for item in self.visual_assets)
        if len(asset_ids) != len(set(asset_ids)) or len(asset_selectors) != len(set(asset_selectors)):
            raise ValueError("baseline visual assets are ambiguous")
        if any(selector not in set(visible_selectors) for selector in asset_selectors):
            raise ValueError("baseline visual asset selector is unknown")
        for value in (
            self.build_record,
            self.creative_plan_record,
            self.design_plan_record,
        ):
            encoded = canonical_json_bytes(value)
            if not encoded or len(encoded) > MAX_JSON_INPUT_BYTES:
                raise ValueError("baseline record is unavailable or oversized")
        return self


class DeckCandidateBaselineLoader(Protocol):
    def load(
        self,
        request: DeckCandidateCompileRequest,
    ) -> DeckCandidateBaseline | Awaitable[DeckCandidateBaseline]: ...


def baseline_from_authenticated_snapshot(
    authenticated: Any,
    *,
    instrument: DeckQualityRuntimeInstrument,
    render_contents: Mapping[str, bytes],
    visual_assets: tuple[BaselineVisualAsset, ...] = (),
) -> DeckCandidateBaseline:
    """Project a verified quality-adapter result into the compiler boundary.

    ``DurableDeckQualityEvidenceAdapter`` authenticates the row, manifest,
    evidence bundle, stages, and object inventory before returning its snapshot.
    Its object reader intentionally keeps render bytes separate, so the wrapper
    supplies those exact bytes here after reading each manifest-addressed render.
    This function rechecks selector and hash coverage before constructing the
    frozen baseline consumed by the compiler.
    """

    try:
        row = authenticated.row
        manifest = authenticated.manifest
        evidence_manifest = authenticated.evidence_manifest
        bundle = authenticated.evidence_bundle
        snapshot = bundle.snapshot
        selectors = tuple(str(selector) for selector in evidence_manifest.selectors)
        if tuple(render_contents) != selectors:
            raise ValueError
        renders = tuple(
            BaselineDeckRender(
                selector=selector,
                content=render_contents[selector],
                sha256=evidence_manifest.render_hashes[selector],
            )
            for selector in selectors
        )
        if (
            row.quality_run_id != evidence_manifest.quality_run_id
            or row.build_id != manifest.build_id
            or snapshot.build_id != manifest.build_id
            or row.user_id != manifest.user_id
            or manifest.thread_id not in {row.thread_id, row.task_id}
            or row.task_id != evidence_manifest.task_id
            or row.thread_id != evidence_manifest.thread_id
            or row.logical_artifact_id != manifest.logical_artifact_id
            or row.artifact_version_id != manifest.current_artifact_version_id
            or row.manifest_revision != manifest.manifest_revision
        ):
            raise ValueError
        return DeckCandidateBaseline(
            build_id=row.build_id,
            user_id=row.user_id,
            thread_id=manifest.thread_id,
            task_id=row.task_id or evidence_manifest.task_id,
            builder_run_id=row.builder_run_id or evidence_manifest.builder_run_id,
            parent_builder_trace_id=(row.parent_builder_trace_id or evidence_manifest.parent_builder_trace_id),
            initial_quality_run_id=row.quality_run_id,
            logical_artifact_id=row.logical_artifact_id,
            initial_artifact_version_id=row.artifact_version_id,
            initial_manifest_revision=row.manifest_revision,
            task_brief=snapshot.brief.request,
            build_record=bundle.build_record,
            creative_plan_record=snapshot.creative_plan,
            design_plan_record=snapshot.design_plan,
            instrument=instrument,
            visible_text=snapshot.visible_text,
            renders=renders,
            visual_assets=visual_assets,
        )
    except DeckCandidateCompilationError:
        raise
    except Exception:
        raise DeckCandidateCompilationError("baseline_invalid") from None


class CandidateDq1Publisher(Protocol):
    def publish(
        self,
        *,
        prepared: PreparedDeckQualityPublication,
        instrument: DeckQualityRuntimeInstrument,
        pptx_bytes: bytes,
        source_pack: DeckQualitySourcePack,
        source_pack_bytes: bytes,
    ) -> DeckQualityProducerBundleReceipt | Awaitable[DeckQualityProducerBundleReceipt]: ...


class _CandidatePublicationStore(Protocol):
    def create_if_absent(
        self,
        object_path: str,
        content: bytes,
        *,
        content_type: str,
    ) -> Literal["created", "exists"]: ...

    def read(self, object_path: str) -> bytes | None: ...

    def read_bounded(self, object_path: str, *, max_bytes: int) -> bytes | None: ...


class DurableCandidateDq1Publisher:
    """Create the candidate snapshot and exact-source-bound DQ-1 outbox."""

    def __init__(
        self,
        *,
        store_factory: Callable[[], _CandidatePublicationStore] | None = None,
    ) -> None:
        self._store_factory = (
            store_factory
            if store_factory is not None
            else SupabaseImmutableObjectStore
        )
        # Injected synchronous stores are useful for deterministic tests.  The
        # live producer bundle must use its native async, deadline-bounded
        # protocol instead of inheriting the snapshot store's per-request
        # synchronous timeouts.
        self._reuse_snapshot_store_for_bundle = store_factory is not None

    @staticmethod
    def _persist_snapshot(
        *,
        store: _CandidatePublicationStore,
        object_path: str,
        pptx_bytes: bytes,
    ) -> None:
        """Create one immutable snapshot with exact ambiguity reconciliation."""

        create_failed = False
        try:
            outcome = store.create_if_absent(
                object_path,
                pptx_bytes,
                content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        except Exception:
            # A create-only POST may commit before its response is lost.  The
            # exact bounded GET below is the ambiguity fence.
            create_failed = True
            outcome = None
        if not create_failed and outcome not in {"created", "exists"}:
            logger.warning(
                "DQ2 candidate publication failed "
                "stage=immutable_snapshot code=create_outcome_invalid"
            )
            raise DeckCandidateCompilationError("publication_failed")

        first_read_failed = False
        try:
            stored = store.read_bounded(
                object_path,
                max_bytes=MAX_CANDIDATE_PPTX_BYTES,
            )
        except Exception:
            first_read_failed = True
            stored = None
        if stored is not None:
            if stored != pptx_bytes:
                logger.warning(
                    "DQ2 candidate publication failed "
                    "stage=immutable_snapshot code=content_conflict"
                )
                raise DeckCandidateCompilationError("publication_failed")
            if create_failed:
                logger.info(
                    "DQ2 candidate publication reconciled "
                    "stage=immutable_snapshot attempt=1"
                )
            return

        if not create_failed:
            code = "readback_failed" if first_read_failed else "missing_after_create"
            logger.warning(
                "DQ2 candidate publication failed "
                "stage=immutable_snapshot code=%s",
                code,
            )
            raise DeckCandidateCompilationError("publication_failed")

        # The first create response is ambiguous and no exact object can be
        # confirmed.  One create-only retry is safe even if the first POST
        # committed: it can only return a conflict for the same key.
        retry_create_failed = False
        try:
            retry_outcome = store.create_if_absent(
                object_path,
                pptx_bytes,
                content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        except Exception:
            retry_create_failed = True
            retry_outcome = None
        if not retry_create_failed and retry_outcome not in {"created", "exists"}:
            logger.warning(
                "DQ2 candidate publication failed "
                "stage=immutable_snapshot code=retry_outcome_invalid"
            )
            raise DeckCandidateCompilationError("publication_failed")

        retry_read_failed = False
        try:
            stored = store.read_bounded(
                object_path,
                max_bytes=MAX_CANDIDATE_PPTX_BYTES,
            )
        except Exception:
            retry_read_failed = True
            stored = None
        if stored is not None:
            if stored != pptx_bytes:
                logger.warning(
                    "DQ2 candidate publication failed "
                    "stage=immutable_snapshot code=content_conflict"
                )
                raise DeckCandidateCompilationError("publication_failed")
            logger.info(
                "DQ2 candidate publication reconciled "
                "stage=immutable_snapshot attempt=2"
            )
            return

        if retry_read_failed:
            code = "retry_readback_failed"
        elif retry_create_failed:
            code = "retry_create_failed"
        else:
            code = "missing_after_retry"
        logger.warning(
            "DQ2 candidate publication failed "
            "stage=immutable_snapshot code=%s",
            code,
        )
        raise DeckCandidateCompilationError("publication_failed")

    async def publish(
        self,
        *,
        prepared: PreparedDeckQualityPublication,
        instrument: DeckQualityRuntimeInstrument,
        pptx_bytes: bytes,
        source_pack: DeckQualitySourcePack,
        source_pack_bytes: bytes,
    ) -> DeckQualityProducerBundleReceipt:
        return await anyio.to_thread.run_sync(
            partial(
                self._publish_sync,
                prepared=prepared,
                instrument=instrument,
                pptx_bytes=pptx_bytes,
                source_pack=source_pack,
                source_pack_bytes=source_pack_bytes,
            ),
            abandon_on_cancel=False,
        )

    def _publish_sync(
        self,
        *,
        prepared: PreparedDeckQualityPublication,
        instrument: DeckQualityRuntimeInstrument,
        pptx_bytes: bytes,
        source_pack: DeckQualitySourcePack,
        source_pack_bytes: bytes,
    ) -> DeckQualityProducerBundleReceipt:
        if not 0 < len(pptx_bytes) <= MAX_CANDIDATE_PPTX_BYTES or _sha256(pptx_bytes) != prepared.artifact_sha256:
            raise DeckCandidateCompilationError("candidate_artifact_invalid")
        if (
            not isinstance(source_pack, DeckQualitySourcePack)
            or not isinstance(source_pack_bytes, bytes)
            or not source_pack_bytes
            or canonical_json_bytes(source_pack) != source_pack_bytes
            or source_pack.artifact_sha256 != prepared.artifact_sha256
            or source_pack.artifact_version_id != prepared.artifact_version_id
            or source_pack.manifest_revision != prepared.manifest_revision
        ):
            raise DeckCandidateCompilationError("publication_failed")
        expected_path = deck_quality_immutable_artifact_snapshot_path(
            user_id=prepared.user_id,
            thread_id=prepared.thread_id,
            build_id=prepared.build_id,
            logical_artifact_id=prepared.logical_artifact_id,
            artifact_version_id=prepared.artifact_version_id,
            artifact_sha256=prepared.artifact_sha256,
            artifact_virtual_path=prepared.artifact_virtual_path,
        )
        if prepared.artifact_storage_object_path != expected_path:
            raise DeckCandidateCompilationError("identity_mismatch")
        try:
            store = self._store_factory()
        except Exception:
            logger.warning(
                "DQ2 candidate publication failed "
                "stage=immutable_snapshot code=store_unavailable"
            )
            raise DeckCandidateCompilationError("publication_failed") from None

        self._persist_snapshot(
            store=store,
            object_path=expected_path,
            pptx_bytes=pptx_bytes,
        )
        try:
            publication_kwargs: dict[str, Any] = {
                "prepared": prepared,
                "instrument": instrument,
                "source_pack": source_pack,
                "source_pack_bytes": source_pack_bytes,
            }
            if self._reuse_snapshot_store_for_bundle:
                publication_kwargs["object_store"] = store
            return persist_deck_quality_producer_bundle(
                **publication_kwargs,
            )
        except DeckCandidateCompilationError:
            raise
        except Exception as error:
            logger.warning(
                "DQ2 candidate publication failed "
                "stage=producer_bundle code=%s",
                safe_deck_quality_publication_error_code(error),
            )
            raise DeckCandidateCompilationError("publication_failed") from None


DeckBuildServiceFactory = Callable[
    [Callable[[str, Any], dict[str, Any]], Callable[[Any, Any, int], dict[str, Any]]],
    DeckBuildService,
]


def _default_service_factory(
    batch_runner: Callable[[str, Any], dict[str, Any]],
    single_runner: Callable[[Any, Any, int], dict[str, Any]],
) -> DeckBuildService:
    return DeckBuildService(
        image_batch_runner=batch_runner,
        image_single_runner=single_runner,
    )


async def _run_sync[ValueT](
    function: Callable[..., ValueT],
    /,
    *args: Any,
    **kwargs: Any,
) -> ValueT:
    """Run blocking work without ever abandoning its worker on cancellation."""

    return await anyio.to_thread.run_sync(
        partial(function, *args, **kwargs),
        abandon_on_cancel=False,
    )


async def _call_maybe_async[ValueT](
    function: Callable[..., ValueT | Awaitable[ValueT]],
    /,
    *args: Any,
    **kwargs: Any,
) -> ValueT:
    """Invoke native async adapters directly and isolate synchronous adapters."""

    if inspect.iscoroutinefunction(function):
        return await function(*args, **kwargs)
    value = await _run_sync(function, *args, **kwargs)
    if inspect.isawaitable(value):
        return await value
    return value


def _json_object(path: Path, *, code: CandidateCompilationErrorCode) -> dict[str, Any]:
    try:
        if not path.is_file() or path.is_symlink() or not 0 < path.stat().st_size <= MAX_JSON_INPUT_BYTES:
            raise ValueError
        raw = path.read_bytes()
        value = json.loads(raw.decode("utf-8"), parse_constant=lambda _value: (_ for _ in ()).throw(ValueError()))
        if not isinstance(value, dict) or canonical_json_bytes(value) == b"":
            raise ValueError
        return value
    except Exception:
        raise DeckCandidateCompilationError(code) from None


def _safe_utf8(content: bytes) -> str:
    try:
        return content.decode("utf-8")
    except UnicodeError:
        raise DeckCandidateCompilationError("source_decode_failed") from None


def _slide_selectors(request: DeckCandidateCompileRequest) -> tuple[str, ...]:
    selectors = tuple(component.selector for component in request.baseline_manifest.components if component.selector != DECK_STYLE_ROOT_SELECTOR)
    if not selectors or any(_SLIDE_SELECTOR_RE.fullmatch(selector) is None for selector in selectors):
        raise DeckCandidateCompilationError("source_graph_invalid")
    expected = tuple(f"slide:{index}" for index in range(1, len(selectors) + 1))
    if selectors != expected:
        raise DeckCandidateCompilationError("source_graph_invalid")
    return selectors


def _validate_request_sources(
    request: DeckCandidateCompileRequest,
) -> tuple[dict[tuple[str, str], str], tuple[str, ...], tuple[str, ...]]:
    components = {component.selector: component for component in request.baseline_manifest.components}
    expected_keys = {(component.selector, role) for component in request.baseline_manifest.components for role in component.source_roles}
    values: dict[tuple[str, str], str] = {}
    flags: dict[str, set[bool]] = {selector: set() for selector in components}
    model_authored: set[tuple[str, str]] = set()
    for source in request.sources:
        key = (source.selector, source.source_role)
        if key in values or key not in expected_keys:
            raise DeckCandidateCompilationError("source_graph_invalid")
        if _sha256(source.content) != source.source_hash:
            raise DeckCandidateCompilationError("source_hash_mismatch")
        values[key] = _safe_utf8(source.content)
        flags[source.selector].add(source.component_version_changed)
        if source.model_authored:
            model_authored.add(key)
    if set(values) != expected_keys or any(len(value) != 1 for value in flags.values()):
        raise DeckCandidateCompilationError("source_graph_invalid")

    expected_model_authored = {(selector, role) for selector in request.program.authorized_selectors for role in request.program.authorized_source_roles[selector]}
    if model_authored != expected_model_authored:
        raise DeckCandidateCompilationError("source_graph_invalid")
    changed = tuple(component.selector for component in request.baseline_manifest.components if flags[component.selector] == {True})
    unchanged = tuple(component.selector for component in request.baseline_manifest.components if flags[component.selector] == {False})
    if not changed or not set(request.program.authorized_selectors).issubset(changed):
        raise DeckCandidateCompilationError("source_graph_invalid")

    slide_selectors = _slide_selectors(request)
    if DECK_STYLE_ROOT_SELECTOR in request.program.authorized_selectors:
        if changed != tuple(component.selector for component in request.baseline_manifest.components):
            raise DeckCandidateCompilationError("source_graph_invalid")
        expected_derived = tuple((selector, "assembled") for selector in slide_selectors)
    else:
        if set(changed) != set(request.program.authorized_selectors):
            raise DeckCandidateCompilationError("source_graph_invalid")
        expected_derived = tuple((selector, "assembled") for selector in slide_selectors if selector in request.program.authorized_selectors and set(request.program.authorized_source_roles[selector]).intersection({"body", "slide_css"}))
    if request.derived_source_targets != expected_derived:
        raise DeckCandidateCompilationError("source_graph_invalid")
    return values, changed, unchanged


def _baseline_slide_records(
    baseline: DeckCandidateBaseline,
    selectors: tuple[str, ...],
) -> dict[str, dict[str, Any]]:
    raw_slides = baseline.build_record.get("slides")
    if not isinstance(raw_slides, list):
        raise DeckCandidateCompilationError("baseline_invalid")
    result: dict[str, dict[str, Any]] = {}
    for raw in raw_slides:
        if not isinstance(raw, dict):
            raise DeckCandidateCompilationError("baseline_invalid")
        selector = raw.get("selector")
        if not isinstance(selector, str) or selector in result:
            raise DeckCandidateCompilationError("baseline_invalid")
        result[selector] = raw
    if tuple(result) != selectors:
        raise DeckCandidateCompilationError("baseline_invalid")
    return result


def _validate_baseline(
    request: DeckCandidateCompileRequest,
    baseline: DeckCandidateBaseline,
    selectors: tuple[str, ...],
) -> None:
    manifest = request.baseline_manifest
    expected_instrument_hash = canonical_sha256(baseline.instrument.lock)
    invalid_identity = any(
        (
            baseline.build_id != request.build_id,
            baseline.user_id != request.user_id,
            baseline.thread_id != request.thread_id,
            baseline.build_id != manifest.build_id,
            baseline.user_id != manifest.user_id,
            baseline.thread_id != manifest.thread_id,
            baseline.initial_manifest_revision != manifest.manifest_revision,
            baseline.initial_manifest_revision != request.program.initial_manifest_revision,
            baseline.initial_artifact_version_id != manifest.current_artifact_version_id,
            baseline.logical_artifact_id != manifest.logical_artifact_id,
            baseline.initial_quality_run_id != request.program.initial_quality_run_id,
            request.candidate_manifest_revision != manifest.manifest_revision + 1,
        )
    )
    if invalid_identity:
        raise DeckCandidateCompilationError("identity_mismatch")
    if expected_instrument_hash != request.program.instrument_hash or baseline.instrument.lock.rubric_version != request.program.rubric_version:
        raise DeckCandidateCompilationError("instrument_mismatch")
    if request.program.plan_revision_allowed:
        raise DeckCandidateCompilationError("plan_revision_forbidden")
    if tuple(item.selector for item in baseline.visible_text) != selectors or tuple(item.selector for item in baseline.renders) != selectors:
        raise DeckCandidateCompilationError("baseline_invalid")
    if baseline.build_record.get("build_id") != request.build_id:
        raise DeckCandidateCompilationError("baseline_invalid")
    if baseline.build_record.get("deck_authoring_contract") not in {
        "compact_model_html_v1",
        "compact_model_html_v2",
    }:
        raise DeckCandidateCompilationError("baseline_invalid")


def _validate_visual_asset_inventory(baseline: DeckCandidateBaseline) -> None:
    raw_assets = baseline.creative_plan_record.get("image_assets")
    if not isinstance(raw_assets, list):
        raise DeckCandidateCompilationError("baseline_invalid")
    planned: set[tuple[str, str]] = set()
    for item in raw_assets:
        if not isinstance(item, dict):
            raise DeckCandidateCompilationError("baseline_invalid")
        asset_id = item.get("asset_id")
        selector = item.get("slide_selector")
        if not isinstance(asset_id, str) or not isinstance(selector, str):
            raise DeckCandidateCompilationError("baseline_invalid")
        planned.add((asset_id, selector))
    if len(planned) != len(raw_assets):
        raise DeckCandidateCompilationError("baseline_invalid")
    available = {(asset.asset_id, asset.selector) for asset in baseline.visual_assets}
    if available != planned:
        raise DeckCandidateCompilationError("baseline_asset_missing")


def _slide_ir(
    *,
    selectors: tuple[str, ...],
    records: Mapping[str, Mapping[str, Any]],
    sources: Mapping[tuple[str, str], str],
) -> list[dict[str, Any]]:
    slides: list[dict[str, Any]] = []
    for selector in selectors:
        record = records[selector]
        title = record.get("title")
        narrative = record.get("narrative")
        role = record.get("role")
        layout = record.get("layout_kind")
        if not all(isinstance(value, str) and value.strip() for value in (title, narrative, role, layout)):
            raise DeckCandidateCompilationError("baseline_invalid")
        baseline_notes = record.get("speaker_notes")
        expected_notes = baseline_notes if isinstance(baseline_notes, str) else ""
        if sources[(selector, "notes")] != expected_notes:
            raise DeckCandidateCompilationError("content_changed")
        slides.append(
            {
                "title": title,
                "narrative": narrative,
                "role": role,
                "layout_kind": layout,
                "claim": record.get("claim") if isinstance(record.get("claim"), str) else None,
                "visual_prompt": record.get("visual_prompt") if isinstance(record.get("visual_prompt"), str) else None,
                "html_body": sources[(selector, "body")],
                "slide_css": sources[(selector, "slide_css")],
                "speaker_notes": sources[(selector, "notes")],
            }
        )
    return slides


def _runtime(
    *,
    root: Path,
    request: DeckCandidateCompileRequest,
    baseline: DeckCandidateBaseline,
    slide_count: int,
    output_path: str,
) -> SimpleNamespace:
    outputs = root / "outputs"
    workspace = root / "workspace"
    uploads = root / "uploads"
    for path in (outputs, workspace, uploads):
        path.mkdir(parents=True, exist_ok=True)
    return SimpleNamespace(
        state={
            "builder_build_id": request.build_id,
            "build_id": request.build_id,
            "thread_id": request.thread_id,
            "parent_thread_id": request.thread_id,
            "user_id": request.user_id,
            "task_id": baseline.task_id,
            "run_id": request.operation_id,
            "builder_pptx_requested_slide_count": slide_count,
            "builder_artifact_target_path": output_path,
            "delegation_context": {"request": baseline.task_brief},
            "thread_data": {
                "outputs_path": str(outputs),
                "workspace_path": str(workspace),
                "uploads_path": str(uploads),
            },
        },
        context={
            "thread_id": request.thread_id,
            "build_foundation_config": BuildFoundationConfig(
                enabled=False,
                manifest_mode="off",
                persist_event_journal=False,
                enable_mutation_transactions=False,
            ),
        },
        config={"configurable": {"thread_id": request.thread_id}},
    )


class _ReuseOnlyImageRunner:
    def __init__(self, assets: tuple[BaselineVisualAsset, ...]) -> None:
        self._by_index = {int(_SLIDE_SELECTOR_RE.fullmatch(asset.selector).group(1)): asset for asset in assets if _SLIDE_SELECTOR_RE.fullmatch(asset.selector) is not None}
        self.calls = 0

    @staticmethod
    def _decode_asset(asset: BaselineVisualAsset) -> None:
        try:
            with Image.open(io.BytesIO(asset.content)) as image:
                image.verify()
            with Image.open(io.BytesIO(asset.content)) as image:
                if image.format != "PNG" or image.width < 1 or image.height < 1:
                    raise ValueError
        except Exception:
            raise DeckCandidateCompilationError("baseline_asset_invalid") from None

    def run_batch(self, manifest_path: str, runtime: Any) -> dict[str, Any]:
        self.calls += 1
        try:
            host = Path(replace_virtual_path(manifest_path, runtime.state["thread_data"]))
            manifest = json.loads(host.read_text(encoding="utf-8"))
            raw_items = manifest.get("items")
            if not isinstance(raw_items, list):
                raise ValueError
            requested_indices = tuple(item.get("slide_index") for item in raw_items if isinstance(item, dict))
            if len(requested_indices) != len(raw_items) or set(requested_indices) != set(self._by_index):
                raise DeckCandidateCompilationError("baseline_asset_missing")
            items: list[dict[str, Any]] = []
            outputs_root = Path(runtime.state["thread_data"]["outputs_path"]).resolve()
            for item in raw_items:
                slide_index = int(item["slide_index"])
                output_file = item.get("output_file")
                if not isinstance(output_file, str):
                    raise ValueError
                asset = self._by_index[slide_index]
                self._decode_asset(asset)
                output_host = Path(replace_virtual_path(output_file, runtime.state["thread_data"]))
                resolved_parent = output_host.parent.resolve()
                if outputs_root != resolved_parent and outputs_root not in resolved_parent.parents:
                    raise ValueError
                output_host.parent.mkdir(parents=True, exist_ok=True)
                output_host.write_bytes(asset.content)
                items.append(
                    {
                        "item_index": slide_index,
                        "output_file": output_file,
                        "success": True,
                        "bytes": len(asset.content),
                        "error_class": None,
                        "reused_asset_hash": asset.sha256,
                    }
                )
            return {
                "summary_present": True,
                "batch_attempted": True,
                "complete": True,
                "requested": len(items),
                "images_generated": len(items),
                "failed": 0,
                "items": items,
                "error_class_histogram": {},
                "source": "immutable_baseline_reuse",
            }
        except DeckCandidateCompilationError:
            raise
        except Exception:
            raise DeckCandidateCompilationError("baseline_asset_invalid") from None

    @staticmethod
    def forbid_single(_slide: Any, _runtime: Any, _attempt_no: int) -> dict[str, Any]:
        raise DeckCandidateCompilationError("baseline_asset_missing")


def _validate_service_result(
    result: DeckBuildResult,
    *,
    request: DeckCandidateCompileRequest,
    slide_count: int,
    output_path: str,
    authoring_contract: str,
) -> None:
    if not isinstance(result, DeckBuildResult):
        raise DeckCandidateCompilationError("service_result_invalid")
    if not result.success:
        if "deck_mechanical_gate_failed" in {
            result.failure_code,
            result.root_failure_code,
        }:
            raise DeckCandidateCompilationError("mechanical_gate_failed")
        raise DeckCandidateCompilationError("service_failed")
    if (
        result.build_id != request.build_id
        or result.pptx_path != output_path
        or result.slide_count != slide_count
        or result.deck_compile_mode != NATIVE_DECK_COMPILE_MODE
        or result.deck_compile_mode in FORBIDDEN_SCREENSHOT_COMPILE_MODES
        or not result.native_required
        or result.legacy_screenshot_debug
        or result.full_slide_picture_count != 0
        or (result.native_editability_score or 0.0) < 0.60
        or result.mechanical_gate_results.get("passed") is not True
        or result.source_retention_report.get("passed") is not True
        or result.native_contrast_report.get("passed") is not True
        or result.missing_visual_count != 0
        or result.expected_visual_count != result.successful_visual_count
        or result.deck_authoring_contract != authoring_contract
    ):
        raise DeckCandidateCompilationError("mechanical_gate_failed")
    native = result.native_mechanical_report
    if native.get("render_success") is not True or native.get("lint_fix_success") is not True or native.get("lint_residue_count") != 0:
        raise DeckCandidateCompilationError("mechanical_gate_failed")


def _read_candidate_artifact(path: Path) -> bytes:
    try:
        if not path.is_file() or path.is_symlink() or not 0 < path.stat().st_size <= MAX_CANDIDATE_PPTX_BYTES:
            raise ValueError
        content = path.read_bytes()
        if not content.startswith(b"PK\x03\x04"):
            raise ValueError
    except Exception:
        raise DeckCandidateCompilationError("candidate_artifact_invalid") from None
    return content


def _safe_package_member_name(name: str) -> bool:
    if not name or len(name) > 1_024 or name.startswith(("/", "\\")) or "\\" in name or "\x00" in name:
        return False
    parts = name.split("/")
    return bool(parts) and all(part not in {"", ".", ".."} for part in parts)


def _canonicalize_pptx_package(content: bytes) -> bytes:
    """Repack OOXML with deterministic ZIP metadata and member order.

    ``python-pptx`` writes the current wall-clock time into every ``ZipInfo``.
    Rebuilding the same native deck later therefore changes the artifact hash
    even when every OOXML/media member is byte-identical.  The canonical form
    intentionally leaves all package member bytes untouched; only the ZIP
    container's order, compression settings, timestamps, permissions, extras,
    and comments are fixed.
    """

    if not 0 < len(content) <= MAX_CANDIDATE_PPTX_BYTES:
        raise DeckCandidateCompilationError("candidate_artifact_invalid")
    try:
        members: dict[str, bytes] = {}
        total_uncompressed = 0
        with zipfile.ZipFile(io.BytesIO(content), mode="r") as source:
            infos = source.infolist()
            if not 0 < len(infos) <= MAX_CANDIDATE_PACKAGE_MEMBERS:
                raise ValueError
            for info in infos:
                if info.filename in members or not _safe_package_member_name(info.filename) or info.is_dir() or info.flag_bits & 0x1 or not 0 <= info.file_size <= MAX_CANDIDATE_PACKAGE_MEMBER_BYTES:
                    raise ValueError
                total_uncompressed += info.file_size
                if total_uncompressed > MAX_CANDIDATE_PACKAGE_UNCOMPRESSED_BYTES:
                    raise ValueError
                with source.open(info, mode="r") as member_file:
                    member = member_file.read(MAX_CANDIDATE_PACKAGE_MEMBER_BYTES + 1)
                if len(member) != info.file_size or len(member) > MAX_CANDIDATE_PACKAGE_MEMBER_BYTES:
                    raise ValueError
                members[info.filename] = member

        output = io.BytesIO()
        with zipfile.ZipFile(
            output,
            mode="w",
            compression=zipfile.ZIP_DEFLATED,
            compresslevel=9,
            strict_timestamps=True,
        ) as destination:
            destination.comment = b""
            for name in sorted(members):
                info = zipfile.ZipInfo(filename=name, date_time=_CANONICAL_ZIP_DATETIME)
                info.compress_type = zipfile.ZIP_DEFLATED
                info.create_system = 3
                info.create_version = 20
                info.extract_version = 20
                info.flag_bits = 0
                info.internal_attr = 0
                info.external_attr = 0o600 << 16
                info.extra = b""
                info.comment = b""
                destination.writestr(info, members[name], compress_type=zipfile.ZIP_DEFLATED, compresslevel=9)
        canonical = output.getvalue()
        if not 0 < len(canonical) <= MAX_CANDIDATE_PPTX_BYTES:
            raise ValueError
        with zipfile.ZipFile(io.BytesIO(canonical), mode="r") as verified:
            if verified.testzip() is not None or verified.namelist() != sorted(members):
                raise ValueError
            if any(info.date_time != _CANONICAL_ZIP_DATETIME or info.compress_type != zipfile.ZIP_DEFLATED or info.extra or info.comment for info in verified.infolist()):
                raise ValueError
        Presentation(io.BytesIO(canonical))
        return canonical
    except DeckCandidateCompilationError:
        raise
    except Exception:
        raise DeckCandidateCompilationError("candidate_artifact_invalid") from None


def _canonicalize_candidate_artifact(path: Path) -> bytes:
    canonical = _canonicalize_pptx_package(_read_candidate_artifact(path))
    try:
        written = path.write_bytes(canonical)
        if written != len(canonical) or path.read_bytes() != canonical:
            raise ValueError
    except Exception:
        raise DeckCandidateCompilationError("candidate_artifact_invalid") from None
    return canonical


def _stable_generated_build_record(record: Mapping[str, Any]) -> dict[str, Any]:
    """Remove only DeckBuild execution telemetry before immutable publication."""

    stable = dict(record)
    for key in _VOLATILE_BUILD_RECORD_KEYS:
        stable.pop(key, None)
    try:
        # Round-trip through the canonical encoder to reject non-JSON values
        # and detach the returned object from the service-owned mapping.
        encoded = canonical_json_bytes(stable)
        if not 0 < len(encoded) <= MAX_JSON_INPUT_BYTES:
            raise ValueError
        decoded = json.loads(encoded)
        if not isinstance(decoded, dict):
            raise ValueError
        return decoded
    except Exception:
        raise DeckCandidateCompilationError("service_result_invalid") from None


def _write_stable_build_record(path: Path, record: Mapping[str, Any]) -> None:
    try:
        encoded = canonical_json_bytes(record)
        if not 0 < len(encoded) <= MAX_JSON_INPUT_BYTES or path.is_symlink():
            raise ValueError
        written = path.write_bytes(encoded)
        if written != len(encoded) or path.read_bytes() != encoded:
            raise ValueError
    except Exception:
        raise DeckCandidateCompilationError("service_result_invalid") from None


def _plan_record(
    baseline: Mapping[str, Any],
    generated: Mapping[str, Any],
    *,
    common: Mapping[str, Any],
    plan_kind: Literal["creative", "design"],
) -> dict[str, Any]:
    if canonical_sha256(baseline) != canonical_sha256(generated):
        raise DeckCandidateCompilationError("plan_changed")
    summary: dict[str, Any] = {
        "schema_version": "sophia-deck-candidate-plan/v1",
        "plan_kind": plan_kind,
        **dict(common),
        "plan_revision_changed": False,
        "frozen_plan_hash": canonical_sha256(baseline),
    }
    for key in (
        "style_lane",
        "signature",
        "rhythm",
        "image_strategy",
    ):
        value = baseline.get(key)
        if isinstance(value, str) and value:
            summary[key] = value
    slide_compositions = baseline.get("slide_compositions")
    image_assets = baseline.get("image_assets")
    if isinstance(slide_compositions, list):
        summary["slide_composition_count"] = len(slide_compositions)
    if isinstance(image_assets, list):
        summary["image_asset_count"] = len(image_assets)
    # The complete, unredacted plans remain in the immutable DQ-1 source pack.
    # The candidate manifest stores only this hash-bound projection so provider
    # prompts or other model-facing material cannot leak into durable metadata.
    return summary


def _generated_sources(
    *,
    build_record: Mapping[str, Any],
    request: DeckCandidateCompileRequest,
    sources: Mapping[tuple[str, str], str],
    selectors: tuple[str, ...],
) -> tuple[DerivedDeckSource, ...]:
    raw_slides = build_record.get("slides")
    if not isinstance(raw_slides, list):
        raise DeckCandidateCompilationError("derived_source_invalid")
    generated: dict[str, str] = {}
    for item in raw_slides:
        if not isinstance(item, dict):
            raise DeckCandidateCompilationError("derived_source_invalid")
        selector = item.get("selector")
        assembled = item.get("html_source")
        if not isinstance(selector, str) or selector in generated or not isinstance(assembled, str) or not assembled:
            raise DeckCandidateCompilationError("derived_source_invalid")
        generated[selector] = assembled
    if tuple(generated) != selectors:
        raise DeckCandidateCompilationError("derived_source_invalid")
    targets = set(request.derived_source_targets)
    for selector in selectors:
        if (selector, "assembled") not in targets and generated[selector] != sources[(selector, "assembled")]:
            raise DeckCandidateCompilationError("derived_source_invalid")
    return tuple(
        DerivedDeckSource(
            selector=selector,
            source_role="assembled",
            content=generated[selector],
        )
        for selector, _role in request.derived_source_targets
    )


def _shape_text_fragments(shape: Any) -> tuple[str, ...]:
    children = getattr(shape, "shapes", None)
    if children is not None:
        return tuple(fragment for child in children for fragment in _shape_text_fragments(child))
    if bool(getattr(shape, "has_table", False)):
        return tuple(paragraph.text for row in shape.table.rows for cell in row.cells for paragraph in cell.text_frame.paragraphs if paragraph.text.strip())
    if bool(getattr(shape, "has_text_frame", False)):
        return tuple(paragraph.text for paragraph in shape.text_frame.paragraphs if paragraph.text.strip())
    return ()


def _candidate_visible_text(content: bytes) -> tuple[VisibleTextSlide, ...]:
    try:
        presentation = Presentation(io.BytesIO(content))
        return tuple(
            VisibleTextSlide(
                selector=f"slide:{index}",
                text="\n".join(fragment.strip() for shape in slide.shapes for fragment in _shape_text_fragments(shape) if fragment.strip()),
                source_hash=canonical_sha256(
                    {
                        "selector": f"slide:{index}",
                        "text": "\n".join(fragment.strip() for shape in slide.shapes for fragment in _shape_text_fragments(shape) if fragment.strip()),
                    }
                ),
            )
            for index, slide in enumerate(presentation.slides, start=1)
        )
    except Exception:
        raise DeckCandidateCompilationError("candidate_artifact_invalid") from None


def _normalized_text_token_sequence(value: str) -> tuple[str, ...]:
    """Return ordered visible-text tokens while ignoring layout whitespace."""

    normalized = unicodedata.normalize("NFKC", value)
    return tuple(_VISIBLE_TEXT_TOKEN_PATTERN.findall(normalized))


def _content_proof(
    *,
    baseline: DeckCandidateBaseline,
    candidate_pptx: bytes,
    selectors: tuple[str, ...],
    candidate_editability_score: float,
) -> ContentPreservationProof:
    candidate = _candidate_visible_text(candidate_pptx)
    baseline_by_selector = {item.selector: item for item in baseline.visible_text}
    candidate_by_selector = {item.selector: item for item in candidate}
    text_preserved = tuple(candidate_by_selector) == selectors and all(
        _normalized_text_token_sequence(candidate_by_selector[selector].text)
        == _normalized_text_token_sequence(baseline_by_selector[selector].text)
        for selector in selectors
    )
    raw_baseline_score = baseline.build_record.get("native_editability_score")
    baseline_score = float(raw_baseline_score) if isinstance(raw_baseline_score, (int, float)) and not isinstance(raw_baseline_score, bool) else 0.0
    native_preserved = candidate_editability_score + MAX_NATIVE_EDITABILITY_DROP >= baseline_score
    if not text_preserved or not native_preserved:
        raise DeckCandidateCompilationError("content_changed")
    return ContentPreservationProof(
        brief_preserved=True,
        initial_slide_count=len(selectors),
        candidate_slide_count=len(candidate),
        required_content_preserved=True,
        factual_content_preserved=True,
        native_editability_preserved=True,
    )


def _inventory_slide_projection(value: object) -> dict[str, Any]:
    if not isinstance(value, dict):
        raise DeckCandidateCompilationError("baseline_invalid")
    shapes = value.get("shapes")
    if not isinstance(shapes, list):
        raise DeckCandidateCompilationError("baseline_invalid")
    compact: list[dict[str, Any]] = []
    for shape in shapes:
        if not isinstance(shape, dict):
            raise DeckCandidateCompilationError("baseline_invalid")
        compact.append(
            {
                "name": shape.get("name"),
                "type": shape.get("type"),
                "text_preview": shape.get("text_preview"),
                "full_slide": bool(shape.get("full_slide")),
            }
        )
    compact.sort(key=lambda item: canonical_json_bytes(item))
    return {
        "shape_count": value.get("shape_count", len(shapes)),
        "full_slide_picture_count": value.get("full_slide_picture_count", 0),
        "shapes": compact,
    }


def _native_inventory_proof(
    *,
    baseline: DeckCandidateBaseline,
    candidate_build: Mapping[str, Any],
    unchanged_selectors: tuple[str, ...],
) -> tuple[str, str]:
    baseline_inventory = baseline.build_record.get("native_shape_inventory")
    candidate_inventory = candidate_build.get("native_shape_inventory")
    if not isinstance(baseline_inventory, dict) or not isinstance(candidate_inventory, dict):
        raise DeckCandidateCompilationError("baseline_invalid")
    baseline_projection = {selector: _inventory_slide_projection(baseline_inventory.get(selector)) for selector in unchanged_selectors if selector.startswith("slide:")}
    candidate_projection = {selector: _inventory_slide_projection(candidate_inventory.get(selector)) for selector in unchanged_selectors if selector.startswith("slide:")}
    if baseline_projection != candidate_projection:
        raise DeckCandidateCompilationError("native_inventory_changed")
    return canonical_sha256(candidate_inventory), canonical_sha256(candidate_projection)


def _render_files(outputs: Path, selectors: tuple[str, ...]) -> dict[str, bytes]:
    render_root = outputs / ".builder" / "deck_native" / "rendered"
    try:
        paths = sorted(
            render_root.glob("slide-*.jpg"),
            key=lambda path: int(path.stem.removeprefix("slide-")),
        )
        if len(paths) != len(selectors) or any(path.is_symlink() for path in paths):
            raise ValueError
        result = {selector: path.read_bytes() for selector, path in zip(selectors, paths, strict=True)}
        if any(not content or len(content) > MAX_BASELINE_RENDER_BYTES for content in result.values()):
            raise ValueError
        return result
    except Exception:
        raise DeckCandidateCompilationError("service_result_invalid") from None


def _normalized_render(content: bytes) -> Image.Image:
    try:
        with Image.open(io.BytesIO(content)) as opened:
            image = opened.convert("RGB")
        if image.width < 1 or image.height < 1:
            raise ValueError
        return image.resize((RENDER_COMPARE_WIDTH, RENDER_COMPARE_HEIGHT), Image.Resampling.LANCZOS)
    except Exception:
        raise DeckCandidateCompilationError("service_result_invalid") from None


def _render_pixel_hash(content: bytes) -> str:
    try:
        with Image.open(io.BytesIO(content)) as opened:
            image = opened.convert("RGB")
        if image.width < 1 or image.height < 1:
            raise ValueError
        return canonical_sha256(
            {
                "schema_version": "sophia-render-pixels/v1",
                "width": image.width,
                "height": image.height,
                "rgb_sha256": _sha256(image.tobytes()),
            }
        )
    except Exception:
        raise DeckCandidateCompilationError("service_result_invalid") from None


def _render_proof(
    *,
    baseline: DeckCandidateBaseline,
    candidate_renders: Mapping[str, bytes],
    selectors: tuple[str, ...],
    changed_selectors: tuple[str, ...],
) -> tuple[dict[str, float], tuple[str, ...]]:
    baseline_by_selector = {item.selector: item.content for item in baseline.renders}
    if tuple(candidate_renders) != selectors or tuple(baseline_by_selector) != selectors:
        raise DeckCandidateCompilationError("baseline_invalid")
    changed_slides = {selector for selector in changed_selectors if selector.startswith("slide:")}
    compared = tuple(selector for selector in selectors if selector not in changed_slides)
    deltas: dict[str, float] = {}
    for selector in selectors:
        candidate = _normalized_render(candidate_renders[selector])
        baseline_image = _normalized_render(baseline_by_selector[selector])
        if selector not in compared:
            continue
        difference = ImageChops.difference(baseline_image, candidate)
        mean = sum(ImageStat.Stat(difference).mean) / 3.0
        deltas[selector] = round(float(mean), 4)
        if not math.isfinite(mean) or mean > MAX_UNCHANGED_RENDER_MEAN_DELTA:
            raise DeckCandidateCompilationError("render_collateral_changed")
    return deltas, compared


def _mechanical_projection(record: Mapping[str, Any]) -> MechanicalProjection:
    raw_checks = record.get("checks")
    if not isinstance(raw_checks, dict) or set(raw_checks) != set(_MECHANICAL_CHECK_IDS):
        raise DeckCandidateCompilationError("mechanical_gate_failed")
    if any(raw_checks[check_id] is not True for check_id in _MECHANICAL_CHECK_IDS):
        raise DeckCandidateCompilationError("mechanical_gate_failed")
    checks = tuple(
        MechanicalCheck(
            check_id=check_id,  # type: ignore[arg-type]
            status="passed",
        )
        for check_id in _MECHANICAL_CHECK_IDS
    )
    return MechanicalProjection(
        status="passed",
        checks=checks,
        authoritative_record_hash=canonical_sha256(record),
    )


def _common_identity(request: DeckCandidateCompileRequest) -> dict[str, Any]:
    return {
        "build_id": request.build_id,
        "transaction_id": request.transaction_id,
        "artifact_version_id": request.artifact_version_id,
        "manifest_revision": request.candidate_manifest_revision,
    }


@dataclass(frozen=True, slots=True)
class _CandidateBuildInputs:
    baseline: DeckCandidateBaseline
    sources: dict[tuple[str, str], str]
    changed_selectors: tuple[str, ...]
    unchanged_selectors: tuple[str, ...]
    selectors: tuple[str, ...]
    slides: list[dict[str, Any]]
    output_path: str


@dataclass(frozen=True, slots=True)
class _PreparedCandidateState:
    request: DeckCandidateCompileRequest
    baseline: DeckCandidateBaseline
    result: DeckBuildResult
    selectors: tuple[str, ...]
    changed_selectors: tuple[str, ...]
    unchanged_selectors: tuple[str, ...]
    pptx_bytes: bytes
    derived_sources: tuple[DerivedDeckSource, ...]
    creative_record: dict[str, Any]
    design_record: dict[str, Any]
    inventory_hash: str
    unchanged_inventory_hash: str
    candidate_renders: dict[str, bytes]
    render_deltas: dict[str, float]
    compared_selectors: tuple[str, ...]
    candidate_score: float
    content: ContentPreservationProof
    prepared: PreparedDeckQualityPublication
    source_pack: DeckQualitySourcePack
    source_pack_bytes: bytes
    mechanical_record: dict[str, Any]
    mechanical: MechanicalProjection


def _prepare_build_inputs(
    request: DeckCandidateCompileRequest,
    baseline_value: object,
) -> _CandidateBuildInputs:
    baseline = DeckCandidateBaseline.model_validate(baseline_value)
    sources, changed_selectors, unchanged_selectors = _validate_request_sources(request)
    selectors = _slide_selectors(request)
    _validate_baseline(request, baseline, selectors)
    _validate_visual_asset_inventory(baseline)
    baseline_slide_records = _baseline_slide_records(baseline, selectors)
    slides = _slide_ir(
        selectors=selectors,
        records=baseline_slide_records,
        sources=sources,
    )
    output_path = f"/mnt/user-data/outputs/.builder/builds/{request.build_id}/artifacts/{request.artifact_version_id}/candidate.pptx"
    return _CandidateBuildInputs(
        baseline=baseline,
        sources=sources,
        changed_selectors=changed_selectors,
        unchanged_selectors=unchanged_selectors,
        selectors=selectors,
        slides=slides,
        output_path=output_path,
    )


def _run_candidate_build(
    *,
    root: Path,
    request: DeckCandidateCompileRequest,
    inputs: _CandidateBuildInputs,
    service_factory: DeckBuildServiceFactory,
) -> tuple[SimpleNamespace, DeckBuildResult]:
    runtime = _runtime(
        root=root,
        request=request,
        baseline=inputs.baseline,
        slide_count=len(inputs.selectors),
        output_path=inputs.output_path,
    )
    asset_runner = _ReuseOnlyImageRunner(inputs.baseline.visual_assets)
    service = service_factory(asset_runner.run_batch, asset_runner.forbid_single)
    with langsmith_tracing_disabled():
        result = service.prepare_and_build(
            runtime=runtime,
            deck_title=str(inputs.baseline.build_record.get("deck_title") or ""),
            slides=inputs.slides,
            output_path=inputs.output_path,
            register=str(inputs.baseline.build_record.get("register") or ""),
            visual_policy=str(inputs.baseline.build_record.get("visual_policy") or ""),
            deck_stylesheet=inputs.sources[(DECK_STYLE_ROOT_SELECTOR, "deck_css")],
            authoring_contract=str(inputs.baseline.build_record.get("deck_authoring_contract") or ""),
            style_profile=(inputs.baseline.build_record.get("style_profile") if isinstance(inputs.baseline.build_record.get("style_profile"), dict) else {}),
            design_plan=inputs.baseline.design_plan_record,
            creative_plan=inputs.baseline.creative_plan_record,
            native_lint_slide_indices=tuple(
                int(_SLIDE_SELECTOR_RE.fullmatch(selector).group(1)) - 1
                for selector in inputs.changed_selectors
                if selector != DECK_STYLE_ROOT_SELECTOR
            ),
        )
    return runtime, result


def _prepare_candidate_state(
    *,
    root: Path,
    runtime: SimpleNamespace,
    request: DeckCandidateCompileRequest,
    inputs: _CandidateBuildInputs,
    result: DeckBuildResult,
) -> _PreparedCandidateState:
    baseline = inputs.baseline
    selectors = inputs.selectors
    _validate_service_result(
        result,
        request=request,
        slide_count=len(selectors),
        output_path=inputs.output_path,
        authoring_contract=str(baseline.build_record["deck_authoring_contract"]),
    )
    output_host = Path(replace_virtual_path(inputs.output_path, runtime.state["thread_data"]))
    pptx_bytes = _canonicalize_candidate_artifact(output_host)
    build_host = Path(replace_virtual_path(result.deck_build_path, runtime.state["thread_data"]))
    generated_build = _stable_generated_build_record(_json_object(build_host, code="service_result_invalid"))
    # DQ-1 captures this file into its immutable source pack. Removing
    # execution telemetry keeps retry publication and PPTX hashes stable.
    _write_stable_build_record(build_host, generated_build)
    generated_creative = _json_object(
        root / "outputs" / "deck_build" / "creative_plan.json",
        code="service_result_invalid",
    )
    generated_design = _json_object(
        root / "outputs" / "deck_build" / "design_plan.json",
        code="service_result_invalid",
    )
    common = _common_identity(request)
    creative_record = _plan_record(
        baseline.creative_plan_record,
        generated_creative,
        common=common,
        plan_kind="creative",
    )
    design_record = _plan_record(
        baseline.design_plan_record,
        generated_design,
        common=common,
        plan_kind="design",
    )
    derived_sources = _generated_sources(
        build_record=generated_build,
        request=request,
        sources=inputs.sources,
        selectors=selectors,
    )
    inventory_hash, unchanged_inventory_hash = _native_inventory_proof(
        baseline=baseline,
        candidate_build=generated_build,
        unchanged_selectors=inputs.unchanged_selectors,
    )
    candidate_renders = _render_files(root / "outputs", selectors)
    render_deltas, compared_selectors = _render_proof(
        baseline=baseline,
        candidate_renders=candidate_renders,
        selectors=selectors,
        changed_selectors=inputs.changed_selectors,
    )
    candidate_score = float(result.native_editability_score or 0.0)
    content = _content_proof(
        baseline=baseline,
        candidate_pptx=pptx_bytes,
        selectors=selectors,
        candidate_editability_score=candidate_score,
    )

    artifact_hash = _sha256(pptx_bytes)
    snapshot_path = deck_quality_immutable_artifact_snapshot_path(
        user_id=request.user_id,
        thread_id=request.thread_id,
        build_id=request.build_id,
        logical_artifact_id=baseline.logical_artifact_id,
        artifact_version_id=request.artifact_version_id,
        artifact_sha256=artifact_hash,
        artifact_virtual_path=inputs.output_path,
    )
    prepared = PreparedDeckQualityPublication(
        outputs_root=root / "outputs",
        artifact_virtual_path=inputs.output_path,
        artifact_storage_object_path=snapshot_path,
        artifact_sha256=artifact_hash,
        artifact_id=request.artifact_version_id,
        logical_artifact_id=baseline.logical_artifact_id,
        artifact_version_id=request.artifact_version_id,
        manifest_revision=request.candidate_manifest_revision,
        build_id=request.build_id,
        user_id=request.user_id,
        thread_id=request.thread_id,
        task_id=baseline.task_id,
        builder_run_id=request.operation_id,
        parent_builder_trace_id=baseline.parent_builder_trace_id,
        task_brief=baseline.task_brief,
        mechanical_gate_results=result.mechanical_gate_results,
        source_retention_report=result.source_retention_report,
        native_contrast_report=result.native_contrast_report,
        native_mechanical_report=result.native_mechanical_report,
        native_editability_score=result.native_editability_score,
        missing_expected_visual_count=result.missing_visual_count,
    )
    source_pack, source_pack_bytes = capture_deck_quality_source_pack(
        prepared=prepared,
        instrument=baseline.instrument,
    )
    mechanical_record = source_pack.mechanical_record
    mechanical = _mechanical_projection(mechanical_record)
    return _PreparedCandidateState(
        request=request,
        baseline=baseline,
        result=result,
        selectors=selectors,
        changed_selectors=inputs.changed_selectors,
        unchanged_selectors=inputs.unchanged_selectors,
        pptx_bytes=pptx_bytes,
        derived_sources=derived_sources,
        creative_record=creative_record,
        design_record=design_record,
        inventory_hash=inventory_hash,
        unchanged_inventory_hash=unchanged_inventory_hash,
        candidate_renders=candidate_renders,
        render_deltas=render_deltas,
        compared_selectors=compared_selectors,
        candidate_score=candidate_score,
        content=content,
        prepared=prepared,
        source_pack=source_pack,
        source_pack_bytes=source_pack_bytes,
        mechanical_record=mechanical_record,
        mechanical=mechanical,
    )


def _finish_candidate_state(
    state: _PreparedCandidateState,
    receipt: DeckQualityProducerBundleReceipt,
) -> DeckCandidateCompilation:
    request = state.request
    baseline = state.baseline
    result = state.result
    common = _common_identity(request)
    bundle_path = deck_quality_producer_bundle_path(receipt.quality_run_id)
    bundle_archive_path = deck_quality_producer_archive_path(receipt.quality_run_id)
    if receipt.bundle_object_path not in {bundle_path, bundle_archive_path}:
        raise DeckCandidateCompilationError("publication_failed")

    candidate_build_record = {
        "schema_version": "sophia-deck-candidate-build/v1",
        **common,
        "slide_count": len(state.selectors),
        "candidate_hash": request.candidate_hash,
        "deck_route": result.deck_route,
        "deck_compile_mode": result.deck_compile_mode,
        "deck_authoring_contract": result.deck_authoring_contract,
        "native_required": result.native_required,
        "native_editability_score": state.candidate_score,
        "full_slide_picture_count": result.full_slide_picture_count,
        "expected_visual_count": result.expected_visual_count,
        "successful_visual_count": result.successful_visual_count,
        "missing_visual_count": result.missing_visual_count,
        "authoritative_build_record_hash": canonical_sha256(state.source_pack.build_record),
        "source_graph_hash": canonical_sha256(
            {
                "sources": [
                    {
                        "selector": source.selector,
                        "source_role": source.source_role,
                        "source_hash": source.source_hash,
                        "model_authored": source.model_authored,
                        "component_version_changed": source.component_version_changed,
                    }
                    for source in request.sources
                ],
                "derived_targets": request.derived_source_targets,
            }
        ),
        "foundation_status": "disabled_for_dq2_candidate",
    }
    native_record = {
        "schema_version": "sophia-deck-candidate-native/v1",
        **common,
        "verified": True,
        "native_editable": True,
        "slide_count": len(state.selectors),
        "native_editability_score": state.candidate_score,
        "full_slide_picture_count": 0,
        "native_shape_inventory_hash": state.inventory_hash,
        "unchanged_native_inventory_hash": state.unchanged_inventory_hash,
    }
    render_record = {
        "schema_version": "sophia-deck-candidate-render-collateral/v1",
        **common,
        "verified": True,
        "within_tolerance": True,
        "expected_selectors": list(state.selectors),
        "rendered_selectors": list(state.selectors),
        "compared_selectors": list(state.compared_selectors),
        "mean_pixel_delta_by_selector": state.render_deltas,
        "tolerance_mean_delta": MAX_UNCHANGED_RENDER_MEAN_DELTA,
        "render_hash_kind": "decoded-rgb-v1",
        "render_hashes": {selector: _render_pixel_hash(state.candidate_renders[selector]) for selector in state.selectors},
    }
    locality = LocalityProof(
        authorized_selectors=request.program.authorized_selectors,
        changed_component_versions=state.changed_selectors,
        unchanged_component_versions=state.unchanged_selectors,
        unexpected_changes=(),
        shared_dependency_changed=DECK_STYLE_ROOT_SELECTOR in state.changed_selectors,
        native_inventory_preserved=True,
        render_collateral_within_tolerance=True,
    )
    publication = {
        "schema_version": "sophia-deck-candidate-dq1-publication/v1",
        **common,
        "quality_run_id": receipt.quality_run_id,
        "bundle_object_path": bundle_path,
        "bundle_archive_object_path": bundle_archive_path,
        "bundle_hash": receipt.bundle_hash,
        "bundle_size_bytes": receipt.bundle_size_bytes,
        "artifact_snapshot_object_path": state.prepared.artifact_storage_object_path,
        "artifact_hash": state.prepared.artifact_sha256,
        "instrument_identity_hash": canonical_sha256(baseline.instrument.lock),
    }
    for record in (candidate_build_record, native_record, render_record, publication):
        if len(canonical_json_bytes(record)) > MAX_NATIVE_RECORD_BYTES:
            raise DeckCandidateCompilationError("service_result_invalid")
    return DeckCandidateCompilation(
        pptx_bytes=state.pptx_bytes,
        derived_sources=state.derived_sources,
        build_record=candidate_build_record,
        creative_plan_record=state.creative_record,
        design_plan_record=state.design_record,
        mechanical_record=state.mechanical_record,
        mechanical=state.mechanical,
        native_record=native_record,
        render_collateral_record=render_record,
        locality=locality,
        content=state.content,
        dq1_publication_metadata=publication,
    )


class ProductionDeckCandidateCompiler:
    """Compile and publish a single deterministic DQ-2 deck candidate."""

    def __init__(
        self,
        *,
        baseline_loader: DeckCandidateBaselineLoader,
        publisher: CandidateDq1Publisher | None = None,
        service_factory: DeckBuildServiceFactory = _default_service_factory,
    ) -> None:
        self._baseline_loader = baseline_loader
        self._publisher = publisher or DurableCandidateDq1Publisher()
        self._service_factory = service_factory

    async def compile(self, request: DeckCandidateCompileRequest) -> DeckCandidateCompilation:
        try:
            baseline_value = await _call_maybe_async(
                self._baseline_loader.load,
                request,
            )
        except DeckCandidateCompilationError:
            raise
        except Exception:
            raise DeckCandidateCompilationError("baseline_unavailable") from None

        try:
            inputs = await _run_sync(
                _prepare_build_inputs,
                request,
                baseline_value,
            )
        except DeckCandidateCompilationError:
            raise
        except (TypeError, ValueError, ValidationError):
            raise DeckCandidateCompilationError("baseline_invalid") from None

        temporary = await _run_sync(
            tempfile.TemporaryDirectory,
            prefix="sophia-dq2-candidate-",
        )
        try:
            root = Path(temporary.name)
            try:
                runtime, result = await _run_sync(
                    _run_candidate_build,
                    root=root,
                    request=request,
                    inputs=inputs,
                    service_factory=self._service_factory,
                )
            except DeckCandidateCompilationError:
                raise
            except Exception:
                raise DeckCandidateCompilationError("service_failed") from None

            try:
                state = await _run_sync(
                    _prepare_candidate_state,
                    root=root,
                    runtime=runtime,
                    request=request,
                    inputs=inputs,
                    result=result,
                )
            except DeckCandidateCompilationError:
                raise
            except Exception:
                raise DeckCandidateCompilationError("service_result_invalid") from None

            try:
                receipt = await _call_maybe_async(
                    self._publisher.publish,
                    prepared=state.prepared,
                    instrument=state.baseline.instrument,
                    pptx_bytes=state.pptx_bytes,
                    source_pack=state.source_pack,
                    source_pack_bytes=state.source_pack_bytes,
                )
            except DeckCandidateCompilationError:
                raise
            except Exception:
                raise DeckCandidateCompilationError("publication_failed") from None
            if not isinstance(receipt, DeckQualityProducerBundleReceipt) or receipt.quality_run_id != state.source_pack.quality_run_id:
                raise DeckCandidateCompilationError("publication_failed")
            return await _run_sync(_finish_candidate_state, state, receipt)
        finally:
            # Temp files are still inputs to the durable publication. Shield
            # cleanup so cancellation cannot leave a worker using a deleted
            # workspace or return while filesystem cleanup is still running.
            with anyio.CancelScope(shield=True):
                await _run_sync(temporary.cleanup)


__all__ = [
    "BaselineDeckRender",
    "BaselineVisualAsset",
    "CandidateDq1Publisher",
    "DeckCandidateBaseline",
    "DeckCandidateBaselineLoader",
    "DeckCandidateCompilationError",
    "DurableCandidateDq1Publisher",
    "ProductionDeckCandidateCompiler",
    "baseline_from_authenticated_snapshot",
]
