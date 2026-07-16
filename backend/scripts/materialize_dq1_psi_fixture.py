from __future__ import annotations

import argparse
import hashlib
import json
import shutil
from pathlib import Path
from typing import Any

from pptx import Presentation

from deerflow.sophia.deck_quality.brief import sanitize_current_request
from deerflow.sophia.deck_quality.canonical import canonical_sha256, file_sha256
from deerflow.sophia.deck_quality.visible_text import visible_text_sidecar


def _write_json(path: Path, value: Any) -> None:
    path.write_text(json.dumps(value, indent=2, sort_keys=True, ensure_ascii=False) + "\n", encoding="utf-8")


def _runs_named(runs: list[dict[str, Any]], name: str) -> list[dict[str, Any]]:
    return sorted(
        (run for run in runs if run.get("name") == name),
        key=lambda run: str(run.get("start_time") or ""),
    )


def _passed(value: Any) -> bool:
    if isinstance(value, bool):
        return value
    if isinstance(value, dict):
        if "passed" in value:
            return bool(value["passed"])
        return str(value.get("status") or "").casefold() in {"pass", "passed", "success", "ok"}
    return False


def _extract_visible_text(pptx_path: Path) -> list[dict[str, Any]]:
    presentation = Presentation(str(pptx_path))
    slides: list[tuple[str, list[str]]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        fragments = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if isinstance(text, str) and text.strip():
                fragments.append(text)
        slides.append((f"slide:{index}", fragments))
    return [item.model_dump(mode="json") for item in visible_text_sidecar(slides)]


def materialize(*, trace_path: Path, pptx_path: Path, render_dir: Path, output_dir: Path) -> None:
    trace = json.loads(trace_path.read_text(encoding="utf-8"))
    runs = trace.get("runs") or []
    root = next(run for run in runs if not run.get("parent_run_id"))
    prepare = _runs_named(runs, "prepare_deck_build")[-1]
    mechanical_run = _runs_named(runs, "deck.native.mechanical_report")[-1]
    gates_run = _runs_named(runs, "deck.mechanical_gates.evaluate")[-1]
    terminal_run = _runs_named(runs, "deck.terminal")[-1]
    emit_run = _runs_named(runs, "deck.emit.decision")[-1]

    prepare_inputs = prepare["inputs"]
    creative_plan = prepare_inputs["creative_plan"]
    design_plan = creative_plan["design_plan"]
    user_messages = root.get("inputs", {}).get("messages") or []
    request = sanitize_current_request(str(user_messages[0].get("content") or ""))
    brief = {
        "request": request,
        "subject": creative_plan["subject"],
        "audience": creative_plan["audience"],
        "goal": creative_plan["goal"],
        "viewing_context": creative_plan["viewing_context"],
        "explicit_brand_style_constraints": [
            "restrained editorial system",
            "warm ivory, ink, muted cobalt, and one ember accent",
            "native editable text and shapes",
        ],
    }

    output_dir.mkdir(parents=True, exist_ok=True)
    fixture_pptx = output_dir / "artifact.pptx"
    if pptx_path.resolve() != fixture_pptx.resolve():
        shutil.copyfile(pptx_path, fixture_pptx)
    fixture_renders = output_dir / "renders"
    fixture_renders.mkdir(exist_ok=True)
    for source in sorted(render_dir.glob("*.png")):
        target = fixture_renders / source.name
        if source.resolve() != target.resolve():
            shutil.copyfile(source, target)

    _write_json(output_dir / "brief.json", brief)
    _write_json(output_dir / "creative_plan.json", creative_plan)
    _write_json(output_dir / "design_plan.json", design_plan)
    _write_json(output_dir / "visible_text.json", _extract_visible_text(fixture_pptx))

    mechanical_outputs = mechanical_run.get("outputs") or {}
    gate_outputs = gates_run.get("outputs") or {}
    terminal_outputs = terminal_run.get("outputs") or {}
    emit_outputs = emit_run.get("outputs") or {}
    gate_passed = _passed(gate_outputs)
    terminal_passed = str(terminal_outputs.get("status") or "").casefold() in {
        "success",
        "completed",
    }
    lint_passed = bool(mechanical_outputs.get("lint_fix_success")) and int(mechanical_outputs.get("lint_residue_count") or 0) == 0
    checks = {
        "authoritative_gate": gate_passed and terminal_passed,
        "source_retention": _passed(mechanical_outputs.get("source_retention")),
        "native_editability": gate_passed and float(mechanical_outputs.get("native_editability_score") or 0) > 0,
        "contrast": _passed(mechanical_outputs.get("contrast")),
        "native_lint": lint_passed,
        "overflow_collision_clipping": gate_passed,
        "render_success": bool(mechanical_outputs.get("render_success")),
        "visual_asset_completeness": int(emit_outputs.get("missing_visual_count") or 0) == 0,
        "artifact_identity": file_sha256(fixture_pptx) == file_sha256(pptx_path),
    }
    mechanical_report = {
        "schema_version": "deck-quality-fixture-mechanical/v1",
        "checks": checks,
        "authoritative_mechanical_report": mechanical_outputs,
        "authoritative_gate_result": gate_outputs,
        "terminal_status": terminal_outputs.get("status"),
        "source_trace_run_started_at": mechanical_run.get("start_time"),
    }
    _write_json(output_dir / "mechanical_report.json", mechanical_report)

    render_hashes = {path.name: file_sha256(path) for path in sorted(fixture_renders.glob("*.png"))}
    source_hashes = {
        "artifact.pptx": file_sha256(fixture_pptx),
        "brief.json": canonical_sha256(brief),
        "creative_plan.json": canonical_sha256(creative_plan),
        "design_plan.json": canonical_sha256(design_plan),
        "mechanical_report.json": canonical_sha256(mechanical_report),
        "deck_stylesheet": hashlib.sha256(str(prepare_inputs["deck_stylesheet"]).encode()).hexdigest(),
        "slides_input": canonical_sha256(prepare_inputs["slides"]),
    }
    manifest = {
        "schema_version": "deck-quality-fixture-bundle/v1",
        "fixture_id": "clean_underdesigned_psi_v1",
        "source": "synthetic production canary",
        "trace_id": trace.get("trace_id"),
        "trace_run_count": trace.get("run_count"),
        "build_id": "build_01KXKNNQ5Z9N198VCMJPDWSBJ0",
        "artifact_id": "artifact_6726b07cedb246eb14a5eabf",
        "artifact_version_id": "artifact_version_01KXKNV1DJ7B4ABS00FN20G1SK",
        "source_hashes": source_hashes,
        "render_hashes": render_hashes,
    }
    _write_json(output_dir / "manifest.json", manifest)


def main() -> int:
    parser = argparse.ArgumentParser(description="Materialize the synthetic DQ-1 PSI fixture")
    parser.add_argument("--trace", type=Path, required=True)
    parser.add_argument("--pptx", type=Path, required=True)
    parser.add_argument("--render-dir", type=Path, required=True)
    parser.add_argument("--output-dir", type=Path, required=True)
    args = parser.parse_args()
    materialize(
        trace_path=args.trace,
        pptx_path=args.pptx,
        render_dir=args.render_dir,
        output_dir=args.output_dir,
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
